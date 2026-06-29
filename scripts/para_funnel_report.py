"""Para Performance Funnels.

USOPC-style "performance funnel" for para triathlon: x-axis = years of experience
(since first para race in our DB), y-axis = Elo rating. Gray band = medalist
benchmark cohort; blue lines = top USA athletes. One chart per category.

Output: standalone PPTX (Para_Performance_Funnels.pptx) with 4 category slides + notes.

Usage:
    python scripts/para_funnel_report.py
"""
from __future__ import annotations

import datetime as dt
import os
import sys
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from sqlalchemy import text
from sqlalchemy.engine import Engine

project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

from tri_analysis.database import get_engine
from tri_analysis.elo_trajectory import (
    PARA_HISTORY_FLOOR,
    build_pool_trajectory,
)

# ---------- Style (mirrors scripts/para_progression_report.py) ----------
USAT_RED   = "#C8102E"
USAT_NAVY  = "#002868"
USAT_GOLD  = "#B8860B"
GREY_BAND  = "#9E9E9E"
GREY_LINE  = "#B5B5B5"
GREY_DARK  = "#555555"

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.titleweight": "bold",
    "figure.dpi": 150,
})

USA_COUNTRY_ALIASES = ["USA", "United States", "United States of America", "U.S.A", "US"]

# Slide order
CATEGORIES = ["PTVI Women", "PTVI Men", "PTS3 Women", "PTWC Women"]
USA_TOP_N = 3

# Per-category benchmark cohort definition.
# Paralympic categories: Tokyo 2020 + Paris 2024 + Wollongong 2025 Worlds medalists.
# PTS3 Women: never a Paralympic event -> all World Championship medalists since 2021.
PTS3_WORLDS_COHORT = {"PTS3 Women"}


# ---------- Cohort / USA selection ----------
def benchmark_cohort_ids(engine: Engine, category: str) -> list[int]:
    """Top-3 finisher athlete_ids defining the medal standard for a category."""
    if category in PTS3_WORLDS_COHORT:
        q = text(
            """
            SELECT DISTINCT r.athlete_id
            FROM events e JOIN race_results r
              ON e.event_id = r.event_id AND e.prog_id = r.prog_id
            WHERE e.is_para = TRUE AND e.prog_name = :cat
              AND r.finish_status = 'FINISH' AND r.finish_position IN (1, 2, 3)
              AND e.cat_name ILIKE '%World Championships%'
              AND e.event_date >= '2021-01-01'
            """
        )
    else:
        q = text(
            """
            SELECT DISTINCT r.athlete_id
            FROM events e JOIN race_results r
              ON e.event_id = r.event_id AND e.prog_id = r.prog_id
            WHERE e.is_para = TRUE AND e.prog_name = :cat
              AND r.finish_status = 'FINISH' AND r.finish_position IN (1, 2, 3)
              AND (
                (e.cat_name ILIKE '%Major Games%' AND e.event_date >= '2016-01-01')
                OR (e.event_name ILIKE '%Wollongong%' AND e.event_date >= '2025-01-01')
              )
            """
        )
    with engine.begin() as c:
        df = pd.read_sql(q, c, params={"cat": category})
    return [int(x) for x in df["athlete_id"].tolist()]


def usa_top_ids(engine: Engine, category: str, n: int = USA_TOP_N) -> list[int]:
    """Top-N most active USA athletes in the category (by distinct event count)."""
    aliases = ", ".join(f"'{a}'" for a in USA_COUNTRY_ALIASES)
    q = text(
        f"""
        SELECT a.athlete_id, COUNT(DISTINCT e.event_id) AS n_events
        FROM race_results r
        JOIN events e ON e.event_id = r.event_id AND e.prog_id = r.prog_id
        JOIN athlete a ON a.athlete_id = r.athlete_id
        WHERE e.is_para = TRUE AND e.prog_name = :cat
          AND a.country IN ({aliases})
          AND r.finish_status = 'FINISH'
        GROUP BY a.athlete_id
        ORDER BY n_events DESC
        LIMIT :n
        """
    )
    with engine.begin() as c:
        df = pd.read_sql(q, c, params={"cat": category, "n": n})
    return [int(x) for x in df["athlete_id"].tolist()]


def name_map(engine: Engine, ids: list[int]) -> dict[int, str]:
    if not ids:
        return {}
    id_list = ", ".join(str(int(a)) for a in ids)
    q = text(f"SELECT athlete_id, full_name FROM athlete WHERE athlete_id IN ({id_list})")
    with engine.begin() as c:
        df = pd.read_sql(q, c)
    return {int(r["athlete_id"]): r["full_name"] for _, r in df.iterrows() if pd.notna(r["full_name"])}


def short_name(full: str) -> str:
    """Strip sub-class suffix (B1/B3/H2) and keep first + last name."""
    import re
    n = re.sub(r"\s+(B[123]|H[12])$", "", str(full).strip())
    return n


# ---------- Band ----------
def _rolling3(s: pd.Series) -> pd.Series:
    return s.rolling(window=3, min_periods=1, center=True).mean()


def build_band(traj_cohort: pd.DataFrame) -> pd.DataFrame:
    """Per-athlete 3-race smoothing, then p25/p50/p75 of Elo per integer experience-year bin."""
    if traj_cohort.empty:
        return pd.DataFrame(columns=["exp_bin", "p25", "p50", "p75", "n"])
    df = traj_cohort.sort_values(["athlete_id", "event_date"]).copy()
    df["elo_smooth"] = df.groupby("athlete_id")["elo_after"].transform(_rolling3)
    df["exp_bin"] = np.floor(df["years_experience"]).astype(int)
    # one value per athlete per bin (latest in bin) to avoid over-weighting busy seasons
    per = df.sort_values("event_date").groupby(["athlete_id", "exp_bin"]).tail(1)
    agg = per.groupby("exp_bin")["elo_smooth"].agg(
        p25=lambda s: s.quantile(0.25),
        p50="median",
        p75=lambda s: s.quantile(0.75),
        n="count",
    ).reset_index()
    return agg


# ---------- Plot ----------
def plot_funnel(traj: pd.DataFrame, cohort_ids: list[int], usa_ids: list[int],
                category: str, names: dict[int, str], out_png: Path):
    cohort = traj[traj["athlete_id"].isin(cohort_ids)].copy()
    band = build_band(cohort)

    fig, ax = plt.subplots(figsize=(13, 7))

    # Band: p25-p75 envelope + median, both clipped to the "reliable" region (>=3 athletes/bin).
    # Beyond that the cohort thins to 1-2 athletes, so a "median standard" there is just one person.
    if not band.empty:
        solid = band[band["n"] >= 3]
        if len(solid) >= 2:
            ax.fill_between(solid["exp_bin"], solid["p25"], solid["p75"],
                            color=GREY_BAND, alpha=0.25, zorder=1,
                            label="Medalist 25th–75th pct")
            ax.plot(solid["exp_bin"], solid["p50"], color=GREY_DARK, linewidth=2.4,
                    linestyle="--", zorder=4, label="Medalist median")
            # faint dotted continuation where the cohort thins (informational only)
            thin = band[band["exp_bin"] >= solid["exp_bin"].max()]
            if len(thin) >= 2:
                ax.plot(thin["exp_bin"], thin["p50"], color=GREY_DARK, linewidth=1.2,
                        linestyle=":", alpha=0.5, zorder=4)

    # Individual cohort trajectories (thin gray)
    for aid, g in cohort.groupby("athlete_id"):
        g = g.sort_values("event_date")
        ax.plot(g["years_experience"], _rolling3(g["elo_after"]),
                color=GREY_LINE, linewidth=1.0, alpha=0.45, zorder=2)

    # USA highlighted trajectories (blue), labeled at line end
    usa_colors = [USAT_NAVY, USAT_RED, USAT_GOLD, "#1f6feb", "#8B0000"]
    for i, aid in enumerate(usa_ids):
        g = traj[traj["athlete_id"] == aid].sort_values("event_date")
        if g.empty:
            continue
        color = usa_colors[i % len(usa_colors)]
        y = _rolling3(g["elo_after"])
        ax.plot(g["years_experience"], y, color=color, linewidth=2.6,
                marker="o", markersize=4, zorder=6,
                label=short_name(names.get(aid, str(aid))))
        # left-censored marker on first point
        first = g.iloc[0]
        if bool(first["left_censored"]):
            ax.scatter([first["years_experience"]], [y.iloc[0]], s=110,
                       facecolor="white", edgecolor=color, linewidth=2.0, zorder=7)
            ax.annotate("≥", (first["years_experience"], y.iloc[0]),
                        textcoords="offset points", xytext=(-12, -4),
                        fontsize=12, color=color, fontweight="bold")

    ax.set_xlabel("Years of experience (since first para race in database)", fontsize=11)
    ax.set_ylabel("Elo rating (within-category pool)", fontsize=11)
    ax.set_title(f"{category} — Performance Funnel vs Medalist Standard",
                 fontsize=15, pad=12)
    ax.grid(True, linestyle=":", alpha=0.4)
    ax.legend(frameon=False, fontsize=9, loc="lower right", ncol=2)
    ax.set_xlim(left=-0.3)
    fig.tight_layout()
    fig.savefig(out_png, bbox_inches="tight")
    plt.close(fig)
    return band, cohort


# ---------- PPTX ----------
def build_pptx(figs: list[tuple[str, Path, str]], notes: list[str], outdir: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00, 0x28, 0x68)
    RED  = RGBColor(0xC8, 0x10, 0x2E)
    GREY = RGBColor(0x55, 0x55, 0x55)
    BLANK = prs.slide_layouts[6]

    def title_bar(slide, title, subtitle=""):
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.runs[0].font.size = Pt(26); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
        if subtitle:
            p2 = tf.add_paragraph(); p2.text = subtitle
            p2.runs[0].font.size = Pt(13); p2.runs[0].font.color.rgb = GREY

    # Title slide
    slide = prs.slides.add_slide(BLANK)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.1), Inches(2.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Para Performance Funnels"
    p.runs[0].font.size = Pt(48); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = "USA athletes vs the medalist standard — Elo rating by years of experience"
    p2.runs[0].font.size = Pt(22); p2.runs[0].font.color.rgb = RED
    p3 = tf.add_paragraph()
    p3.text = "Benchmark = Paralympic (Tokyo 2020 + Paris 2024) + World Championship medalists. PTVI · PTWC · PTS3."
    p3.runs[0].font.size = Pt(13); p3.runs[0].font.color.rgb = GREY

    # Category slides
    for title, png, caption in figs:
        slide = prs.slides.add_slide(BLANK)
        title_bar(slide, title)
        slide.shapes.add_picture(str(png), Inches(0.5), Inches(1.15), width=Inches(12.3))
        if caption:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = caption
            p.runs[0].font.size = Pt(10); p.runs[0].font.color.rgb = GREY

    # Notes slide
    slide = prs.slides.add_slide(BLANK)
    title_bar(slide, "How to Read These Funnels — Notes & Caveats")
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + n
        p.runs[0].font.size = Pt(13); p.runs[0].font.color.rgb = GREY
        p.space_after = Pt(6)

    out = outdir / "Para_Performance_Funnels.pptx"
    prs.save(str(out))
    return out


# ---------- Main ----------
def main():
    engine = get_engine()
    stamp = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    outdir = Path("para_triathlon_analysis/output") / f"funnels_{stamp}"
    outdir.mkdir(parents=True, exist_ok=True)

    figs = []
    cat_label = {"PTVI Women": "Women's PTVI", "PTVI Men": "Men's PTVI",
                 "PTS3 Women": "Women's PTS3", "PTWC Women": "Women's PTWC"}

    for category in CATEGORIES:
        traj = build_pool_trajectory(engine, category)
        cohort_ids = benchmark_cohort_ids(engine, category)
        usa_ids = usa_top_ids(engine, category)
        ids_for_names = list(set(cohort_ids) | set(usa_ids))
        names = name_map(engine, ids_for_names)

        png = outdir / f"funnel_{category.replace(' ', '_')}.png"
        band, cohort = plot_funnel(traj, cohort_ids, usa_ids, category, names, png)

        n_cohort = cohort["athlete_id"].nunique()
        usa_names = ", ".join(short_name(names.get(a, str(a))) for a in usa_ids)
        bench_src = ("World Championship medalists since 2021"
                     if category in PTS3_WORLDS_COHORT
                     else "Tokyo 2020 + Paris 2024 + Wollongong 2025 medalists")
        caption = (f"Gray = {n_cohort} benchmark athletes ({bench_src}). "
                   f"Blue = top USA by event count: {usa_names}. "
                   f"Open marker (≥) = athlete's career predates our 2017 data floor.")
        figs.append((f"{cat_label[category]} — Performance Funnel", png, caption))
        print(f"[{category}] cohort={n_cohort} athletes, USA={usa_names}")

    notes = [
        "X-axis is YEARS OF EXPERIENCE, not age — para athletes enter the sport at very different ages, "
        "so experience (time since first para race) is the fairer development clock.",
        "Y-axis is Elo rating computed WITHIN each category pool (PTVI, PTWC, PTS3 are separate competition "
        "pools). Ratings are not comparable across categories. All athletes start at 1500.",
        "The gray band is the 25th–75th percentile of medalist trajectories; the dashed line is the median — "
        "this is the 'medal-capable' development corridor. Blue USA lines below the band are developing toward it.",
        "DATA FLOOR: our para history begins 2017-03-11. Athletes who competed before then are left-censored "
        "(open '≥' marker) — their true experience is higher than plotted, so their curve should be read as "
        "shifted right.",
        "AGE: a true age overlay was not feasible — only ~19% of benchmark medalists have a birth year on file. "
        "Qualitatively, PTVI/PTWC medalists tend to peak (highest Elo) around 6–10 years of experience, with "
        "several still climbing past year 10; the band's downturn at the far right reflects late-career medalists "
        "rather than a hard experience ceiling.",
        "A USA athlete who is also a medalist (e.g. Kendall Gretsch) appears BOTH as a blue line and inside the "
        "gray band — they are part of the standard they are measured against.",
        "Source: World Triathlon official results (race_results, events). Elo: pairwise, field-size-scaled K, "
        "tier-weighted (Worlds/Paralympics weighted highest).",
    ]
    pptx = build_pptx(figs, notes, outdir)
    print(f"\nDeck: {pptx}")
    print(f"All outputs in: {outdir}")


if __name__ == "__main__":
    main()
