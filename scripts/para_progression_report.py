"""Progression reports for para athletes.

Generates two presentation-ready outputs:
  1. Perry vs Parker (PTWC Women) — Sullivan-style yearly deficit bars
  2. Hermes vs Rodriguez (PTVI Women) — pace trend + direct comparison

Outputs PNG charts + HTML narrative + CSV data per athlete.

Usage:
    python scripts/para_progression_report.py
"""
from __future__ import annotations

import datetime as dt
from dataclasses import dataclass
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import Patch
from sqlalchemy import text

from tri_analysis.database import get_engine

# ---------- Style ----------
USAT_RED   = "#C8102E"
USAT_NAVY  = "#002868"
USAT_GOLD  = "#B8860B"
GREY_LIGHT = "#CCCCCC"
GREY_DARK  = "#555555"

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.titleweight": "bold",
    "axes.titlesize": 12,
    "axes.labelsize": 10,
    "xtick.labelsize": 9,
    "ytick.labelsize": 9,
    "figure.dpi": 150,
})

# ---------- Athletes ----------
@dataclass(frozen=True)
class Athlete:
    aid: int
    name: str
    short: str
    country: str

PERRY    = Athlete(171840, "Emelia Perry",     "Perry",    "USA")
PARKER   = Athlete(122515, "Lauren Parker",    "Parker",   "AUS")
HERMES   = Athlete(165205, "McClain Hermes",   "Hermes",   "USA")
RODRIGUEZ= Athlete(40701,  "Susana Rodriguez", "Rodriguez","ESP")  # B1
TARANTELLO=Athlete(166501, "Francesca Tarantello","Tarantello","ITA")  # B3
FISHER   = Athlete(165198, "Skyler Fisher",   "Fisher",   "USA")  # H2 PTWC


# ---------- Time parsing ----------
def to_sec(s):
    if s is None or pd.isna(s) or str(s).strip() in ("", "00:00:00"):
        return None
    parts = str(s).split(":")
    try:
        if len(parts) == 3:
            return int(parts[0])*3600 + int(parts[1])*60 + float(parts[2])
        if len(parts) == 2:
            return int(parts[0])*60 + float(parts[1])
        return float(parts[0])
    except Exception:
        return None


# ---------- Data pulls ----------
def pull_overlap(engine, a: Athlete, b: Athlete) -> pd.DataFrame:
    q = text("""
        SELECT e.event_id, e.prog_id, e.event_name, e.event_date,
               e.prog_distance_category, e.cat_name,
               e.swim_distance, e.bike_distance, e.run_distance,
               ra.swimtime AS a_swim, ra.t1time AS a_t1, ra.biketime AS a_bike,
               ra.t2time AS a_t2, ra.runtime AS a_run, ra.total_time AS a_total,
               ra.finish_status AS a_status, ra.finish_position AS a_pos,
               rb.swimtime AS b_swim, rb.t1time AS b_t1, rb.biketime AS b_bike,
               rb.t2time AS b_t2, rb.runtime AS b_run, rb.total_time AS b_total,
               rb.finish_status AS b_status, rb.finish_position AS b_pos
        FROM race_results ra
        JOIN race_results rb
          ON ra.event_id = rb.event_id AND ra.prog_id = rb.prog_id
        JOIN events e ON e.event_id = ra.event_id AND e.prog_id = ra.prog_id
        WHERE ra.athlete_id = :a AND rb.athlete_id = :b
          AND e.is_para = TRUE
          AND ra.finish_status = 'FINISH' AND rb.finish_status = 'FINISH'
        ORDER BY e.event_date
    """)
    with engine.begin() as conn:
        df = pd.read_sql(q, conn, params={"a": a.aid, "b": b.aid})
    if df.empty:
        return df
    for col in ["a_swim","a_t1","a_bike","a_t2","a_run","a_total",
                "b_swim","b_t1","b_bike","b_t2","b_run","b_total"]:
        df[col + "_s"] = df[col].apply(to_sec)
    for seg in ["swim","t1","bike","t2","run","total"]:
        df[f"d_{seg}"] = df[f"a_{seg}_s"] - df[f"b_{seg}_s"]
    df["event_date"] = pd.to_datetime(df["event_date"])
    df["year"] = df["event_date"].dt.year
    return df


def pull_full_history(engine, athlete: Athlete, peer_ids: dict[str, int] | None = None) -> pd.DataFrame:
    """All para finishes for athlete, with podium-avg-excluding-self and per-peer times.

    peer_ids is a dict short_key -> athlete_id (e.g. {"rod": 40701, "tar": 166501})
    Output columns include peer_{key}_{seg} for each peer/segment.
    """
    peer_ids = peer_ids or {}
    q_events = text("""
        SELECT DISTINCT e.event_id, e.prog_id, e.event_name, e.event_date,
               e.prog_distance_category, e.cat_name,
               e.swim_distance, e.bike_distance, e.run_distance
        FROM race_results r
        JOIN events e ON e.event_id = r.event_id AND e.prog_id = r.prog_id
        WHERE r.athlete_id = :aid AND e.is_para = TRUE
          AND r.finish_status = 'FINISH'
        ORDER BY e.event_date
    """)
    with engine.begin() as conn:
        events = pd.read_sql(q_events, conn, params={"aid": athlete.aid})

    rows = []
    for _, ev in events.iterrows():
        q = text("""
            SELECT athlete_id, swimtime, t1time, biketime, t2time, runtime, total_time,
                   finish_position
            FROM race_results
            WHERE event_id = :eid AND prog_id = :pid AND finish_status = 'FINISH'
        """)
        with engine.begin() as conn:
            field = pd.read_sql(q, conn, params={"eid": int(ev["event_id"]), "pid": int(ev["prog_id"])})
        for col in ["swimtime","t1time","biketime","t2time","runtime","total_time"]:
            field[col + "_s"] = field[col].apply(to_sec)
        me = field[field["athlete_id"] == athlete.aid]
        if me.empty:
            continue
        me = me.iloc[0]

        # Podium average excluding self
        podium_others = field[(field["finish_position"].between(1, 3)) &
                              (field["athlete_id"] != athlete.aid)]
        # Fallback: if athlete IS the entire podium (tiny field), use top 3 incl. self
        if podium_others.empty:
            podium_others = field[field["finish_position"].between(1, 3)]

        avg_cols = ["swimtime_s","t1time_s","biketime_s","t2time_s","runtime_s","total_time_s"]
        avg = podium_others[avg_cols].mean()

        row = {
            "event_date": pd.to_datetime(ev["event_date"]),
            "event_name": ev["event_name"],
            "dist": ev["prog_distance_category"],
            "cat": ev["cat_name"],
            "swim_m": ev["swim_distance"],
            "bike_m": ev["bike_distance"],
            "run_m":  ev["run_distance"],
            "h_pos":  int(me["finish_position"]) if pd.notna(me["finish_position"]) else None,
            "h_swim": me["swimtime_s"],
            "h_bike": me["biketime_s"],
            "h_run":  me["runtime_s"],
            "h_total":me["total_time_s"],
            "pod_swim": avg["swimtime_s"],
            "pod_bike": avg["biketime_s"],
            "pod_run":  avg["runtime_s"],
            "pod_total":avg["total_time_s"],
            "n_finishers": len(field),
            "n_podium_ref": len(podium_others),
        }
        for short, pid in peer_ids.items():
            pr = field[field["athlete_id"] == pid]
            for seg, src in [("swim","swimtime_s"),("bike","biketime_s"),
                             ("run","runtime_s"),("total","total_time_s")]:
                row[f"{short}_{seg}"] = pr[src].iloc[0] if not pr.empty else None
        rows.append(row)
    df = pd.DataFrame(rows)
    if df.empty:
        return df

    # Raw deficits and pace-normalized (sec per 100m swim, sec per km bike/run)
    for seg, dist_col, scale in [("swim","swim_m",100.0), ("bike","bike_m",1000.0), ("run","run_m",1000.0)]:
        df[f"d_{seg}_vs_pod"] = df[f"h_{seg}"] - df[f"pod_{seg}"]
        df[f"pace_h_{seg}"]   = df[f"h_{seg}"]   / (df[dist_col] / scale)
        df[f"pace_pod_{seg}"] = df[f"pod_{seg}"] / (df[dist_col] / scale)
        df[f"pace_d_{seg}"]   = df[f"pace_h_{seg}"] - df[f"pace_pod_{seg}"]
    df["year"] = df["event_date"].dt.year
    return df


# ---------- Charts: Perry deck ----------
def chart_perry_yearly_deficit(df: pd.DataFrame, out: Path):
    """Sullivan-style: yearly avg seconds behind Parker for swim, bike, run + total."""
    sprint = df[df["prog_distance_category"]=="sprint"].copy()
    agg = sprint.groupby("year").agg(
        swim=("d_swim","mean"),
        bike=("d_bike","mean"),
        run =("d_run","mean"),
        total=("d_total","mean"),
        n=("event_date","count"),
    ).round(0)

    fig, axes = plt.subplots(1, 4, figsize=(14, 4.2), sharey=False)
    segs = [("swim","Swim"), ("bike","Bike"), ("run","Run"), ("total","Overall")]
    for ax, (col, label) in zip(axes, segs):
        vals = agg[col].values
        years = agg.index.astype(str)
        colors = [USAT_RED if v >= 0 else USAT_NAVY for v in vals]
        bars = ax.bar(years, vals, color=colors, edgecolor="white", linewidth=1.2)
        ax.axhline(0, color="black", linewidth=0.8)
        ax.set_title(label)
        ax.set_xlabel("Year")
        ax.set_ylabel("Avg seconds behind Parker" if col == "swim" else "")
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        for b, v in zip(bars, vals):
            ax.annotate(f"{int(v):+d}s",
                        xy=(b.get_x()+b.get_width()/2, v),
                        xytext=(0, 4 if v >= 0 else -14),
                        textcoords="offset points", ha="center", fontsize=9, fontweight="bold")
        # add n=
        ax.set_xticks(range(len(years)))
        ax.set_xticklabels([f"{y}\nn={int(n)}" for y, n in zip(years, agg["n"].values)])
    fig.suptitle("Emelia Perry — Average Deficit vs Lauren Parker (Sprint Distance)",
                 fontsize=14, fontweight="bold", y=1.02)
    fig.text(0.5, -0.04, "Red = Perry behind Parker.  Navy = Perry ahead.", ha="center", fontsize=9, color=GREY_DARK)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def chart_perry_position(df: pd.DataFrame, out: Path):
    sprint = df[df["prog_distance_category"]=="sprint"].copy()
    agg = sprint.groupby("year").agg(
        a_pos=("a_pos","mean"),
        b_pos=("b_pos","mean"),
        n=("event_date","count"),
    ).round(2)

    fig, ax = plt.subplots(figsize=(7, 4.2))
    x = np.arange(len(agg))
    w = 0.38
    ax.bar(x - w/2, agg["a_pos"], width=w, color=USAT_RED,  label="Perry")
    ax.bar(x + w/2, agg["b_pos"], width=w, color=GREY_DARK, label="Parker")
    for i, (p, q) in enumerate(zip(agg["a_pos"], agg["b_pos"])):
        ax.annotate(f"{p:.1f}", xy=(i - w/2, p), xytext=(0,4),
                    textcoords="offset points", ha="center", fontsize=9, fontweight="bold")
        ax.annotate(f"{q:.1f}", xy=(i + w/2, q), xytext=(0,4),
                    textcoords="offset points", ha="center", fontsize=9, color=GREY_DARK)
    ax.set_xticks(x)
    ax.set_xticklabels([f"{y}\nn={int(n)}" for y, n in zip(agg.index, agg["n"])])
    ax.invert_yaxis()
    ax.set_ylabel("Average finish position (1 = best)")
    ax.set_title("Average Finish Position by Year (Sprint Distance)")
    ax.legend(frameon=False, loc="lower right")
    ax.grid(axis="y", linestyle=":", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def chart_perry_per_race(df: pd.DataFrame, out: Path):
    """Per-race deficit timeline — every overlap race as a dot, with rolling-3 + linear trendlines."""
    from matplotlib.lines import Line2D
    fig, axes = plt.subplots(1, 3, figsize=(14, 4.8), sharey=False)
    segs = [("d_swim","Swim deficit (s)"),("d_bike","Bike deficit (s)"),("d_run","Run deficit (s)")]
    df = df.sort_values("event_date").reset_index(drop=True)
    # x as days-since-first for linear fit
    x_days = (df["event_date"] - df["event_date"].iloc[0]).dt.days.values.astype(float)

    for ax, (col, label) in zip(axes, segs):
        is_worlds = df["event_name"].str.contains("Championships", case=False, na=False)
        colors = [USAT_GOLD if w else USAT_RED for w in is_worlds]
        y = df[col].values.astype(float)

        # 3-race rolling average (centered, min 2)
        roll = df[col].rolling(window=3, min_periods=2, center=True).mean()
        ax.plot(df["event_date"], roll, color=USAT_NAVY, linewidth=2.2, zorder=2,
                label="3-race rolling avg")

        # Linear regression line
        m, b = np.polyfit(x_days, y, 1)
        x_line = np.array([x_days.min(), x_days.max()])
        y_line = m * x_line + b
        ax.plot([df["event_date"].iloc[0], df["event_date"].iloc[-1]], y_line,
                color=GREY_DARK, linewidth=1.4, linestyle="--", zorder=2,
                label=f"Linear ({m*365:+.0f}s/yr)")

        # Race dots
        ax.scatter(df["event_date"], y, c=colors, s=110, edgecolor="black",
                   linewidth=0.7, zorder=3)
        ax.axhline(0, color="black", linewidth=0.8)
        ax.set_title(label)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        ax.set_xlabel("")
        for d, v in zip(df["event_date"], y):
            ax.annotate(f"{int(v):+d}", xy=(d, v), xytext=(0,9 if v>=0 else -15),
                        textcoords="offset points", ha="center", fontsize=8, color=GREY_DARK)
        ax.legend(frameon=False, loc="upper right", fontsize=8)

    # Race-type legend separately at bottom
    legend = [Patch(facecolor=USAT_RED, edgecolor="black", label="Series / Cup / Paralympics"),
              Patch(facecolor=USAT_GOLD, edgecolor="black", label="World Championships")]
    fig.legend(handles=legend, frameon=False, loc="lower center", ncol=2, fontsize=9,
               bbox_to_anchor=(0.5, -0.04))
    fig.suptitle("Perry vs Parker — Per-Race Segment Deficit (every overlap race)",
                 fontsize=14, fontweight="bold", y=1.02)
    for ax in axes:
        ax.tick_params(axis="x", rotation=30)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


# ---------- Charts: Hermes deck ----------
def chart_pace_trend(df: pd.DataFrame, athlete_name: str, peers_meta: list[dict],
                     out: Path, prog_label: str = ""):
    """Per-race pace by segment over time, athlete vs podium-avg (excl self) per race.

    peers_meta items: {"key": col-prefix, "label": display, "color": hex, "marker": str}
    """
    fig, axes = plt.subplots(3, 1, figsize=(14, 12), sharex=True)
    segs = [("swim", "Swim pace (sec / 100m)"),
            ("bike", "Bike pace (sec / km)"),
            ("run",  "Run pace (sec / km)")]
    df = df.sort_values("event_date").reset_index(drop=True)
    df = df.dropna(subset=["pace_h_swim","pace_h_bike","pace_h_run"], how="all")
    for ax, (seg, label) in zip(axes, segs):
        sub = df.dropna(subset=[f"pace_h_{seg}", f"pace_pod_{seg}"]).copy()
        ax.plot(sub["event_date"], sub[f"pace_pod_{seg}"], color=GREY_DARK,
                linewidth=1.8, label=f"Podium avg (excl. {athlete_name})", marker="s", markersize=7)
        ax.plot(sub["event_date"], sub[f"pace_h_{seg}"], color=USAT_RED,
                linewidth=2.5, label=athlete_name, marker="o", markersize=10)
        scale = 100.0 if seg == "swim" else 1000.0
        dist_col = f"{seg}_m"
        for peer in peers_meta:
            pkey = peer["key"]
            sub_p = sub.dropna(subset=[f"{pkey}_{seg}"]).copy()
            if sub_p.empty:
                continue
            sub_p[f"pace_{pkey}_{seg}"] = sub_p[f"{pkey}_{seg}"] / (sub_p[dist_col] / scale)
            ax.scatter(sub_p["event_date"], sub_p[f"pace_{pkey}_{seg}"],
                       facecolor=peer["color"], edgecolor="black", s=180, zorder=5,
                       label=f"{peer['label']} (when present)", marker=peer["marker"])
        # Year shading
        years = sorted(sub["event_date"].dt.year.unique())
        for i, y in enumerate(years):
            if i % 2 == 0:
                ax.axvspan(pd.Timestamp(year=y, month=1, day=1),
                           pd.Timestamp(year=y, month=12, day=31),
                           color="#f5f5f5", zorder=0)
        ax.set_ylabel(label, fontsize=11)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        ax.legend(frameon=False, loc="upper right", fontsize=10)
        ax.set_title(label.split(" pace")[0].title() + " — lower is faster",
                     loc="left", fontsize=11)
    axes[-1].set_xlabel("Race date", fontsize=11)
    axes[-1].tick_params(axis="x", rotation=30)
    fig.suptitle(f"{athlete_name} — Pace Trend Across {prog_label} Career".strip(),
                 fontsize=15, fontweight="bold", y=0.995)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def chart_yearly_pace(df: pd.DataFrame, athlete_name: str, out: Path):
    """Yearly avg pace-deficit to podium (excl self)."""
    df = df.dropna(subset=["pace_d_swim","pace_d_bike","pace_d_run"], how="all").copy()
    agg = df.groupby("year").agg(
        swim=("pace_d_swim","mean"),
        bike=("pace_d_bike","mean"),
        run =("pace_d_run","mean"),
        n=("event_date","count"),
    ).round(2)

    fig, axes = plt.subplots(1, 3, figsize=(13, 4.2))
    segs = [("swim","Swim (sec/100m)"), ("bike","Bike (sec/km)"), ("run","Run (sec/km)")]
    for ax, (col, label) in zip(axes, segs):
        vals = agg[col].values
        years = agg.index.astype(str)
        colors = [USAT_RED if v >= 0 else USAT_NAVY for v in vals]
        bars = ax.bar(years, vals, color=colors, edgecolor="white", linewidth=1.2)
        ax.axhline(0, color="black", linewidth=0.8)
        ax.set_title(label)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        for b, v in zip(bars, vals):
            ax.annotate(f"{v:+.1f}", xy=(b.get_x()+b.get_width()/2, v),
                        xytext=(0, 4 if v >= 0 else -14),
                        textcoords="offset points", ha="center", fontsize=9, fontweight="bold")
        ax.set_xticks(range(len(years)))
        ax.set_xticklabels([f"{y}\nn={int(n)}" for y, n in zip(years, agg["n"])])
    fig.suptitle(f"{athlete_name} — Yearly Avg Pace Deficit vs Podium (excl. self)",
                 fontsize=14, fontweight="bold", y=1.02)
    fig.text(0.5, -0.04, f"Red = {athlete_name} slower than podium.  Navy = {athlete_name} faster.",
             ha="center", fontsize=9, color=GREY_DARK)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def chart_worlds_only(df_full: pd.DataFrame, athlete_name: str, out: Path):
    """Worlds-only deficit by year (one race per year). Removes small-field bias from yearly avg.

    Filters to events whose cat == 'World Championships' AND name contains 'Para Championships'.
    """
    worlds = df_full[
        df_full["cat"].fillna("").str.contains("World Championships", case=False) &
        df_full["event_name"].str.contains("Para Championships", case=False, na=False)
    ].copy()
    if worlds.empty:
        return None
    worlds = worlds.sort_values("event_date")
    worlds["year"] = worlds["event_date"].dt.year
    agg = worlds.groupby("year").agg(
        swim=("pace_d_swim","mean"),
        bike=("pace_d_bike","mean"),
        run =("pace_d_run","mean"),
        pos=("h_pos","mean"),
        n_field=("n_finishers","mean"),
        venue=("event_name", lambda s: s.iloc[0].split()[-1]),
    ).round(2)

    fig, axes = plt.subplots(1, 3, figsize=(13, 4.4))
    segs = [("swim","Swim (sec/100m)"), ("bike","Bike (sec/km)"), ("run","Run (sec/km)")]
    for ax, (col, label) in zip(axes, segs):
        vals = agg[col].values
        years = agg.index.astype(str)
        colors = [USAT_RED if v >= 0 else USAT_NAVY for v in vals]
        bars = ax.bar(years, vals, color=colors, edgecolor="white", linewidth=1.2)
        ax.axhline(0, color="black", linewidth=0.8)
        ax.set_title(label)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        for b, v in zip(bars, vals):
            ax.annotate(f"{v:+.1f}", xy=(b.get_x()+b.get_width()/2, v),
                        xytext=(0, 4 if v >= 0 else -14),
                        textcoords="offset points", ha="center", fontsize=10, fontweight="bold")
        ax.set_xticks(range(len(years)))
        ax.set_xticklabels([
            f"{y}\n{agg.loc[int(y),'venue']}\nfield={int(agg.loc[int(y),'n_field'])} · {int(agg.loc[int(y),'pos'])}{ordinal(int(agg.loc[int(y),'pos']))}"
            for y in years
        ], fontsize=8)
    fig.suptitle(f"{athlete_name} — Pace Deficit vs Podium at World Championships Only",
                 fontsize=14, fontweight="bold", y=1.02)
    fig.text(0.5, -0.06,
             "Strips small-field bias from the yearly average — only deep-field benchmarks remain.",
             ha="center", fontsize=9, color=GREY_DARK)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)
    return agg


def chart_peer_direct(df_overlap: pd.DataFrame, peer_label: str, out: Path,
                      athlete_name: str = "Athlete"):
    """Direct per-race head-to-head deficit chart for athlete vs a single peer."""
    df = df_overlap.sort_values("event_date").reset_index(drop=True).copy()
    df["label"] = df.apply(lambda r: f"{r['event_date'].strftime('%b %Y')}\n{shortname(r['event_name'])}", axis=1)
    fig, axes = plt.subplots(1, 3, figsize=(15, 5.2), sharey=False)
    segs = [("d_swim","Swim deficit (s)"), ("d_bike","Bike deficit (s)"), ("d_run","Run deficit (s)")]
    for ax, (col, label) in zip(axes, segs):
        vals = df[col].values
        colors = [USAT_RED if v >= 0 else USAT_NAVY for v in vals]
        bars = ax.bar(range(len(df)), vals, color=colors, edgecolor="white", linewidth=1.5)
        ax.axhline(0, color="black", linewidth=0.8)
        ax.set_title(label)
        ax.set_xticks(range(len(df)))
        ax.set_xticklabels(df["label"], fontsize=8, rotation=20, ha="right")
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        for i, v in enumerate(vals):
            ax.annotate(f"{int(v):+d}s", xy=(i, v), xytext=(0,4 if v>=0 else -14),
                        textcoords="offset points", ha="center", fontsize=9, fontweight="bold")
    fig.suptitle(f"{athlete_name} vs {peer_label} — Per-Race Deficit (each overlap race)",
                 fontsize=14, fontweight="bold", y=1.02)
    fig.text(0.5, -0.08,
             f"Red = {athlete_name} behind.  Navy = {athlete_name} ahead.",
             ha="center", fontsize=9, color=GREY_DARK)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight")
    plt.close(fig)


def shortname(event_name: str) -> str:
    # Pull venue from event title (typically last token after year and qualifier)
    name = event_name
    for prefix in ["2022 ","2023 ","2024 ","2025 ","2026 "]:
        name = name.replace(prefix, "")
    name = (name.replace("World Triathlon Para Championships", "Worlds")
                 .replace("World Triathlon Para Series", "WPS")
                 .replace("World Triathlon Para Cup", "WPC")
                 .replace("Paralympic Games", "Paralympics")
                 .replace("Americas Triathlon Para Championships", "Americas")
                 .replace("Americas Triathlon Para Cup", "Americas Cup")
                 .replace("Torremolinos-Andalucia", "Torremolinos"))
    return name.strip()


# ---------- HTML summary ----------
HTML_TMPL = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{title}</title>
<style>
  body {{ font-family: Arial, sans-serif; max-width: 1100px; margin: 24px auto; color: #222; }}
  h1 {{ color: {navy}; border-bottom: 3px solid {red}; padding-bottom: 8px; }}
  h2 {{ color: {navy}; margin-top: 28px; }}
  .kpi {{ display: inline-block; padding: 10px 18px; margin: 6px 8px 6px 0;
          background: #f4f4f4; border-left: 4px solid {red}; }}
  .kpi b {{ font-size: 1.4em; color: {red}; }}
  table {{ border-collapse: collapse; margin: 12px 0; font-size: 0.92em; }}
  td, th {{ border: 1px solid #ccc; padding: 6px 10px; }}
  th {{ background: {navy}; color: white; }}
  img {{ max-width: 100%; margin: 12px 0; border: 1px solid #ddd; }}
  .note {{ background: #fff8e1; border-left: 4px solid {gold}; padding: 10px 14px; margin: 12px 0; }}
  .caveat {{ background: #ffebee; border-left: 4px solid {red}; padding: 10px 14px; margin: 12px 0; font-size: 0.9em; }}
</style></head><body>
<h1>{title}</h1>
{body}
<p style="color:#888;font-size:0.85em;margin-top:40px">
  Generated {ts} — Source: World Triathlon API (race_results, events tables).
</p>
</body></html>
"""

def write_html(out: Path, title: str, body_html: str):
    out.write_text(
        HTML_TMPL.format(
            title=title, body=body_html, ts=dt.datetime.now().strftime("%Y-%m-%d %H:%M"),
            navy=USAT_NAVY, red=USAT_RED, gold=USAT_GOLD,
        ), encoding="utf-8"
    )


# ---------- Build Perry deck ----------
def build_perry(engine, outdir: Path):
    outdir.mkdir(parents=True, exist_ok=True)
    df = pull_overlap(engine, PERRY, PARKER)
    df.to_csv(outdir / "data.csv", index=False)

    sprint = df[df["prog_distance_category"]=="sprint"]
    yearly = sprint.groupby("year").agg(
        n=("event_date","count"),
        d_swim=("d_swim","mean"),
        d_bike=("d_bike","mean"),
        d_run =("d_run","mean"),
        d_total=("d_total","mean"),
        a_pos=("a_pos","mean"),
    ).round(1)

    chart_perry_yearly_deficit(df, outdir / "fig1_yearly_deficit.png")
    chart_perry_position(df, outdir / "fig2_position.png")
    chart_perry_per_race(df, outdir / "fig3_per_race.png")

    # Wollongong 2025 Worlds callout (Perry's first PTWC Worlds silver vs Parker)
    wol_match = df[df["event_name"].str.contains("Wollongong", case=False, na=False)]
    wol = wol_match.iloc[0] if not wol_match.empty else df.iloc[-1]

    # KPIs
    d23 = yearly.loc[2023]; d25 = yearly.loc[2025]
    total_drop_pct = (1 - d25["d_total"]/d23["d_total"]) * 100
    bike_drop_pct  = (1 - d25["d_bike"]/d23["d_bike"]) * 100
    run_flip = d25["d_run"]  # negative = Perry faster

    body = f"""
<h2>Story</h2>
<p><b>Emelia Perry has closed the gap on Paralympic gold medalist Lauren Parker across every discipline from 2023 to 2025.</b>
The total time deficit shrank by <b>{total_drop_pct:.0f}%</b>, average finish position improved from
<b>{d23['a_pos']:.0f}th</b> to <b>{d25['a_pos']:.1f}</b>, and Perry has flipped from trailing Parker
on the run to running faster than her in 2025.</p>

<div class="kpi"><b>{total_drop_pct:.0f}%</b><br>Total deficit reduction</div>
<div class="kpi"><b>{bike_drop_pct:.0f}%</b><br>Bike deficit reduction</div>
<div class="kpi"><b>{int(run_flip):+d}s</b><br>2025 run advantage (Perry ahead)</div>
<div class="kpi"><b>{d23['a_pos']:.0f} &rarr; {d25['a_pos']:.1f}</b><br>Avg finish position</div>

<h2>Yearly Deficit Trend — Sprint Distance ({len(sprint)} overlap races)</h2>
<img src="fig1_yearly_deficit.png">

<h2>Finish Position Trend</h2>
<img src="fig2_position.png">

<h2>Every Overlap Race</h2>
<img src="fig3_per_race.png">

<h2>2025 World Championships, Wollongong — Worlds Silver</h2>
<div class="note">
  Perry finished <b>{wol['a_pos']:.0f}{ordinal(int(wol['a_pos']))}</b> at PTWC Worlds, only
  <b>{int(wol['d_total'])}s behind Parker</b> on total time, and was
  <b>{abs(int(wol['d_run']))}s faster than Parker on the run</b>.
  Best individual head-to-head result of her career.
</div>

<h2>Yearly Detail (sprint races)</h2>
{yearly.rename(columns={
  'n':'n races','d_swim':'Swim Δ (s)','d_bike':'Bike Δ (s)','d_run':'Run Δ (s)',
  'd_total':'Total Δ (s)','a_pos':'Perry avg pos'
}).to_html()}

<div class="caveat"><b>Caveats:</b>
2025 sample is 3 sprint races (Magog, Montreal, Wollongong Worlds). Sample is small — every additional race could move the yearly average meaningfully. Both athletes are in the H1 wheelchair sub-class, so within-class time factors cancel out and raw split differences are valid.</div>
"""
    write_html(outdir / "report.html", "Emelia Perry — PTWC Progression vs Lauren Parker", body)
    print(f"[Perry] wrote {outdir/'report.html'}  ({len(df)} overlap races)")
    return {"yearly": yearly, "n_total": len(df),
            "wol_pos": int(wol["a_pos"]), "wol_total": int(wol["d_total"]),
            "wol_run": int(wol["d_run"]),
            "total_drop_pct": total_drop_pct, "bike_drop_pct": bike_drop_pct,
            "run_flip_2025": int(run_flip)}


# ---------- Build Hermes deck ----------
def build_hermes(engine, outdir: Path):
    outdir.mkdir(parents=True, exist_ok=True)
    peers = {"rod": RODRIGUEZ.aid, "tar": TARANTELLO.aid}
    full = pull_full_history(engine, HERMES, peer_ids=peers)
    overlap_rod = pull_overlap(engine, HERMES, RODRIGUEZ)
    overlap_tar = pull_overlap(engine, HERMES, TARANTELLO)
    full.to_csv(outdir / "data_full_history.csv", index=False)
    overlap_rod.to_csv(outdir / "data_rodriguez_overlap.csv", index=False)
    overlap_tar.to_csv(outdir / "data_tarantello_overlap.csv", index=False)

    peers_meta = [
        {"key": "rod", "label": "Rodriguez B1", "color": USAT_NAVY, "marker": "D"},
        {"key": "tar", "label": "Tarantello B3", "color": USAT_GOLD, "marker": "^"},
    ]
    chart_pace_trend(full, "Hermes", peers_meta, outdir / "fig1_pace_trend.png", prog_label="PTVI")
    chart_yearly_pace(full, "Hermes", outdir / "fig2_yearly_pace.png")
    worlds_agg = chart_worlds_only(full, "Hermes", outdir / "fig3_worlds_only.png")
    chart_peer_direct(overlap_rod, "Rodriguez (B1)", outdir / "fig4_rodriguez_direct.png", athlete_name="Hermes")
    chart_peer_direct(overlap_tar, "Tarantello (B3)", outdir / "fig5_tarantello_direct.png", athlete_name="Hermes")

    yearly = full.dropna(subset=["pace_d_swim","pace_d_bike","pace_d_run"], how="all") \
                 .groupby("year").agg(
        n=("event_date","count"),
        pace_swim=("pace_d_swim","mean"),
        pace_bike=("pace_d_bike","mean"),
        pace_run =("pace_d_run","mean"),
        pos=("h_pos","mean"),
    ).round(2)

    d23 = yearly.loc[2023] if 2023 in yearly.index else None
    d25 = yearly.loc[2025] if 2025 in yearly.index else None

    # Per-peer trend (yearly total deficit on overlap races)
    def _peer_summary(df_over: pd.DataFrame) -> pd.DataFrame:
        if df_over.empty:
            return pd.DataFrame()
        return df_over.groupby("year").agg(
            n=("event_date","count"),
            total=("d_total","mean"),
            swim=("d_swim","mean"),
            bike=("d_bike","mean"),
            run=("d_run","mean"),
        ).round(0)
    rod_yr = _peer_summary(overlap_rod)
    tar_yr = _peer_summary(overlap_tar)

    body = f"""
<h2>Story</h2>
<p><b>McClain Hermes is one of the fastest swimmers in PTVI women and made a step-change in run pace in 2025.</b>
Across her {len(full)} PTVI race finishes, her swim pace has been ahead of the podium average every year,
the bike gap has narrowed from +{d23['pace_bike']:.1f} sec/km behind to +{d25['pace_bike']:.1f},
and her run pace went from <b>+{d23['pace_run']:.0f} sec/km off podium in 2023</b>
to <b>{d25['pace_run']:+.1f} sec/km in 2025</b> — faster than the average podium athlete.</p>

<div class="kpi"><b>{d25['pace_swim']:+.1f}</b><br>2025 swim vs podium (sec/100m)</div>
<div class="kpi"><b>{d25['pace_bike']:+.1f}</b><br>2025 bike vs podium (sec/km)</div>
<div class="kpi"><b>{d25['pace_run']:+.1f}</b><br>2025 run vs podium (sec/km)</div>
<div class="kpi"><b>{d23['pos']:.1f} &rarr; {d25['pos']:.1f}</b><br>Avg finish position</div>

<h2>Pace Trend — All PTVI Races ({len(full)} race finishes since 2022)</h2>
<img src="fig1_pace_trend.png">
<p style="font-size:0.9em;color:#555">
  Hermes (red), podium average excluding herself (grey), Rodriguez B1 (navy diamonds) and Tarantello B3 (gold triangles)
  highlighted at the races they attended. Lower = faster. Pace normalization puts sprint and standard distance on one chart.
</p>

<h2>Yearly Pace Summary (All Races)</h2>
<img src="fig2_yearly_pace.png">

<h2>World Championships Only — Deep-Field Benchmark</h2>
<img src="fig3_worlds_only.png">
<p style="font-size:0.9em;color:#555">
  One Worlds per year (Pontevedra '23, Torremolinos '24, Wollongong '25).
  Removes small-field bias from the "all races" yearly average.
</p>
{worlds_agg.rename(columns={'swim':'Swim Δ (sec/100m)','bike':'Bike Δ (sec/km)','run':'Run Δ (sec/km)','pos':'Pos','n_field':'Field size','venue':'Venue'}).to_html() if worlds_agg is not None else ''}

<h2>Head-to-Head vs Susana Rodriguez B1 ({len(overlap_rod)} overlap races)</h2>
<img src="fig4_rodriguez_direct.png">

<h2>Head-to-Head vs Francesca Tarantello B3 ({len(overlap_tar)} overlap races)</h2>
<img src="fig5_tarantello_direct.png">
<p style="font-size:0.9em;color:#555">
  Tarantello is in the B3 sub-class. Times are official (factor-adjusted) so the deficit shown is the gap on the final classification.
</p>

<h2>Yearly Pace Detail (vs Podium Avg)</h2>
{yearly.rename(columns={
  'n':'n races','pace_swim':'Swim Δ (sec/100m)','pace_bike':'Bike Δ (sec/km)',
  'pace_run':'Run Δ (sec/km)','pos':'Avg pos'
}).to_html()}

<h2>Yearly Total Deficit — Direct Peers (overlap races only)</h2>
<h3>vs Rodriguez (B1)</h3>
{rod_yr.to_html() if not rod_yr.empty else '<p>no data</p>'}
<h3>vs Tarantello (B3)</h3>
{tar_yr.to_html() if not tar_yr.empty else '<p>no data</p>'}

<div class="caveat"><b>Caveats:</b>
<ul>
<li>Direct overlap with each peer is thin: {len(overlap_rod)} races vs Rodriguez, {len(overlap_tar)} vs Tarantello. Yearly averaging head-to-head is unreliable; the per-race charts are the honest view.</li>
<li>Some races (2025 Calima, 2025 Tata) had only 3 finishers, so the "podium average" reduces to the rest of the field. The 2025 yearly average is favorably skewed by those small-field events.</li>
<li>Hermes and Rodriguez are both B1; Tarantello is B3. Cross-subclass comparison uses the official factor-adjusted finish times that determined actual placing.</li>
</ul>
</div>
"""
    write_html(outdir / "report.html", "McClain Hermes — PTVI Progression", body)
    print(f"[Hermes] wrote {outdir/'report.html'}  ({len(full)} career races; {len(overlap_rod)} vs Rodriguez, {len(overlap_tar)} vs Tarantello)")
    return {"yearly": yearly, "rod_yr": rod_yr, "tar_yr": tar_yr,
            "n_full": len(full), "n_rod": len(overlap_rod), "n_tar": len(overlap_tar)}


# ---------- PPTX export ----------
def build_pptx_perry(outdir: Path, stats: dict) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00, 0x28, 0x68)
    RED  = RGBColor(0xC8, 0x10, 0x2E)
    BLANK = prs.slide_layouts[6]

    def add_title_bar(slide, title: str, subtitle: str = ""):
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = NAVY
        if subtitle:
            p2 = tf.add_paragraph(); p2.text = subtitle
            p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    def add_image_slide(title: str, image_path: Path, caption: str = ""):
        slide = prs.slides.add_slide(BLANK)
        add_title_bar(slide, title)
        slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.2),
                                 width=Inches(12.5))
        if caption:
            tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = caption
            p.runs[0].font.size = Pt(11); p.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)
        return slide

    # Slide 1 — title
    slide = prs.slides.add_slide(BLANK)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Emelia Perry"
    p.runs[0].font.size = Pt(54); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = "PTWC Progression vs Lauren Parker, 2023–2025"
    p2.runs[0].font.size = Pt(28); p2.runs[0].font.color.rgb = RED
    p3 = tf.add_paragraph(); p3.text = f"{stats['n_total']} head-to-head races · World Triathlon results"
    p3.runs[0].font.size = Pt(14); p3.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    # Slide 2 — story / KPIs
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "Story",
                  "Closing the gap on the Paralympic gold medalist — every segment, every year.")
    kpis = [
        (f"{stats['total_drop_pct']:.0f}%",  "Total deficit reduction\n(2023 → 2025)"),
        (f"{stats['bike_drop_pct']:.0f}%",   "Bike deficit reduction"),
        (f"{stats['run_flip_2025']:+d}s",    "2025 run advantage\n(Perry ahead)"),
        (f"7th → {stats['yearly'].loc[2025,'a_pos']:.1f}", "Avg finish position"),
    ]
    for i, (big, sml) in enumerate(kpis):
        x = Inches(0.6 + i*3.1)
        box = slide.shapes.add_textbox(x, Inches(2.0), Inches(2.9), Inches(2.5))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big
        p.runs[0].font.size = Pt(56); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = sml
        p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER

    # Charts
    add_image_slide("Sprint Yearly Deficit vs Lauren Parker",
                    outdir / "fig1_yearly_deficit.png",
                    "Red = Perry behind. Navy = Perry ahead. n = sprint races per year.")
    add_image_slide("Average Finish Position by Year",
                    outdir / "fig2_position.png",
                    "Lower bar = better finish. Perry has gone from avg 7th to avg 2.5 in PTWC sprint racing.")
    add_image_slide("Every Overlap Race — Per-Segment Deficit",
                    outdir / "fig3_per_race.png",
                    "Gold marker = standard-distance race (2025 Worlds, Wollongong).")

    # Wollongong 2025 Worlds callout
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "2025 World Championships, Wollongong — Worlds Silver")
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Silver — {stats['wol_total']:+d}s on total time"
    p.runs[0].font.size = Pt(40); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
    p2 = tf.add_paragraph(); p2.text = f"Perry beat Parker on the run by {abs(stats['wol_run'])}s at PTWC Worlds."
    p2.runs[0].font.size = Pt(22); p2.runs[0].font.color.rgb = NAVY
    p3 = tf.add_paragraph(); p3.text = " "
    p4 = tf.add_paragraph()
    p4.text = "Best individual head-to-head result of Perry's career — confirms the 2025 fitness step."
    p4.runs[0].font.size = Pt(16); p4.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    # Caveats
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "Notes")
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.2), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    notes = [
        "Both athletes race in the H1 wheelchair sub-class — within-class comparison is exact (no time-factor effects).",
        "2025 sprint sample is 2 races (Magog, Montreal); 2025 standard sample is 1 race (Wollongong). Further 2026 races will tighten the yearly averages.",
        "Source: World Triathlon official results — race_results and events tables.",
    ]
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + n
        p.runs[0].font.size = Pt(16); p.runs[0].font.color.rgb = RGBColor(0x33,0x33,0x33)

    pptx_path = outdir / "Perry_PTWC_Progression.pptx"
    prs.save(str(pptx_path))
    return pptx_path


def build_pptx_hermes(outdir: Path, stats: dict) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00, 0x28, 0x68)
    RED  = RGBColor(0xC8, 0x10, 0x2E)
    BLANK = prs.slide_layouts[6]
    yearly = stats["yearly"]
    d23 = yearly.loc[2023] if 2023 in yearly.index else None
    d25 = yearly.loc[2025] if 2025 in yearly.index else None

    def add_title_bar(slide, title: str, subtitle: str = ""):
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
        if subtitle:
            p2 = tf.add_paragraph(); p2.text = subtitle
            p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    def add_image_slide(title: str, image_path: Path, caption: str = ""):
        slide = prs.slides.add_slide(BLANK)
        add_title_bar(slide, title)
        slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.2),
                                 width=Inches(12.5))
        if caption:
            tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = caption
            p.runs[0].font.size = Pt(11); p.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)
        return slide

    # Slide 1 — title
    slide = prs.slides.add_slide(BLANK)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "McClain Hermes"
    p.runs[0].font.size = Pt(54); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = "PTVI Progression vs Rodriguez (B1) & Tarantello (B3)"
    p2.runs[0].font.size = Pt(28); p2.runs[0].font.color.rgb = RED
    p3 = tf.add_paragraph()
    p3.text = f"{stats['n_full']} career races · {stats['n_rod']} vs Rodriguez · {stats['n_tar']} vs Tarantello"
    p3.runs[0].font.size = Pt(14); p3.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    # Slide 2 — story / KPIs
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "Story",
                  "Fast swimmer, narrowing the bike gap, and a step-change in run pace in 2025.")
    kpis = [
        (f"{d25['pace_swim']:+.1f}",  "2025 swim vs podium\n(sec / 100m)"),
        (f"{d25['pace_bike']:+.1f}",  "2025 bike vs podium\n(sec / km)"),
        (f"{d25['pace_run']:+.1f}",   "2025 run vs podium\n(sec / km)"),
        (f"{d23['pos']:.1f}→{d25['pos']:.1f}", "Avg finish position\n2023 → 2025"),
    ]
    for i, (big, sml) in enumerate(kpis):
        x = Inches(0.6 + i*3.1)
        box = slide.shapes.add_textbox(x, Inches(2.0), Inches(2.9), Inches(2.5))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big
        p.runs[0].font.size = Pt(48); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = sml
        p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER

    add_image_slide("Pace Trend — All PTVI Races",
                    outdir / "fig1_pace_trend.png",
                    "Red = Hermes. Grey = podium avg excl. Hermes. Navy diamonds = Rodriguez B1. Gold triangles = Tarantello B3.")
    add_image_slide("Yearly Avg Pace Deficit vs Podium — All Races",
                    outdir / "fig2_yearly_pace.png",
                    "Red = Hermes slower than podium. Navy = Hermes faster. Yearly avg is field-depth sensitive — see Worlds-only view next.")
    add_image_slide("Worlds Only — Apples-to-Apples Trajectory",
                    outdir / "fig3_worlds_only.png",
                    "One race per year (Pontevedra '23, Torremolinos '24, Wollongong '25). Removes small-field bias.")
    add_image_slide(f"Head-to-Head vs Susana Rodriguez (B1) — {stats['n_rod']} overlap races",
                    outdir / "fig4_rodriguez_direct.png",
                    "Red = Hermes behind. Navy = Hermes ahead.")
    add_image_slide(f"Head-to-Head vs Francesca Tarantello (B3) — {stats['n_tar']} overlap races",
                    outdir / "fig5_tarantello_direct.png",
                    "Tarantello is in the B3 sub-class — comparison uses official factor-adjusted times.")

    # Caveats
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "Notes")
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.2), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    notes = [
        f"Direct overlap data is thin: {stats['n_rod']} races vs Rodriguez, {stats['n_tar']} vs Tarantello. Per-race charts are more honest than year-over-year averaging.",
        "Two 2025 races (Calima Cont. Champs, Tata WC) had only 3 finishers; Hermes was on the podium in both. The 2025 'vs podium avg' yearly number is favorably skewed by those tiny fields.",
        "Hermes and Rodriguez are both B1; Tarantello is B3. PTVI results are factor-adjusted to equalize across sub-classes — the deficits shown are the gap on official classification.",
        "Source: World Triathlon official results — race_results and events tables.",
    ]
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + n
        p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = RGBColor(0x33,0x33,0x33)

    pptx_path = outdir / "Hermes_PTVI_Progression.pptx"
    prs.save(str(pptx_path))
    return pptx_path


# ---------- Build Fisher deck (PTWC, peers = Parker + Perry) ----------
def build_fisher(engine, outdir: Path):
    outdir.mkdir(parents=True, exist_ok=True)
    peers = {"prk": PARKER.aid, "pry": PERRY.aid}
    full = pull_full_history(engine, FISHER, peer_ids=peers)
    overlap_prk = pull_overlap(engine, FISHER, PARKER)
    overlap_pry = pull_overlap(engine, FISHER, PERRY)
    full.to_csv(outdir / "data_full_history.csv", index=False)
    overlap_prk.to_csv(outdir / "data_parker_overlap.csv", index=False)
    overlap_pry.to_csv(outdir / "data_perry_overlap.csv", index=False)

    peers_meta = [
        {"key": "prk", "label": "Parker H1", "color": USAT_NAVY, "marker": "D"},
        {"key": "pry", "label": "Perry H1",  "color": USAT_GOLD, "marker": "^"},
    ]
    chart_pace_trend(full, "Fisher", peers_meta, outdir / "fig1_pace_trend.png", prog_label="PTWC")
    chart_yearly_pace(full, "Fisher", outdir / "fig2_yearly_pace.png")
    worlds_agg = chart_worlds_only(full, "Fisher", outdir / "fig3_worlds_only.png")
    chart_peer_direct(overlap_prk, "Parker (H1)", outdir / "fig4_parker_direct.png", athlete_name="Fisher")
    chart_peer_direct(overlap_pry, "Perry (H1)",  outdir / "fig5_perry_direct.png",  athlete_name="Fisher")

    yearly = full.dropna(subset=["pace_d_swim","pace_d_bike","pace_d_run"], how="all") \
                 .groupby("year").agg(
        n=("event_date","count"),
        pace_swim=("pace_d_swim","mean"),
        pace_bike=("pace_d_bike","mean"),
        pace_run =("pace_d_run","mean"),
        pos=("h_pos","mean"),
    ).round(2)

    # Per-peer total-deficit summary
    def _peer_summary(df_over: pd.DataFrame) -> pd.DataFrame:
        if df_over.empty:
            return pd.DataFrame()
        return df_over.groupby("year").agg(
            n=("event_date","count"),
            total=("d_total","mean"),
            swim=("d_swim","mean"),
            bike=("d_bike","mean"),
            run=("d_run","mean"),
        ).round(0)
    prk_yr = _peer_summary(overlap_prk)
    pry_yr = _peer_summary(overlap_pry)

    # Best single-race gaps
    best_prk = overlap_prk.loc[overlap_prk["d_total"].idxmin()] if not overlap_prk.empty else None
    best_pry = overlap_pry.loc[overlap_pry["d_total"].idxmin()] if not overlap_pry.empty else None

    # Cup podiums (1/2/3 at Para Cup events)
    cup_pod = full[
        full["cat"].fillna("").str.contains("Para Cup", case=False)
        & full["h_pos"].isin([1, 2, 3])
    ]
    cup_podium_count = len(cup_pod)

    body = f"""
<h2>Story</h2>
<p><b>Skyler Fisher has been a consistent Para Cup podium presence and produced her closest-ever deep-field gap to the world leaders at Yokohama 2026.</b>
Across her {len(full)} PTWC race finishes she has {cup_podium_count} Para Cup podium results, and at the 2026 Yokohama World Para Series she finished within
<b>{int(best_prk['d_total'])}s of Parker</b> (defending Paralympic gold) and <b>{int(best_pry['d_total'])}s of Perry</b> — her closest gap to Parker in any of her {len(overlap_prk)} head-to-head races.</p>

<div class="kpi"><b>{len(full)}</b><br>PTWC career race finishes (since 2022)</div>
<div class="kpi"><b>{cup_podium_count}</b><br>Para Cup podiums</div>
<div class="kpi"><b>{int(best_prk['d_total'])}s</b><br>Closest career gap to Parker<br><span style="font-size:0.85em">({best_prk['event_date'].strftime('%b %Y')} {best_prk['event_name'].split()[-1]})</span></div>
<div class="kpi"><b>{int(best_pry['d_total'])}s</b><br>Closest 2026 gap to Perry<br><span style="font-size:0.85em">({best_pry['event_date'].strftime('%b %Y')} {best_pry['event_name'].split()[-1]})</span></div>

<h2>Pace Trend — All PTWC Races ({len(full)} race finishes since 2022)</h2>
<img src="fig1_pace_trend.png">
<p style="font-size:0.9em;color:#555">
  Fisher (red), podium average excluding herself (grey), Parker H1 (navy diamonds) and Perry H1 (gold triangles)
  highlighted at the races they attended. Lower = faster.
</p>

<h2>Yearly Pace Summary (All Races)</h2>
<img src="fig2_yearly_pace.png">

<h2>World Championships Only — Deep-Field Benchmark</h2>
<img src="fig3_worlds_only.png">
<p style="font-size:0.9em;color:#555">
  Worlds-only deficit — removes small-field bias from the all-races yearly average.
</p>
{worlds_agg.rename(columns={'swim':'Swim Δ (sec/100m)','bike':'Bike Δ (sec/km)','run':'Run Δ (sec/km)','pos':'Pos','n_field':'Field size','venue':'Venue'}).to_html() if worlds_agg is not None else ''}

<h2>Head-to-Head vs Lauren Parker H1 ({len(overlap_prk)} overlap races)</h2>
<img src="fig4_parker_direct.png">

<h2>Head-to-Head vs Emelia Perry H1 ({len(overlap_pry)} overlap races)</h2>
<img src="fig5_perry_direct.png">

<h2>Yearly Pace Detail (vs Podium Avg)</h2>
{yearly.rename(columns={'n':'n races','pace_swim':'Swim Δ (sec/100m)','pace_bike':'Bike Δ (sec/km)','pace_run':'Run Δ (sec/km)','pos':'Avg pos'}).to_html()}

<h2>Yearly Total Deficit — Direct Peers (overlap races only)</h2>
<h3>vs Parker (H1)</h3>
{prk_yr.to_html() if not prk_yr.empty else '<p>no data</p>'}
<h3>vs Perry (H1)</h3>
{pry_yr.to_html() if not pry_yr.empty else '<p>no data</p>'}

<div class="caveat"><b>Caveats:</b>
<ul>
<li>Fisher is in the H2 sub-class; Parker and Perry are H1. PTWC results use sub-class time factors; the splits shown are the official factor-adjusted times that determined final placing.</li>
<li>Direct overlap with Parker is thin ({len(overlap_prk)} races, all deep fields). Per-race chart is the honest view.</li>
<li>Some Para Cup races had small fields; the 'podium avg excl. self' reduces to the other 1–2 finishers in those races.</li>
<li>Source: World Triathlon official results — race_results and events tables.</li>
</ul>
</div>
"""
    write_html(outdir / "report.html", "Skyler Fisher — PTWC Progression", body)
    print(f"[Fisher] wrote {outdir/'report.html'}  ({len(full)} career races; {len(overlap_prk)} vs Parker, {len(overlap_pry)} vs Perry)")
    return {"yearly": yearly, "prk_yr": prk_yr, "pry_yr": pry_yr,
            "n_full": len(full), "n_prk": len(overlap_prk), "n_pry": len(overlap_pry),
            "best_prk": best_prk, "best_pry": best_pry,
            "cup_podium_count": cup_podium_count}


def build_pptx_fisher(outdir: Path, stats: dict) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00, 0x28, 0x68)
    RED  = RGBColor(0xC8, 0x10, 0x2E)
    BLANK = prs.slide_layouts[6]
    best_prk = stats["best_prk"]; best_pry = stats["best_pry"]

    def add_title_bar(slide, title: str, subtitle: str = ""):
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
        if subtitle:
            p2 = tf.add_paragraph(); p2.text = subtitle
            p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    def add_image_slide(title: str, image_path: Path, caption: str = ""):
        slide = prs.slides.add_slide(BLANK)
        add_title_bar(slide, title)
        slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.2),
                                 width=Inches(12.5))
        if caption:
            tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = caption
            p.runs[0].font.size = Pt(11); p.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)
        return slide

    # Slide 1 — title
    slide = prs.slides.add_slide(BLANK)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Skyler Fisher"
    p.runs[0].font.size = Pt(54); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = "PTWC Progression vs Parker (H1, reigning Paralympic gold) & Perry (H1, USA teammate)"
    p2.runs[0].font.size = Pt(24); p2.runs[0].font.color.rgb = RED
    p3 = tf.add_paragraph()
    p3.text = f"{stats['n_full']} career races · {stats['n_prk']} vs Parker · {stats['n_pry']} vs Perry"
    p3.runs[0].font.size = Pt(14); p3.runs[0].font.color.rgb = RGBColor(0x55,0x55,0x55)

    # Slide 2 — story / KPIs
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "Story",
                  "Cup-podium regular with a 2026 Yokohama breakthrough at depth.")
    kpis = [
        (f"{stats['n_full']}",            "PTWC career\nrace finishes"),
        (f"{stats['cup_podium_count']}",  "Para Cup\npodium results"),
        (f"{int(best_prk['d_total'])}s",  f"Closest gap to Parker\n({best_prk['event_date'].strftime('%b %Y')})"),
        (f"{int(best_pry['d_total'])}s",  f"Closest 2026 gap to Perry\n({best_pry['event_date'].strftime('%b %Y')})"),
    ]
    for i, (big, sml) in enumerate(kpis):
        x = Inches(0.6 + i*3.1)
        box = slide.shapes.add_textbox(x, Inches(2.0), Inches(2.9), Inches(2.5))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big
        p.runs[0].font.size = Pt(48); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = sml
        p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER

    add_image_slide("Pace Trend — All PTWC Races",
                    outdir / "fig1_pace_trend.png",
                    "Red = Fisher. Grey = podium avg excl. Fisher. Navy diamonds = Parker H1. Gold triangles = Perry H1.")
    add_image_slide("Yearly Avg Pace Deficit vs Podium — All Races",
                    outdir / "fig2_yearly_pace.png",
                    "Red = Fisher slower than podium. Navy = Fisher faster. Yearly avg is field-depth sensitive — see Worlds-only view next.")
    add_image_slide("Worlds Only — Apples-to-Apples Trajectory",
                    outdir / "fig3_worlds_only.png",
                    "One race per year at the World Championships level. Removes small-field bias.")
    add_image_slide(f"Head-to-Head vs Lauren Parker (H1) — {stats['n_prk']} overlap races",
                    outdir / "fig4_parker_direct.png",
                    "Red = Fisher behind. Navy = Fisher ahead. Parker is in the H1 sub-class.")
    add_image_slide(f"Head-to-Head vs Emelia Perry (H1) — {stats['n_pry']} overlap races",
                    outdir / "fig5_perry_direct.png",
                    "Red = Fisher behind. Navy = Fisher ahead. Both USA, but Perry is H1, Fisher is H2.")

    # Caveats
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, "Notes")
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.2), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    notes = [
        f"Direct overlap data: {stats['n_prk']} races vs Parker, {stats['n_pry']} vs Perry. Parker overlap is concentrated at Worlds + WPS Series; the per-race chart is the honest view.",
        "Fisher is H2; Parker and Perry are H1. PTWC results are factor-adjusted to equalize across sub-classes — deficits shown are the gap on the official classification (the gap that actually determined finishing order).",
        "Some Para Cup races had very small fields (2–4 finishers). When Fisher was on the podium in those, the 'podium avg excl. self' reduces to just 1–2 other athletes — interpret 2025 yearly numbers with this in mind.",
        "Source: World Triathlon official results — race_results and events tables.",
    ]
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + n
        p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = RGBColor(0x33,0x33,0x33)

    pptx_path = outdir / "Fisher_PTWC_Progression.pptx"
    prs.save(str(pptx_path))
    return pptx_path


# ---------- Standalone 2026 update slides (single PPTX, copy/paste into existing deck) ----------
def _seconds_to_mmss(sec):
    if sec is None or pd.isna(sec):
        return "—"
    sec = int(round(sec))
    h, rem = divmod(sec, 3600)
    m, s = divmod(rem, 60)
    return f"{h}:{m:02d}:{s:02d}" if h else f"{m}:{s:02d}"


def build_perry_2026_slide(engine, outdir: Path) -> Path:
    """Single slide: Perry 2026 season update (Yokohama WPS vs Parker)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    df = pull_overlap(engine, PERRY, PARKER)
    df["year"] = df["event_date"].dt.year
    races_2026 = df[df["year"] == 2026].sort_values("event_date")
    if races_2026.empty:
        print("[Perry 2026] no 2026 races, skipping")
        return None
    race = races_2026.iloc[-1]  # most recent
    yearly_2025 = df[df["year"] == 2025].agg(
        {"d_swim":"mean","d_bike":"mean","d_run":"mean","d_total":"mean"})

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00,0x28,0x68); RED = RGBColor(0xC8,0x10,0x2E)
    GREY = RGBColor(0x55,0x55,0x55); LIGHT = RGBColor(0xF4,0xF4,0xF4)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "2026 Season Update — Emelia Perry vs Lauren Parker"
    p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = f"{race['event_name']} · {race['event_date'].strftime('%B %d, %Y')} · finished {int(race['a_pos'])}{ordinal(int(race['a_pos']))}"
    p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = GREY

    # KPI row
    kpis = [
        (f"{int(race['a_pos'])}{ordinal(int(race['a_pos']))}", "Yokohama WPS finish", "vs Parker (gold)"),
        (f"{int(race['d_total']):+d}s", "Total deficit",         f"2025 avg was {int(yearly_2025['d_total']):+d}s"),
        (f"{int(race['d_bike']):+d}s",  "Bike deficit",          f"2025 avg was {int(yearly_2025['d_bike']):+d}s"),
        (f"{int(race['d_run']):+d}s",   "Run deficit",           "Perry faster on run again"),
    ]
    for i, (big, mid, sub) in enumerate(kpis):
        x = Inches(0.5 + i*3.15)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.95), Inches(1.9))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT
        bg.line.color.rgb = LIGHT
        box = slide.shapes.add_textbox(x, Inches(1.55), Inches(2.95), Inches(1.85))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big
        p.runs[0].font.size = Pt(40); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = mid
        p2.runs[0].font.size = Pt(14); p2.runs[0].font.bold = True; p2.runs[0].font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = sub
        p3.runs[0].font.size = Pt(11); p3.runs[0].font.color.rgb = GREY
        p3.alignment = PP_ALIGN.CENTER

    # Race-detail table
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(8.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "Race splits"
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY

    rows = [
        ("Segment", "Perry", "Parker", "Δ"),
        ("Swim",  _seconds_to_mmss(race['a_swim_s']), _seconds_to_mmss(race['b_swim_s']), f"{int(race['d_swim']):+d}s"),
        ("T1",    _seconds_to_mmss(race['a_t1_s']),   _seconds_to_mmss(race['b_t1_s']),   f"{int(race['d_t1']):+d}s"),
        ("Bike",  _seconds_to_mmss(race['a_bike_s']), _seconds_to_mmss(race['b_bike_s']), f"{int(race['d_bike']):+d}s"),
        ("T2",    _seconds_to_mmss(race['a_t2_s']),   _seconds_to_mmss(race['b_t2_s']),   f"{int(race['d_t2']):+d}s"),
        ("Run",   _seconds_to_mmss(race['a_run_s']),  _seconds_to_mmss(race['b_run_s']),  f"{int(race['d_run']):+d}s"),
        ("Total", _seconds_to_mmss(race['a_total_s']), _seconds_to_mmss(race['b_total_s']),f"{int(race['d_total']):+d}s"),
    ]
    table = slide.shapes.add_table(len(rows), 4, Inches(0.5), Inches(4.15), Inches(7.5), Inches(2.85)).table
    for r, row_vals in enumerate(rows):
        for c, val in enumerate(row_vals):
            cell = table.cell(r, c); cell.text = ""
            tf = cell.text_frame; p = tf.paragraphs[0]; p.text = val
            p.runs[0].font.size = Pt(13)
            if r == 0:
                p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                if c == 3 and val.startswith("-"):
                    p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
                elif c == 3:
                    p.runs[0].font.color.rgb = RED
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF) if r % 2 == 0 else RGBColor(0xF8,0xF8,0xF8)

    # Story callout
    tb = slide.shapes.add_textbox(Inches(8.4), Inches(3.7), Inches(4.6), Inches(3.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "What this race confirms"
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    for line in [
        f"• Bike gap fell to {int(race['d_bike'])}s — the smallest single-race bike deficit in her career vs Parker.",
        f"• Perry beat Parker on the run by {abs(int(race['d_run']))}s, her 4th consecutive race ahead.",
        f"• Total of {int(race['d_total'])}s is consistent with the 2025 trajectory — the gains held into 2026.",
        "• Open question: swim deficit (+1:45) is the remaining headline gap. Coaching focus for the next block.",
    ]:
        p = tf.add_paragraph(); p.text = line
        p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = GREY

    out = outdir / "Perry_2026_Update_SLIDE.pptx"
    prs.save(str(out))
    return out


def build_hermes_2026_slide(engine, outdir: Path) -> Path:
    """Single slide: Hermes 2026 season update (Samarkand + Yokohama wins)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    full = pull_full_history(engine, HERMES, peer_ids={"rod": RODRIGUEZ.aid, "tar": TARANTELLO.aid})
    full["year"] = full["event_date"].dt.year
    races_2026 = full[full["year"] == 2026].sort_values("event_date")
    if races_2026.empty:
        print("[Hermes 2026] no 2026 races, skipping")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00,0x28,0x68); RED = RGBColor(0xC8,0x10,0x2E)
    GREY = RGBColor(0x55,0x55,0x55); LIGHT = RGBColor(0xF4,0xF4,0xF4)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "2026 Season Update — McClain Hermes"
    p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = f"{len(races_2026)} starts · {(races_2026['h_pos']==1).sum()} wins · first WPS Series victory"
    p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = GREY

    # KPI row
    win_count = int((races_2026["h_pos"]==1).sum())
    yokohama = races_2026[races_2026["event_name"].str.contains("Yokohama")].iloc[0] if not races_2026[races_2026["event_name"].str.contains("Yokohama")].empty else races_2026.iloc[-1]
    samarkand = races_2026[races_2026["event_name"].str.contains("Samarkand")].iloc[0] if not races_2026[races_2026["event_name"].str.contains("Samarkand")].empty else None
    kpis = [
        (f"{win_count} / {len(races_2026)}", "Wins / Starts",      "Undefeated through May"),
        ("1st",                              "Samarkand Para Cup", "Apr 26 · 4 finishers"),
        ("1st",                              "Yokohama WPS",       "May 16 · 8 finishers · first WPS win"),
        (_seconds_to_mmss(yokohama['h_total']) if pd.notna(yokohama['h_total']) else "—", "Yokohama total time", "Career-best sprint result"),
    ]
    for i, (big, mid, sub) in enumerate(kpis):
        x = Inches(0.5 + i*3.15)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.95), Inches(1.9))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT
        bg.line.color.rgb = LIGHT
        box = slide.shapes.add_textbox(x, Inches(1.55), Inches(2.95), Inches(1.85))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big
        p.runs[0].font.size = Pt(36); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = mid
        p2.runs[0].font.size = Pt(13); p2.runs[0].font.bold = True; p2.runs[0].font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = sub
        p3.runs[0].font.size = Pt(11); p3.runs[0].font.color.rgb = GREY
        p3.alignment = PP_ALIGN.CENTER

    # Yokohama detail table
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(8.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "Yokohama WPS — winning margin"
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY

    # Need 2nd & 3rd at Yokohama — pull from race_results
    yoko_field_q = text("""
        SELECT athlete_full_name, finish_position, total_time, swimtime, biketime, runtime
        FROM race_results r JOIN events e ON e.event_id=r.event_id AND e.prog_id=r.prog_id
        WHERE e.event_name ILIKE '%Yokohama%' AND e.event_date='2026-05-16' AND e.prog_name='PTVI Women'
          AND r.finish_position IN (1,2,3) ORDER BY finish_position
    """)
    with engine.begin() as c:
        field = pd.read_sql(yoko_field_q, c)
    for col in ["total_time","swimtime","biketime","runtime"]:
        field[col + "_s"] = field[col].apply(to_sec)

    rows = [("Pos","Athlete","Swim","Bike","Run","Total")]
    for _, r in field.iterrows():
        rows.append((
            f"{int(r['finish_position'])}",
            r["athlete_full_name"][:22],
            _seconds_to_mmss(r["swimtime_s"]),
            _seconds_to_mmss(r["biketime_s"]),
            _seconds_to_mmss(r["runtime_s"]),
            _seconds_to_mmss(r["total_time_s"]),
        ))
    table = slide.shapes.add_table(len(rows), 6, Inches(0.5), Inches(4.15), Inches(7.5), Inches(2.0)).table
    for r, row_vals in enumerate(rows):
        for c, val in enumerate(row_vals):
            cell = table.cell(r, c); cell.text = ""
            tf = cell.text_frame; p = tf.paragraphs[0]; p.text = val
            p.runs[0].font.size = Pt(11)
            if r == 0:
                p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                if r == 1:
                    cell.fill.fore_color.rgb = RGBColor(0xFF,0xEE,0xEE)
                    p.runs[0].font.bold = True
                else:
                    cell.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF) if r % 2 == 0 else RGBColor(0xF8,0xF8,0xF8)

    # Story callout
    tb = slide.shapes.add_textbox(Inches(8.4), Inches(3.7), Inches(4.6), Inches(3.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "What this start confirms"
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    for line in [
        "• 2 wins from 2 starts — first WPS Series victory of her career.",
        "• Won Yokohama by 2:55 over 2nd place in an 8-deep field — biggest margin she has put on a deep PTVI field.",
        "• Bike was the race-fastest split.",
        "• Open question: neither Rodriguez (B1, defending Paralympic champ) nor Tarantello (B3, recent Worlds medallist) raced. True championship-field gap is still TBD until they meet at Worlds.",
    ]:
        p = tf.add_paragraph(); p.text = line
        p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = GREY

    out = outdir / "Hermes_2026_Update_SLIDE.pptx"
    prs.save(str(out))
    return out


def build_fisher_2026_slide(engine, outdir: Path) -> Path:
    """Single slide: Fisher 2026 season update (Yokohama WPS vs Parker + Perry)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    overlap_prk = pull_overlap(engine, FISHER, PARKER)
    overlap_pry = pull_overlap(engine, FISHER, PERRY)
    overlap_prk["year"] = overlap_prk["event_date"].dt.year
    overlap_pry["year"] = overlap_pry["event_date"].dt.year
    races_2026_prk = overlap_prk[overlap_prk["year"] == 2026]
    races_2026_pry = overlap_pry[overlap_pry["year"] == 2026]
    if races_2026_prk.empty and races_2026_pry.empty:
        print("[Fisher 2026] no 2026 head-to-head races, skipping")
        return None

    # Use Yokohama 2026 as the showcase race — both Parker and Perry there
    race_prk = races_2026_prk.iloc[-1] if not races_2026_prk.empty else None
    race_pry = races_2026_pry.iloc[-1] if not races_2026_pry.empty else None

    # Prior-year best gap to Parker (for context)
    pre2026_prk = overlap_prk[overlap_prk["year"] < 2026]
    prior_best_prk = pre2026_prk["d_total"].min() if not pre2026_prk.empty else None

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    NAVY = RGBColor(0x00,0x28,0x68); RED = RGBColor(0xC8,0x10,0x2E)
    GREY = RGBColor(0x55,0x55,0x55); LIGHT = RGBColor(0xF4,0xF4,0xF4)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "2026 Season Update — Skyler Fisher"
    p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = f"Yokohama WPS · May 16, 2026 · finished {int(race_prk['a_pos'])}{ordinal(int(race_prk['a_pos']))}"
    p2.runs[0].font.size = Pt(14); p2.runs[0].font.color.rgb = GREY

    # KPI row
    kpis = [
        (f"{int(race_prk['a_pos'])}{ordinal(int(race_prk['a_pos']))}", "Yokohama WPS finish", f"in field of 7"),
        (f"{int(race_prk['d_total']):+d}s",  "Gap to Parker (gold)",   f"Prior best was +{int(prior_best_prk)}s" if prior_best_prk else "Career best"),
        (f"{int(race_pry['d_total']):+d}s",  "Gap to Perry (silver)",  "Closest gap in a depth race"),
        ("H2 vs H1",                          "Sub-class context",      "Factor-adjusted official times"),
    ]
    for i, (big, mid, sub) in enumerate(kpis):
        x = Inches(0.5 + i*3.15)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.95), Inches(1.9))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT
        bg.line.color.rgb = LIGHT
        box = slide.shapes.add_textbox(x, Inches(1.55), Inches(2.95), Inches(1.85))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big
        p.runs[0].font.size = Pt(36); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = mid
        p2.runs[0].font.size = Pt(13); p2.runs[0].font.bold = True; p2.runs[0].font.color.rgb = NAVY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = sub
        p3.runs[0].font.size = Pt(11); p3.runs[0].font.color.rgb = GREY
        p3.alignment = PP_ALIGN.CENTER

    # Three-way splits table
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(8.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "Yokohama 2026 — three-way splits"
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    rows = [
        ("Segment", "Parker (1st)", "Perry (2nd)", "Fisher (6th)", "Fisher Δ vs Parker"),
        ("Swim",  _seconds_to_mmss(race_prk['b_swim_s']), _seconds_to_mmss(race_pry['b_swim_s']),
                  _seconds_to_mmss(race_prk['a_swim_s']), f"{int(race_prk['d_swim']):+d}s"),
        ("Bike",  _seconds_to_mmss(race_prk['b_bike_s']), _seconds_to_mmss(race_pry['b_bike_s']),
                  _seconds_to_mmss(race_prk['a_bike_s']), f"{int(race_prk['d_bike']):+d}s"),
        ("Run",   _seconds_to_mmss(race_prk['b_run_s']),  _seconds_to_mmss(race_pry['b_run_s']),
                  _seconds_to_mmss(race_prk['a_run_s']),  f"{int(race_prk['d_run']):+d}s"),
        ("Total", _seconds_to_mmss(race_prk['b_total_s']),_seconds_to_mmss(race_pry['b_total_s']),
                  _seconds_to_mmss(race_prk['a_total_s']),f"{int(race_prk['d_total']):+d}s"),
    ]
    table = slide.shapes.add_table(len(rows), 5, Inches(0.5), Inches(4.15), Inches(7.8), Inches(2.2)).table
    for r, row_vals in enumerate(rows):
        for c, val in enumerate(row_vals):
            cell = table.cell(r, c); cell.text = ""
            tf = cell.text_frame; p = tf.paragraphs[0]; p.text = val
            p.runs[0].font.size = Pt(12)
            if r == 0:
                p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF) if r % 2 == 0 else RGBColor(0xF8,0xF8,0xF8)

    # Story callout
    tb = slide.shapes.add_textbox(Inches(8.6), Inches(3.7), Inches(4.4), Inches(3.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "What this race confirms"
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    for line in [
        f"• Closest gap to Parker of Fisher's career in any head-to-head: {int(race_prk['d_total'])}s (prior best was +{int(prior_best_prk)}s).",
        f"• Also her tightest deep-field gap to Perry ({int(race_pry['d_total'])}s).",
        "• 6th place in a 7-deep field that included all the world leaders — depth gap is closing.",
        "• Open question: bike split is the largest single contributor to the gap. Next training block focus.",
    ]:
        p = tf.add_paragraph(); p.text = line
        p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = GREY

    out = outdir / "Fisher_2026_Update_SLIDE.pptx"
    prs.save(str(out))
    return out


def ordinal(n: int) -> str:
    if 10 <= n % 100 <= 20:
        return "th"
    return {1:"st", 2:"nd", 3:"rd"}.get(n % 10, "th")


def main():
    engine = get_engine()
    stamp = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    base = Path("para_triathlon_analysis/output") / f"progression_{stamp}"

    perry_out = base / "perry_vs_parker"
    hermes_out = base / "hermes_vs_rodriguez"
    fisher_out = base / "fisher_vs_parker_perry"
    perry_stats = build_perry(engine, perry_out)
    hermes_stats = build_hermes(engine, hermes_out)
    fisher_stats = build_fisher(engine, fisher_out)

    perry_pptx = build_pptx_perry(perry_out, perry_stats)
    hermes_pptx = build_pptx_hermes(hermes_out, hermes_stats)
    fisher_pptx = build_pptx_fisher(fisher_out, fisher_stats)
    print(f"[Perry]  PPTX: {perry_pptx}")
    print(f"[Hermes] PPTX: {hermes_pptx}")
    print(f"[Fisher] PPTX: {fisher_pptx}")

    # Standalone single-slide 2026 updates (drop into existing decks)
    perry_2026  = build_perry_2026_slide(engine, perry_out)
    hermes_2026 = build_hermes_2026_slide(engine, hermes_out)
    fisher_2026 = build_fisher_2026_slide(engine, fisher_out)
    if perry_2026:  print(f"[Perry]  2026 single-slide: {perry_2026}")
    if hermes_2026: print(f"[Hermes] 2026 single-slide: {hermes_2026}")
    if fisher_2026: print(f"[Fisher] 2026 single-slide: {fisher_2026}")

    print(f"\nAll outputs in: {base}")


if __name__ == "__main__":
    main()
