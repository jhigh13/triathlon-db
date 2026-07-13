#!/usr/bin/env python
"""
relay_analysis.py — World Triathlon Mixed Relay Championships historical analysis.

Produces a PowerPoint deck covering nation performance, leg split profiles,
race dynamics (gap evolution), and competitive density.

Usage:
    python scripts/relay_analysis.py --venue Hamburg --years 6
    python scripts/relay_analysis.py --venue Hamburg --output relay_deck.pptx
"""

import argparse
import io
import os
import sys
from datetime import date, timedelta

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd

_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(_SCRIPT_DIR, ".."))
if REPO_ROOT not in sys.path:
    sys.path.insert(0, REPO_ROOT)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import text

from tri_analysis.database import get_engine

# ── Paths ─────────────────────────────────────────────────────────────────────
DEFAULT_OUTPUT_DIR = os.path.join(REPO_ROOT, "ppt files")
USAT_LOGO_PATH = os.path.join(REPO_ROOT, "docs", "power_bi_files", "USA_Triathlon_Logo.jpg")

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY          = RGBColor(0x00, 0x20, 0x60)
RED           = RGBColor(0xC0, 0x00, 0x00)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY      = RGBColor(0xBF, 0xBF, 0xBF)
DARK_GRAY     = RGBColor(0x26, 0x26, 0x26)
SUBTITLE_BLUE = RGBColor(0xAA, 0xBB, 0xDD)

C_NAVY       = "#002060"
C_RED        = "#C00000"
C_LIGHT_BLUE = "#4472C4"
C_GRAY       = "#808080"
C_ORANGE     = "#E76F51"
C_GREEN      = "#2A9D8F"
C_GOLD       = "#F4A261"
C_VIOLET     = "#9B5DE5"

FONT    = "Arial"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

LEG_LABELS = {1: "W1", 2: "M1", 3: "W2", 4: "M2"}

YEAR_PALETTE = [C_NAVY, C_LIGHT_BLUE, C_GREEN, C_ORANGE, C_RED, C_VIOLET,
                "#9B5DE5", "#F72585", "#4CC9F0", "#7209B7"]

plt.rcParams.update({
    "font.family": "Arial",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.titleweight": "bold",
})


# ── Utilities ─────────────────────────────────────────────────────────────────

def parse_time_to_seconds(t) -> float | None:
    if t is None or str(t).strip() in ("", "None", "nan"):
        return None
    t = str(t).strip()
    parts = t.split(":")
    try:
        if len(parts) == 3:
            return int(parts[0]) * 3600 + int(parts[1]) * 60 + float(parts[2])
        elif len(parts) == 2:
            return int(parts[0]) * 60 + float(parts[1])
        else:
            return float(t)
    except (ValueError, IndexError):
        return None


def seconds_to_mmss(s, always_hours: bool = False) -> str:
    if s is None or (isinstance(s, float) and np.isnan(s)):
        return "—"
    s = int(round(float(s)))
    h = s // 3600
    m = (s % 3600) // 60
    sec = s % 60
    if h > 0 or always_hours:
        return f"{h}:{m:02d}:{sec:02d}"
    return f"{m}:{sec:02d}"


def spearman_r(x, y) -> float:
    n = len(x)
    if n < 3:
        return float("nan")
    rx = pd.Series(x).rank().values
    ry = pd.Series(y).rank().values
    d = rx - ry
    return 1.0 - 6.0 * float((d ** 2).sum()) / (n * (n * n - 1))


# ── pptx helpers ──────────────────────────────────────────────────────────────

def _add_textbox(slide, text: str, left, top, width, height,
                 font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
                 italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    run.font.color.rgb = color or DARK_GRAY
    return txb


def add_slide_chrome(slide, title: str, subtitle: str = "", show_logo: bool = True):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    _add_textbox(slide, title, Inches(0.25), Inches(0.05), Inches(10.5), Inches(0.65),
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, subtitle, Inches(0.25), Inches(0.68), Inches(10.5), Inches(0.35),
                     font_size=12, color=SUBTITLE_BLUE)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(1.05), SLIDE_W, Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()

    if show_logo and os.path.exists(USAT_LOGO_PATH):
        slide.shapes.add_picture(USAT_LOGO_PATH, Inches(11.9), Inches(0.03),
                                 Inches(1.35), Inches(1.0))

    footer = slide.shapes.add_shape(1, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15))
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY
    footer.line.fill.background()


def fig_to_image(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── DB queries ────────────────────────────────────────────────────────────────

def query_relay_events(engine, venue: str, years_back: int) -> pd.DataFrame:
    cutoff = date.today() - timedelta(days=years_back * 365)
    sql = text("""
        SELECT e.event_id, e.prog_id, e.event_name, e.event_date, e.event_venue
        FROM events e
        WHERE (e.event_venue ILIKE :venue_pat OR e.event_name ILIKE :venue_pat)
          AND e.prog_name ILIKE '%relay%'
          AND e.event_date >= :cutoff
          AND e.event_date <= :today
        ORDER BY e.event_date ASC
    """)
    with engine.connect() as conn:
        return pd.read_sql(sql, conn, params={
            "venue_pat": f"%{venue}%",
            "cutoff": cutoff,
            "today": date.today(),
        })


def load_relay_legs(engine, events_df: pd.DataFrame) -> pd.DataFrame:
    """One row per athlete-leg with splits in seconds and country from athlete table."""
    if events_df.empty:
        return pd.DataFrame()

    chunks = []
    for _, ev in events_df.iterrows():
        eid, pid = int(ev.event_id), int(ev.prog_id)
        year = pd.to_datetime(ev.event_date).year

        sql = text("""
            SELECT
                rr.athlete_id,
                rr.athlete_full_name,
                rr.finish_position,
                rr.start_num,
                rr.swimtime,
                rr.t1time,
                rr.biketime,
                rr.t2time,
                rr.runtime,
                rr.total_time,
                COALESCE(NULLIF(a.country, ''), '???') AS country
            FROM race_results rr
            LEFT JOIN athlete a ON rr.athlete_id = a.athlete_id
            WHERE rr.event_id = :eid AND rr.prog_id = :pid
              AND rr.total_time IS NOT NULL
        """)
        with engine.connect() as conn:
            df = pd.read_sql(sql, conn, params={"eid": eid, "pid": pid})

        if df.empty:
            continue

        df["leg"] = pd.to_numeric(df["start_num"], errors="coerce")
        df = df[df["leg"].isin([1.0, 2.0, 3.0, 4.0])].copy()
        df["leg"] = df["leg"].astype(int)

        df["team_rank"] = pd.to_numeric(df["finish_position"], errors="coerce")
        df = df[df["team_rank"].notna()].copy()
        df["team_rank"] = df["team_rank"].astype(int)

        for src, dst in [
            ("swimtime", "swim_s"), ("biketime", "bike_s"), ("runtime", "run_s"),
            ("t1time", "t1_s"),    ("t2time", "t2_s"),    ("total_time", "leg_total_s"),
        ]:
            df[dst] = df[src].apply(parse_time_to_seconds)

        df["year"] = year
        df["event_id"] = eid
        df["event_name"] = str(ev.event_name)
        df["leg_label"] = df["leg"].map(LEG_LABELS)
        chunks.append(df)

    if not chunks:
        return pd.DataFrame()
    return pd.concat(chunks, ignore_index=True)


def build_team_cumulative(legs_df: pd.DataFrame) -> pd.DataFrame:
    """Cumulative team time and gap to leader at each leg handoff."""
    if legs_df.empty:
        return pd.DataFrame()

    records = []
    for (year, country, team_rank), grp in legs_df.groupby(
            ["year", "country", "team_rank"], sort=False):
        grp_s = grp.sort_values("leg").drop_duplicates("leg", keep="first")
        if not {1, 2, 3, 4}.issubset(set(grp_s["leg"])):
            continue
        cum = 0.0
        valid = True
        team_records = []
        for _, row in grp_s.iterrows():
            if pd.isna(row["leg_total_s"]):
                valid = False
                break
            cum += row["leg_total_s"]
            team_records.append({
                "year": year, "country": country, "team_rank": team_rank,
                "leg": int(row["leg"]), "leg_label": row["leg_label"],
                "cumulative_s": cum,
            })
        if valid:
            records.extend(team_records)

    if not records:
        return pd.DataFrame()

    cum_df = pd.DataFrame(records)
    leader = (cum_df.groupby(["year", "leg"])["cumulative_s"]
              .min().reset_index().rename(columns={"cumulative_s": "leader_s"}))
    cum_df = cum_df.merge(leader, on=["year", "leg"])
    cum_df["gap_to_leader_s"] = cum_df["cumulative_s"] - cum_df["leader_s"]
    return cum_df


def build_team_summary(legs_df: pd.DataFrame) -> pd.DataFrame:
    """One row per (year, country, team_rank) with total team time."""
    if legs_df.empty:
        return pd.DataFrame()
    agg = (legs_df.groupby(["year", "country", "team_rank"])
           .agg(team_time=("leg_total_s", "sum"),
                leg_count=("leg", "count"),
                event_name=("event_name", "first"))
           .reset_index())
    return agg[agg["leg_count"] >= 4].copy()


# ── Slide builders ────────────────────────────────────────────────────────────

def add_relay_title_slide(prs: Presentation, venue: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    acc = slide.shapes.add_shape(1, Inches(12.5), Inches(0), Inches(0.83), SLIDE_H)
    acc.fill.solid()
    acc.fill.fore_color.rgb = RED
    acc.line.fill.background()

    rule = slide.shapes.add_shape(1, Inches(0.5), Inches(4.1), Inches(11.8), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = WHITE
    rule.line.fill.background()

    if os.path.exists(USAT_LOGO_PATH):
        slide.shapes.add_picture(USAT_LOGO_PATH, Inches(0.6), Inches(0.4),
                                 Inches(2.1), Inches(1.6))

    _add_textbox(slide, f"{venue.upper()} MIXED RELAY",
                 Inches(0.5), Inches(1.9), Inches(11.8), Inches(1.0),
                 font_size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "WORLD CHAMPIONSHIPS",
                 Inches(0.5), Inches(2.85), Inches(11.8), Inches(0.7),
                 font_size=32, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _add_textbox(slide,
                 f"Mixed Relay Historical Analysis  \u00b7  {date.today().strftime('%B %Y')}",
                 Inches(0.5), Inches(4.3), Inches(11.8), Inches(0.6),
                 font_size=18, color=SUBTITLE_BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "USA Triathlon High Performance",
                 Inches(0.5), Inches(5.0), Inches(11.8), Inches(0.5),
                 font_size=14, color=RGBColor(0x80, 0x90, 0xB0), align=PP_ALIGN.CENTER)


def add_relay_overview_slide(prs: Presentation, events_df: pd.DataFrame,
                              legs_df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Event History  \u2014  World Mixed Relay Championships",
                     subtitle="Recent editions at Hamburg; individual-leg splits available 2020\u20132025")

    teams_per_year: dict[int, int] = {}
    if not legs_df.empty:
        tpy = legs_df.groupby(["year", "team_rank"]).size().reset_index()
        tpy = tpy.groupby("year").size()
        teams_per_year = tpy.to_dict()

    col_x = [Inches(0.4), Inches(1.7), Inches(6.6), Inches(11.9)]
    col_w = [Inches(1.1), Inches(4.7), Inches(5.1), Inches(1.1)]
    headers = ["Year", "Date", "Event", "Teams"]
    header_y = Inches(1.3)
    header_h = Inches(0.36)

    hdr = slide.shapes.add_shape(1, Inches(0.3), header_y, Inches(12.5), header_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()
    for txt, x, w in zip(headers, col_x, col_w):
        _add_textbox(slide, txt, x, header_y + Inches(0.04), w, header_h,
                     font_size=11, bold=True, color=WHITE)

    row_y = Inches(1.66)
    row_h = Inches(0.38)
    for i, (_, ev) in enumerate(events_df.iterrows()):
        yr = pd.to_datetime(ev.event_date).year
        dt = pd.to_datetime(ev.event_date).strftime("%d %b %Y")
        teams_str = str(teams_per_year.get(yr, "\u2014"))
        bg_clr = LIGHT_GRAY if i % 2 == 0 else WHITE

        bg = slide.shapes.add_shape(1, Inches(0.3), row_y, Inches(12.5), row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_clr
        bg.line.fill.background()

        for txt, x, w in zip([str(yr), dt, str(ev.event_name), teams_str],
                               col_x, col_w):
            _add_textbox(slide, txt, x, row_y + Inches(0.06), w, row_h,
                         font_size=11, color=DARK_GRAY)
        row_y += row_h

    _add_textbox(slide,
                 "Note: 2013\u20132019 editions are in the events table; "
                 "individual-leg timing data not available for those years.",
                 Inches(0.4), Inches(6.88), Inches(12.0), Inches(0.38),
                 font_size=9, italic=True, color=MID_GRAY)


def add_nation_podiums_slide(prs: Presentation, team_summary: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Nation Podium Record  \u2014  Mixed Relay World Championships",
                     subtitle="Gold / silver / bronze finishes by nation, 2020\u20132025")

    medals = team_summary[team_summary["team_rank"] <= 3].copy()
    medals["medal"] = medals["team_rank"].map({1: "Gold", 2: "Silver", 3: "Bronze"})
    medal_counts = (medals.groupby(["country", "medal"])
                    .size().unstack(fill_value=0)
                    .reindex(columns=["Gold", "Silver", "Bronze"], fill_value=0))
    medal_counts["pts"] = (medal_counts["Gold"] * 3
                           + medal_counts["Silver"] * 2
                           + medal_counts["Bronze"])
    medal_counts = medal_counts.sort_values(["pts", "Gold"], ascending=False)

    fig, ax = plt.subplots(figsize=(12.5, 5.2))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    countries = medal_counts.index.tolist()
    x = np.arange(len(countries))
    w = 0.22
    bars_g = ax.bar(x - w, medal_counts["Gold"],   width=w,
                    color="#FFD700", label="Gold",   edgecolor="white", linewidth=0.5)
    bars_s = ax.bar(x,     medal_counts["Silver"], width=w,
                    color="#C0C0C0", label="Silver", edgecolor="white", linewidth=0.5)
    bars_b = ax.bar(x + w, medal_counts["Bronze"], width=w,
                    color="#CD7F32", label="Bronze", edgecolor="white", linewidth=0.5)

    for bars in (bars_g, bars_s, bars_b):
        for b in bars:
            h = b.get_height()
            if h:
                ax.text(b.get_x() + b.get_width() / 2, h + 0.05, str(int(h)),
                        ha="center", va="bottom", fontsize=10, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(countries, fontsize=12)
    ax.set_ylabel("Count", fontsize=11)
    ax.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax.legend(fontsize=11, loc="upper right")
    ax.set_title("Medals by Nation  (2020\u20132025)", fontsize=14, pad=10)
    mx = medal_counts[["Gold", "Silver", "Bronze"]].values.max()
    ax.set_ylim(0, mx + 1.5)
    plt.tight_layout(pad=1.5)

    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.25),
                              Inches(12.73), Inches(5.9))


def add_results_table_slide(prs: Presentation, team_summary: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Year-by-Year Results  \u2014  Final Standings",
                     subtitle="Top-4 finish order and winning team time for each edition")

    col_x   = [Inches(0.4),  Inches(1.9),  Inches(4.3),
               Inches(6.7),  Inches(9.1),  Inches(11.5)]
    col_w   = [Inches(1.3),  Inches(2.2),  Inches(2.2),
               Inches(2.2),  Inches(2.2),  Inches(1.5)]
    headers = ["Year", "1st", "2nd", "3rd", "4th", "Win Time"]
    header_y = Inches(1.3)
    header_h = Inches(0.38)

    hdr = slide.shapes.add_shape(1, Inches(0.3), header_y, Inches(12.8), header_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()
    for txt, x, w in zip(headers, col_x, col_w):
        _add_textbox(slide, txt, x, header_y + Inches(0.04), w, header_h,
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    row_y = Inches(1.68)
    row_h = Inches(0.44)
    for i, year in enumerate(sorted(team_summary["year"].unique())):
        yr_df = team_summary[team_summary["year"] == year].sort_values("team_rank")

        def _get(rank):
            row = yr_df[yr_df["team_rank"] == rank]
            return row.iloc[0]["country"] if not row.empty else "\u2014"

        win_row = yr_df[yr_df["team_rank"] == 1]
        win_time = (seconds_to_mmss(win_row.iloc[0]["team_time"], always_hours=True)
                    if not win_row.empty else "\u2014")
        cells = [str(year), _get(1), _get(2), _get(3), _get(4), win_time]

        bg_clr = LIGHT_GRAY if i % 2 == 0 else WHITE
        bg = slide.shapes.add_shape(1, Inches(0.3), row_y, Inches(12.8), row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_clr
        bg.line.fill.background()

        for txt, x, w in zip(cells, col_x, col_w):
            _add_textbox(slide, txt, x, row_y + Inches(0.07), w, row_h,
                         font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        row_y += row_h


def add_leg_split_profile_slide(prs: Presentation, legs_df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Leg Split Profiles  \u2014  Average Time by Segment",
                     subtitle="Mean swim / bike / run per relay leg; whiskers show field min\u2013max range")

    fig, axes = plt.subplots(1, 4, figsize=(13, 5.0), sharey=False)
    fig.patch.set_facecolor("white")

    split_cols   = ["swim_s", "bike_s", "run_s"]
    split_labels = ["Swim", "Bike", "Run"]
    split_colors = [C_LIGHT_BLUE, C_NAVY, C_RED]

    for ax, leg_num in zip(axes, [1, 2, 3, 4]):
        ax.set_facecolor("white")
        sub = legs_df[legs_df["leg"] == leg_num].dropna(subset=split_cols)
        if sub.empty:
            ax.set_title(LEG_LABELS[leg_num], fontsize=12)
            continue

        means = [sub[c].mean() for c in split_cols]
        mins  = [sub[c].min()  for c in split_cols]
        maxs  = [sub[c].max()  for c in split_cols]

        x = np.arange(len(split_cols))
        bars = ax.bar(x, means, color=split_colors, edgecolor="white",
                      linewidth=0.5, width=0.55)
        for xi, mn, lo, hi in zip(x, means, mins, maxs):
            ax.vlines(xi, lo, hi, color=C_GRAY, linewidth=1.5, zorder=3)
            ax.scatter([xi, xi], [lo, hi], color=C_GRAY, s=22, zorder=4)

        for b, m in zip(bars, means):
            ax.text(b.get_x() + b.get_width() / 2, m + 1,
                    seconds_to_mmss(m), ha="center", va="bottom",
                    fontsize=8, fontweight="bold", color="black")

        ax.set_xticks(x)
        ax.set_xticklabels(split_labels, fontsize=9)
        ax.set_title(f"Leg {LEG_LABELS[leg_num]}", fontsize=12, fontweight="bold")
        ax.yaxis.set_major_formatter(
            mticker.FuncFormatter(lambda v, _: f"{int(v//60)}:{int(round(v%60)):02d}"))
        ax.set_ylabel("Time" if leg_num == 1 else "", fontsize=9)
        ax.grid(axis="y", alpha=0.25, linestyle="--")

    plt.tight_layout(pad=1.2)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.25),
                              Inches(12.73), Inches(5.9))


def add_country_leg_strengths_slide(prs: Presentation, legs_df: pd.DataFrame,
                                     top_n: int = 10):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Country Leg Strengths  \u2014  Relative Performance by Leg",
                     subtitle="Z-score of mean leg time per slot. Green = faster than field average, red = slower.")

    leg_order = [1, 2, 3, 4]

    counts = legs_df.groupby("country")["leg"].count()
    top_countries = counts.nlargest(top_n).index.tolist()
    sub = legs_df[legs_df["country"].isin(top_countries)].copy()

    means = (sub.groupby(["country", "leg"])["leg_total_s"]
             .mean().unstack(level="leg").reindex(columns=leg_order))

    zscores = (means - means.mean()) / means.std()

    means["avg"] = means.mean(axis=1)
    ordered = means.sort_values("avg").index.tolist()
    means.drop(columns="avg", inplace=True)
    zscores = zscores.loc[ordered]
    means = means.loc[ordered]

    z_matrix   = zscores.values.astype(float)
    mean_matrix = means.values.astype(float)

    n_rows = len(ordered)
    fig_h  = max(4.0, n_rows * 0.45 + 0.6)
    fig, ax = plt.subplots(figsize=(8, fig_h))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    vabs = max(float(np.nanmax(np.abs(z_matrix))), 1.5)
    im = ax.imshow(z_matrix, cmap="RdYlGn_r", vmin=-vabs, vmax=vabs, aspect="auto")

    ax.set_xticks(range(4))
    ax.set_xticklabels([LEG_LABELS[l] for l in leg_order], fontsize=12, fontweight="bold")
    ax.set_yticks(range(n_rows))
    ax.set_yticklabels(ordered, fontsize=11)

    for ri in range(n_rows):
        for ci in range(4):
            m_val = mean_matrix[ri, ci]
            z_val = z_matrix[ri, ci]
            if not np.isnan(m_val):
                txt_color = "black" if abs(z_val) < 1.2 else "white"
                ax.text(ci, ri, seconds_to_mmss(m_val),
                        ha="center", va="center", fontsize=9,
                        color=txt_color, fontweight="bold")

    plt.colorbar(im, ax=ax, label="Z-score (negative = faster)",
                 fraction=0.03, pad=0.02)
    plt.tight_layout(pad=1.2)

    img_h = min(Inches(5.8), Inches(0.8 + n_rows * 0.45))
    slide.shapes.add_picture(fig_to_image(fig), Inches(2.0), Inches(1.3),
                              Inches(9.5), img_h)


def add_gap_evolution_slide(prs: Presentation, cum_df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Race Dynamics  \u2014  Gap to Leader at Each Handoff",
                     subtitle="Cumulative gap in seconds at end of each leg. 0 = leading team. "
                              "Top 8 teams shown per year.")

    years  = sorted(cum_df["year"].unique())
    n      = len(years)
    ncols  = min(3, n)
    nrows  = (n + ncols - 1) // ncols

    fig, axes = plt.subplots(nrows, ncols, figsize=(13, 5.0), squeeze=False)
    fig.patch.set_facecolor("white")

    leg_nums    = [1, 2, 3, 4]
    leg_xlabels = [LEG_LABELS[l] for l in leg_nums]

    for idx, year in enumerate(years):
        row_i, col_i = divmod(idx, ncols)
        ax = axes[row_i][col_i]
        ax.set_facecolor("white")

        yr_df   = cum_df[cum_df["year"] == year].sort_values(["team_rank", "leg"])
        max_gap = 0

        for rank in sorted(yr_df["team_rank"].unique())[:8]:
            team = yr_df[yr_df["team_rank"] == rank].sort_values("leg")
            gaps = team["gap_to_leader_s"].tolist()
            legs = team["leg"].tolist()
            country = team.iloc[-1]["country"] if not team.empty else ""

            if rank == 1:
                color, lw, alpha, zorder = "#FFD700", 2.5, 1.0, 5
            elif rank == 2:
                color, lw, alpha, zorder = "#C0C0C0", 2.0, 0.9, 4
            elif rank == 3:
                color, lw, alpha, zorder = "#CD7F32", 1.8, 0.8, 3
            else:
                color, lw, alpha, zorder = C_GRAY, 1.0, 0.35, 2

            ax.plot(legs, gaps, color=color, linewidth=lw, alpha=alpha,
                    zorder=zorder, marker="o", markersize=4)
            if gaps:
                max_gap = max(max_gap, max(gaps))
                if rank <= 5 and country:
                    ax.annotate(country, (legs[-1], gaps[-1]),
                                xytext=(4, 0), textcoords="offset points",
                                fontsize=7, va="center",
                                color=color if rank <= 3 else C_GRAY)

        ax.set_xticks(leg_nums)
        ax.set_xticklabels(leg_xlabels, fontsize=8)
        ax.set_ylabel("Gap (s)", fontsize=8)
        ax.set_title(str(year), fontsize=11, fontweight="bold")
        ax.set_ylim(-2, max(max_gap + 8, 20))
        ax.axhline(0, color=C_NAVY, linewidth=0.7, linestyle="--", alpha=0.4)
        ax.grid(axis="y", alpha=0.2, linestyle=":")

    for idx in range(n, nrows * ncols):
        row_i, col_i = divmod(idx, ncols)
        axes[row_i][col_i].set_visible(False)

    plt.tight_layout(pad=1.2)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.25),
                              Inches(12.73), Inches(5.9))


def add_competitive_density_slide(prs: Presentation, cum_df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Competitive Density  \u2014  Teams Within Threshold at Each Handoff",
                     subtitle="Average number of teams (excl. leader) within 5 / 10 / 20 seconds "
                              "of the leader, averaged across all editions")

    leg_nums      = [1, 2, 3, 4]
    leg_xlabels   = [f"After {LEG_LABELS[l]}" for l in leg_nums]
    thresholds    = [5, 10, 20]
    thresh_colors = [C_GREEN, C_LIGHT_BLUE, C_ORANGE]
    thresh_labels = ["\u22645 s", "\u226410 s", "\u226420 s"]

    density = {t: [] for t in thresholds}
    for leg in leg_nums:
        sub    = cum_df[(cum_df["leg"] == leg) & (cum_df["gap_to_leader_s"] > 0)]
        years  = sub["year"].unique()
        for t in thresholds:
            counts = [(sub[(sub["year"] == yr) & (sub["gap_to_leader_s"] <= t)].shape[0])
                      for yr in years]
            density[t].append(float(np.mean(counts)) if counts else 0.0)

    fig, ax = plt.subplots(figsize=(10, 5.0))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    x = np.arange(4)
    w = 0.25
    for i, (t, color, label) in enumerate(zip(thresholds, thresh_colors, thresh_labels)):
        bars = ax.bar(x + (i - 1) * w, density[t], width=w,
                      color=color, edgecolor="white", linewidth=0.5, label=label)
        for b, v in zip(bars, density[t]):
            if v > 0.05:
                ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.05,
                        f"{v:.1f}", ha="center", va="bottom",
                        fontsize=9, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(leg_xlabels, fontsize=11)
    ax.set_ylabel("Avg. teams within threshold", fontsize=11)
    ax.set_title("Race Compactness by Handoff  (avg across editions)", fontsize=13, pad=8)
    ax.legend(fontsize=11, loc="upper right")
    ax.yaxis.set_major_locator(mticker.MaxNLocator(integer=False))
    ax.grid(axis="y", alpha=0.25, linestyle="--")
    all_vals = [v for vals in density.values() for v in vals]
    ax.set_ylim(0, max(all_vals) * 1.3 if all_vals else 5)
    plt.tight_layout(pad=1.5)

    slide.shapes.add_picture(fig_to_image(fig), Inches(1.5), Inches(1.25),
                              Inches(10.0), Inches(5.9))


def add_leg1_predictive_slide(prs: Presentation, cum_df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide,
                     "Predictive Power of Leg 1  \u2014  Does W1 Decide the Race?",
                     subtitle="W1 (first swimmer) exit rank vs final team finish. "
                              "Points on the diagonal = perfect prediction.")

    leg1 = cum_df[cum_df["leg"] == 1][["year", "country", "cumulative_s"]].copy()
    leg1["leg1_rank"] = (leg1.groupby("year")["cumulative_s"]
                         .rank(method="min").astype(int))

    final = cum_df[cum_df["leg"] == 4][["year", "country", "team_rank"]].copy()
    merged = leg1.merge(final, on=["year", "country"], how="inner")

    fig, (ax_l, ax_r) = plt.subplots(1, 2, figsize=(13, 5.0),
                                      gridspec_kw={"width_ratios": [2, 1]})
    fig.patch.set_facecolor("white")

    ax_l.set_facecolor("white")
    years      = sorted(merged["year"].unique())
    pal        = YEAR_PALETTE[:len(years)]

    for year, color in zip(years, pal):
        sub = merged[merged["year"] == year]
        ax_l.scatter(sub["leg1_rank"], sub["team_rank"],
                     color=color, s=60, alpha=0.8, label=str(year),
                     edgecolors="white", linewidth=0.5, zorder=3)
        for _, row in sub.iterrows():
            if row["team_rank"] <= 4 or row["leg1_rank"] <= 4:
                ax_l.annotate(row["country"],
                              (row["leg1_rank"], row["team_rank"]),
                              xytext=(3, 2), textcoords="offset points",
                              fontsize=7, color=color, alpha=0.85)

    max_rank = max(int(merged["leg1_rank"].max()), int(merged["team_rank"].max())) + 1
    ax_l.plot([1, max_rank], [1, max_rank], color=C_GRAY, linewidth=1,
              linestyle="--", alpha=0.5, zorder=1)
    ax_l.set_xlabel("Position after W1 (first swim exit)", fontsize=11)
    ax_l.set_ylabel("Final team finish position", fontsize=11)
    ax_l.set_title("W1 Exit Position vs Final Result", fontsize=12, pad=8)
    ax_l.legend(fontsize=9, loc="upper left", title="Year")
    ax_l.set_xlim(0.5, max_rank + 0.5)
    ax_l.set_ylim(0.5, max_rank + 0.5)
    ax_l.xaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax_l.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax_l.grid(alpha=0.2, linestyle=":")

    ax_r.set_facecolor("white")
    ax_r.axis("off")

    rho          = spearman_r(merged["leg1_rank"].tolist(), merged["team_rank"].tolist())
    led          = merged[merged["leg1_rank"] == 1]
    led_won      = led[led["team_rank"] == 1]
    pct_won      = 100.0 * len(led_won) / len(led) if len(led) > 0 else 0.0

    stats = (f"Spearman \u03c1 = {rho:.2f}\n\n"
             f"W1 leader won: {len(led_won)}/{len(led)} races\n"
             f"({pct_won:.0f}%)\n\n"
             f"n = {len(merged)} team\u2013years")
    ax_r.text(0.08, 0.72, stats, transform=ax_r.transAxes,
              fontsize=13, va="top", color=C_NAVY,
              bbox=dict(boxstyle="round,pad=0.5", facecolor="#EEF2FA",
                        edgecolor=C_NAVY, alpha=0.9))
    ax_r.text(0.08, 0.22,
              "\u03c1 > 0.6 = swim order\ndominates result\n"
              "\u03c1 < 0.4 = race decided\nafter the swim",
              transform=ax_r.transAxes, fontsize=10, va="top",
              color=C_GRAY, style="italic")

    plt.tight_layout(pad=1.2)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.25),
                              Inches(12.73), Inches(5.9))


def add_summary_slide(prs: Presentation, legs_df: pd.DataFrame,
                      team_summary: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Key Takeaways  \u2014  Hamburg Mixed Relay",
                     subtitle="2020\u20132025 analysis summary")

    if team_summary.empty or legs_df.empty:
        _add_textbox(slide, "Insufficient data for summary.",
                     Inches(0.5), Inches(2.0), Inches(12.0), Inches(1.0),
                     font_size=14, color=DARK_GRAY)
        return

    gold = (team_summary[team_summary["team_rank"] == 1]
            .groupby("country").size().sort_values(ascending=False))
    top_nation = gold.index[0] if not gold.empty else "\u2014"
    top_wins   = int(gold.iloc[0]) if not gold.empty else 0

    w1_swim = legs_df[legs_df["leg"] == 1].groupby("country")["swim_s"].mean().sort_values()
    top_w1_swim      = w1_swim.index[0] if not w1_swim.empty else "\u2014"
    top_w1_swim_time = seconds_to_mmss(w1_swim.iloc[0]) if not w1_swim.empty else "\u2014"

    m1_bike = legs_df[legs_df["leg"] == 2].groupby("country")["bike_s"].mean().sort_values()
    top_m1_bike      = m1_bike.index[0] if not m1_bike.empty else "\u2014"
    top_m1_bike_time = seconds_to_mmss(m1_bike.iloc[0]) if not m1_bike.empty else "\u2014"

    win_times = (team_summary[team_summary["team_rank"] == 1]
                 .sort_values("year")[["year", "team_time"]])
    if len(win_times) >= 2:
        delta = win_times.iloc[-1]["team_time"] - win_times.iloc[0]["team_time"]
        y0, y1 = int(win_times.iloc[0]["year"]), int(win_times.iloc[-1]["year"])
        trend = (f"{abs(delta):.0f}s faster" if delta < 0
                 else f"{abs(delta):.0f}s slower" if delta > 0 else "unchanged")
        trend_line = f"Winning time trend: {trend} from {y0} to {y1}"
    else:
        trend_line = ""

    bullets = [
        f"Most decorated nation: {top_nation} ({top_wins} gold medal{'s' if top_wins != 1 else ''})",
        f"Fastest W1 swim: {top_w1_swim} (avg {top_w1_swim_time})",
        f"Fastest M1 bike: {top_m1_bike} (avg {top_m1_bike_time})",
        trend_line,
        "Race format: W\u2013M\u2013W\u2013M leg order, 750m swim / 20km bike / 5km run per leg",
        "Split data coverage: 2020\u20132025  (6 editions with full individual-leg timing)",
    ]
    bullets = [b for b in bullets if b]

    y = Inches(1.55)
    for bullet in bullets:
        dot = slide.shapes.add_shape(1, Inches(0.45), y + Inches(0.13),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RED
        dot.line.fill.background()
        _add_textbox(slide, bullet, Inches(0.72), y, Inches(12.0), Inches(0.5),
                     font_size=14, color=DARK_GRAY)
        y += Inches(0.65)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate a World Mixed Relay Championships analysis PowerPoint.")
    parser.add_argument("--venue",  default="Hamburg",
                        help="Venue name (default: Hamburg)")
    parser.add_argument("--years",  type=int, default=6,
                        help="Years to look back (default: 6)")
    parser.add_argument("--output", default=None,
                        help="Output .pptx path (default: auto-named in 'ppt files/')")
    parser.add_argument("--top-n",  type=int, default=10,
                        help="Nations shown in leg-strength heatmap (default: 10)")
    args = parser.parse_args()

    fname = args.output or f"{args.venue.replace(' ', '_')}_Relay_{date.today()}.pptx"
    if not os.path.isabs(fname):
        if os.path.dirname(fname):
            os.makedirs(os.path.dirname(fname), exist_ok=True)
        else:
            os.makedirs(DEFAULT_OUTPUT_DIR, exist_ok=True)
            fname = os.path.join(DEFAULT_OUTPUT_DIR, fname)

    print("Connecting to database...")
    engine = get_engine()

    print(f"Querying relay events at '{args.venue}' (past {args.years} years)...")
    events_df = query_relay_events(engine, args.venue, args.years)
    if events_df.empty:
        print(f"No relay events found for venue '{args.venue}'. Exiting.")
        sys.exit(1)
    for _, ev in events_df.iterrows():
        print(f"  {ev.event_date}  {ev.event_name}")

    print("Loading individual leg results...")
    legs_df = load_relay_legs(engine, events_df)
    if legs_df.empty:
        print("No leg-level result data found. Exiting.")
        sys.exit(1)
    print(f"  {len(legs_df)} athlete-leg rows across "
          f"{legs_df['year'].nunique()} edition(s).")

    print("Building cumulative team times...")
    cum_df       = build_team_cumulative(legs_df)
    team_summary = build_team_summary(legs_df)
    print(f"  {len(team_summary)} complete team-year records.")

    print("Building PowerPoint...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_relay_title_slide(prs, args.venue);           print("  1/10  Title")
    add_relay_overview_slide(prs, events_df, legs_df);print("  2/10  Overview")
    add_nation_podiums_slide(prs, team_summary);      print("  3/10  Nation podiums")
    add_results_table_slide(prs, team_summary);       print("  4/10  Results table")
    add_leg_split_profile_slide(prs, legs_df);        print("  5/10  Leg split profile")
    add_country_leg_strengths_slide(prs, legs_df,
                                    top_n=args.top_n);print("  6/10  Country leg strengths")
    add_gap_evolution_slide(prs, cum_df);             print("  7/10  Gap evolution")
    add_competitive_density_slide(prs, cum_df);       print("  8/10  Competitive density")
    add_leg1_predictive_slide(prs, cum_df);           print("  9/10  Leg-1 predictive power")
    add_summary_slide(prs, legs_df, team_summary);    print(" 10/10  Summary")

    prs.save(fname)
    print(f"\nDone! Saved to: {fname}")


if __name__ == "__main__":
    main()
