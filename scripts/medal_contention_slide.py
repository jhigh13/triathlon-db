#!/usr/bin/env python
"""
medal_contention_slide.py -- Branded USAT deck on WTCS medal contention.

Builds a 2-slide PowerPoint:
  1. "A Global Rarity" -- how few athletes have ever medalled at a WTCS, KPI
     cards, distinct-medalists-by-nation bar chart, and the USA roster.
  2. "The Pipeline Behind the Podium" -- podium vs. top-10 depth by nation,
     showing the USA has the 2nd-deepest top-10 pool in the world.

Reuses the brand chrome/helpers from venue_analysis.py.

Usage:
    python scripts/medal_contention_slide.py                 # wtcs, rolling 4y
    python scripts/medal_contention_slide.py --scope majors
    python scripts/medal_contention_slide.py --output "ppt files/medal_contention.pptx"
"""
import argparse
import os
import sys
from datetime import date

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

_repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _repo_root not in sys.path:
    sys.path.insert(0, _repo_root)
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

# Reuse brand constants + chrome helpers from the venue deck builder.
import venue_analysis as va  # noqa: E402
from venue_analysis import (  # noqa: E402
    NAVY, RED, WHITE, DARK_GRAY, LIGHT_GRAY, MID_GRAY,
    C_NAVY, C_RED, SLIDE_W, SLIDE_H, FONT,
    _add_textbox, _set_cell, add_slide_chrome, fig_to_image,
)
from medal_contention import fetch_medal_data, country_medalists, country_depth  # noqa: E402

HIGHLIGHT = "United States"
C_HL = C_RED
C_BASE = C_NAVY
C_LIGHT = "#8FA8D4"  # light navy for secondary series
DARK_GRAY_HEX = "#262626"


def _ordinal(n: int) -> str:
    suffix = "th" if 10 <= n % 100 <= 20 else {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")
    return f"{n}{suffix}"


# ── KPI cards ───────────────────────────────────────────────────────────────
def _kpi_card(slide, left, top, width, big, label, accent=NAVY):
    card = slide.shapes.add_shape(1, left, top, width, Inches(1.15))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    card.shadow.inherit = False
    _add_textbox(slide, str(big), left, top + Inches(0.06), width, Inches(0.62),
                 font_size=34, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _add_textbox(slide, label, left + Inches(0.08), top + Inches(0.7),
                 width - Inches(0.16), Inches(0.4),
                 font_size=10.5, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── Charts ──────────────────────────────────────────────────────────────────
def _bar_medalists(cm: pd.DataFrame, top_n: int = 10):
    """Horizontal bar: distinct medalists by nation, USA highlighted."""
    d = cm.head(top_n).iloc[::-1]  # biggest on top
    colors = [C_HL if c == HIGHLIGHT else C_BASE for c in d.index]
    fig, ax = plt.subplots(figsize=(5.6, 4.05))
    bars = ax.barh(d.index, d["total"], color=colors, edgecolor="white", height=0.72)
    for bar, (men, women) in zip(bars, d[["men", "women"]].values):
        ax.text(bar.get_width() + 0.12, bar.get_y() + bar.get_height() / 2,
                f"{int(bar.get_width())}", va="center", ha="left",
                fontsize=10, fontweight="bold", color=DARK_GRAY_HEX)
    ax.set_xlabel("Distinct athletes who have medalled", fontsize=9.5)
    ax.set_xlim(0, d["total"].max() + 1.4)
    ax.tick_params(axis="y", labelsize=10)
    ax.tick_params(axis="x", labelsize=8)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    fig.tight_layout()
    return fig_to_image(fig)


def _bar_depth(depth: pd.DataFrame, top_n: int = 10):
    """Grouped horizontal bars: podium medalists vs. top-10 depth per nation."""
    d = depth.head(top_n).iloc[::-1]
    y = range(len(d))
    fig, ax = plt.subplots(figsize=(6.2, 4.15))
    h = 0.38
    ax.barh([i + h / 2 for i in y], d["top10"], height=h, color=C_LIGHT,
            edgecolor="white", label="Top-10 finishers")
    med_colors = [C_HL if c == HIGHLIGHT else C_BASE for c in d.index]
    ax.barh([i - h / 2 for i in y], d["medalists"], height=h, color=med_colors,
            edgecolor="white", label="Medallists (podium)")
    for i, (m, t) in enumerate(zip(d["medalists"], d["top10"])):
        ax.text(t + 0.2, i + h / 2, str(int(t)), va="center", fontsize=8.5, color=DARK_GRAY_HEX)
        ax.text(m + 0.2, i - h / 2, str(int(m)), va="center", fontsize=8.5,
                fontweight="bold", color=DARK_GRAY_HEX)
    ax.set_yticks(list(y))
    ax.set_yticklabels(d.index, fontsize=10)
    ax.set_xlabel("Distinct athletes", fontsize=9.5)
    ax.set_xlim(0, d["top10"].max() + 2)
    ax.tick_params(axis="x", labelsize=8)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    ax.legend(loc="lower right", fontsize=8.5, frameon=False)
    fig.tight_layout()
    return fig_to_image(fig)


# ── USA roster table ─────────────────────────────────────────────────────────
def _usa_roster_table(slide, medalists: pd.DataFrame, left, top, width):
    usa = (medalists[medalists.country == HIGHLIGHT]
           .sort_values(["gender", "medals"], ascending=[True, False]))
    rows = len(usa) + 1
    tbl = slide.shapes.add_table(rows, 3, left, top, width, Inches(0.35 * rows)).table
    tbl.columns[0].width = Inches(2.15)
    tbl.columns[1].width = Inches(1.05)
    tbl.columns[2].width = Inches(1.55)
    for j, h in enumerate(["USA medallist", "Medals", "Last WTCS medal"]):
        _set_cell(tbl.cell(0, j), h, bold=True, font_size=10.5, color=WHITE,
                  bg_color="002060", align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, (_, r) in enumerate(usa.iterrows(), start=1):
        shade = "FFFFFF" if i % 2 else "EEF1F7"
        last = pd.to_datetime(r["last_medal_date"]).strftime("%b %Y") if pd.notna(r["last_medal_date"]) else "-"
        gtag = "W" if r["gender"] == "female" else "M"
        _set_cell(tbl.cell(i, 0), f"{r['athlete_full_name']} ({gtag})", font_size=10,
                  bg_color=shade, align=PP_ALIGN.LEFT)
        _set_cell(tbl.cell(i, 1), int(r["medals"]), font_size=10, bg_color=shade)
        _set_cell(tbl.cell(i, 2), last, font_size=10, bg_color=shade)


# ── Scope-dependent wording ──────────────────────────────────────────────────
def _labels(scope: str) -> dict:
    if scope == "majors":
        return {
            "title": "Medal Contention  —  A Global Rarity",
            "subtitle": "Distinct athletes who have won a WTCS or Olympic medal",
            "medal": "WTCS or Olympic",
            "chart": "Distinct medallists by nation (top 10)",
            "nations_card": "nations have ANY WTCS/Olympic medallist",
            "roster_hdr": "USA's proven medallists (WTCS + Olympic)",
            "funnel_hdr": "The global funnel (4 yrs, WTCS + Olympic)",
            "footer": ("MAJORS = World Triathlon Championship Series/Finals rounds + individual Olympic "
                       "medals (Elite Men/Women). Excludes Mixed Relay (team), T100, and regional/distance "
                       "championships. Source: World Triathlon results database."),
        }
    return {
        "title": "WTCS Medal Contention  —  A Global Rarity",
        "subtitle": "Distinct athletes who have won a World Triathlon Championship Series medal",
        "medal": "WTCS",
        "chart": "Distinct WTCS medallists by nation (top 10)",
        "nations_card": "nations have ANY WTCS medallist",
        "roster_hdr": "USA's proven WTCS medallists",
        "funnel_hdr": "The global funnel (4 yrs, WTCS)",
        "footer": ("WTCS = World Triathlon Championship Series rounds + season Championship Finals "
                   "(world-title race). Excludes T100, Olympics, and regional/distance championships. "
                   "Source: World Triathlon results database."),
    }


# ── Slides ───────────────────────────────────────────────────────────────────
def slide_rarity(prs, primary, wtcs, majors, scope, since, today):
    lab = _labels(scope)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(
        slide, lab["title"],
        f"{lab['subtitle']}  ·  {since.strftime('%b %Y')}–{today.strftime('%b %Y')} (4 yrs)")

    m = primary[primary.medals > 0]
    n_all = m.athlete_id.nunique()
    n_women = m[m.gender == "female"].athlete_id.nunique()
    n_men = m[m.gender == "male"].athlete_id.nunique()
    cm = country_medalists(m)
    n_nations = cm.shape[0]
    usa_total = int(cm.loc[HIGHLIGHT, "total"]) if HIGHLIGHT in cm.index else 0
    usa_women = int(cm.loc[HIGHLIGHT, "women"]) if HIGHLIGHT in cm.index else 0
    # "n-th most" = how many nations have strictly more, +1 (ties share a rank)
    usa_rank = int((cm["total"] > usa_total).sum()) + 1
    usa_w_rank = int((cm["women"] > usa_women).sum()) + 1
    usa_w_tied = int((cm["women"] == usa_women).sum()) > 1  # another nation shares the count
    usa_w_phrase = f"{'tied for the ' if usa_w_tied else 'the '}{_ordinal(usa_w_rank)}-most"
    n_wtcs = wtcs[wtcs.medals > 0].athlete_id.nunique()
    n_mj = majors[majors.medals > 0].athlete_id.nunique()

    # KPI cards
    cw, gap = Inches(3.02), Inches(0.18)
    x0 = Inches(0.35)
    kind = "medallists" if scope == "majors" else "WTCS medallists"
    _kpi_card(slide, x0, Inches(1.28), cw, n_all,
              f"{kind} in 4 yrs\n{n_men} men · {n_women} women", NAVY)
    _kpi_card(slide, x0 + cw + gap, Inches(1.28), cw, n_nations,
              lab["nations_card"], NAVY)
    _kpi_card(slide, x0 + 2 * (cw + gap), Inches(1.28), cw, usa_total,
              f"USA medallists — {_ordinal(usa_rank)}-most of any nation", RED)
    if scope == "majors":
        added = n_mj - n_wtcs
        _kpi_card(slide, x0 + 3 * (cw + gap), Inches(1.28), cw, n_wtcs,
                  f"also hold a WTCS medal — the Olympics adds just "
                  f"{added} new name{'s' if added != 1 else ''}", MID_GRAY)
    else:
        _kpi_card(slide, x0 + 3 * (cw + gap), Inches(1.28), cw, f"+{n_mj - n_wtcs}",
                  f"more with Olympic medals ({n_mj} total 'majors')", MID_GRAY)

    # Bar chart (left)
    _add_textbox(slide, lab["chart"],
                 Inches(0.35), Inches(2.62), Inches(5.6), Inches(0.3),
                 font_size=13, bold=True, color=NAVY)
    slide.shapes.add_picture(_bar_medalists(cm), Inches(0.3), Inches(2.95),
                             height=Inches(4.05))

    # Right column: callout + USA roster
    _add_textbox(
        slide,
        f"Only {n_women} women on the planet have won a {lab['medal']} medal in 4 years — "
        f"and the USA has {usa_women} of them, {usa_w_phrase} of any nation. "
        f"\"Medal contention\" is a club of a few dozen, not a roster you can simply assemble.",
        Inches(6.15), Inches(2.72), Inches(6.9), Inches(1.35),
        font_size=13, bold=False, color=DARK_GRAY)
    _add_textbox(slide, lab["roster_hdr"], Inches(6.15), Inches(4.15),
                 Inches(6.9), Inches(0.3), font_size=13, bold=True, color=NAVY)
    _usa_roster_table(slide, m, Inches(6.15), Inches(4.5), Inches(4.75))

    _add_textbox(slide, lab["footer"], Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.32),
                 font_size=8, italic=True, color=MID_GRAY)


def slide_depth(prs, primary, scope, since, today):
    lab = _labels(scope)
    depth = country_depth(primary)
    usa_t10 = int(depth.loc[HIGHLIGHT, "top10"]) if HIGHLIGHT in depth.index else 0
    usa_t10_rank = int((depth["top10"] > usa_t10).sum()) + 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(
        slide, "The Pipeline Behind the Podium",
        f"Podiums are scarce — but the USA has the {_ordinal(usa_t10_rank)}-deepest pool "
        f"of top-10 athletes in the world")

    g_med = int((primary.medals > 0).sum())
    g_t6 = int((primary.top6 > 0).sum())
    g_t10 = int((primary.top10 > 0).sum())

    _add_textbox(slide, "Podium medallists vs. top-10 finishers, by nation (top 10)",
                 Inches(0.35), Inches(1.3), Inches(6.4), Inches(0.3),
                 font_size=13, bold=True, color=NAVY)
    slide.shapes.add_picture(_bar_depth(depth), Inches(0.25), Inches(1.68),
                             height=Inches(4.15))

    # Funnel KPI cards (right)
    fx = Inches(7.15)
    _add_textbox(slide, lab["funnel_hdr"], fx, Inches(1.35),
                 Inches(5.9), Inches(0.3), font_size=13, bold=True, color=NAVY)
    _kpi_card(slide, fx, Inches(1.75), Inches(1.85), g_t10, "reached a\ntop-10", NAVY)
    _kpi_card(slide, fx + Inches(2.0), Inches(1.75), Inches(1.85), g_t6, "reached a\ntop-6", NAVY)
    _kpi_card(slide, fx + Inches(4.0), Inches(1.75), Inches(1.85), g_med, "won a\nmedal", RED)

    # Contrast callout
    usa_row = depth.loc[HIGHLIGHT] if HIGHLIGHT in depth.index else None
    fr_row = depth.loc["France"] if "France" in depth.index else None
    if usa_row is not None and fr_row is not None:
        box = slide.shapes.add_shape(1, fx, Inches(3.35), Inches(5.85), Inches(2.15))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY
        box.line.width = Pt(1.25)
        box.shadow.inherit = False
        _add_textbox(
            slide,
            "Depth, not just peaks",
            fx + Inches(0.2), Inches(3.5), Inches(5.5), Inches(0.35),
            font_size=14, bold=True, color=RED)
        _add_textbox(
            slide,
            f"France: {int(fr_row['medalists'])} medallists but only {int(fr_row['top10'])} top-10 "
            f"finishers — elite but top-heavy.\n\n"
            f"USA: {int(usa_row['medalists'])} medallists and {int(usa_row['top10'])} top-10 finishers "
            f"— the {_ordinal(usa_t10_rank)}-deepest pipeline in the world, with a wave of athletes "
            f"knocking on the podium door.",
            fx + Inches(0.2), Inches(3.9), Inches(5.5), Inches(1.5),
            font_size=12.5, color=DARK_GRAY)

    _add_textbox(
        slide,
        f"Top-6 / top-10 counts include the medallists. {lab['footer']}",
        Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.32),
        font_size=8, italic=True, color=MID_GRAY)


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--scope", choices=("wtcs", "majors"), default="majors",
                    help="majors (WTCS + individual Olympic medals, default) or wtcs.")
    ap.add_argument("--years", type=int, default=4)
    ap.add_argument("--since", type=str, default=None)
    ap.add_argument("--output", type=str, default=None)
    args = ap.parse_args()

    today = date.today()
    since_str = args.since or f"{today.year - args.years}-{today.month:02d}-{today.day:02d}"
    since = pd.to_datetime(since_str).date()

    engine = va.get_engine()
    wtcs = fetch_medal_data(engine, since_str, "wtcs")
    majors = fetch_medal_data(engine, since_str, "majors")
    primary = majors if args.scope == "majors" else wtcs

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_rarity(prs, primary, wtcs, majors, args.scope, since, today)
    slide_depth(prs, primary, args.scope, since, today)

    out = args.output or os.path.join(
        va.DEFAULT_OUTPUT_DIR, f"medal_contention_{args.scope}_{since}_to_{today}.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"Saved deck: {out}")


if __name__ == "__main__":
    main()
