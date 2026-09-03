#!/usr/bin/env python
"""
venue_analysis.py — Historical race analysis for a given venue.

Queries the triathlon-db PostgreSQL database and produces a PowerPoint report
with swim/bike/run split analysis, pack dynamics, and weather history.

Usage:
    python scripts/venue_analysis.py --venue Yokohama --years 8
    python scripts/venue_analysis.py --venue Yokohama --years 8 --gender men
    python scripts/venue_analysis.py --venue Yokohama --years 8 --output report.pptx
"""

import argparse
import io
import os
import re
import sys
from collections import Counter
from datetime import date, timedelta

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd
import requests
from dotenv import load_dotenv
from sqlalchemy import create_engine, text

# Allow running from repo root or scripts/ subdirectory
_repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _repo_root not in sys.path:
    sys.path.insert(0, _repo_root)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(_SCRIPT_DIR, ".."))
DEFAULT_OUTPUT_DIR = os.path.join(REPO_ROOT, "ppt files")
USAT_LOGO_PATH = os.path.join(REPO_ROOT, "docs", "power_bi_files", "USA_Triathlon_Logo.jpg")

# ── Brand colours (matching USAT reference deck) ──────────────────────────────
NAVY = RGBColor(0x00, 0x20, 0x60)       # #002060
RED = RGBColor(0xC0, 0x00, 0x00)        # #C00000
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0xBF, 0xBF, 0xBF)
DARK_GRAY = RGBColor(0x26, 0x26, 0x26)
LIGHT_NAVY = RGBColor(0x00, 0x47, 0xAB)  # accent blue for charts
GREEN = RGBColor(0x1B, 0x7F, 0x3A)      # #1B7F3A on-pace / positive verdict
SUBTITLE_BLUE = RGBColor(0xAA, 0xBB, 0xDD)  # header subtitle / banner sub-text

# Chart colour palette
C_NAVY = "#002060"
C_RED = "#C00000"
C_LIGHT_BLUE = "#4472C4"
C_LIGHT_RED = "#FF6B6B"
C_GRAY = "#808080"
C_ORANGE = "#E76F51"
C_GREEN = "#2A9D8F"
C_GOLD = "#F4A261"
C_VIOLET = "#9B5DE5"

FONT = "Arial"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Matplotlib global style
plt.rcParams.update({
    "font.family": "Arial",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.titleweight": "bold",
})


# ── DB helpers ─────────────────────────────────────────────────────────────────

def get_engine():
    load_dotenv(override=True)
    uri = os.environ.get(
        "DB_URI",
        "postgresql+psycopg://postgres:Bc020406%21@localhost:5432/triathlon_results",
    )
    return create_engine(uri)


def parse_time_to_seconds(t) -> float | None:
    """Convert 'H:MM:SS', 'MM:SS', or raw-second strings to float seconds."""
    if t is None or str(t).strip() in ("", "None", "nan"):
        return None
    t = str(t).strip()
    parts = t.split(":")
    try:
        if len(parts) == 3:
            return int(parts[0]) * 3600 + int(parts[1]) * 60 + float(parts[2])
        elif len(parts) == 2:
            return int(parts[0]) * 60 + float(parts[1])
        else:
            return float(t)
    except (ValueError, IndexError):
        return None


def adjust_outlier(series: pd.Series, threshold: float = 2.0) -> pd.Series:
    """Mark implausibly small split times as NaN so they don't poison min/mean.

    Three filters, applied in order:
      1. Non-positive values (timing-system gaps recorded as ``00:00:00``) → NaN.
      2. Values below 50% of the median split → NaN (clearly bad data — e.g.
         a 9-minute swim when the field ran 17 minutes).
      3. Iteratively, if the smallest remaining value is more than ``threshold``×
         faster than the second-smallest, drop it as a residual outlier.
    """
    s = pd.to_numeric(series, errors="coerce")
    s = s.where(s > 0, np.nan)
    valid = s.dropna()
    if len(valid) < 2:
        return s
    median_val = float(valid.median())
    if median_val > 0:
        s = s.where(s.isna() | (s >= median_val * 0.5), np.nan)
    for _ in range(5):
        valid = s.dropna()
        if len(valid) < 2:
            break
        sorted_vals = valid.sort_values()
        if sorted_vals.iloc[0] * threshold < sorted_vals.iloc[1]:
            s = s.mask(s == sorted_vals.iloc[0], np.nan)
        else:
            break
    return s


# World Triathlon appends a sport-class sub-code to Para athlete names:
#   PTWC -> "... H1" / "... H2" (handcycle sub-class)
#   PTVI -> "... B1" / "... B2" / "... B3" (vision sub-class)
# Strip it so the displayed surname is the athlete, not the code.
_PARA_SUBCLASS_RE = re.compile(r"\s+[HBT]\d\s*$", re.IGNORECASE)


def _para_surname(full_name: str | None) -> str:
    """Return a display surname for a Para athlete, stripping any trailing
    sport-class sub-code (H1/H2/B1/B2/B3)."""
    if not full_name:
        return "—"
    cleaned = _PARA_SUBCLASS_RE.sub("", str(full_name)).strip()
    parts = cleaned.split()
    return parts[-1] if parts else "—"


def seconds_to_mmss(s: float | None, always_hours: bool = False) -> str:
    if s is None:
        return "—"
    s = int(round(s))
    h = s // 3600
    m = (s % 3600) // 60
    sec = s % 60
    if h > 0 or always_hours:
        return f"{h}:{m:02d}:{sec:02d}"
    return f"{m}:{sec:02d}"


def query_venue_events(engine, venue: str, years_back: int, gender_filter: str) -> pd.DataFrame:
    cutoff = date.today() - timedelta(days=years_back * 365)
    gender_clauses = []
    if gender_filter in ("men", "both"):
        gender_clauses.append("(e.prog_name ILIKE '%Elite Men%' AND e.prog_name NOT ILIKE '%Women%')")
    if gender_filter in ("women", "both"):
        gender_clauses.append("e.prog_name ILIKE '%Elite Women%'")
    gender_sql = " OR ".join(gender_clauses)

    sql = text(f"""
        SELECT
            e.event_id, e.prog_id,
            e.event_name, e.event_date, e.event_venue,
            e.cat_name, e.prog_name, e.prog_distance_category,
            e.swim_distance, e.bike_distance, e.run_distance,
            e.swim_laps, e.bike_laps, e.run_laps,
            e.temperature_air, e.temperature_water,
            e.wind, e.wetsuit, e.weather,
            e.wind_speed_kmh, e.wind_gust_kmh,
            e.apparent_temp, e.precipitation_mm
        FROM events e
        WHERE e.event_venue ILIKE :venue_pat
          AND e.event_date >= :cutoff
          AND e.event_date <= :today
          AND ({gender_sql})
          AND (e.is_para IS NULL OR e.is_para IS NOT TRUE)
        ORDER BY e.event_date ASC
    """)
    with engine.connect() as conn:
        df = pd.read_sql(sql, conn, params={
            "venue_pat": f"%{venue}%",
            "cutoff": cutoff,
            "today": date.today(),
        })
    return df


def query_venue_para_events(engine, venue: str, years_back: int,
                             para_classes: list[str] | None = None) -> pd.DataFrame:
    """Query Para Series events at a venue.

    para_classes: optional list of class names to filter, e.g. ['PTS5', 'PTVI'].
    If None, returns all Para sport classes.
    Excludes future events (event_date <= today).
    """
    cutoff = date.today() - timedelta(days=years_back * 365)
    class_clause = ""
    if para_classes:
        # Build an OR clause: prog_name ILIKE '%PTS5%' OR prog_name ILIKE '%PTVI%' ...
        parts = " OR ".join(f"e.prog_name ILIKE '%{c}%'" for c in para_classes)
        class_clause = f"AND ({parts})"

    sql = text(f"""
        SELECT
            e.event_id, e.prog_id,
            e.event_name, e.event_date, e.event_venue,
            e.cat_name, e.prog_name, e.prog_distance_category,
            e.swim_distance, e.bike_distance, e.run_distance,
            e.swim_laps, e.bike_laps, e.run_laps,
            e.temperature_air, e.temperature_water,
            e.wind, e.wetsuit, e.weather,
            e.wind_speed_kmh, e.wind_gust_kmh,
            e.apparent_temp, e.precipitation_mm
        FROM events e
        WHERE (e.event_venue ILIKE :venue_pat OR e.event_name ILIKE :venue_pat)
          AND e.event_date >= :cutoff
          AND e.event_date <= :today
          AND e.is_para = TRUE
          {class_clause}
        ORDER BY e.event_date ASC, e.prog_name ASC
    """)
    with engine.connect() as conn:
        df = pd.read_sql(sql, conn, params={
            "venue_pat": f"%{venue}%",
            "cutoff": cutoff,
            "today": date.today(),
        })
    return df


def query_splits(engine, event_id: int, prog_id: int) -> dict:
    sql = text("""
        SELECT swimtime, biketime, runtime, total_time
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid AND finish_status = 'FINISH'
    """)
    with engine.connect() as conn:
        df = pd.read_sql(sql, conn, params={"eid": event_id, "pid": prog_id})

    result = {}
    for col, key in [("swimtime", "swim"), ("biketime", "bike"), ("runtime", "run"), ("total_time", "total")]:
        secs = df[col].apply(parse_time_to_seconds).dropna()
        secs = adjust_outlier(secs, threshold=1.5)
        secs = secs.dropna()
        result[f"{key}_fastest"] = float(secs.min()) if len(secs) else None
        result[f"{key}_avg"] = float(secs.mean()) if len(secs) else None
        result[f"{key}_n"] = len(secs)
    return result


def query_winner(engine, event_id: int, prog_id: int) -> dict:
    sql = text("""
        SELECT athlete_full_name, total_time
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid AND finish_position = 1
        LIMIT 1
    """)
    with engine.connect() as conn:
        row = conn.execute(sql, {"eid": event_id, "pid": prog_id}).fetchone()
    if row:
        return {"winner_name": row[0], "winner_time": row[1]}
    return {"winner_name": None, "winner_time": None}


def query_swim_exit_group(engine, event_id: int, prog_id: int) -> int | None:
    sql = text("""
        SELECT COUNT(*)
        FROM position_metrics pm
        JOIN race_results rr USING (event_id, prog_id, athlete_id)
        WHERE pm.event_id = :eid AND pm.prog_id = :pid
          AND rr.finish_status = 'FINISH'
          AND pm.behindswim <= 15
    """)
    with engine.connect() as conn:
        row = conn.execute(sql, {"eid": event_id, "pid": prog_id}).fetchone()
    return int(row[0]) if row and row[0] else None


def query_bike_evolution(engine, event_id: int, prog_id: int) -> dict:
    """Lead pack = athletes within 15s of leader at T1/T2.
    Chase group = athletes within 15s of the first athlete outside the lead pack.
    All data from position_metrics for consistency across all race tiers."""
    sql = text("""
        SELECT
            COUNT(CASE WHEN pm.behindt1   <= 15 THEN 1 END)               AS t1_lead_n,
            MAX(CASE WHEN pm.behindt1     <= 15 THEN pm.behindt1   END)   AS t1_lead_tail,
            COUNT(CASE WHEN pm.behindbike <= 15 THEN 1 END)               AS t2_lead_n,
            MAX(CASE WHEN pm.behindbike   <= 15 THEN pm.behindbike END)   AS t2_lead_tail,
            MIN(CASE WHEN pm.behindbike   >  15 THEN pm.behindbike END)   AS t2_chase_head
        FROM position_metrics pm
        JOIN race_results rr USING (event_id, prog_id, athlete_id)
        WHERE pm.event_id = :eid AND pm.prog_id = :pid
          AND rr.finish_status = 'FINISH'
    """)
    result: dict = {
        "t1_lead_pack": None, "t2_lead_pack": None,
        "t2_chase_gap": None, "t2_chase_pack": None,
    }

    with engine.connect() as conn:
        row = conn.execute(sql, {"eid": event_id, "pid": prog_id}).fetchone()

    if not row or row[0] is None:
        return result

    result["t1_lead_pack"] = int(row[0]) if row[0] else None
    result["t2_lead_pack"] = int(row[2]) if row[2] else None

    # Chase gap: seconds from last lead-pack athlete to first chase athlete at T2
    if row[4] is not None and row[3] is not None:
        result["t2_chase_gap"] = max(0, int(row[4]) - int(row[3]))

    # Chase group: athletes within 15s of the first chase athlete (symmetric window)
    if row[4] is not None:
        chase_head = float(row[4])
        sql2 = text("""
            SELECT COUNT(*)
            FROM position_metrics pm
            JOIN race_results rr USING (event_id, prog_id, athlete_id)
            WHERE pm.event_id = :eid AND pm.prog_id = :pid
              AND rr.finish_status = 'FINISH'
              AND pm.behindbike > 15
              AND pm.behindbike <= :chase_window
        """)
        with engine.connect() as conn:
            row2 = conn.execute(sql2, {
                "eid": event_id, "pid": prog_id,
                "chase_window": chase_head + 15,
            }).fetchone()
        result["t2_chase_pack"] = int(row2[0]) if row2 and row2[0] else None

    return result


def query_dropout_profile(engine, event_id: int, prog_id: int) -> list[dict]:
    """For each DNF/LAP athlete, determine their last valid checkpoint and
    whether they were in the lead group at that point.

    Valid = non-zero split AND not an outlier (>2× faster than 2nd fastest).
    Mirrors the logic in tri_analysis/metrics.py:adjust_outlier.
    """
    sql = text("""
        SELECT athlete_id, athlete_full_name, finish_status,
               swimtime, t1time, biketime, t2time, runtime
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid
          AND finish_status IN ('DNF', 'LAP', 'DNS')
    """)
    sql_lead = text("""
        SELECT swimtime, t1time, biketime, t2time, runtime
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid
          AND finish_status = 'FINISH'
    """)

    with engine.connect() as conn:
        dnf_df = pd.read_sql(sql, conn, params={"eid": event_id, "pid": prog_id})
        fin_df = pd.read_sql(sql_lead, conn, params={"eid": event_id, "pid": prog_id})

    if dnf_df.empty or fin_df.empty:
        return []

    def _to_secs(series: pd.Series) -> pd.Series:
        return series.apply(lambda t: parse_time_to_seconds(t) or 0.0)

    # Build cumulative elapsed for finishers to get leader time at each checkpoint
    fs = pd.DataFrame()
    fs["swim"]  = _to_secs(fin_df["swimtime"])
    fs["t1"]    = fs["swim"] + _to_secs(fin_df["t1time"])
    fs["bike"]  = fs["t1"]  + _to_secs(fin_df["biketime"])

    # Apply outlier filter per checkpoint (same threshold as metrics.py)
    for col in ["swim", "t1", "bike"]:
        fs[col] = adjust_outlier(fs[col][fs[col] > 0].reindex(fs.index, fill_value=0))

    leader_swim = float(fs["swim"][fs["swim"] > 0].min())
    leader_t1   = float(fs["t1"][fs["t1"] > 0].min())
    leader_bike = float(fs["bike"][fs["bike"] > 0].min())

    results = []
    for _, row in dnf_df.iterrows():
        swim_s  = parse_time_to_seconds(row["swimtime"])  or 0.0
        t1_s    = parse_time_to_seconds(row["t1time"])    or 0.0
        bike_s  = parse_time_to_seconds(row["biketime"])  or 0.0

        elapsed_swim = swim_s
        elapsed_t1   = swim_s + t1_s
        elapsed_bike = elapsed_t1 + bike_s

        # Determine last valid checkpoint (non-zero, non-outlier)
        # A split is suspicious if it's >2× faster than the fastest finisher split
        def _is_outlier_vs_leader(elapsed: float, leader: float) -> bool:
            return leader > 0 and elapsed > 0 and elapsed * 2 < leader

        last_cp   = None
        behind    = None

        if elapsed_swim > 0 and not _is_outlier_vs_leader(elapsed_swim, leader_swim):
            last_cp = "swim"
            behind  = round(elapsed_swim - leader_swim)
        if elapsed_t1 > 0 and t1_s > 0 and not _is_outlier_vs_leader(elapsed_t1, leader_t1):
            last_cp = "t1"
            behind  = round(elapsed_t1 - leader_t1)
        if elapsed_bike > 0 and bike_s > 0 and not _is_outlier_vs_leader(elapsed_bike, leader_bike):
            last_cp = "bike"
            behind  = round(elapsed_bike - leader_bike)

        if last_cp is None:
            continue  # no valid data

        results.append({
            "name":        row["athlete_full_name"],
            "status":      row["finish_status"],
            "last_cp":     last_cp,
            "behind":      behind,
            "in_lead_grp": behind is not None and behind <= 15,
        })

    return results


def query_field_size(engine, event_id: int, prog_id: int) -> tuple[int | None, int | None]:
    """Returns (finishers, non_finishers) where non_finishers = DNS+DNF+LAP."""
    sql = text("""
        SELECT
            COUNT(*) FILTER (WHERE finish_status = 'FINISH')     AS finishers,
            COUNT(*) FILTER (WHERE finish_status != 'FINISH')    AS non_finishers
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid
    """)
    with engine.connect() as conn:
        row = conn.execute(sql, {"eid": event_id, "pid": prog_id}).fetchone()
    if row:
        return (int(row[0]) or None, int(row[1]) or None)
    return (None, None)


def query_position_times(engine, event_id: int, prog_id: int) -> dict:
    """Return total_time in seconds for positions 1, 2, 3, 5, 10, 20."""
    positions = [1, 2, 3, 5, 10, 20]
    sql = text("""
        SELECT finish_position, total_time, athlete_full_name
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid
          AND finish_position = ANY(:positions)
          AND finish_status = 'FINISH'
        ORDER BY finish_position
    """)
    with engine.connect() as conn:
        df = pd.read_sql(sql, conn, params={"eid": event_id, "pid": prog_id, "positions": positions})

    result = {}
    winner_secs = None
    for _, row in df.iterrows():
        pos = int(row.finish_position)
        secs = parse_time_to_seconds(row.total_time)
        result[pos] = {"secs": secs, "name": row.athlete_full_name, "time_str": row.total_time}
        if pos == 1:
            winner_secs = secs

    # Compute gaps from winner
    for pos in positions:
        if pos in result and winner_secs and result[pos]["secs"]:
            result[pos]["gap"] = result[pos]["secs"] - winner_secs
        elif pos in result:
            result[pos]["gap"] = None

    return result


def collect_race_data(engine, events_df: pd.DataFrame) -> list[dict]:
    rows = []
    for _, ev in events_df.iterrows():
        eid, pid = int(ev.event_id), int(ev.prog_id)
        splits = query_splits(engine, eid, pid)
        if splits.get("total_n", 0) == 0:
            print(f"  Skipping {ev.event_date} {ev.prog_name} — no results (cancelled or unavailable)")
            continue
        winner = query_winner(engine, eid, pid)
        swim_group = query_swim_exit_group(engine, eid, pid)
        bike = query_bike_evolution(engine, eid, pid)
        dropouts = query_dropout_profile(engine, eid, pid)
        field_size, non_finishers = query_field_size(engine, eid, pid)
        field_size = splits.get("total_n") or field_size or None
        pos_times = query_position_times(engine, eid, pid)

        rows.append({
            "event_id": eid,
            "prog_id": pid,
            "year": pd.to_datetime(ev.event_date).year,
            "date": ev.event_date,
            "event_name": ev.event_name,
            "cat_name": ev.cat_name or "",
            "prog_name": ev.prog_name or "",
            "prog_distance_category": ev.prog_distance_category or "",
            "swim_km": ev.swim_distance,
            "bike_km": ev.bike_distance,
            "run_km": ev.run_distance,
            "bike_laps": ev.bike_laps,
            "run_laps": ev.run_laps,
            "field_size": field_size,
            "non_finishers": non_finishers,
            "swim_exit_group": swim_group,
            "t1_lead_pack":   bike["t1_lead_pack"],
            "t2_lead_pack":   bike["t2_lead_pack"],
            "t2_chase_gap":   bike["t2_chase_gap"],
            "t2_chase_pack":  bike["t2_chase_pack"],
            "dropouts":       dropouts,
            "temp_air": ev.temperature_air,
            "temp_water": ev.temperature_water,
            "wind_raw": ev.wind,
            "wind_kmh": ev.wind_speed_kmh,
            "wetsuit": ev.wetsuit,
            "weather": ev.weather,
            # Open-Meteo enriched fields (filled later by enrich_rows_with_openmeteo)
            "humidity": None,
            "precip": None,
            "pm25": None,
            "aqi": None,
            "uv_index": None,
            "winner_name": winner["winner_name"],
            "winner_time": winner["winner_time"],
            "pos_times": pos_times,
            **splits,
        })
    return rows


# ── Deep-dive helpers (detailed splits, season norms, startlist) ──────────────

def load_detailed_splits(xlsx_path: str) -> dict[str, pd.DataFrame]:
    """Load detailed per-lap splits Excel for a single race.

    Returns {'men': df, 'women': df} with columns:
        Rank, Bib, Name, Nat,
        swim_sec, bike_sec, run_sec, total_sec,
        end_swim_elapsed, end_bike_elapsed, end_run_elapsed,
        dnf (bool — True if Total contained DNF/LAP)

    Handles both Male/female sheet naming variants and both
    T1.1/T2.1 (women) and TA1/TA2 (men) discipline-column variants.
    """
    xl = pd.ExcelFile(xlsx_path)
    sheets_lc = {n.lower(): n for n in xl.sheet_names}
    out: dict[str, pd.DataFrame] = {}
    for gender, candidates in [
        ("men",   ["male", "men", "m"]),
        ("women", ["female", "women", "f"]),
    ]:
        sheet = next((sheets_lc[c] for c in candidates if c in sheets_lc), None)
        if not sheet:
            continue
        df = pd.read_excel(xl, sheet_name=sheet)
        cols = list(df.columns)

        t1_col = "T1.1" if "T1.1" in cols else ("TA1" if "TA1" in cols else None)
        t2_col = "T2.1" if "T2.1" in cols else ("TA2" if "TA2" in cols else None)

        for col, target in [("Swim", "swim_sec"), ("Bike", "bike_sec"), ("Run", "run_sec")]:
            df[target] = df[col].apply(parse_time_to_seconds) if col in cols else None

        df["t1_sec"] = df[t1_col].apply(parse_time_to_seconds) if t1_col else 0
        df["t2_sec"] = df[t2_col].apply(parse_time_to_seconds) if t2_col else 0

        # Non-finishers are marked DNF/DNS/LAP in the Rank column (Total is 'n.a.')
        rank_str = df["Rank"].astype(str)
        total_str = df["Total"].astype(str)
        df["dnf"] = (
            rank_str.str.contains(r"DNF|DNS|LAP", case=False, na=False)
            | total_str.str.contains(r"DNF|DNS|LAP", case=False, na=False)
        )
        df["total_sec"] = df["Total"].apply(parse_time_to_seconds)

        df["end_swim_elapsed"] = df["swim_sec"]
        bike_cumulative = (
            df["swim_sec"].fillna(0)
            + df["t1_sec"].fillna(0)
            + df["bike_sec"].fillna(0)
        )
        df["end_bike_elapsed"] = bike_cumulative.where(
            df["swim_sec"].notna() & df["bike_sec"].notna(),
            pd.NA,
        )
        df["end_run_elapsed"] = df["total_sec"]

        keep = ["Rank", "Bib", "Name", "Nat",
                "swim_sec", "bike_sec", "run_sec", "total_sec",
                "end_swim_elapsed", "end_bike_elapsed", "end_run_elapsed",
                "dnf"]
        out[gender] = df[[c for c in keep if c in df.columns]].copy()
    return out


def compute_pack_scatter(df: pd.DataFrame) -> dict[str, dict]:
    """For each leg, compute scatter data: gap-to-leader vs placement-at-checkpoint.

    Returns dict keyed by leg label ('Swim', 'Bike', 'Run') with:
        {"points": [{"name": ..., "gap": ..., "placement": ..., "in_lead_swim": bool}],
         "leader_time": float | None}

    Swim-exit lead pack = athletes within 15s of the swim leader.
    Run leg includes finishers only (DNF/LAP excluded).
    """
    out: dict[str, dict] = {}

    swim_valid = df[df["end_swim_elapsed"].notna()].copy()
    if swim_valid.empty:
        lead_swim_names: set = set()
    else:
        swim_leader = float(swim_valid["end_swim_elapsed"].min())
        lead_swim_names = set(
            swim_valid.loc[swim_valid["end_swim_elapsed"] - swim_leader <= 15, "Name"]
        )

    for leg_col, label, finishers_only in [
        ("end_swim_elapsed", "Swim", False),
        ("end_bike_elapsed", "Bike", False),
        ("end_run_elapsed",  "Run",  True),
    ]:
        sub = df[df[leg_col].notna()].copy()
        if finishers_only:
            sub = sub[~sub["dnf"]]
        if sub.empty:
            out[label] = {"points": [], "leader_time": None}
            continue
        sub = sub.sort_values(leg_col).reset_index(drop=True)
        leader_time = float(sub[leg_col].iloc[0])
        sub["gap"] = sub[leg_col].astype(float) - leader_time
        sub["placement"] = sub.index + 1
        out[label] = {
            "points": [
                {
                    "name": row["Name"],
                    "gap": float(row["gap"]),
                    "placement": int(row["placement"]),
                    "in_lead_swim": row["Name"] in lead_swim_names,
                }
                for _, row in sub.iterrows()
            ],
            "leader_time": leader_time,
        }
    return out


def map_excel_names_to_athlete_ids(engine, event_id: int, prog_id: int, names: list[str]) -> dict[str, int]:
    """Match Excel athlete names to athlete_ids via race_results for the same historical race."""
    sql = text("""
        SELECT athlete_id, athlete_full_name
        FROM race_results
        WHERE event_id = :eid AND prog_id = :pid AND athlete_id IS NOT NULL
    """)
    df = pd.read_sql(sql, engine, params={"eid": event_id, "pid": prog_id})
    lookup = {str(n).lower().strip(): int(aid) for aid, n in zip(df.athlete_id, df.athlete_full_name)}
    return {n: lookup[str(n).lower().strip()] for n in names if str(n).lower().strip() in lookup}


def query_wtcs_season_splits(engine, athlete_ids: list[int], venue_date,
                             gender: str) -> pd.DataFrame:
    """WTCS Standard-distance splits for the given athletes in the 12 months
    ending at venue_date. Excludes DNF/DNS/LAP/DSQ rows and the venue race itself.

    Returns columns: athlete_id, event_date, event_name, swim_sec, bike_sec, run_sec.
    """
    if not athlete_ids:
        return pd.DataFrame(columns=["athlete_id", "event_date", "event_name",
                                     "swim_sec", "bike_sec", "run_sec"])
    gender_clause = (
        "(e.prog_name ILIKE '%Elite Men%' AND e.prog_name NOT ILIKE '%Women%')"
        if gender.lower().startswith("m")
        else "e.prog_name ILIKE '%Elite Women%'"
    )
    sql = text(f"""
        SELECT rr.athlete_id, e.event_date, e.event_name, e.event_venue,
               rr.swimtime, rr.biketime, rr.runtime
        FROM race_results rr
        JOIN events e USING (event_id, prog_id)
        WHERE rr.athlete_id = ANY(:athlete_ids)
          AND e.cat_name ILIKE '%Championship Series%'
          AND (e.prog_distance_category IS NULL OR e.prog_distance_category ILIKE 'standard')
          AND e.event_date BETWEEN (:venue_date - INTERVAL '12 months') AND :venue_date
          AND {gender_clause}
          AND (rr.finish_status IS NULL OR rr.finish_status NOT IN ('DNF', 'DNS', 'DSQ', 'LAP'))
    """)
    df = pd.read_sql(sql, engine, params={"athlete_ids": list(athlete_ids),
                                          "venue_date": venue_date})
    for src, dst in [("swimtime", "swim_sec"), ("biketime", "bike_sec"), ("runtime", "run_sec")]:
        df[dst] = adjust_outlier(df[src].apply(parse_time_to_seconds))
    return df


def find_upcoming_event(engine, venue: str) -> dict | None:
    """Find the upcoming/most-recent-future event_id at this venue (men's prog as anchor).
    Includes cat_name so downstream tier detection uses the UPCOMING event's tier
    (e.g. 'World Cup') rather than falling back to an old historical event's tier."""
    sql = text("""
        SELECT event_id,
               MIN(event_date) AS event_date,
               MIN(event_name) AS event_name,
               MIN(cat_name)   AS cat_name
        FROM events
        WHERE (event_venue ILIKE :v OR event_name ILIKE :v)
          AND event_date >= CURRENT_DATE
        GROUP BY event_id
        ORDER BY event_date ASC
        LIMIT 1
    """)
    df = pd.read_sql(sql, engine, params={"v": f"%{venue}%"})
    if df.empty:
        return None
    return df.iloc[0].to_dict()


def query_top_by_world_ranking(engine, gender: str, on_startlist_event_id: int | None,
                               limit: int = 3) -> pd.DataFrame:
    """Top N athletes by current world ranking. If on_startlist_event_id is given,
    restrict to athletes on that event's startlist."""
    cat_id = 13 if gender.lower().startswith("m") else 14
    if on_startlist_event_id:
        sql = text("""
            WITH latest AS (
              SELECT MAX(retrieved_at) AS dt FROM athlete_rankings WHERE ranking_cat_id = :cat
            )
            SELECT DISTINCT ar.athlete_id, ar.athlete_name,
                   ar.rank_position, ar.total_points
            FROM athlete_rankings ar
            JOIN latest l ON ar.retrieved_at = l.dt
            JOIN program_entries pe
              ON pe.athlete_id = ar.athlete_id
             AND pe.event_id = :eid
             AND pe.is_active = TRUE
             AND pe.entry_type = 'start'
            WHERE ar.ranking_cat_id = :cat
            ORDER BY ar.rank_position ASC
            LIMIT :lim
        """)
        return pd.read_sql(sql, engine, params={"cat": cat_id, "eid": on_startlist_event_id,
                                                "lim": limit})
    sql = text("""
        SELECT ar.athlete_id, ar.athlete_name, ar.rank_position, ar.total_points
        FROM athlete_rankings ar
        WHERE ar.ranking_cat_id = :cat
          AND ar.retrieved_at = (SELECT MAX(retrieved_at) FROM athlete_rankings WHERE ranking_cat_id = :cat)
        ORDER BY ar.rank_position ASC
        LIMIT :lim
    """)
    return pd.read_sql(sql, engine, params={"cat": cat_id, "lim": limit})


def query_prior_podium_on_startlist(engine, prior_event_id: int, prior_prog_id: int,
                                    upcoming_event_id: int) -> pd.DataFrame:
    """Athletes on the upcoming startlist who finished top-3 at the prior race."""
    sql = text("""
        SELECT rr.athlete_id, rr.athlete_full_name,
               rr.position_sort AS position, rr.total_time
        FROM race_results rr
        JOIN program_entries pe
          ON pe.athlete_id = rr.athlete_id
         AND pe.event_id = :upcoming
         AND pe.is_active = TRUE
         AND pe.entry_type = 'start'
        WHERE rr.event_id = :prior_eid AND rr.prog_id = :prior_pid
          AND rr.position_sort IS NOT NULL
          AND rr.position_sort <= 3
        ORDER BY rr.position_sort
    """)
    return pd.read_sql(sql, engine, params={"upcoming": upcoming_event_id,
                                            "prior_eid": prior_event_id,
                                            "prior_pid": prior_prog_id})


# ── Open-Meteo / geocoding helpers ────────────────────────────────────────────

_GEOCODE_CACHE: dict = {}
_WEATHER_CACHE: dict = {}
_AQI_CACHE:     dict = {}

# Weather sampling window (start_hour, end_hour) inclusive, in local time.
# Defaults to a morning window but is widened to cover the actual race start
# times via set_weather_window() so afternoon/evening races (e.g. WTCS London
# at 14:30 / 16:15) are captured — not just morning races.
_WEATHER_WINDOW: tuple[int, int] = (8, 12)


def set_weather_window(race_starts) -> tuple[int, int]:
    """Set the global weather sampling window from a venue's race start times.
    Window spans the earliest start to ~2 h after the latest start (to cover the
    race duration). Falls back to (8, 12) when no start times are available."""
    global _WEATHER_WINDOW
    hours = []
    for _lbl, hhmm in (race_starts or []):
        try:
            hours.append(int(str(hhmm)[:2]))
        except (ValueError, IndexError):
            continue
    if hours:
        lo = max(0, min(hours))
        hi = min(23, max(hours) + 2)
        _WEATHER_WINDOW = (lo, hi)
    else:
        _WEATHER_WINDOW = (8, 12)
    return _WEATHER_WINDOW


def _in_weather_window(hh: str) -> bool:
    """True if a 2-char hour string falls in the current weather window."""
    try:
        h = int(hh)
    except (ValueError, TypeError):
        return False
    return _WEATHER_WINDOW[0] <= h <= _WEATHER_WINDOW[1]


# Hardcoded fallback for venues that have been geocoded historically. Avoids
# relying on Nominatim when local SSL certs / network are unavailable.
VENUE_COORDS_FALLBACK: dict[str, tuple[float, float]] = {
    "karlovy vary": (50.226, 12.855),  # Lake Rolava, Karlovy Vary
    "huatulco":  (15.831, -96.320),
    "alghero":   (40.564,   8.319),
    "yokohama":  (35.444, 139.638),
    "abu dhabi": (24.466,  54.367),
    "quiberon":  (47.485,  -3.114),
    "montreal":  (45.503, -73.534),   # Parc Jean-Drapeau, Notre-Dame Island
    "hamburg":   (53.554,   9.994),   # Binnenalster / Rathausmarkt, city center
    "edmonton":  (53.527, -113.547),  # William Hawrelak Park, 9330 Groat Rd NW (race venue)
    "asuncion":  (-25.252, -57.642),  # La Isla del Delta, Nueva Asunción (Presidente Hayes, across the river)
    "weihai":    (37.5043, 122.1290), # Wave to Wonderland swim venue (from athlete-guide water-quality sample coords)
    "london":    (51.508,   0.029),   # ExCeL London, Royal Victoria Dock
    "rio de janeiro": (-22.971, -43.185),  # Copacabana Beach, Avenida Atlântica
    "rio":       (-22.971, -43.185),  # Copacabana Beach (alias)
}


def geocode_venue(venue: str) -> tuple | None:
    """Geocode a venue string via Nominatim OSM. Returns (lat, lon) or None.
    Falls back to VENUE_COORDS_FALLBACK on network/SSL failure."""
    if venue in _GEOCODE_CACHE:
        return _GEOCODE_CACHE[venue]
    try:
        r = requests.get(
            "https://nominatim.openstreetmap.org/search",
            params={"q": venue, "format": "json", "limit": 1},
            headers={"User-Agent": "triathlon-venue-analysis/1.0"},
            timeout=10,
        )
        if r.ok and r.json():
            d = r.json()[0]
            result = (float(d["lat"]), float(d["lon"]))
            print(f"  [Geocode] '{venue}' -> {result[0]:.3f}, {result[1]:.3f}")
            _GEOCODE_CACHE[venue] = result
            return result
    except Exception as exc:
        print(f"  [Geocode] Nominatim failed for '{venue}': {exc}")
    # Fallback to hardcoded coords for known venues
    fb = VENUE_COORDS_FALLBACK.get(venue.lower().strip())
    if fb:
        print(f"  [Geocode] Using fallback coords for '{venue}' -> {fb[0]:.3f}, {fb[1]:.3f}")
        _GEOCODE_CACHE[venue] = fb
        return fb
    _GEOCODE_CACHE[venue] = None
    return None


def _ometa_race_avg(hourly: dict, times: list, col: str) -> float | None:
    """Average an Open-Meteo hourly field over the current race window
    (set via set_weather_window() — covers the actual race start times)."""
    idxs = [i for i, t in enumerate(times) if len(t) >= 13 and _in_weather_window(t[11:13])]
    if not idxs:
        idxs = list(range(_WEATHER_WINDOW[0], _WEATHER_WINDOW[1] + 1))
    vals = [hourly[col][i] for i in idxs
            if i < len(hourly.get(col, [])) and hourly[col][i] is not None]
    return round(sum(vals) / len(vals), 1) if vals else None


def fetch_openmeteo_weather(lat: float, lon: float, race_date: str) -> dict:
    """Fetch race-window weather from Open-Meteo archive (hours 08–11 local)."""
    key = (round(lat, 2), round(lon, 2), race_date)
    if key in _WEATHER_CACHE:
        return _WEATHER_CACHE[key]
    result: dict = {}
    try:
        r = requests.get(
            "https://archive-api.open-meteo.com/v1/archive",
            params={
                "latitude": lat, "longitude": lon,
                "start_date": race_date, "end_date": race_date,
                "hourly": "temperature_2m,relative_humidity_2m,wind_speed_10m,precipitation",
                "wind_speed_unit": "kmh",
                "timezone": "auto",
            },
            timeout=15,
        )
        if r.ok:
            hrs = r.json().get("hourly", {})
            times = hrs.get("time", [])
            result = {
                "temp_air_om":  _ometa_race_avg(hrs, times, "temperature_2m"),
                "humidity_om":  _ometa_race_avg(hrs, times, "relative_humidity_2m"),
                "wind_kmh_om":  _ometa_race_avg(hrs, times, "wind_speed_10m"),
                "precip_om":    _ometa_race_avg(hrs, times, "precipitation"),
            }
    except Exception as exc:
        print(f"  [OpenMeteo weather] {race_date}: {exc}")
    _WEATHER_CACHE[key] = result
    return result


def fetch_openmeteo_airquality(lat: float, lon: float, race_date: str) -> dict:
    """Fetch race-window air quality from Open-Meteo AQ archive."""
    key = (round(lat, 2), round(lon, 2), race_date)
    if key in _AQI_CACHE:
        return _AQI_CACHE[key]
    result: dict = {}
    try:
        r = requests.get(
            "https://air-quality-api.open-meteo.com/v1/air-quality",
            params={
                "latitude": lat, "longitude": lon,
                "start_date": race_date, "end_date": race_date,
                "hourly": "pm2_5,european_aqi,uv_index",
                "timezone": "auto",
            },
            timeout=15,
        )
        if r.ok:
            hrs = r.json().get("hourly", {})
            times = hrs.get("time", [])
            result = {
                "pm25": _ometa_race_avg(hrs, times, "pm2_5"),
                "aqi":  _ometa_race_avg(hrs, times, "european_aqi"),
                "uv":   _ometa_race_avg(hrs, times, "uv_index"),
            }
    except Exception as exc:
        print(f"  [OpenMeteo AQI] {race_date}: {exc}")
    _AQI_CACHE[key] = result
    return result


_MARINE_CACHE: dict = {}


def fetch_openmeteo_marine(lat: float, lon: float, race_date: str) -> dict:
    """Fetch sea surface temperature from Open-Meteo Marine API (global coverage)."""
    key = (round(lat, 2), round(lon, 2), race_date)
    if key in _MARINE_CACHE:
        return _MARINE_CACHE[key]
    result: dict = {}
    try:
        r = requests.get(
            "https://marine-api.open-meteo.com/v1/marine",
            params={
                "latitude": lat, "longitude": lon,
                "start_date": race_date, "end_date": race_date,
                "hourly": "sea_surface_temperature",
                "timezone": "auto",
            },
            timeout=15,
        )
        if r.ok:
            hrs = r.json().get("hourly", {})
            times = hrs.get("time", [])
            sst = _ometa_race_avg(hrs, times, "sea_surface_temperature")
            if sst is not None:
                result["sst"] = sst
    except Exception as exc:
        print(f"  [OpenMeteo Marine] {race_date}: {exc}")
    _MARINE_CACHE[key] = result
    return result


def fetch_sst_for_race_day(lat: float, lon: float, race_date: str,
                           years_back: int = 7) -> tuple[float | None, str]:
    """Best-effort sea-surface temperature for race_date.

    Tries the marine API for the actual race_date first (forecast window).
    Falls back to averaging the same calendar day across prior years.
    Returns (sst_celsius, source) where source is 'forecast' / 'climatology' / 'unavailable'.
    """
    direct = fetch_openmeteo_marine(lat, lon, race_date)
    if direct.get("sst") is not None:
        return float(direct["sst"]), "forecast"

    target = pd.to_datetime(race_date).date()
    current_year = date.today().year
    vals: list[float] = []
    for yr_off in range(1, years_back + 1):
        yr = current_year - yr_off
        d = date(yr, target.month, target.day).isoformat()
        result = fetch_openmeteo_marine(lat, lon, d)
        if result.get("sst") is not None:
            vals.append(float(result["sst"]))
    if vals:
        return sum(vals) / len(vals), "climatology"
    return None, "unavailable"


# ── Forward forecast + climatology (used by Race-Day Forecast slide) ─────────

def fetch_openmeteo_forecast(lat: float, lon: float, race_date: str) -> dict:
    """Forward forecast for race_date. Returns hourly arrays for 04:00–12:00 local,
    or empty dict if race_date is outside the 16-day forecast window.

    Result keys: 'time', 'temperature_2m', 'apparent_temperature',
    'relative_humidity_2m', 'wind_speed_10m', 'precipitation_probability', 'uv_index'.
    """
    try:
        r = requests.get(
            "https://api.open-meteo.com/v1/forecast",
            params={
                "latitude": lat, "longitude": lon,
                "start_date": race_date, "end_date": race_date,
                "hourly": ",".join([
                    "temperature_2m", "apparent_temperature",
                    "relative_humidity_2m", "wind_speed_10m",
                    "precipitation_probability", "uv_index",
                ]),
                "wind_speed_unit": "kmh",
                "timezone": "auto",
                # NB: Open-Meteo rejects `forecast_days` alongside start_date/end_date
                # with a 400, which silently forced every deck onto the climatology
                # fallback. The explicit date range already bounds the request.
            },
            timeout=15,
        )
        if not r.ok:
            return {}
        hrs = r.json().get("hourly", {})
        times = hrs.get("time", [])
        # Confirm race_date is actually present in the returned series
        if not any(t.startswith(race_date) for t in times):
            return {}
        # Chart window: one hour of lead-in before the earliest start, through
        # the end of the race window (covers afternoon/evening starts too).
        lo = max(0, _WEATHER_WINDOW[0] - 1)
        hi = _WEATHER_WINDOW[1]
        keep_idx = [i for i, t in enumerate(times)
                    if t.startswith(race_date) and lo <= int(t[11:13]) <= hi]
        if not keep_idx:
            return {}
        out: dict = {"time": [times[i][11:16] for i in keep_idx]}
        for col in ["temperature_2m", "apparent_temperature",
                    "relative_humidity_2m", "wind_speed_10m",
                    "precipitation_probability", "uv_index"]:
            vals = hrs.get(col, [])
            out[col] = [vals[i] if i < len(vals) else None for i in keep_idx]
        return out
    except Exception as exc:
        print(f"  [OpenMeteo forecast] {race_date}: {exc}")
        return {}


def fetch_openmeteo_climatology(lat: float, lon: float, race_date: str,
                                years_back: int = 7) -> dict:
    """Average each hour 04:00–12:00 across the same calendar day for the last
    `years_back` years. Used as a fallback when the race date is outside the
    forward-forecast window."""
    target = pd.to_datetime(race_date).date()
    current_year = date.today().year
    rows_per_hour: dict[str, dict[str, list]] = {}
    fields = ["temperature_2m", "apparent_temperature",
              "relative_humidity_2m", "wind_speed_10m",
              "precipitation", "uv_index"]
    for yr_off in range(1, years_back + 1):
        yr = current_year - yr_off
        d = date(yr, target.month, target.day).isoformat()
        try:
            r = requests.get(
                "https://archive-api.open-meteo.com/v1/archive",
                params={
                    "latitude": lat, "longitude": lon,
                    "start_date": d, "end_date": d,
                    "hourly": ",".join(fields),
                    "wind_speed_unit": "kmh",
                    "timezone": "auto",
                },
                timeout=15,
            )
            if not r.ok:
                continue
            hrs = r.json().get("hourly", {})
            times = hrs.get("time", [])
            lo = max(0, _WEATHER_WINDOW[0] - 1)
            hi = _WEATHER_WINDOW[1]
            for i, t in enumerate(times):
                hour = t[11:16]
                if not (lo <= int(t[11:13]) <= hi):
                    continue
                bucket = rows_per_hour.setdefault(hour, {f: [] for f in fields})
                for f in fields:
                    arr = hrs.get(f, [])
                    if i < len(arr) and arr[i] is not None:
                        bucket[f].append(arr[i])
        except Exception as exc:
            print(f"  [OpenMeteo climatology {d}]: {exc}")
            continue
    if not rows_per_hour:
        return {}
    times_sorted = sorted(rows_per_hour.keys())
    out: dict = {"time": times_sorted}
    out_field_map = {
        "temperature_2m": "temperature_2m",
        "apparent_temperature": "apparent_temperature",
        "relative_humidity_2m": "relative_humidity_2m",
        "wind_speed_10m": "wind_speed_10m",
        "precipitation": "precipitation",
        "uv_index": "uv_index",
    }
    for src, dst in out_field_map.items():
        out[dst] = [
            (sum(rows_per_hour[h][src]) / len(rows_per_hour[h][src]))
            if rows_per_hour[h][src] else None
            for h in times_sorted
        ]
    # Climatology has no precip-probability — synthesise from mean precipitation > 0.1mm
    out["precipitation_probability"] = [
        100.0 if (p is not None and p > 0.1) else 0.0 for p in out["precipitation"]
    ]
    return out


def enrich_rows_with_openmeteo(rows: list[dict], coords: tuple | None) -> None:
    """Fill weather + AQ + water temp fields in-place using Open-Meteo where WT API left NULLs."""
    if coords is None:
        return
    lat, lon = coords
    for row in rows:
        race_date = str(pd.to_datetime(row["date"]).date())
        wx     = fetch_openmeteo_weather(lat, lon, race_date)
        aq     = fetch_openmeteo_airquality(lat, lon, race_date)
        marine = fetch_openmeteo_marine(lat, lon, race_date)
        # Fill temp/wind only if WT API returned nothing (mark with * so slide notes source)
        if not row.get("temp_air") and wx.get("temp_air_om") is not None:
            row["temp_air"] = f"{wx['temp_air_om']}°*"
        if not row.get("wind_kmh") and wx.get("wind_kmh_om") is not None:
            row["wind_kmh"] = wx["wind_kmh_om"]
        # Water temp from Marine API where WT API is absent
        if not row.get("temp_water") and marine.get("sst") is not None:
            row["temp_water"] = f"{marine['sst']:.1f}°*"
        # Humidity and precip always from Open-Meteo
        row["humidity"] = wx.get("humidity_om")
        row["precip"]   = wx.get("precip_om")
        # Air quality always from Open-Meteo
        row["pm25"]     = aq.get("pm25")
        row["aqi"]      = aq.get("aqi")
        row["uv_index"] = aq.get("uv")


# ── Venue content helpers ──────────────────────────────────────────────────────

def _build_wt_event_url(event_name: str) -> str:
    """Convert an event name to a triathlon.org event URL slug."""
    slug = event_name.lower()
    slug = re.sub(r"[^\w\s]", "", slug)
    slug = re.sub(r"\s+", "_", slug.strip())
    return f"https://www.triathlon.org/events/event/{slug}"


def fetch_venue_content(venue: str, events_df: pd.DataFrame) -> dict:
    """Build WT event URLs from DB data; try the API for the upcoming event date."""
    content: dict = {
        "wt_event_urls": {},   # {year: url}
        "last_year_url": None,
        "upcoming_event_url": None,
    }

    # Build WT event URLs from DB event names (deduped by event_id)
    for _, row in events_df.drop_duplicates(subset=["event_id"]).iterrows():
        yr = pd.to_datetime(row.event_date).year
        content["wt_event_urls"][yr] = _build_wt_event_url(str(row.event_name))

    if content["wt_event_urls"]:
        last_yr = max(content["wt_event_urls"])
        content["last_year_url"] = content["wt_event_urls"][last_yr]

    this_year = date.today().year
    content["upcoming_event_url"] = (
        content["wt_event_urls"].get(this_year)
        or content["wt_event_urls"].get(this_year + 1)
    )

    # Try the WT API for the upcoming event slug (more reliable than constructed URL)
    load_dotenv(override=True)
    api_key = os.getenv("TRI_API_KEY")
    if api_key:
        try:
            r = requests.get(
                "https://api.triathlon.org/v1/events",
                headers={"apikey": api_key},
                params={"q": venue, "limit": 20},
                timeout=10,
            )
            if r.ok:
                body = r.json()
                # API may return {"data": {"events": [...]}} or {"data": [...]}
                data_block = body.get("data", {})
                if isinstance(data_block, list):
                    event_list = data_block
                elif isinstance(data_block, dict):
                    event_list = data_block.get("events", [])
                else:
                    event_list = []
                for ev in event_list:
                    ev_date = str(ev.get("event_date", ""))
                    ev_name = str(ev.get("event_name", ""))
                    if ev_date >= str(date.today()) and "championship series" in ev_name.lower():
                        slug = ev.get("event_slug") or _build_wt_event_url(ev_name).split("/")[-1]
                        content["upcoming_event_url"] = (
                            f"https://www.triathlon.org/events/event/{slug}"
                        )
                        break
        except Exception as exc:
            print(f"  [Info] WT API fetch failed: {exc}")

    return content


def _derive_keys_to_success(rows: list[dict]) -> list[str]:
    """Analyze race data patterns and return 3 key success factors."""
    keys = []

    # 1. Swim exit importance
    swim_groups = [r["swim_exit_group"] for r in rows if r.get("swim_exit_group")]
    if swim_groups:
        avg = sum(swim_groups) / len(swim_groups)
        if avg >= 15:
            keys.append(
                f"Swim positioning matters — an average of {int(avg)} athletes exit the water "
                "within 15 sec of the leader. Entering the bike in the front pack is critical; "
                "conserve energy with a disciplined swim line rather than sprinting alone."
            )
        else:
            keys.append(
                f"Fast swimming is decisive — only ~{int(avg)} athletes typically leave the water "
                "together. A strong swim directly builds a race-winning gap before T1."
            )
    else:
        keys.append(
            f"[PLACEHOLDER: Key 1 — Describe swim strategy and how the swim exit shapes the race at {rows[0]['event_name'] if rows else 'this venue'}, "
            "e.g. whether a large group exits together or swim speed is decisive.]"
        )

    # 2. Bike / breakaway tendency
    packs = [r["lead_bike_pack"] for r in rows if r.get("lead_bike_pack")]
    if packs:
        avg = sum(packs) / len(packs)
        small_count = sum(1 for p in packs if p < 10)
        if small_count >= len(packs) / 2:
            keys.append(
                f"Breakaway bike legs are common — {small_count} of {len(packs)} races had fewer than 10 "
                "riders in the lead pack. Strong cyclists can create decisive gaps; athletes need "
                "the engine to push a high pace from the gun."
            )
        elif avg >= 20:
            keys.append(
                f"Group bike racing dominates (avg lead pack: {int(avg)} athletes). Surviving the bike "
                "in the front group is non-negotiable — the race is ultimately decided on the run."
            )
        else:
            keys.append(
                f"Mixed bike dynamics — the lead group averages {int(avg)} riders. Bike endurance keeps "
                "athletes in contention, but the run determines who stands on the podium."
            )
    else:
        keys.append(
            "[PLACEHOLDER: Key 2 — Describe bike strategy: is this a pack race or do breakaways form? "
            "What power/tactics are needed to stay in the lead group?]"
        )

    # 3. Run
    run_fastest = [r["run_fastest"] for r in rows if r.get("run_fastest")]
    run_avg = [r["run_avg"] for r in rows if r.get("run_avg")]
    if run_fastest and run_avg:
        avg_fast = sum(run_fastest) / len(run_fastest)
        avg_mean = sum(run_avg) / len(run_avg)
        pace_km = avg_fast / 10  # assumes 10 km run
        spread = int(avg_mean - avg_fast)
        keys.append(
            f"Run speed separates the podium — the fastest run pace averages ~{seconds_to_mmss(pace_km)}/km "
            f"with a ~{spread}s spread between fastest and average finisher. "
            "Either a pre-built lead from the bike or exceptional run speed is required to win."
        )
    else:
        keys.append(
            "[PLACEHOLDER: Key 3 — Describe run dynamics: is this a sprint finish or does the winner "
            "typically build a lead? What run pace is required to podium?]"
        )

    return keys


# ── PPTX helpers ───────────────────────────────────────────────────────────────

def _set_cell(cell, text_val: str, bold=False, font_size=10, color=None,
              align=PP_ALIGN.CENTER, bg_color=None, italic=False):
    """Set text, font, and background of a table cell."""
    tf = cell.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = align
    # Clear any existing runs
    for run in para.runs:
        run.text = ""
    run = para.runs[0] if para.runs else para.add_run()
    run.text = str(text_val)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    run.font.color.rgb = color or DARK_GRAY

    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # Remove existing fills
        for existing in tcPr.findall(qn("a:solidFill")):
            tcPr.remove(existing)
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", str(bg_color))


def _add_textbox(slide, text: str, left, top, width, height,
                 font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
                 italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    run.font.color.rgb = color or DARK_GRAY
    return txb


def add_slide_chrome(slide, title: str, subtitle: str = "", show_logo: bool = True):
    """Add the standard chrome: navy header bar, red accent line, USAT logo."""
    # Navy header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title text in header
    _add_textbox(slide, title, Inches(0.25), Inches(0.05), Inches(10.5), Inches(0.65),
                 font_size=28, bold=True, color=WHITE)

    if subtitle:
        _add_textbox(slide, subtitle, Inches(0.25), Inches(0.68), Inches(10.5), Inches(0.35),
                     font_size=12, color=RGBColor(0xAA, 0xBB, 0xDD))

    # Red accent line
    accent = slide.shapes.add_shape(1, Inches(0), Inches(1.05), SLIDE_W, Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()

    # USAT logo top-right
    if show_logo and os.path.exists(USAT_LOGO_PATH):
        slide.shapes.add_picture(USAT_LOGO_PATH, Inches(11.9), Inches(0.03), Inches(1.35), Inches(1.0))

    # Thin navy footer strip
    footer = slide.shapes.add_shape(1, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15))
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY
    footer.line.fill.background()


def fig_to_image(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── Slide builders ─────────────────────────────────────────────────────────────

def add_title_slide(prs: Presentation, venue: str, years_back: int,
                    analysis_label: str = "Elite Analysis"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full navy background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Red diagonal accent bar at right
    accent = slide.shapes.add_shape(1, Inches(12.5), Inches(0), Inches(0.83), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()

    # White horizontal rule
    rule = slide.shapes.add_shape(1, Inches(0.5), Inches(4.1), Inches(11.8), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = WHITE
    rule.line.fill.background()

    # USAT logo centered
    if os.path.exists(USAT_LOGO_PATH):
        slide.shapes.add_picture(USAT_LOGO_PATH, Inches(0.6), Inches(0.4), Inches(2.1), Inches(1.6))

    # Main title
    _add_textbox(slide, f"{venue.upper()} RACE PREVIEW",
                 Inches(0.5), Inches(2.0), Inches(11.8), Inches(1.4),
                 font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle — pull the upcoming race month from EVENT_SCHEDULES when available
    sched = EVENT_SCHEDULES.get(venue.lower().strip())
    when_str = sched.get("date_range", "").split(",")[-1].strip() if sched else ""
    if not when_str:
        when_str = date.today().strftime("%B %Y")
    elif when_str.isdigit():  # year-only
        when_str = sched.get("date_range", date.today().strftime("%B %Y"))
    _add_textbox(slide, f"{analysis_label}  ·  {when_str}",
                 Inches(0.5), Inches(4.3), Inches(11.8), Inches(0.6),
                 font_size=18, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

    _add_textbox(slide, "USA Triathlon High Performance",
                 Inches(0.5), Inches(5.0), Inches(11.8), Inches(0.5),
                 font_size=14, color=RGBColor(0x80, 0x90, 0xB0), align=PP_ALIGN.CENTER)


def add_section_divider(prs: Presentation, gender: str, prefix: str = "ELITE "):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    red_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
    red_bar.fill.solid()
    red_bar.fill.fore_color.rgb = RED
    red_bar.line.fill.background()

    rule = slide.shapes.add_shape(1, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = WHITE
    rule.line.fill.background()

    _add_textbox(slide, f"{prefix}{gender.upper()}",
                 Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.8),
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    if os.path.exists(USAT_LOGO_PATH):
        slide.shapes.add_picture(USAT_LOGO_PATH, Inches(11.0), Inches(6.3), Inches(2.0), Inches(1.1))


def add_overview_slide(prs: Presentation, rows: list[dict], gender: str, venue: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Results Summary  —  Elite {gender.title()}", venue)

    # Column definitions: (header, width_inches, align)
    cols = [
        ("Year",          0.65,  PP_ALIGN.CENTER),
        ("Date",          0.95,  PP_ALIGN.CENTER),
        ("Category",      1.9,   PP_ALIGN.LEFT),
        ("Field",         0.55,  PP_ALIGN.CENTER),
        ("DNF/LAP/DNS",   0.75,  PP_ALIGN.CENTER),
        ("Winner",        2.0,   PP_ALIGN.LEFT),
        ("Winner Time",   1.0,   PP_ALIGN.CENTER),
        ("Air Temp",      0.8,   PP_ALIGN.CENTER),
        ("Water Temp",    0.85,  PP_ALIGN.CENTER),
        ("Swim km",       0.7,   PP_ALIGN.CENTER),
        ("Bike km",       0.7,   PP_ALIGN.CENTER),
        ("Run km",        0.7,   PP_ALIGN.CENTER),
    ]
    n_rows = len(rows) + 1
    n_cols = len(cols)
    row_h = min(0.56, 5.9 / n_rows)
    total_w = sum(c[1] for c in cols)
    scale = 12.73 / total_w  # fit within slide width
    scaled_widths = [c[1] * scale for c in cols]

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.3), Inches(1.25),
        Inches(12.73), Inches(row_h * n_rows)
    )
    tbl = tbl_shape.table
    for i, w in enumerate(scaled_widths):
        tbl.columns[i].width = Inches(w)

    # Header row
    for ci, (col_name, _, align) in enumerate(cols):
        _set_cell(tbl.cell(0, ci), col_name, bold=True, font_size=14,
                  color=WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY)

    def _fmt_dist(val) -> str:
        """Format a distance value stored in either km or metres."""
        if val is None or (isinstance(val, float) and np.isnan(val)):
            return "—"
        # Values > 100 are assumed to be stored in metres; convert to km
        if val > 100:
            val = val / 1000.0
        return f"{val:.2f}" if val < 2 else f"{val:.1f}"

    def _dims_equal(a, b) -> bool:
        """Compare two distance tuples treating NaN as equal."""
        for x, y in zip(a, b):
            x_nan = x is None or (isinstance(x, float) and np.isnan(x))
            y_nan = y is None or (isinstance(y, float) and np.isnan(y))
            if x_nan and y_nan:
                continue
            if x_nan != y_nan or x != y:
                return False
        return True

    prev_dims = None
    footnote_needed = False
    for ri, row in enumerate(rows, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE

        current_dims = (row["swim_km"], row["bike_km"], row["run_km"])
        course_flag = ""
        if prev_dims is not None and not _dims_equal(prev_dims, current_dims):
            course_flag = " *"
            footnote_needed = True
        prev_dims = current_dims

        cat = (row["cat_name"] or row["prog_distance_category"] or "—").replace("World Triathlon Championship Series", "WTCS")
        winner_last = ""
        if row["winner_name"]:
            parts = row["winner_name"].split()
            winner_last = parts[-1] if parts else row["winner_name"]

        air   = str(row.get("temp_air")   or "—").strip()
        water = str(row.get("temp_water") or "—").strip()

        fs  = row.get("field_size") or 0
        dnf = row.get("non_finishers") or 0
        total_field = (fs + dnf) if (fs or dnf) else None

        vals = [
            str(row["year"]),
            pd.to_datetime(row["date"]).strftime("%b %d"),
            cat + course_flag,
            str(total_field) if total_field else "—",
            str(dnf) if dnf else "—",
            winner_last,
            seconds_to_mmss(parse_time_to_seconds(row["winner_time"])),
            air,
            water,
            _fmt_dist(row["swim_km"]),
            _fmt_dist(row["bike_km"]),
            _fmt_dist(row["run_km"]),
        ]
        for ci, (v, (col_name, _, align)) in enumerate(zip(vals, cols)):
            cell_font = 12 if col_name == "Category" else 14
            _set_cell(tbl.cell(ri, ci), v, font_size=cell_font, align=align, bg_color=bg)

    if footnote_needed:
        _add_textbox(slide, "* Course distance changed vs. prior year",
                     Inches(0.3), Inches(7.1), Inches(6), Inches(0.28),
                     font_size=8, color=MID_GRAY, italic=True)


def add_swim_slide(prs: Presentation, rows: list[dict], gender: str, venue: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Swim Exit Dynamics  —  Elite {gender.title()}", venue)

    years = [str(r["year"]) for r in rows]
    group_sizes = [r["swim_exit_group"] or 0 for r in rows]
    fastest = [r["swim_fastest"] for r in rows]
    avg_swim = [r["swim_avg"] for r in rows]

    fig, axes = plt.subplots(1, 2, figsize=(13, 5.2))
    fig.patch.set_facecolor("white")

    # Left: swim exit group size
    ax1 = axes[0]
    bar_colors = [C_NAVY if g >= 20 else C_RED for g in group_sizes]
    bars = ax1.bar(years, group_sizes, color=bar_colors, edgecolor="white", linewidth=0.5, width=0.6)
    for bar, val in zip(bars, group_sizes):
        if val:
            ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                     str(val), ha="center", va="bottom", fontsize=10, fontweight="bold", color="#333333")
    ax1.set_title("Athletes Within 15 sec of Swim Leader", fontsize=13, pad=8)
    ax1.set_ylabel("# Athletes", fontsize=11)
    ax1.set_ylim(0, max(group_sizes + [1]) * 1.3)
    ax1.tick_params(axis="x", rotation=45, labelsize=10)
    ax1.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax1.set_facecolor("white")

    # Right: fastest vs avg swim time — dot/lollipop chart, zoomed y-axis
    ax2 = axes[1]
    x = np.arange(len(years))
    fast_secs = [f for f in fastest]
    avg_secs = [a for a in avg_swim]
    all_secs = [v for v in fast_secs + avg_secs if v]

    ax2.vlines(x, fast_secs, avg_secs, color=C_GRAY, linewidth=1.5, alpha=0.5, zorder=1)
    ax2.scatter(x, fast_secs, s=110, color=C_NAVY, zorder=3, label="Fastest", edgecolor="white", linewidth=1.2)
    ax2.scatter(x, avg_secs, s=110, color=C_LIGHT_BLUE, zorder=3, label="Average", edgecolor="white", linewidth=1.2)

    for xi, secs in zip(x, fast_secs):
        if secs:
            ax2.annotate(seconds_to_mmss(secs), (xi, secs), xytext=(0, -14),
                         textcoords="offset points", ha="center", fontsize=9,
                         fontweight="bold", color=C_NAVY)
    for xi, secs in zip(x, avg_secs):
        if secs:
            ax2.annotate(seconds_to_mmss(secs), (xi, secs), xytext=(0, 9),
                         textcoords="offset points", ha="center", fontsize=9, color=C_LIGHT_BLUE)

    ax2.set_xticks(x)
    ax2.set_xticklabels(years, rotation=45, fontsize=10)
    ax2.set_title("Fastest vs. Average Swim Time", fontsize=13, pad=8)
    ax2.set_ylabel("Time (min:sec)", fontsize=11)
    ax2.legend(fontsize=10, loc="upper right")
    ax2.set_facecolor("white")
    ax2.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v, _: f"{int(v // 60)}:{int(round(v % 60)):02d}")
    )
    if all_secs:
        pad = max(8, (max(all_secs) - min(all_secs)) * 0.25)
        ax2.set_ylim(min(all_secs) - pad, max(all_secs) + pad)
    ax2.grid(axis="y", alpha=0.3, linestyle="--")

    plt.tight_layout(pad=1.5)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.25), Inches(12.73), Inches(5.9))


def add_bike_evolution_slide(prs: Presentation, rows: list[dict], gender: str, venue: str):
    """Bike pack evolution: lead group size at T1 vs T2, plus chase gap per year."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Bike Pack Evolution  —  Elite {gender.title()}", venue)

    _add_textbox(
        slide,
        "Lead pack = athletes within 15 sec of leader. "
        "T1 = lead group entering the bike; T2 = lead group exiting the bike.",
        Inches(0.3), Inches(1.18), Inches(12.73), Inches(0.38),
        font_size=10.5, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )

    years       = [str(r["year"]) for r in rows]
    t1_sizes    = [r.get("t1_lead_pack") or 0 for r in rows]
    t2_sizes    = [r.get("t2_lead_pack") or 0 for r in rows]
    chase_gaps  = [r.get("t2_chase_gap")       for r in rows]  # seconds, may be None

    fig, (ax_left, ax_right) = plt.subplots(1, 2, figsize=(12.5, 4.5))
    fig.patch.set_facecolor("white")
    x = np.arange(len(years))

    # ── Left: T1 vs T2 pack size per year ─────────────────────────────────────
    # Draw connecting line (consolidation arrow) first
    for xi, (t1, t2) in enumerate(zip(t1_sizes, t2_sizes)):
        if t1 and t2:
            clr = C_GREEN if t2 > t1 else (C_RED if t2 < t1 else C_GRAY)
            ax_left.plot([xi, xi], [t1, t2], color=clr, linewidth=2, alpha=0.6, zorder=1)

    sc1 = ax_left.scatter(x, t1_sizes, s=120, color=C_LIGHT_BLUE, zorder=3,
                          label="T1 (swim exit)", edgecolor="white", linewidth=1.2)
    sc2 = ax_left.scatter(x, t2_sizes, s=120, color=C_NAVY, zorder=3,
                          label="T2 (bike exit)", edgecolor="white", linewidth=1.2)

    for xi, (t1, t2) in enumerate(zip(t1_sizes, t2_sizes)):
        if t1:
            ax_left.annotate(str(t1), (xi, t1), xytext=(0, 8),
                             textcoords="offset points", ha="center",
                             fontsize=9, color=C_LIGHT_BLUE, fontweight="bold")
        if t2:
            ax_left.annotate(str(t2), (xi, t2), xytext=(0, -14),
                             textcoords="offset points", ha="center",
                             fontsize=9, color=C_NAVY, fontweight="bold")

    ax_left.set_xticks(x)
    ax_left.set_xticklabels(years, rotation=45, fontsize=10)
    ax_left.set_title("Lead Pack Size: Entering T1 vs Exiting T2", fontsize=12, pad=8)
    ax_left.set_ylabel("Athletes in lead group", fontsize=11)
    ax_left.legend(fontsize=10, loc="upper left")
    ax_left.set_facecolor("white")
    ax_left.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    valid = [v for v in t1_sizes + t2_sizes if v]
    if valid:
        ax_left.set_ylim(0, max(valid) * 1.35)

    # ── Right: chase group gap at T2 ──────────────────────────────────────────
    # If lead pack > 75% of field, treat as one main group (show 7s sentinel)
    MAIN_GROUP_SENTINEL = 7
    adjusted_gaps = []
    is_main_group = []
    for r, g in zip(rows, chase_gaps):
        fs = r.get("field_size") or 0
        t2 = r.get("t2_lead_pack") or 0
        if fs and t2 and t2 >= 0.75 * fs:
            adjusted_gaps.append(MAIN_GROUP_SENTINEL)
            is_main_group.append(True)
        else:
            adjusted_gaps.append(g)
            is_main_group.append(False)

    gap_vals = [g if g is not None else 0 for g in adjusted_gaps]
    bar_colors = []
    for g, main in zip(adjusted_gaps, is_main_group):
        if g is None:    bar_colors.append(C_GRAY)
        elif main:       bar_colors.append(C_NAVY)   # one main group
        elif g <= 30:    bar_colors.append(C_GREEN)
        elif g <= 60:    bar_colors.append(C_GOLD)
        else:            bar_colors.append(C_RED)

    bars = ax_right.bar(years, gap_vals, color=bar_colors, edgecolor="white",
                        linewidth=0.5, width=0.55)
    for bar, g, main in zip(bars, adjusted_gaps, is_main_group):
        if g:
            label = "Main group" if main else f"{g}s"
            ax_right.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                          label, ha="center", va="bottom", fontsize=10, fontweight="bold")

    ax_right.set_title("Gap: Lead Pack → Chase Group at T2", fontsize=12, pad=8)
    ax_right.set_ylabel("Seconds behind lead pack", fontsize=11)
    ax_right.tick_params(axis="x", rotation=45, labelsize=10)
    ax_right.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))
    ax_right.set_facecolor("white")

    from matplotlib.patches import Patch
    ax_right.legend(handles=[
        Patch(facecolor=C_NAVY,  label="One main group (>75% of field)"),
        Patch(facecolor=C_GREEN, label="≤30s (tight)"),
        Patch(facecolor=C_GOLD,  label="31–60s (gap)"),
        Patch(facecolor=C_RED,   label=">60s (breakaway)"),
        Patch(facecolor=C_GRAY,  label="No data"),
    ], fontsize=9, loc="upper right")

    plt.tight_layout(pad=1.5)
    slide.shapes.add_picture(fig_to_image(fig),
                             Inches(0.3), Inches(1.62), Inches(12.73), Inches(4.1))

    # ── Summary table ──────────────────────────────────────────────────────────
    tbl_cols = ["Year", "Lead Size @ T1", "Lead Size @ T2", "Change",
                "Chase Group @ T2", "Chase Gap*"]
    n_rows = 1 + len(rows)
    tbl_shape = slide.shapes.add_table(
        n_rows, len(tbl_cols),
        Inches(0.3), Inches(5.78), Inches(12.73), Inches(1.55),
    )
    tbl = tbl_shape.table
    col_w = [Inches(0.75), Inches(2.0), Inches(2.0), Inches(1.3),
             Inches(2.5), Inches(4.18)]
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = w

    for ci, hdr in enumerate(tbl_cols):
        _set_cell(tbl.cell(0, ci), hdr, bold=True, color=WHITE, bg_color=NAVY,
                  font_size=10)

    for ri, row in enumerate(rows, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        t1 = row.get("t1_lead_pack")
        t2 = row.get("t2_lead_pack")
        chg = ""
        if t1 and t2:
            diff = t2 - t1
            chg = f"+{diff}" if diff > 0 else str(diff)
        cg  = str(row.get("t2_chase_pack") or "—")
        gap = f"{row['t2_chase_gap']}s" if row.get("t2_chase_gap") is not None else "—"
        vals = [str(row["year"]), str(t1 or "—"), str(t2 or "—"),
                chg or "—", cg, gap]
        for ci, v in enumerate(vals):
            _set_cell(tbl.cell(ri, ci), v, font_size=9.5, bg_color=bg,
                      align=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 "Lead pack = athletes within 15 sec of leader. "
                 "Chase group = athletes within 15 sec of the first athlete outside the lead pack. "
                 "* Chase Gap = seconds from the last lead-pack athlete to the first chase-group athlete at T2.",
                 Inches(0.3), Inches(7.18), Inches(12.73), Inches(0.22),
                 font_size=8, color=MID_GRAY, italic=True)

    # Dropout annotation: summarise DNF/LAP athletes per year that were in the lead group
    dropout_notes = []
    for r in rows:
        d_list = r.get("dropouts") or []
        if not d_list:
            continue
        in_lead = [d for d in d_list if d["in_lead_grp"]]
        total   = len(d_list)
        note = f"{r['year']}: {total} DNF/LAP"
        if in_lead:
            cps = ", ".join(sorted({d['last_cp'] for d in in_lead}))
            note += f" ({len(in_lead)} in lead group at {cps})"
        dropout_notes.append(note)

    if dropout_notes:
        _add_textbox(slide,
                     "Dropouts: " + "  |  ".join(dropout_notes),
                     Inches(0.3), Inches(7.35), Inches(12.73), Inches(0.18),
                     font_size=7.5, color=MID_GRAY, italic=True)


def add_run_slide(prs: Presentation, rows: list[dict], gender: str, venue: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Run Split Analysis  —  Elite {gender.title()}", venue)

    years = [str(r["year"]) for r in rows]
    fastest = [r["run_fastest"] for r in rows]
    avg_run  = [r["run_avg"]     for r in rows]

    fig, ax = plt.subplots(figsize=(12, 5.3))
    fig.patch.set_facecolor("white")
    x = np.arange(len(years))

    fast_secs = [f for f in fastest]
    avg_secs = [a for a in avg_run]
    all_secs = [v for v in fast_secs + avg_secs if v]

    ax.vlines(x, fast_secs, avg_secs, color=C_GRAY, linewidth=1.5, alpha=0.5, zorder=1)
    ax.scatter(x, fast_secs, s=130, color=C_NAVY, zorder=3, label="Fastest", edgecolor="white", linewidth=1.2)
    ax.scatter(x, avg_secs, s=130, color=C_LIGHT_BLUE, zorder=3, label="Average", edgecolor="white", linewidth=1.2)

    for xi, secs in zip(x, fast_secs):
        if secs:
            ax.annotate(seconds_to_mmss(secs), (xi, secs), xytext=(0, -16),
                        textcoords="offset points", ha="center", fontsize=10,
                        fontweight="bold", color=C_NAVY)
    for xi, secs in zip(x, avg_secs):
        if secs:
            ax.annotate(seconds_to_mmss(secs), (xi, secs), xytext=(0, 10),
                        textcoords="offset points", ha="center", fontsize=10, color=C_LIGHT_BLUE)

    ax.set_xticks(x)
    ax.set_xticklabels(years, rotation=45, fontsize=11)
    ax.set_title("Fastest vs. Average Run Split by Year", fontsize=14, pad=10)
    ax.set_ylabel("Time (min:sec)", fontsize=12)
    ax.legend(fontsize=11, loc="upper right")
    ax.set_facecolor("white")
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v, _: f"{int(v // 60)}:{int(round(v % 60)):02d}")
    )
    if all_secs:
        pad = max(15, (max(all_secs) - min(all_secs)) * 0.25)
        ax.set_ylim(min(all_secs) - pad, max(all_secs) + pad)
    ax.grid(axis="y", alpha=0.3, linestyle="--")

    plt.tight_layout(pad=1.5)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.5), Inches(1.22), Inches(12.3), Inches(5.95))


def add_position_times_slide(prs: Presentation, rows: list[dict], gender: str, venue: str):
    """Line chart + table: times for 1st, 2nd, 3rd, 5th, 10th, 20th place per year."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Key Position Times  —  Elite {gender.title()}", venue)

    positions = [1, 2, 3, 5, 10, 20]
    pos_labels = {1: "1st", 2: "2nd", 3: "3rd", 5: "5th", 10: "10th", 20: "20th"}
    colors_pos = {1: C_NAVY, 2: C_RED, 3: C_ORANGE, 5: C_LIGHT_BLUE, 10: C_GREEN, 20: C_GRAY}

    years = [str(r["year"]) for r in rows]

    # ── Left: line chart of gap from winner (seconds) ──
    fig, (ax_left, ax_right) = plt.subplots(1, 2, figsize=(12.5, 3.4))
    fig.patch.set_facecolor("white")

    for pos in [2, 3, 5, 10, 20]:
        gaps = []
        for row in rows:
            pt = row["pos_times"].get(pos)
            gaps.append(pt["gap"] if pt and pt.get("gap") is not None else np.nan)
        ax_left.plot(years, gaps, marker="o", label=pos_labels[pos],
                     color=colors_pos[pos], linewidth=2, markersize=6)
        # Label endpoint
        last_valid = [(y, g) for y, g in zip(years, gaps) if not np.isnan(g)]
        if last_valid:
            ly, lg = last_valid[-1]
            ax_left.annotate(f"{int(round(lg))}s", xy=(ly, lg),
                             xytext=(4, 0), textcoords="offset points",
                             fontsize=7.5, color=colors_pos[pos], fontweight="bold")

    ax_left.set_title("Gap from Winner (seconds)", fontsize=11, pad=8)
    ax_left.set_ylabel("Seconds behind 1st place", fontsize=10)
    ax_left.tick_params(axis="x", rotation=45, labelsize=9)
    ax_left.set_facecolor("white")
    ax_left.legend(fontsize=9, loc="upper left")
    ax_left.axhline(0, color=C_NAVY, linewidth=0.8, alpha=0.4)

    # ── Right: bar chart of average gaps per position ──
    avg_gaps = []
    for pos in [2, 3, 5, 10, 20]:
        all_gaps = []
        for row in rows:
            pt = row["pos_times"].get(pos)
            if pt and pt.get("gap") is not None:
                all_gaps.append(pt["gap"])
        avg_gaps.append(np.mean(all_gaps) if all_gaps else 0)

    bar_labels = [pos_labels[p] for p in [2, 3, 5, 10, 20]]
    bar_clrs = [colors_pos[p] for p in [2, 3, 5, 10, 20]]
    bars = ax_right.bar(bar_labels, avg_gaps, color=bar_clrs, edgecolor="white", width=0.55)
    for bar, val in zip(bars, avg_gaps):
        if val:
            ax_right.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                          f"{int(round(val))}s", ha="center", va="bottom", fontsize=9.5, fontweight="bold")
    ax_right.set_title("Avg Gap from 1st Place (all years)", fontsize=11, pad=8)
    ax_right.set_ylabel("Seconds behind winner", fontsize=10)
    ax_right.set_facecolor("white")
    ax_right.tick_params(axis="x", labelsize=10)

    plt.tight_layout(pad=1.2)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.18), Inches(12.73), Inches(3.3))

    # ── Summary table: actual times per position per year ──
    tbl_rows = 1 + len(rows)
    tbl_cols = 1 + len(positions)
    tbl_top = Inches(4.58)
    row_h = min(0.38, 2.7 / tbl_rows)
    tbl_h = Inches(row_h * tbl_rows)
    tbl_shape = slide.shapes.add_table(tbl_rows, tbl_cols,
                                       Inches(0.3), tbl_top, Inches(12.73), tbl_h)
    tbl = tbl_shape.table
    col_widths = [Inches(0.7)] + [Inches(12.03 / len(positions))] * len(positions)
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    _set_cell(tbl.cell(0, 0), "Year", bold=True, color=WHITE, bg_color=NAVY)
    for ci, pos in enumerate(positions, start=1):
        _set_cell(tbl.cell(0, ci), pos_labels[pos], bold=True, color=WHITE, bg_color=NAVY)

    for ri, row in enumerate(rows, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        _set_cell(tbl.cell(ri, 0), str(row["year"]), bold=True, color=DARK_GRAY, bg_color=bg)
        for ci, pos in enumerate(positions, start=1):
            pt = row["pos_times"].get(pos)
            if pt and pt.get("secs"):
                time_str = seconds_to_mmss(pt["secs"])
                gap_str = f"+{int(round(pt['gap']))}s" if pt.get("gap") else ""
                val = f"{time_str}\n{gap_str}" if gap_str else time_str
            else:
                val = "—"
            _set_cell(tbl.cell(ri, ci), val, font_size=8.5, bg_color=bg)


def add_weather_slide(prs: Presentation, rows: list[dict], gender: str, venue: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Race Conditions  —  Elite {gender.title()}", venue)

    cols = [
        ("Year",        0.65, PP_ALIGN.CENTER),
        ("Date",        0.95, PP_ALIGN.CENTER),
        ("Air Temp",    1.0,  PP_ALIGN.CENTER),
        ("Humidity",    0.9,  PP_ALIGN.CENTER),
        ("Wind (km/h)", 1.0,  PP_ALIGN.CENTER),
        ("Precip (mm)", 0.95, PP_ALIGN.CENTER),
        ("Water Temp",  0.95, PP_ALIGN.CENTER),
        ("Wetsuit",     0.9,  PP_ALIGN.CENTER),
        ("Conditions",  5.49, PP_ALIGN.LEFT),
    ]
    n_rows = len(rows) + 1
    row_h = min(0.50, 5.0 / n_rows)
    tbl_shape = slide.shapes.add_table(
        n_rows, len(cols),
        Inches(0.3), Inches(1.25),
        Inches(12.73), Inches(row_h * n_rows)
    )
    tbl = tbl_shape.table
    total_w = sum(c[1] for c in cols)
    scale = 12.73 / total_w
    for ci, (_, w, _) in enumerate(cols):
        tbl.columns[ci].width = Inches(w * scale)

    for ci, (name, _, _) in enumerate(cols):
        _set_cell(tbl.cell(0, ci), name, bold=True, font_size=12,
                  color=WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY)

    for ri, row in enumerate(rows, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        wind_val = (f"{row['wind_kmh']:.0f}" if row.get("wind_kmh") else (row.get("wind_raw") or "—"))
        air    = str(row.get("temp_air")   or "—").strip()
        water  = str(row.get("temp_water") or "—").strip()
        wetsuit = str(row.get("wetsuit")   or "—").strip()
        weather = str(row.get("weather")   or "—").strip()
        humid  = f"{row['humidity']:.0f}%" if row.get("humidity") is not None else "—"
        precip = f"{row['precip']:.1f}" if row.get("precip") is not None else "—"
        if len(weather) > 80:
            weather = weather[:77] + "..."

        vals = [str(row["year"]), pd.to_datetime(row["date"]).strftime("%b %d"),
                air, humid, str(wind_val), precip, water, wetsuit, weather]
        for ci, (v, (_, _, align)) in enumerate(zip(vals, cols)):
            _set_cell(tbl.cell(ri, ci), v, font_size=10, align=align, bg_color=bg)

    _add_textbox(slide,
                 "Air temp, humidity, wind, and precipitation sourced from Open-Meteo historical archive "
                 f"(race window {_WEATHER_WINDOW[0]:02d}:00–{_WEATHER_WINDOW[1]:02d}:00 local) where WT API data is absent. "
                 "Values marked * are Open-Meteo estimates. "
                 "Water temp and wetsuit ruling are from the World Triathlon API per program.",
                 Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.72),
                 font_size=9, color=MID_GRAY, italic=True)


# Sport-class display order for Para tables/charts.
_PARA_CLASS_ORDER = ["PTWC", "PTS2", "PTS3", "PTS4", "PTS5", "PTVI"]


def _para_class_sort_key(label: str) -> int:
    up = label.upper()
    for i, c in enumerate(_PARA_CLASS_ORDER):
        if c in up:
            return i
    return 99


def _para_gender_progs(para_events_df: pd.DataFrame, gender: str,
                       para_classes: list[str] | None) -> list[tuple[str, str]]:
    """Ordered list of (class_label, full_prog_name) for a single gender.

    Matches the gender on the final program-name token so 'Women' is not caught
    by an 'endswith("men")' test. Honours an optional para_classes filter.
    """
    progs: set[tuple[str, str]] = set()
    for pn in para_events_df["prog_name"].dropna().unique():
        parts = pn.rsplit(" ", 1)
        if len(parts) != 2 or parts[1].lower() != gender.lower():
            continue
        label = parts[0].strip()
        if para_classes and not any(pc.upper() in label.upper() for pc in para_classes):
            continue
        progs.add((label, pn))
    return sorted(progs, key=lambda t: _para_class_sort_key(t[0]))


def add_para_results_summary_slide(prs: Presentation, para_events_df: pd.DataFrame,
                                   engine, venue: str, gender: str,
                                   para_classes: list[str] | None):
    """Para Series results summary for one gender — one row per event-year, one
    column per sport class, cell = winning athlete surname + total time."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Para Series Results  —  {gender}", venue)

    progs = _para_gender_progs(para_events_df, gender, para_classes)
    prog_names = [pn for _, pn in progs]
    sub = para_events_df[para_events_df.prog_name.isin(prog_names)]
    if not progs or sub.empty:
        _add_textbox(slide, f"No Para {gender} results found for this venue in the selected period.",
                     Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.5),
                     font_size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    # Winner per (event_id, prog_name)
    winners: dict[tuple, dict] = {}
    for _, row in sub.iterrows():
        winners[(int(row.event_id), row.prog_name)] = query_winner(
            engine, int(row.event_id), int(row.prog_id))

    # One row per event-year, dropping years where no class produced a winner
    events = (sub.drop_duplicates(subset=["event_id"])
                 .sort_values("event_date")[["event_id", "event_date"]]
                 .to_dict("records"))
    def _ev_has_data(eid: int) -> bool:
        for pn in prog_names:
            w = winners.get((eid, pn)) or {}
            if w.get("winner_name") or w.get("winner_time"):
                return True
        return False

    events = [ev for ev in events if _ev_has_data(int(ev["event_id"]))]
    if not events:
        _add_textbox(slide, f"No completed Para {gender} races found for this venue.",
                     Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.5),
                     font_size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    # Columns: Year | Date | <one per class>
    fixed_cols = [("Year", 0.7), ("Date", 0.95)]
    class_w = (12.73 - sum(w for _, w in fixed_cols)) / max(len(progs), 1)
    all_cols = fixed_cols + [(label, class_w) for label, _ in progs]

    n_rows = 1 + len(events)
    row_h = min(0.65, 5.6 / n_rows)
    total_w = sum(w for _, w in all_cols)
    scale = 12.73 / total_w

    tbl_shape = slide.shapes.add_table(
        n_rows, len(all_cols), Inches(0.3), Inches(1.25),
        Inches(12.73), Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for ci, (_, w) in enumerate(all_cols):
        tbl.columns[ci].width = Inches(w * scale)

    for ci, (name, _) in enumerate(all_cols):
        _set_cell(tbl.cell(0, ci), name, bold=True, font_size=11,
                  color=WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY)

    for ri, ev in enumerate(events, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        eid = int(ev["event_id"])
        _set_cell(tbl.cell(ri, 0), str(pd.to_datetime(ev["event_date"]).year),
                  bold=True, font_size=11, bg_color=bg)
        _set_cell(tbl.cell(ri, 1), pd.to_datetime(ev["event_date"]).strftime("%b %d"),
                  font_size=10.5, bg_color=bg)
        for ci_offset, (_, pn) in enumerate(progs):
            ci = len(fixed_cols) + ci_offset
            w = winners.get((eid, pn)) or {}
            name_raw = w.get("winner_name")
            t_str = str(w.get("winner_time") or "")
            t_mm = seconds_to_mmss(parse_time_to_seconds(t_str)) if t_str else None
            if not name_raw and not t_mm:
                _set_cell(tbl.cell(ri, ci), "—", font_size=9, color=MID_GRAY, bg_color=bg)
                continue
            # Name may be absent in older Para records — still surface the
            # winning time so the course record / trend stays readable.
            surname = _para_surname(name_raw) if name_raw else "n/a"
            _set_cell(tbl.cell(ri, ci), f"{surname}\n{t_mm or '—'}",
                      font_size=9, bg_color=bg)

    _add_textbox(
        slide,
        f"World Para Series {gender} winners by sport class. Surname + winning total time. "
        "Years with no completed racing are omitted.",
        Inches(0.3), Inches(7.08), Inches(12.73), Inches(0.25),
        font_size=9, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER,
    )


def add_para_split_trend_slide(prs: Presentation, para_events_df: pd.DataFrame,
                               engine, venue: str, gender: str,
                               para_classes: list[str] | None,
                               split: str = "total", split_label: str = "Winning Time"):
    """Per-gender trend chart: fastest <split> by sport class across years."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"{split_label} Trends  —  Para {gender}", venue)

    progs = _para_gender_progs(para_events_df, gender, para_classes)
    if not progs:
        _add_textbox(slide, f"No Para {gender} split data available.",
                     Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.5),
                     font_size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    fast_key = f"{split}_fastest"
    # class_label -> {year: fastest_split_seconds}
    class_data: dict[str, dict[int, float]] = {}
    for label, pn in progs:
        match = para_events_df[para_events_df.prog_name == pn]
        per_year: dict[int, float] = {}
        for _, row in match.iterrows():
            splits = query_splits(engine, int(row.event_id), int(row.prog_id))
            val = splits.get(fast_key)
            if val:
                yr = pd.to_datetime(row.event_date).year
                per_year[yr] = min(val, per_year.get(yr, val))
        if per_year:
            class_data[label] = per_year

    all_years = sorted({yr for d in class_data.values() for yr in d})
    if not all_years:
        _add_textbox(slide, f"No {split_label.lower()} data available for Para {gender}.",
                     Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.5),
                     font_size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    ordered_labels = [label for label, _ in progs if label in class_data]
    fig, ax = plt.subplots(figsize=(12, 5.6), dpi=150)
    fig.patch.set_facecolor("white")
    palette = [C_NAVY, C_RED, C_LIGHT_BLUE, C_ORANGE, C_GREEN, C_GOLD, C_VIOLET, C_GRAY]
    x = np.arange(len(all_years))
    for i, label in enumerate(ordered_labels):
        vals = [class_data[label].get(yr) for yr in all_years]
        color = palette[i % len(palette)]
        ax.plot(x, vals, marker="o", label=label, color=color, linewidth=1.8, markersize=6)
        for xi, v in zip(x, vals):
            if v is not None:
                ax.annotate(seconds_to_mmss(v), (xi, v), xytext=(0, 8),
                            textcoords="offset points", ha="center", fontsize=7.5,
                            color=color, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels([str(y) for y in all_years], rotation=45, fontsize=10)
    ax.set_title(f"Fastest {split_label} by Sport Class per Year — {gender}",
                 fontsize=13, pad=8, fontweight="bold")
    ax.set_ylabel("Time (h:mm:ss)" if split == "total" else "Time (min:sec)", fontsize=11)
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v, _: seconds_to_mmss(v))
    )
    ax.legend(fontsize=9, loc="upper right", ncols=min(3, max(1, len(ordered_labels) // 2)))
    ax.grid(axis="y", alpha=0.3, linestyle="--")
    ax.set_facecolor("white")
    plt.tight_layout(pad=1.5)
    slide.shapes.add_picture(fig_to_image(fig), Inches(0.3), Inches(1.25), Inches(12.73), Inches(5.95))


def add_environmental_risk_slide(
    prs: Presentation, venue: str, unique_rows: list[dict]
):
    """Venue-level environmental risk: weather history + air quality + health placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Environmental Risk Profile", venue)

    # ── Top half: weather + AQ tables side by side ────────────────────────────
    # LEFT: weather table  (Year | Air Temp | Humidity | Wind | Precip)
    wx_cols = [
        ("Year",        0.65, PP_ALIGN.CENTER),
        ("Date",        0.85, PP_ALIGN.CENTER),
        ("Air Temp",    0.95, PP_ALIGN.CENTER),
        ("Humidity",    0.85, PP_ALIGN.CENTER),
        ("Wind km/h",   0.85, PP_ALIGN.CENTER),
        ("Precip mm",   0.85, PP_ALIGN.CENTER),
    ]
    n_rows = len(unique_rows) + 1
    row_h = min(0.44, 2.7 / n_rows)

    wx_shape = slide.shapes.add_table(
        n_rows, len(wx_cols),
        Inches(0.3), Inches(1.28),
        Inches(5.15), Inches(row_h * n_rows),
    )
    wx_tbl = wx_shape.table
    total_w = sum(c[1] for c in wx_cols)
    for ci, (_, w, _) in enumerate(wx_cols):
        wx_tbl.columns[ci].width = Inches(w / total_w * 5.15)
    for ci, (hdr, _, _) in enumerate(wx_cols):
        _set_cell(wx_tbl.cell(0, ci), hdr, bold=True, font_size=10,
                  color=WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY)

    for ri, row in enumerate(unique_rows, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        air   = str(row.get("temp_air")  or "—").strip()
        humid = f"{row['humidity']:.0f}%" if row.get("humidity") is not None else "—"
        wind  = f"{row['wind_kmh']:.0f}" if row.get("wind_kmh") else (row.get("wind_raw") or "—")
        prec  = f"{row['precip']:.1f}"   if row.get("precip")   is not None else "—"
        vals  = [str(row["year"]), pd.to_datetime(row["date"]).strftime("%b %d"),
                 air, humid, str(wind), prec]
        for ci, (v, (_, _, align)) in enumerate(zip(vals, wx_cols)):
            _set_cell(wx_tbl.cell(ri, ci), v, font_size=10, align=align, bg_color=bg)

    # Weather section label
    wx_bar = slide.shapes.add_shape(1, Inches(0.3), Inches(1.18), Inches(5.15), Inches(0.1))
    wx_bar.fill.solid(); wx_bar.fill.fore_color.rgb = NAVY; wx_bar.line.fill.background()
    _add_textbox(slide, f"RACE-DAY WEATHER ({_WEATHER_WINDOW[0]:02d}:00–{_WEATHER_WINDOW[1]:02d}:00 local)",
                 Inches(0.35), Inches(1.19), Inches(5.0), Inches(0.1),
                 font_size=9.5, bold=True, color=NAVY)

    # RIGHT: air quality table  (Year | PM2.5 | AQI | UV | Risk)
    def _aqi_label(aqi) -> str:
        if aqi is None: return "—"
        aqi = int(aqi)
        if aqi <= 20:  return "Good"
        if aqi <= 40:  return "Fair"
        if aqi <= 60:  return "Moderate"
        if aqi <= 80:  return "Poor"
        return "Very Poor"

    aq_cols = [
        ("Year",      0.65, PP_ALIGN.CENTER),
        ("PM2.5",     0.85, PP_ALIGN.CENTER),
        ("EU AQI",    0.75, PP_ALIGN.CENTER),
        ("UV Index",  0.85, PP_ALIGN.CENTER),
        ("Risk",      1.1,  PP_ALIGN.CENTER),
    ]
    aq_shape = slide.shapes.add_table(
        n_rows, len(aq_cols),
        Inches(5.65), Inches(1.28),
        Inches(4.45), Inches(row_h * n_rows),
    )
    aq_tbl = aq_shape.table
    total_w2 = sum(c[1] for c in aq_cols)
    for ci, (_, w, _) in enumerate(aq_cols):
        aq_tbl.columns[ci].width = Inches(w / total_w2 * 4.45)
    for ci, (hdr, _, _) in enumerate(aq_cols):
        _set_cell(aq_tbl.cell(0, ci), hdr, bold=True, font_size=10,
                  color=WHITE, align=PP_ALIGN.CENTER, bg_color=NAVY)

    for ri, row in enumerate(unique_rows, start=1):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        pm25 = f"{row['pm25']:.1f}" if row.get("pm25") is not None else "—"
        aqi  = f"{int(row['aqi'])}" if row.get("aqi") is not None else "—"
        uv   = f"{row['uv_index']:.1f}" if row.get("uv_index") is not None else "—"
        risk = _aqi_label(row.get("aqi"))
        vals = [str(row["year"]), pm25, aqi, uv, risk]
        for ci, (v, (_, _, align)) in enumerate(zip(vals, aq_cols)):
            _set_cell(aq_tbl.cell(ri, ci), v, font_size=10, align=align, bg_color=bg)

    _add_textbox(slide, "AIR QUALITY (Open-Meteo archive)",
                 Inches(5.7), Inches(1.19), Inches(4.3), Inches(0.1),
                 font_size=9.5, bold=True, color=NAVY)

    # ── Bottom: health / water quality risk section ───────────────────────────
    tbl_bottom = Inches(1.28 + row_h * n_rows + 0.25)

    risk_bar = slide.shapes.add_shape(1, Inches(0.3), tbl_bottom, Inches(12.73), Inches(0.35))
    risk_bar.fill.solid(); risk_bar.fill.fore_color.rgb = RED; risk_bar.line.fill.background()
    _add_textbox(slide, "HEALTH & WATER QUALITY RISK FACTORS",
                 Inches(0.35), tbl_bottom + Inches(0.02), Inches(12.2), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    risk_items = [
        ("Water Quality",
         f"[PLACEHOLDER: Historical water quality data for {venue} — e.g. bacterial counts, "
         "EU Bathing Water Directive status, known algae / pollution events. "
         "Check local environmental agency reports for May race window.]"),
        ("Athlete Health Incidents",
         "[PLACEHOLDER: Any known GI illness, respiratory, or heat-related incidents at this "
         "or comparable venues. E.g. Sunderland 2023 — reported GI illness post-swim for multiple athletes.]"),
        ("Water Temperature",
         f"[PLACEHOLDER: Historical open-water temperature for {venue} in May "
         "(not available from Open-Meteo; check Copernicus Marine / NOAA or WT historical results).]"),
        ("Heat / Humidity Risk",
         "Auto-derived from above table — flag if humidity >70% or air temp >25°C at race time."),
    ]

    y_offset = tbl_bottom + Inches(0.42)
    item_h = (SLIDE_H - y_offset - Inches(0.2)) / len(risk_items)
    for label, text in risk_items:
        _add_textbox(slide, f"• {label}: {text}",
                     Inches(0.4), y_offset, Inches(12.5), item_h,
                     font_size=10, color=DARK_GRAY)
        y_offset += item_h

    _add_textbox(slide,
                 "Weather and AQ data: Open-Meteo historical archive (archive-api.open-meteo.com). "
                 "EU AQI: Good ≤20, Fair ≤40, Moderate ≤60, Poor ≤80, Very Poor >80.",
                 Inches(0.3), Inches(7.18), Inches(12.73), Inches(0.22),
                 font_size=8, color=MID_GRAY, italic=True)


def add_course_differentiators_slide(
    prs: Presentation, venue: str, all_rows: list[dict],
    events_df: pd.DataFrame, content: dict
):
    """Gender-neutral venue intro: narrative, course features, quick links.

    Prefers structured data from VENUE_PREVIEW + SWIM/BIKE/RUN_COURSE_PROFILES.
    Falls back to historical race rows (`all_rows`) when no profile is defined.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Key Course Differentiators", venue)

    vkey = venue.lower().strip()
    preview = VENUE_PREVIEW.get(vkey, {})
    swim_p  = SWIM_COURSE_PROFILES.get(vkey, {})
    bike_p  = BIKE_COURSE_PROFILES.get(vkey, {})
    run_p   = RUN_COURSE_PROFILES.get(vkey, {})

    # ── Left column: venue name + narrative ───────────────────────────────────
    _add_textbox(slide, venue.upper(),
                 Inches(0.3), Inches(1.2), Inches(6.6), Inches(0.55),
                 font_size=20, bold=True, color=NAVY)

    narrative = preview.get("narrative") or (
        f"[PLACEHOLDER: 2-3 sentences on what makes {venue} unique as a race venue — "
        "geography, setting, and character.]"
    )
    _add_textbox(slide, narrative,
                 Inches(0.3), Inches(1.78), Inches(6.6), Inches(1.35),
                 font_size=11, color=DARK_GRAY)

    # Feature bullets header bar
    feat_bar = slide.shapes.add_shape(1, Inches(0.3), Inches(3.25), Inches(6.6), Inches(0.35))
    feat_bar.fill.solid()
    feat_bar.fill.fore_color.rgb = NAVY
    feat_bar.line.fill.background()
    _add_textbox(slide, "COURSE FEATURES",
                 Inches(0.35), Inches(3.27), Inches(6.2), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    features = preview.get("features") or [
        f"[PLACEHOLDER: Feature 1 — swim character]",
        f"[PLACEHOLDER: Feature 2 — bike character]",
        f"[PLACEHOLDER: Feature 3 — run character]",
        f"[PLACEHOLDER: Feature 4 — conditions]",
    ]
    for i, feat in enumerate(features[:4]):
        _add_textbox(slide, f"•  {feat}",
                     Inches(0.4), Inches(3.7 + i * 0.72), Inches(6.4), Inches(0.7),
                     font_size=11, color=DARK_GRAY)

    # ── Right column: course at a glance ───────────────────────────────────────
    glance_bar = slide.shapes.add_shape(1, Inches(7.2), Inches(1.2), Inches(5.8), Inches(0.35))
    glance_bar.fill.solid()
    glance_bar.fill.fore_color.rgb = NAVY
    glance_bar.line.fill.background()
    _add_textbox(slide, "COURSE AT A GLANCE",
                 Inches(7.25), Inches(1.22), Inches(5.4), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    def _fmt_d(v) -> str:
        if v is None or (isinstance(v, float) and np.isnan(v)):
            return "—"
        if v > 100:
            v = v / 1000.0
        if v < 1:
            return f"{int(round(v * 1000))} m"
        return f"{v:g} km"

    def _fmt_swim_line(p: dict) -> str:
        d = p.get("total_km")
        laps = p.get("laps")
        if not d:
            return "—"
        if laps and laps > 1 and p.get("loop_km"):
            return f"{_fmt_d(d)} – {laps} laps × {_fmt_d(p['loop_km'])}"
        if laps == 1:
            return f"{_fmt_d(d)} – 1 lap"
        return _fmt_d(d)

    def _fmt_multi(p: dict) -> str:
        d = p.get("total_km")
        laps = p.get("loops") or p.get("laps")
        loop_km = p.get("loop_km")
        if not d:
            return "—"
        if laps and loop_km:
            return f"{_fmt_d(d)} – {laps} laps × {loop_km:g} km"
        if laps:
            return f"{_fmt_d(d)} – {laps} laps"
        return _fmt_d(d)

    # Prefer profile data; fall back to historical row for venues without profiles
    ref_row = next(
        (r for r in sorted(all_rows, key=lambda x: x["year"], reverse=True)
         if r.get("swim_km") or r.get("bike_km") or r.get("run_km")),
        None
    )
    swim_label = (_fmt_swim_line(swim_p) if swim_p
                  else (_fmt_d(ref_row.get("swim_km")) if ref_row else "—"))
    bike_label = (_fmt_multi(bike_p) if bike_p
                  else (_fmt_d(ref_row.get("bike_km")) if ref_row else "—"))
    run_label  = (_fmt_multi(run_p) if run_p
                  else (_fmt_d(ref_row.get("run_km")) if ref_row else "—"))

    format_label = preview.get("format_label") or "Sprint Triathlon"
    bike_laps_val = bike_p.get("loops") if bike_p else (
        int(ref_row["bike_laps"]) if ref_row and pd.notna(ref_row.get("bike_laps")) else None)
    run_laps_val  = run_p.get("laps") if run_p else (
        int(ref_row["run_laps"]) if ref_row and pd.notna(ref_row.get("run_laps")) else None)

    glance_items = [
        ("Swim",       swim_label),
        ("Bike",       bike_label),
        ("Run",        run_label),
        ("Format",     format_label),
        ("Bike Laps",  str(bike_laps_val) if bike_laps_val else "—"),
        ("Run Laps",   str(run_laps_val) if run_laps_val else "—"),
    ]

    for i, (label, val) in enumerate(glance_items):
        _add_textbox(slide, f"{label}:",
                     Inches(7.3), Inches(1.65 + i * 0.37), Inches(2.0), Inches(0.35),
                     font_size=12, bold=True, color=NAVY)
        _add_textbox(slide, val,
                     Inches(9.0), Inches(1.65 + i * 0.37), Inches(4.0), Inches(0.35),
                     font_size=12, color=DARK_GRAY)

    # ── Right column: quick links ──────────────────────────────────────────────
    links_y = 4.0
    links_bar = slide.shapes.add_shape(1, Inches(7.2), Inches(links_y), Inches(5.8), Inches(0.35))
    links_bar.fill.solid()
    links_bar.fill.fore_color.rgb = NAVY
    links_bar.line.fill.background()
    _add_textbox(slide, "QUICK LINKS",
                 Inches(7.25), Inches(links_y + 0.02), Inches(5.4), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    race_info_url = preview.get("race_info_url") or content.get("upcoming_event_url")
    race_info_text = preview.get("race_info_text") or (
        race_info_url or f"Race Info | {venue}")

    _add_textbox(slide, "Current Race Info:",
                 Inches(7.3), Inches(links_y + 0.45), Inches(5.6), Inches(0.3),
                 font_size=11, bold=True, color=NAVY)
    _add_textbox(slide, race_info_text or "—",
                 Inches(7.3), Inches(links_y + 0.78), Inches(5.6), Inches(0.42),
                 font_size=10, color=DARK_GRAY, italic=True)


def add_course_map_slide(
    prs: Presentation, venue: str, all_rows: list[dict], content: dict
):
    """Course profile and map slide with placeholder images + lap details table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Course Profile & Map", venue)

    this_year = date.today().year

    for col_i, (label, yr_hint) in enumerate([
        (f"{this_year} Course", "upcoming course map or race photo"),
        (f"{this_year - 1} Course", "previous year course map or race photo"),
    ]):
        left = Inches(0.3 + col_i * 6.55)
        width = Inches(6.2)
        height = Inches(3.65)

        # Gray placeholder box
        ph = slide.shapes.add_shape(1, left, Inches(1.22), width, height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_GRAY
        ph.line.color.rgb = MID_GRAY

        _add_textbox(slide,
                     f"[PLACEHOLDER: Insert {yr_hint} here]\n"
                     f"Right-click → Replace Image, or paste screenshot\n"
                     f"of the World Triathlon course map PDF.",
                     left + Inches(0.2), Inches(2.35), width - Inches(0.4), Inches(1.0),
                     font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

        _add_textbox(slide, label,
                     left, Inches(4.93), width, Inches(0.3),
                     font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        src_url = (
            content.get("upcoming_event_url") if col_i == 0
            else content.get("last_year_url")
        ) or "[PLACEHOLDER: World Triathlon event page URL]"
        _add_textbox(slide, f"Source: {src_url}",
                     left, Inches(5.22), width, Inches(0.27),
                     font_size=9, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Bottom: course details table
    ref_row = next(
        (r for r in sorted(all_rows, key=lambda x: x["year"], reverse=True)
         if r.get("swim_km") or r.get("bike_km") or r.get("run_km")),
        None
    )

    def _fd(v, unit="") -> str:
        if v is None or (isinstance(v, float) and np.isnan(v)):
            return "—"
        if v > 100:
            v = v / 1000.0
        return f"{v:.1f}{unit}"

    def _fl(v) -> str:
        if v is None or (isinstance(v, float) and np.isnan(v)):
            return "—"
        return str(int(v))

    r = ref_row or {}
    tbl_data = [
        ("Discipline", "Distance",           "Laps",              "Terrain / Notes"),
        ("Swim",       _fd(r.get("swim_km"), " km"), _fl(r.get("swim_laps")),
         f"[PLACEHOLDER: e.g. Open-water, Bay of {venue}]"),
        ("Bike",       _fd(r.get("bike_km"), " km"), _fl(r.get("bike_laps")),
         "[PLACEHOLDER: e.g. Flat coastal road circuit, closed course]"),
        ("Run",        _fd(r.get("run_km"),  " km"), _fl(r.get("run_laps")),
         "[PLACEHOLDER: e.g. Flat multi-lap, cobblestone city-centre finish]"),
    ]

    tbl_shape = slide.shapes.add_table(
        4, 4, Inches(0.3), Inches(5.42), Inches(12.73), Inches(1.8)
    )
    tbl = tbl_shape.table
    for ci, w in enumerate([Inches(1.6), Inches(1.7), Inches(1.0), Inches(8.43)]):
        tbl.columns[ci].width = w

    for ri, row_data in enumerate(tbl_data):
        for ci, val in enumerate(row_data):
            is_hdr = ri == 0
            bg = NAVY if is_hdr else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
            clr = WHITE if is_hdr else DARK_GRAY
            align = PP_ALIGN.CENTER if ci < 3 else PP_ALIGN.LEFT
            _set_cell(tbl.cell(ri, ci), val, bold=is_hdr,
                      font_size=11 if is_hdr else 10, color=clr, bg_color=bg, align=align)


# ── Per-venue course profile data ─────────────────────────────────────────────
# Each dict is populated from official organiser course maps + race-info pages.
# All three discipline dicts share parallel structure so the slide builders can
# be lightly templated.

BIKE_COURSE_PROFILES: dict[str, dict] = {
    "karlovy vary": {
        "source":        "2026 WTCS Karlovy Vary Athletes Guide V5 (Lake Rolava to city centre)",
        "loop_km":       4.9,
        "loops":         7,
        "total_km":      40.8,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Sheltered spa-valley setting — wind rarely decisive",
        "surface":        "80% asphalt / 20% pavement (cobbles) — fully closed roads",
        "key_features": [
            "6.5 km approach from Lake Rolava, then 7 technically challenging laps of 4.9 km = 40.8 km total",
            "THE defining feature: roughly 20% of every lap is ascent — the main climb is repeated 7 times and shreds the field",
            "Corrugated terrain, 80% asphalt / 20% pavement — cobbled sections add vibration load; check tyre pressure and bottle security",
            "Lapped athletes are REMOVED from the race — no rule exception here; a bad swim can end the day early",
            "One of the hardest bike courses on the circuit: expect the lead group to shrink substantially rather than a bunch finish",
            "Team wheel station near TA2; neutral wheel station on Nábřeží Osvobození street",
            "Saturday escorted familiarisation covers the approach plus 2 laps, with a regroup at the top of the main climb",
        ],
    },
    "rio de janeiro": {
        "source":        "2026 World Triathlon Cup Rio Athlete's Guide v2.0 (Copacabana)",
        "loop_km":       3.47,
        "loops":         6,
        "total_km":      20.8,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Beachfront — Atlantic sea breeze along Av Atlântica",
        "surface":        "Closed seafront circuit — Avenida Atlântica, Copacabana",
        "key_features": [
            "Fast and flat: 6 counterclockwise laps along the Copacabana seafront (Avenida Atlântica) = 20.8 km",
            "Classic beachfront out-and-back — the same waters/roads that hosted the 2016 Olympic triathlon",
            "Dead-flat, high-speed course built for spectators; expect large lead packs and a bunch finish onto the run",
            "Sea-breeze exposure along the beachfront straights can shape pack dynamics into the turns",
            "NOTE: guide is internally inconsistent — header says 20.8 km / 6 laps, body text says 20 km / 4 × 5 km. Confirm final lap count at briefing.",
        ],
    },
    "london": {
        "source":        "2026 WTCS London Elite Athlete's Guide (ExCeL, Royal Docks)",
        "loop_km":       2.64,
        "loops":         7,
        "total_km":      20.64,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Royal Docks — exposed dockside straights, variable",
        "surface":        "Closed private roads + national highways (contraflow) — Royal Albert Way",
        "key_features": [
            "Flat and fast: 7 × 2.64 km out-and-back laps on Royal Albert Way (+0.75 km bike-out, +1.23 km bike-in) = 20.64 km",
            "Anti-clockwise; short 2.64 km laps mean lapped athletes are allowed to stay on course (WT rule exception)",
            "Dead-flat dockside course — pure power/aero; the race will very likely come down to the run",
            "Out-and-back on a wide highway in contraflow — big lead packs form easily; positioning into the U-turns is key",
            "Course uses private roads/highways only open during official familiarization (Sat 09:00–09:20) — no race-direction preview otherwise",
            "Wheel stations on the Lorry Way / dockside; mark position early on the short laps",
        ],
    },
    "montreal": {
        "source":        "World Triathlon Para Series Montréal Elite Athlete Guide — 27 Jun 2026",
        "loop_km":       4.1,
        "loops":         5,
        "total_km":      21.5,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Open island — variable",
        "surface":        "Closed F1 circuit — Circuit Gilles-Villeneuve, Parc Jean-Drapeau",
        "key_features": [
            "Bike held on the Circuit Gilles-Villeneuve (CGV) F1 track — clockwise 4.1 km loop × 5 = 21.5 km",
            "Pancake-flat, smooth tarmac throughout — no significant climbs; pure power/aero course",
            "Team wheel station at the CGV entrance; Neutral wheel station mid-lap with full disc/rim brake inventory",
            "Bike penalty box sits at the lap/transition split — athletes pass it 5 times across the race",
            "RaceRanger drafting devices are mandatory for PTS/PTVI classes — install at bike racking on Friday if not done at home",
            "Open island setting on Notre-Dame Island — winds off the St. Lawrence can swing crosswinds on the straights",
        ],
    },
    "quiberon": {
        "source":        "World Triathlon Elite Athlete Guide — WTCS Quiberon, 20 Jun 2026",
        "loop_km":       5.5,
        "loops":         4,
        "total_km":      22.0,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Coastal (~20 km/h)",
        "surface":        "Closed coastal/urban circuit — Boulevard René Cassin, Quiberon Peninsula",
        "key_features": [
            "Generally flat profile with no significant climbs — race typically decided on the run, not the bike",
            "Multiple urban changes of direction and turns in town sections — bike-handling premium on lap entries",
            "Coastal sections exposed to Atlantic wind — crosswinds can split the pack hard, especially near Pointe Riberen",
            "Two wheel stations: René Cassin Street (just out of T1) and Fort Neuf Street / Parking Nautique near Pointe Riberen",
            "Slightly longer than standard sprint — 22 km total over 4 × 5.5 km laps",
        ],
    },
    "huatulco": {
        "source":        "asdeporte 2026 Elite Cycling course map",
        "loop_km":       5.0,
        "loops":         4,
        "total_km":      20.0,
        "gain_per_lap_m":  84.4,
        "loss_per_lap_m":  84.6,
        "max_grade_pos":   18.4,
        "max_grade_neg":  -14.9,
        "avg_grade_pct":   3.2,
        "surface":        "Closed urban circuit (Bahía de Santa Cruz → Vialidad 5 → return)",
        "key_features": [
            "Out-and-back along Blvd Santa Cruz with a tight technical hairpin at the neutral wheel station",
            "Single sustained climb up Vialidad 5 — short but punchy (max 18.4%); strong riders can split the field every lap",
            "Fast technical descent back into Bahía de Santa Cruz; max -14.9% grade rewards confident bike-handling",
            "Tight start/finish chicane through transition — drafting packs reshuffle every lap on the climb and descent",
        ],
    },
    "hamburg": {
        "source":        "World Triathlon race info — 2026 WTCS Hamburg",
        "loop_km":       3.3,
        "loops":         6,
        "total_km":      19.8,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "City-center setting — generally sheltered, variable",
        "surface":        "Closed urban circuit — Rathausmarkt, Alsterufer, Jungfernstieg",
        "key_features": [
            "6 × 3.3 km laps (19.8 km total) around Hamburg's city center — one of the most technical WTCS bike courses",
            "Pancake-flat throughout, but high-frequency turns through Rathausmarkt and along the Alster reward bike-handling over raw power",
            "Drafting packs form easily on the long straights; Hamburg historically produces large lead packs of 30+ athletes",
            "Cobblestone sections near the Rathaus and Jungfernstieg add vibration load — check tire pressure pre-race",
            "Wheel station on each lap — mark position early; packs do not slow for late changes",
        ],
    },
    "edmonton": {
        "source":        "2026 World Triathlon Cup Edmonton Athlete Guide v2 (Do North Events) — July 18-19, 2026",
        "loop_km":       6.8,
        "loops":         3,
        "total_km":      20.4,
        "gain_per_lap_m":  85.0,
        "loss_per_lap_m":  87.0,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "River valley setting — generally sheltered, variable on the Saskatchewan Dr exposure",
        "surface":        "Closed road circuit — Groat Rd, Saskatchewan Dr NW, William Hawrelak Park Rd",
        "key_features": [
            "3 × 6.8 km laps (20.4 km total) out of Hawrelak Park up Groat Rd and along Saskatchewan Dr NW",
            "One sustained climb/descent per lap out of the river valley (+85 m / -87 m per lap per the athlete guide profile) — first true elevation test on this venue since 2021",
            "Course map lists overall elevation gain as 76.5 m — legs will feel the repeated valley climb across all 3 laps",
            "Team Wheel station just outside transition (both-side access); Neutral Wheel station at the loop's south end on Saskatchewan Dr, before the right turn back toward the park",
            "Neutral wheels: 700c front/rear, 11/12-speed, 140/160mm rotor options, 28mm tires at 65/70 psi (front/rear)",
            "Newly renovated venue ahead of the 2027 World Triathlon Multisport Championships — this is the first elite race on the rebuilt course",
        ],
    },
    "asuncion": {
        "source":        "2026 World Triathlon Cup Asunción Athlete's Guide + official cycling course flyer — 9 Aug 2026",
        "loop_km":       6.66,
        "loops":         6,
        "total_km":      39.96,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Open riverside / floodplain terrain — little shelter, wind is the main variable on a flat course",
        "surface":        "Closed avenue on La Isla del Delta — smooth tarmac, traffic-free",
        "key_features": [
            "Standard-distance bike: 6 × 6.66 km laps = 39.96 km on a closed main avenue through the Delta development",
            "Flat throughout with no published elevation — this is a pure power and positioning course, not a selection course",
            "Long out-and-back shape with the Zona Operativa (transition) mid-course; few technical features beyond the end turnarounds",
            "Flat + 6 laps + a 40 km distance strongly favours large packs forming and staying together to T2",
            "Single wheel station on course (guide does not specify neutral wheel inventory) — confirm placement at Saturday's briefing",
            "Open floodplain with minimal tree cover: crosswind exposure on the long straights is the realistic pack-splitting mechanism",
        ],
    },
    "weihai": {
        "source":        "2026 WTCS Weihai Athletes' Guide (updated 22 July 2026) + WT race info",
        "loop_km":       5.0,
        "loops":         8,
        "total_km":      40.0,
        "gain_per_lap_m":  None,
        "loss_per_lap_m":  None,
        "max_grade_pos":   None,
        "max_grade_neg":   None,
        "avg_grade_pct":   None,
        "wind":           "Coastal — guide gives August max average 5 m/s (~18 km/h), prevailing SE",
        "surface":        "Closed city roads — Haibin North Rd, Yuhua Rd turnaround, Dongshan Rd",
        "key_features": [
            "8 × 5 km laps = 40 km, the highest lap count of any 2026 preview venue — a turnaround every lap means 16+ hard decelerations and re-accelerations",
            "Route: transition at the Wave to Wonderland running track, north on Haibin North Rd to the Yuhua Rd turnaround, left onto Dongshan Rd, back down Haibin North Rd",
            "Lap-heavy, turnaround-heavy design punishes riders who sit at the back of the pack — repeated surges out of each turn are where the accordion effect bites",
            "Two wheel stations: team wheels at the transition exit, neutral wheels at the bike turnaround. Neutral stock covers rim, 140 mm and 160 mm disc, 11- and 12-speed",
            "Spare wheels must be registered at Registration on race day and labelled with athlete name, race number and country",
            "Coastal seafront exposure on Haibin North Rd — the guide's prevailing SE wind will be a cross/head component on one leg of every lap",
        ],
    },
}

SWIM_COURSE_PROFILES: dict[str, dict] = {
    "karlovy vary": {
        "source":            "2026 WTCS Karlovy Vary Athletes Guide V5 (Lake Rolava)",
        "total_km":          1.5,
        "laps":              2,
        "loop_km":           0.75,
        "layout":            "Lake Rolava — 2 anti-clockwise laps, Australian exit between laps",
        "format":            "Freshwater lake — sheltered, island-lined course",
        "start_type":        "Pontoon start (0.6 m high, 75 cm slots); first turn buoy at 205 m",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (18.0, 20.0),
        "wetsuit_note":      "History says pack both: wetsuit ALLOWED in 2025 (18.0 °C), 2023 (19.0 °C) and 2022 (19.1 °C), but FORBIDDEN in 2024 when the lake hit 23–25 °C. Ruling confirmed at the Sat 12 Sep briefing.",
        "key_features": [
            "2 × 750 m anti-clockwise laps in Lake Rolava — pontoon start, hard 205 m opening leg to the first turn buoy",
            "Australian exit: run the wooden pier at the end of lap 1 and dive back in — rehearse the exit/re-entry at Saturday familiarisation",
            "Navigation-heavy: swim behind the big island, under the bridge, left between the big and small islands, then right to the pier",
            "Sheltered lake water — no current or surf, but the island turns make sighting and pack position matter",
            "Wetsuit ruling has flipped year to year (18–19 °C most editions, 23–25 °C in 2024) — bring both setups",
            "Swim exits into TA1 at Lake Rolava; last-minute gear is transported by the LOC to the finish area in the city",
        ],
        "missing": [
            "Confirmed race-day water temperature + wetsuit ruling (announced at the briefing)",
        ],
    },
    "rio de janeiro": {
        "source":            "2026 World Triathlon Cup Rio Athlete's Guide v2.0 (Copacabana)",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "Copacabana Beach — single counterclockwise ocean loop, buoy-marked",
        "format":            "Open ocean — Atlantic surf, beach start",
        "start_type":        "Beach start; water entry, exit at same point into transition",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (22.0, 24.0),
        "wetsuit_note":      "Guide states August water 22–24 °C — non-wetsuit (WT forbids wetsuits > 22 °C). Temp confirmed at race briefing.",
        "key_features": [
            "Single 750 m counterclockwise ocean loop off Copacabana Beach — the 2016 Olympic swim venue",
            "Beach start: sprint down the sand and dolphin-dive through Atlantic shore break — surf entry is a real skill factor",
            "Ocean swell and chop possible; sighting on bigger surface waves than a sheltered dock/lake course",
            "Warm water (22–24 °C) = non-wetsuit; no cold-water concern, but manage a hot, sunny morning start",
            "Exit at the same point as entry, running straight into transition on the sand",
        ],
        "missing": [
            "Confirmed race-day water temperature + surf/swell forecast (announced at briefing)",
        ],
    },
    "london": {
        "source":            "2026 WTCS London Elite Athlete's Guide (Royal Docks, ExCeL)",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "Royal Victoria Dock — sheltered dockside, 315 m to first buoy",
        "format":            "Enclosed freshwater dock — no current, urban dock water",
        "start_type":        "Pontoon start, anti-clockwise",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (18.0, 21.0),
        "wetsuit_note":      "Late-July London dock water typically ~18–21 °C — borderline; wetsuit likely optional (mandatory < 16 °C, forbidden > 22 °C). Confirm at the pre-race water-temp reading.",
        "key_features": [
            "Single 750 m anti-clockwise loop in the Royal Victoria Dock, off the ExCeL dockside",
            "Long 315 m leg to the first buoy — a hard, honest opening straight that stretches the field before any turn",
            "Enclosed dock = flat water, no current or chop; a pure swim-fitness course with no navigation tricks",
            "Pontoon start; last-minute gear bagged at the start and ferried to the athlete lounge",
            "Dock is closed to swimming outside the Friday 14:00–14:30 familiarization window",
            "Afternoon starts (W 14:30 / M 16:15) — water will be at its warmest of the day",
        ],
        "missing": [
            "Confirmed race-day water temperature (published at the pre-race reading)",
        ],
    },
    "montreal": {
        "source":            "World Triathlon Para Series Montréal Elite Athlete Guide — 27 Jun 2026",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "Olympic Basin — Parc Jean-Drapeau (1976 Olympics rowing venue)",
        "format":            "Sheltered rowing basin — fresh water, calm conditions",
        "start_type":        "In-water start from pontoon",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (22.0, 24.0),
        "wetsuit_note":      "Guide states June water ~23 °C — under WT rules wetsuit is forbidden > 22 °C; expect non-wetsuit race",
        "key_features": [
            "Sheltered Olympic Basin water — no chop, no current, no wave issues; navigation-only swim",
            "PTWC swim warm-up on course 06:15–06:45; PTS/PTVI use warm-up lane 07:00–08:15",
            "Lounge opens 05:30 (PTWC) / 06:15 (PTS/PTVI); athletes' intro 10 min before category start",
            "Water quality monitored by private lab — guide flags it consistently inside WT thresholds",
            "Warm water + warm air at 7 AM start — no cold-water acclimation concerns",
        ],
        "missing": [
            "Buoy layout for individual race (start gathers on pontoon — full layout TBD at familiarization)",
        ],
    },
    "quiberon": {
        "source":            "World Triathlon Elite Athlete Guide — WTCS Quiberon, 20 Jun 2026",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "Atlantic Ocean — beach loop off Boulevard René Cassin",
        "format":            "Open ocean — Atlantic / Bay of Quiberon",
        "start_type":        "Beach start (rue René Cassin, beach access from Place du Doued)",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (14.9, 19.3),
        "wetsuit_note":      "Wetsuit territory — guide states mid-June Atlantic water 14.9–19.3 °C; under WT rules wetsuit is mandatory < 16 °C and optional 16–20 °C",
        "key_features": [
            "Single 750 m loop in the open Atlantic off the Quiberon peninsula — cold water by elite standards",
            "Beach start running entry from rue René Cassin — typical pack-formation chaos through the first 200 m",
            "Atlantic exposure brings chop and swell — sighting on bigger surface waves than typical WTCS venues",
            "Cold-water acclimation matters: 14.9–19.3 °C range means the first 100 m breath-control is a real risk",
            "Wetsuit decision driven by morning measurement — bring both options; men race at 10:00, women at 12:00",
            "Water quality testing site \"Lombardsbrücke\" referenced in guide (text appears templated from Hamburg) — actual Quiberon measuring points TBD pre-race",
        ],
        "missing": [
            "Buoy configuration and turn structure (not in race brief)",
            "Day-of water temperature (Atlantic SST will be filled from Open-Meteo Marine forecast)",
        ],
    },
    "huatulco": {
        "source":            "asdeporte 2026 Elite Swim course map (Natación Elite — Bahía de Santa Cruz)",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "U-shape rectangle — 300 m out, 150 m across, 300 m back",
        "format":            "Ocean — Bahía de Santa Cruz (protected marina)",
        "start_type":        "Beach start (running entry) at Plaza Santa Cruz",
        "water_temp_c":      None,
        "wetsuit_note":      "Non-wetsuit likely — June water temps in Bahía de Santa Cruz typically run 27–29 °C, above the 22 °C threshold",
        "key_features": [
            "Single 750 m loop — 300 m east leg, 150 m across south, 300 m west leg, beach exit by transition",
            "Tight rectangle inside the protected marina — minimal chop, current normally low at 06:30 start",
            "Two turn buoys at the south end — Chief Swim films from outside; angled kayak inside prevents course-cutting",
            "Sun rising due east at 06:30 — bright glare on the outbound 300 m east leg",
            "Extraction zone marked at north dock; short beach run into T1",
        ],
        "missing": [],
    },
    "hamburg": {
        "source":            "World Triathlon race info — 2026 WTCS Hamburg",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "Binnenalster / Kleine Alster — Hamburg city center",
        "format":            "Sheltered urban lake — fresh water, calm conditions",
        "start_type":        "Pontoon start",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (18.0, 22.0),
        "wetsuit_note":      "July Alster water typically 18–22 °C — borderline wetsuit territory; expect non-wetsuit or optional depending on morning measurement",
        "key_features": [
            "Single 750 m loop in the sheltered Binnenalster / Kleine Alster — flat, calm, navigation-friendly",
            "Pontoon start into the Alster canal — fast, mass-start entry with good visibility",
            "Fresh water lake in city center — no ocean chop, swell, or current; pure positioning swim",
            "July morning water mid-to-high teens Celsius — bring both wetsuit options; decision made day-of per WT measurement",
            "Short run into T1 from the Alster exit ramp — transition area directly adjacent to city hall",
        ],
        "missing": [],
    },
    "edmonton": {
        "source":            "2026 World Triathlon Cup Edmonton Athlete Guide v2 (Do North Events) — July 18-19, 2026",
        "total_km":          0.75,
        "laps":              1,
        "loop_km":           0.75,
        "layout":            "Hawrelak Park Lake — single clockwise loop",
        "format":            "Sheltered park lake — fresh water, calm conditions",
        "start_type":        "In-water start off the swim start platform",
        "water_temp_c":      19.0,
        "expected_water_temp_range_c": (17.0, 21.0),
        "wetsuit_note":      "Guide lists average July water ~19 °C (66 °F) — under WT rules (mandatory <16 °C, optional 16–20 °C) this sits in optional-wetsuit territory; final call posted at the Friday pre-race briefing",
        "key_features": [
            "Single 750 m clockwise loop on Hawrelak Park's lake — calm, sheltered, no current or chop",
            "First (yellow) buoy passed on the left shoulder at 200 m; all subsequent orange buoys passed on the right shoulder",
            "Run to T1 is 220 m from the swim exit — longer transition run than most sprint venues",
            "Swim familiarization Friday July 17, 14:00-15:00 only — no other in-water access permitted for safety",
            "Elite Men warm-up 13:00-13:45, Elite Women warm-up 15:00-15:45 on race day",
        ],
        "missing": [],
    },
    "asuncion": {
        "source":            "2026 World Triathlon Cup Asunción Athlete's Guide + official swimming course flyer — 9 Aug 2026",
        "total_km":          1.5,
        "laps":              2,
        "loop_km":           0.75,
        "layout":            "Delta lagoon — 2 × 750 m loops, start off the floating Candock platforms",
        "format":            "Sheltered inland lagoon — fresh water, calm",
        "start_type":        "Floating pontoon (Candock) start",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (18.0, 22.0),
        "wetsuit_note":      "GENUINELY UNCERTAIN — guide forecasts 18–22 °C, which straddles the elite threshold (wetsuit forbidden at 20.0 °C and above, optional 16.0–19.9 °C). Pack both and expect the call at the official race-morning measurement",
        "key_features": [
            "Standard-distance swim: 2 × 750 m loops in the Delta lagoon, not the Paraguay River — sheltered, no current",
            "In-water start from floating Candock platforms rather than a beach or pontoon run-in",
            "Two laps means an Australian-style exit or turn per lap — sighting and re-acceleration off each turn matters more than a single-loop sprint swim",
            "Wetsuit call is a live tactical variable: 18–22 °C forecast sits right on the 20 °C elite cut-off",
            "Women's 07:30 start swims the coldest water of the day; men at 10:30 get three more hours of warming, so the two races will not necessarily get the same wetsuit ruling",
            "Swim familiarization is Saturday 10:45–11:30 only — the venue is closed to athletes on all other days",
        ],
        "missing": [
            "Buoy count and turn configuration (not published in the guide)",
            "Measured water quality data (not published; lagoon is a managed development water body)",
        ],
    },
    "weihai": {
        "source":            "2026 WTCS Weihai Athletes' Guide (updated 22 July 2026)",
        "total_km":          1.5,
        "laps":              2,
        "loop_km":           0.75,
        "layout":            "Yellow Sea — 2 counter-clockwise laps off the south side of Wave to Wonderland",
        "format":            "Open sea bay — salt water, seafront start",
        "start_type":        "Waterfront start off the south shore of Wave to Wonderland",
        "water_temp_c":      None,
        "expected_water_temp_range_c": (24.0, 26.0),
        "wetsuit_note":      "NON-WETSUIT, effectively certain. Water measured 24.2–25.9 °C in every recorded edition (2017, 2023, 2024, 2025) and wetsuits were forbidden all four times. Do not pack a race wetsuit expecting to use it",
        "key_features": [
            "Standard-distance swim: 2 × 750 m counter-clockwise laps in the Yellow Sea off Wave to Wonderland",
            "Salt water and open bay — sighting, swell and chop matter far more here than at the sheltered lagoon venues elsewhere on the calendar",
            "Water has run 24–26 °C in every measured edition, so this is a reliably warm, reliably non-wetsuit swim — overheating in the swim is a bigger risk than cold",
            "Water quality sampled 27 May 2026 at three points: pH 8.12–8.48, E. coli and enterococcus both undetected in all three samples",
            "Two laps means a turn/exit sequence mid-swim — position off the first lap decides which bike pack you land in",
            "Swim familiarization is Friday 28 Aug 10:00–11:00 only; swimming the race course at any other time is prohibited",
        ],
        "missing": [
            "Buoy count and exact turn configuration (not published in the guide)",
            "Tide and current data for race morning (open-bay venue; not published)",
        ],
    },
}

RUN_COURSE_PROFILES: dict[str, dict] = {
    "karlovy vary": {
        "source":            "2026 WTCS Karlovy Vary Athletes Guide V5 (city centre)",
        "total_km":          10.0,
        "laps":              4,
        "loop_km":           2.5,
        "surface":           "Mixed asphalt and pavement — spa-town centre streets",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "LOW",
        "key_features": [
            "4 laps of 2.5 km (10 km) through the Karlovy Vary spa-town centre — out of TA2 at Theatre Square",
            "Mixed asphalt and pavement surface — footing changes underfoot; racing flats over ultra-light spikes",
            "Two aid stations per lap: first ~160 m out of TA2, second at the far turnaround on Zahradní street",
            "Penalty box on Stará Louka street close to TA2 — passed every lap, so a drafting call is served fast",
            "Mid-September Czech conditions: mild and often cool (historic race-day air 14–27 °C) — low heat stress",
            "After a brutally selective bike, the run is usually contested by a much-reduced front group",
        ],
        "missing": [
            "Elevation profile (city-centre laps — no published grade data)",
        ],
    },
    "rio de janeiro": {
        "source":            "2026 World Triathlon Cup Rio Athlete's Guide v2.0 (Copacabana)",
        "total_km":          5.0,
        "laps":              2,
        "loop_km":           2.5,
        "surface":           "Paved seafront promenade — Avenida Atlântica, Copacabana",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "MODERATE",
        "key_features": [
            "Flat and fast: 2 counterclockwise laps of 2.5 km along the Copacabana seafront (5 km total)",
            "Iconic promenade course — the black-and-white Portuguese-wave pavement of Avenida Atlântica",
            "Aid stations positioned throughout each lap; flat profile favours aggressive, even-paced running",
            "August is Rio winter but sun on the exposed beachfront can push perceived heat up — 09:00 men's start runs into rising sun",
            "Race decided here: with a flat swim and flat bike, the run is where the World Cup is won",
        ],
        "missing": [
            "Elevation profile (course is flat — no published grade data)",
        ],
    },
    "london": {
        "source":            "2026 WTCS London Elite Athlete's Guide (ExCeL, Royal Docks)",
        "total_km":          4.92,
        "laps":              3,
        "loop_km":           1.64,
        "surface":           "Mixed indoor/outdoor — through the ExCeL venue hub each lap",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "LOW",
        "key_features": [
            "3 anti-clockwise laps (1.76 + 1.76 + 1.34 km = 4.92 km) routed through the main ExCeL venue hub each lap",
            "Flat and fast; the run runs partly indoors through the venue — crowd-loud and atmospheric, a signature London feature",
            "2 aid stations per lap (sealed water only); litter zones 20 m before / 80 m after each",
            "Penalty box sits just ahead of the finish chute on the right — easy to see, no excuse for a missed stand-down",
            "Lap counter at transition exit; finish chute is a left-hand turn on lap 3 (final 62 m)",
            "Cool UK summer conditions + partly-indoor route = negligible heat stress; expect very fast run splits to decide the race",
        ],
        "missing": [
            "Elevation profile (course is flat — no published grade data)",
        ],
    },
    "montreal": {
        "source":            "World Triathlon Para Series Montréal Elite Athlete Guide — 27 Jun 2026",
        "total_km":          5.0,
        "laps":              2,
        "loop_km":           2.5,
        "surface":           "Asphalt — Circuit Gilles-Villeneuve + Olympic Basin path",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "MODERATE",
        "key_features": [
            "Flat 2.5 km loop × 2 — exit transition north of basin, enter CGV at the hairpin, head west to pit lane",
            "Two left turns at the pit lane drop onto the north-side path along the Olympic Basin",
            "Aid stations / water at 370 m, 1.1 km, 2.3 km, 3.0 km, 4.1 km — five access points per athlete",
            "Run penalty box at ~2.6 km (loop turnaround); passed twice — drafting infractions compound across laps",
            "PTVI Free Leading zones: T2 exit→AS1, AS2 on CGV, paddocks turnaround, AS2 on Chemin Nord, Run Penalty Box turn",
            "Run team wheel station (for racing wheelchairs) co-located with Aid Station 2; both-side access",
        ],
        "missing": [
            "Elevation profile (course described as flat — no published grade data)",
        ],
    },
    "quiberon": {
        "source":            "World Triathlon Elite Athlete Guide — WTCS Quiberon, 20 Jun 2026",
        "total_km":          5.0,
        "laps":              2,
        "loop_km":           2.5,
        "surface":           "Asphalt out-and-back; small stabilized (gravel) section per lap",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "key_features": [
            "Flat out-and-back asphalt course with one short stabilized-surface section per lap",
            "Aid stations at 0.2 km (just after T2) and 1.25 km — pass each twice across the 2 laps",
            "Penalty Box located on the left, ~50 m before the Transition Area — easy to miss if drafting flagged",
            "Course measurement for coaches Friday 19:00 from the Finish Area (Bd René Cassin) — register at Friday package distribution",
            "Cool Atlantic conditions favour aggressive pacing; June daily high ~23 °C, low ~13 °C — no heat-stress concern",
            "Lap structure: 2 × 2.5 km out-and-back (inferred from aid station spacing in the brief)",
        ],
        "missing": [
            "Elevation profile and grade data (course described as \"flat\" — no numbers published)",
            "Exact turnaround coordinates (route map TBD in the guide PDF)",
        ],
    },
    "huatulco": {
        "source":            "asdeporte 2026 Elite Run course map (Carrera Elite — Boulevard Santa Cruz)",
        "total_km":          5.0,
        "laps":              2,
        "loop_km":           2.5,
        "surface":           "Paved urban — Blvd Santa Cruz / Camino a Santa Cruz",
        "gain_per_lap_m":    55.3,
        "loss_per_lap_m":    55.0,
        "max_grade_pos":     17.3,
        "max_grade_neg":    -15.3,
        "avg_grade_pct":     4.2,
        "heat_risk":         "HIGH",
        "elevation_range_m": (6, 32),
        "key_features": [
            "Hilly 2.5 km out-and-back loop — +55.3 m gain per lap (≈110 m total over 5 km)",
            "Max climb 17.3% on Camino a Santa Cruz — short steep pinch that bites harder on lap 2",
            "Steep -15.3% descent back into Bahía de Santa Cruz — light braking + high cadence preserves the legs",
            "08:30 men's start runs into peak surface temps with limited shade on the climb",
            "Aid stations at 1 km and 2 km plus start/finish; penalty box at start/finish — drafting tickets compound across both laps",
            "Course turnaround at the 2 km mark off Blvd Chahue — sharp 180° on a slight uphill",
        ],
        "missing": [],
    },
    "hamburg": {
        "source":            "World Triathlon race info — 2026 WTCS Hamburg",
        "total_km":          10.0,
        "laps":              2,
        "loop_km":           5.0,
        "surface":           "Closed urban — Rathausmarkt, Alsterarkaden, city hall precinct",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     0.0,
        "heat_risk":         "LOW–MODERATE",
        "key_features": [
            "2 × 5 km laps through Hamburg's historic city center — flat, fast, spectator-dense course",
            "Cobblestone sections on Rathausmarkt and adjacent squares add vibration; stay relaxed in the upper body",
            "Wide avenues on most of the loop allow for tactical racing and late-race surges",
            "13:15 / 15:00 summer starts — July Hamburg averages ~22 °C, manageable but warm; aid stations every lap",
            "Penalty box positioned at the finish area — drafting tickets from the bike carry into the run",
        ],
        "missing": [],
    },
    "edmonton": {
        "source":            "2026 World Triathlon Cup Edmonton Athlete Guide v2 (Do North Events) — July 18-19, 2026",
        "total_km":          5.0,
        "laps":              3,
        "loop_km":           1.66,
        "surface":           "Paved park paths — loop sits next to the transition area",
        "gain_per_lap_m":    17.0,
        "loss_per_lap_m":    20.0,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "LOW–MODERATE",
        "key_features": [
            "3 × 1.66 km laps (5 km total), counter-clockwise, run next to the transition area by design to avoid a first-runner/last-biker clash",
            "Aid station passed 6 times total (out and back per lap) — sealed water bottles at 400 m, 850 m, 2000 m, 2450 m, 3600 m, 4050 m, served on the athlete's right",
            "Littering zones marked green (start) / red (end) at each aid and wheel station — time penalties apply outside them",
            "Run penalty box on the left side of the road, 200 m before the finish — passed every lap",
            "July averages a 23 °C high / 12 °C low with a late sunset (21:32) — moderate heat risk for the 14:00/16:00 Elite starts, well short of Huatulco-level heat stress",
        ],
        "missing": [],
    },
    "asuncion": {
        "source":            "2026 World Triathlon Cup Asunción Athlete's Guide + official running course flyer — 9 Aug 2026",
        "total_km":          10.0,
        "laps":              4,
        "loop_km":           2.35,
        "surface":           "Paved closed avenue alongside the lagoon and river frontage",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "LOW — August is Paraguayan winter (see Race-Day Weather slide)",
        "key_features": [
            "Standard-distance run: 4 × 2.35 km laps plus a final 600 m spur = 10 km, described as flat",
            "Long straight out-and-back shape along the avenue — very few turns, so pace is honest and gaps are visible from a long way back",
            "Penalty box sits near transition and is passed every lap — bike drafting tickets are unusually costly over 4 laps",
            "Flat, non-technical and fast: on this profile the race is most likely decided by run legs off a large intact bike pack",
            "August is winter here: both races run in the cool morning and finish by ~12:20, well before the daily maximum. Heat is not the limiter this race normally is in a Cup",
            "Climatology puts the men's 10:30 start a few degrees warmer than the women's 07:30, but the gap varies year to year — the Race-Day Weather slide carries the live forecast and governs",
        ],
        "missing": [
            "Aid station count and placement (not specified in the guide)",
            "Elevation profile (course stated as flat, no published grade data)",
        ],
    },
    "weihai": {
        "source":            "2026 WTCS Weihai Athletes' Guide (updated 22 July 2026)",
        "total_km":          10.0,
        "laps":              4,
        "loop_km":           2.5,
        "surface":           "Closed seafront road — Haibin North Rd out-and-back",
        "gain_per_lap_m":    None,
        "loss_per_lap_m":    None,
        "max_grade_pos":     None,
        "max_grade_neg":     None,
        "avg_grade_pct":     None,
        "heat_risk":         "MODERATE (men 09:00) / HIGHER for women (12:00 sun load)",
        "key_features": [
            "4 × 2.5 km laps = 10 km: north up Haibin North Rd to a turnaround, u-turn, return to transition, then into the finish chute at the main venue",
            "Pure out-and-back on a seafront road — no technical features, so this is an honest 10 km where run legs are fully exposed",
            "Three aid stations with sealed bottled water AND iced sponges — the sponges matter, this is the hottest run on the current preview set",
            "Run penalty box sits near the dismount line and is passed every lap; bike drafting tickets are expensive over 4 laps",
            "HEAT IS NOT SYMMETRIC: WT's recorded race-day air temp has been 3.9–4.7 °C higher for the women in each of the last three editions (2023, 2024, 2025)",
            "Ambient climatology only warms ~1 °C from 09:00 to 12:00 here (coastal sea breeze), so most of that recorded gap is midday SUN and on-course radiant load, not air temperature. Shade is nil on this seafront out-and-back",
            "2026 moves to 29 August, ~4 weeks earlier than the 2024/2025 editions. The closest analogue is 2023 (26 Aug), the hottest on record here: 27.2 °C men / 31.9 °C women",
        ],
        "missing": [
            "Elevation profile (not published; seafront road, assumed flat to gently rolling)",
            "Exact aid-station distances (guide confirms three stations but not their placement)",
        ],
    },
}

# ── Travel & arrival data (per venue, origin = Denver, CO USAT HQ) ───────────
TRAVEL_PROFILES: dict[str, dict] = {
    "karlovy vary": {
        "origin":          "Denver, CO to Karlovy Vary race week",
        "core_read": (
            "Fly Denver to Prague (PRG) — a one-stop overnight via a European hub (LHR, FRA, "
            "AMS, CDG) or an East Coast gateway; roughly 13–15 h door-to-door and a +8 h shift "
            "(MDT to CEST), the biggest jet-lag load of the late-season block. From Prague it is a "
            "~2 h ground transfer to Karlovy Vary (~125 km). The LOC provides a FREE bus transfer "
            "for elite athletes and coaches (bike included) — but you must email flight details to "
            "transfer@citytriathlon.cz before 31 Aug 2026 or it becomes 20 EUR per person each way. "
            "Arrive by Thursday to absorb the time shift before Saturday familiarisation."
        ),
        "stats": [
            ("Primary route",    "DEN → PRG",        "1-stop, overnight"),
            ("Flight time",      "~13–15 h",         "via LHR / FRA / AMS"),
            ("Time zones",       "+8 h",             "Czechia ahead of Denver"),
            ("Airport Transfer", "PRG → Karlovy Vary", "Free LOC bus (register!)"),
            ("Transfer Distance","~125 km / ~2 h",   "Coach transfer"),
        ],
        "hotel_title": "Accommodation / Venue Logistics Read",
        "hotel_bullets": [
            "Briefing is at the Spa Hotel Thermal (Sat 12 Sep, 16:00–16:30) — staying central keeps race week simple.",
            "FREE LOC bus PRG to Karlovy Vary: email arrival time + flight number to transfer@citytriathlon.cz before 31 Aug 2026.",
            "Private LOC transfer if you miss the deadline: 140 EUR for 1–3 people, 180 EUR for 4–7. Extra escort staff 20 EUR per person per journey.",
            "Race-weekend shuttle runs between Divadelní náměstí (T2) and Lake Rolava (T1) — about 20 min, with a funicular link at the Imperial stop.",
            "Two-transition race: plan bag logistics carefully — TA2 checks in first (city), then TA1 at the lake.",
            "Free daily 2 h pool access at Mattoni Aréna and run track at AC Start stadium — bring photo ID.",
        ],
        "food_title": "Food / Grocery Options",
        "food_bullets": [
            "Karlovy Vary is a compact spa town — supermarkets (Billa, Albert, Lidl, Tesco) are an easy walk or short bus from the centre.",
            "Tap water is safe throughout Czechia; the town's famous mineral springs are drinkable but heavily mineralised — not a race-week hydration plan.",
            "Czech restaurant staples run heavy (dumplings, pork, sauces) — identify a couple of lighter pasta-and-rice options early in the week.",
            "Athletes Lounge provides water, fruit and energy bars before and after the race at both the start and finish areas.",
            "Split start times (W 10:00 / M 15:00) mean very different fuelling timelines — men face a full morning of waiting; plan lunch deliberately.",
        ],
    },
    "rio de janeiro": {
        "origin":          "Denver, CO → Rio de Janeiro race week",
        "core_read": (
            "Denver → Rio (GIG) is an overnight, one-stop haul — typically via ATL, IAH, MIA, or "
            "JFK/EWR on United/American/LATAM, ~13–15 h door-to-door. Only +4 h time-zone shift "
            "(MDT → BRT), so jet-lag is moderate — the bigger factor is the long overnight leg. The "
            "LOC runs a FREE airport shuttle GIG → Othon Palace Copacabana (bikes included) on "
            "July 29–31 — registration is mandatory or you self-arrange (~BRL 80–150 by Uber/99). "
            "Land by ~July 30 to settle before the Aug 1 briefing and Aug 2 race."
        ),
        "stats": [
            ("Primary route",    "DEN → GIG",        "1-stop, overnight"),
            ("Flight time",      "~13–15 h",         "via ATL / IAH / MIA"),
            ("Time zones",       "+4 h",             "Rio ahead of Denver"),
            ("Airport Transfer", "GIG → Othon Palace","Free LOC shuttle (register)"),
            ("Transfer Distance","~20 km / 30–60 min","Traffic-dependent"),
        ],
        "hotel_title": "Host Hotel / Accommodation Read",
        "hotel_bullets": [
            "Official host hotel: Othon Palace Copacabana, Av Atlântica 3264 — directly facing Copacabana Beach.",
            "Walking distance to the Triathlon Arena (central median of Av Atlântica, ~height 3806) — no race-morning commute.",
            "Free LOC airport shuttle GIG ↔ Othon (bikes included) on Jul 29–31 / Aug 2–3 — mandatory registration via the LOC form.",
            "Copacabana is dense and walkable; the beachfront promenade is the venue, training ground, and social hub in one.",
            "Brazil entry: check current US-passport visa/ETA status well before travel — requirements have shifted recently.",
        ],
        "food_title": "Food / Grocery Options Near Venue",
        "food_bullets": [
            "Grocery: Zona Sul and Pão de Açúcar supermarkets along/near Av Nossa Senhora de Copacabana for staples, water, fruit, snacks.",
            "Bottled/filtered water only — do not drink tap water in Rio; stock sealed bottles for race week to avoid GI risk.",
            "Athlete-friendly: açaí bowls, grilled meats/fish, rice-and-beans, and juice bars are everywhere — easy carb + protein, but be selective with raw/street food pre-race.",
            "The athletes' lounge provides sealed water + fruit before/after the race; bring your own race-morning fueling from the hotel.",
            "Sun + exposed beachfront: pre-stock electrolytes and high-SPF sunscreen for daily training on the sand.",
        ],
    },
    "london": {
        "origin":          "Denver, CO → London race week",
        "core_read": (
            "Denver → London is a single long-haul leg — British Airways runs a DEN → LHR nonstop, "
            "plus many one-stop options via ORD / EWR / JFK into LHR or LGW. Overnight eastbound, "
            "~9 h flight + 7 h time-zone shift (MDT → BST). Land at least 3–4 days early to adapt to "
            "the +7 h shift AND the unusual afternoon race times (W 14:30 / M 16:15) — body-clock "
            "management is a real factor here. From Heathrow, the ExCeL is ~1 h by taxi or the "
            "Elizabeth Line (change to DLR at Custom House). No LOC airport shuttle — self-arrange."
        ),
        "stats": [
            ("Primary route",    "DEN → LHR",        "BA nonstop / 1-stop"),
            ("Flight time",      "~9 h",             "Overnight eastbound"),
            ("Time zones",       "+7 h",             "London ahead of Denver"),
            ("Airport → Venue",  "LHR → ExCeL",      "Elizabeth Line + DLR"),
            ("Transfer Distance","~40 km / 60–75 min","Custom House stn = venue"),
        ],
        "hotel_title": "Host Hotel / Accommodation Read",
        "hotel_bullets": [
            "Venue is the ExCeL London (Royal Docks) — athlete access via Door S5, South Halls.",
            "On-site / adjacent hotels: Aloft, Novotel, Sunborn Yacht, CopperBox and Premier Inn ExCeL — walkable to transition, ideal for an afternoon race (no early commute).",
            "Custom House (Elizabeth Line + DLR) is the venue station — anything on those lines is a fast, low-stress commute.",
            "Late-July London is peak season — book early; the Docklands fills up around ExCeL events.",
            "UK entry: US athletes need an ETA (Electronic Travel Authorisation) before travel — apply well in advance.",
        ],
        "food_title": "Food / Grocery Options Near Venue",
        "food_bullets": [
            "Grocery: Waitrose / Tesco / Sainsbury's around Canary Wharf (2 DLR stops) for staples, water, and room food.",
            "On-site: the ExCeL boulevard has cafés and chains, but options thin outside event hours — stock the room.",
            "Athlete-friendly: Canary Wharf (Crossrail Place / Jubilee Place) has extensive protein-forward and quick-serve options a short DLR hop away.",
            "Afternoon-race fueling note: with a 14:30 / 16:15 start, plan a full breakfast + measured lunch — this is not a dawn race; the pre-race timeline is longer than usual.",
            "Water is tap-safe throughout London.",
        ],
    },
    "montreal": {
        "origin":          "Denver, CO → Montréal race week",
        "core_read": (
            "Denver → Montréal is a same-day continental hop — multiple nonstop options (Air Canada / "
            "United / WestJet seasonal) plus easy one-stop itineraries through ORD, YYZ, or EWR. "
            "Only +2 h time-zone shift (MDT → EDT), so jet-lag impact is minimal. LOC operates a "
            "scheduled airport shuttle YUL ↔ Alt Hotel Montréal between June 23–28 (60 USD mandatory "
            "per-person fee). Arrive at least Wednesday (Jun 24) to use the official Wed/Thu pool + gym "
            "training windows ahead of Friday's familiarization."
        ),
        "stats": [
            ("Primary route",    "DEN → YUL",       "Nonstop available"),
            ("Flight time",      "~3h 45m",         "Nonstop estimate"),
            ("Time zones",       "+2h",             "Montréal ahead of Denver"),
            ("Airport Transfer", "YUL → Alt Hotel", "LOC shuttle (60 USD)"),
            ("Transfer Distance","19 km / ~20 min", "By car/shuttle"),
        ],
        "hotel_title": "Host Hotel / Official Accommodation Read",
        "hotel_bullets": [
            "Official hotel: Alt Hotel Montréal (Germain Group), 120 Peel Street, H3C 0L8.",
            "Block code \"2606OBTRIA\" — book via the WTPS reservations link, email reservations.altmontreal@germainhotels.com, or call 514.375.0220.",
            "Distance to venue (Parc Jean-Drapeau): 6.3 km — 15 min by car or bike; LOC provides free venue transport on familiarization + race day.",
            "Underground parking on site (Indigo-managed); cash / Visa / Mastercard accepted.",
            "Hotel is downtown (Griffintown adjacent) — easy walk to Old Montréal, restaurants, and metro access.",
        ],
        "food_title": "Food / Grocery Options Near Hotel + Venue",
        "food_bullets": [
            "Best stock-up: Provigo / IGA / Métro grocery stores within walking distance of Alt Hotel (Peel St / Saint-Antoine corridor).",
            "Quick basics: Tim Hortons / Couche-Tard / Starbucks dotted around downtown; pharmacies (Jean Coutu / Pharmaprix) for sports nutrition + electrolytes.",
            "Athlete-friendly restaurants: Marcus, Burgundy Lion, and the Cours Mont-Royal food court are within 10–15 min walk for protein-forward options.",
            "Venue note: Parc Jean-Drapeau is on Notre-Dame Island — bring race-day food/fluid from the hotel; on-island vendor options are limited and seasonal.",
            "Heat consideration: Late June can hit 30 °C+ — pre-stock electrolytes / ice packs at hotel for daily training blocks.",
        ],
    },
    "quiberon": {
        "origin":          "Denver, CO → Quiberon race week",
        "core_read": (
            "There are no direct Denver → Nantes flights. The cleanest athlete route is a one-stop "
            "itinerary through CDG, Dublin, Amsterdam, or another European hub into Nantes, then use "
            "the LOC/Nirvana scheduled shuttle to Quiberon. Many flight options include an overnight "
            "long-haul travel as a part of the leg. Minimize the 8h time zone shift by shifting daily "
            "schedule earlier before travel."
        ),
        "stats": [
            ("Primary route",    "DEN → NTE",       "1 stop likely"),
            ("Flight time",      "~12h",            "Plus layover"),
            ("Time zones",       "+8h",             "France ahead of Denver"),
            ("Airport Transfer", "NTE → Quiberon",  "LOC Scheduled Transfer"),
            ("Transfer Distance","~2 hours",        ""),
        ],
        "hotel_title": "Host Hotel / Official Accommodation Read",
        "hotel_bullets": [
            "Official accommodation is 3 hotels: Sofitel Quiberon Thalassa Sea & Spa, Mercure Quiberon Hotel, and ENV-I2N.",
            "Sofitel = premium recovery/wellness choice, oceanfront thalassotherapy setting.",
            "Mercure = practical team option: renovated rooms, close to beach/thalasso area, free Wi-Fi, bike storage in package.",
            "ENV-I2N = budget-oriented accommodation within a short ride/cycle to the event site.",
        ],
        "food_title": "Food / Grocery Options Near Hotel + Venue",
        "food_bullets": [
            "Best stock-up target: Super U Quiberon (116 Rue du Port de Pêche) for groceries, water, snacks, and room food.",
            "Quick basics: town-center bakeries, small markets, pharmacies, and convenience food around Quiberon centre / beach area.",
            "Restaurant read: seafood/crêperie-heavy resort town; plan simple athlete meals if avoiding rich sauces or unfamiliar seafood.",
            "Logistics note: peninsula access is constrained by one road/rail link. Potential for heavy traffic if coming from outside town center.",
        ],
    },
    "hamburg": {
        "origin":          "Denver, CO → Hamburg race week",
        "core_read": (
            "Denver → Hamburg requires one or two stops — typical routes connect through London (LHR), Amsterdam (AMS), "
            "Frankfurt (FRA), or Paris (CDG) before the short onward hop into Hamburg (HAM). "
            "Lufthansa and United operate the most athlete-friendly itineraries from DEN. "
            "CEST (UTC+2) is 8 hours ahead of MDT, so plan 3–4 days of pre-race arrival to manage jet lag. "
            "Hamburg Airport (HAM) is 10 km from the city center — taxi/rideshare ~20 min, S-Bahn line S1 "
            "links the airport to the Hauptbahnhof in ~25 min. Most WTCS hotels sit within 1–2 km of the Alster."
        ),
        "stats": [
            ("Primary route",    "DEN → HAM",        "1–2 stops (FRA/AMS/LHR typical)"),
            ("Flight time",      "~14–16h",          "Including layover"),
            ("Time zones",       "+8h",              "CEST ahead of MDT"),
            ("Airport Transfer", "HAM → City Center", "S-Bahn S1 ~25 min or taxi ~20 min"),
            ("Transfer Distance","~10 km",            "HAM Airport to Binnenalster"),
        ],
        "hotel_title": "Host Hotel / Official Accommodation",
        "hotel_bullets": [
            "Historic WTCS host hotels cluster along the Außenalster and inner city: Radisson Blu, Fairmont Hotel Vier Jahreszeiten, and NH Hotels are common LOC choices.",
            "Binnenalster / Rathausmarkt area gives walking distance to transition — no shuttle needed for pre-race check-in or familiarization.",
            "City centre parking is limited and expensive — most athletes walk or use S-Bahn from the hotel to the venue.",
            "Check the official 2026 WTCS Hamburg athlete guide (typically released 2–3 weeks prior) for confirmed host hotel and rate codes.",
        ],
        "food_title": "Food / Grocery Options Near Hotel + Venue",
        "food_bullets": [
            "Best stock-up: REWE, Edeka, or Aldi branches within 10 min walk of the Binnenalster area; Hamburger Str / Mönckebergstraße corridor has most options.",
            "Quick basics: dm (pharmacy/snack), Starbucks, Backwerk, and bakeries throughout the Rathausmarkt pedestrian zone.",
            "Athlete-friendly meals: Zum Alten Rathaus, Block House (Hamburg chain, reliable protein), and Superkitchen near the Alster area.",
            "Race-morning note: T1 / transition area is surrounded by city center — café and bakery options open early for race-morning breakfast.",
        ],
    },
    "edmonton": {
        "origin":          "Denver, CO → Edmonton race week",
        "core_read": (
            "Denver → Edmonton is a same-zone trip — both cities run Mountain Time (MDT in July), so there is "
            "zero jet-lag adjustment. No historical nonstop DEN-YEG service; typical itineraries connect through "
            "Calgary (YYC), Vancouver (YVR), or Seattle (SEA) on WestJet, Air Canada, Alaska, or United. LOC-provided "
            "shuttles run Wednesday July 15 through Monday July 20 between the airport and the Westin Hotel only — "
            "submit arrival/departure details by July 5, 2026 or transportation is not guaranteed. A visitor visa or "
            "eTA is required for most nationalities entering Canada; contact Triathlon Canada for a visa support letter."
        ),
        "stats": [
            ("Primary route",    "DEN → YEG",        "1 stop typical (YYC/YVR/SEA)"),
            ("Time zones",       "+0h",              "Both Mountain Time (MDT)"),
            ("Airport Transfer", "YEG → The Westin", "LOC shuttle, ~30 km / 40 min"),
            ("Venue Transfer",   "Westin → Hawrelak Park", "LOC shuttle, 4.5 km / 7-10 min"),
            ("Arrival deadline", "July 5, 2026",     "Submit itinerary form or fees may apply"),
        ],
        "hotel_title": "Host Hotel / Official Accommodation",
        "hotel_bullets": [
            "Official host hotel: The Westin Edmonton, 10135 100 St NW, Edmonton, AB T5J 0N7 — 4.5 km / 7-10 min from the venue by car or bike.",
            "Alternative discounted block: Lister Hall, 11613 87 Ave NW — 2.0 km / 5 min from the venue; book via front desk 780-492-6056, reference \"2026 World Triathlon Cup.\"",
            "LOC free shuttle covers Westin ↔ Hawrelak Park for familiarization and race day only — athletes staying elsewhere arrange their own transport.",
            "Underground parking available at the Westin for a fee, managed directly by the hotel.",
        ],
        "food_title": "Food / Grocery Options Near Hotel + Venue",
        "food_bullets": [
            "Not itemized in the athlete guide — the Westin sits downtown with typical hotel-district dining and grocery access within a few blocks.",
            "Lister Hall option is on the University of Alberta campus — Whyte Ave (109 St / 82 Ave) strip is a short ride away for groceries and casual dining.",
            "Kinsmen Sports Centre (swim training venue) and Hawrelak Park both sit along the river valley bike path connecting the Westin, Lister Hall, and the U of A.",
        ],
    },
    "asuncion": {
        "origin":          "Denver, CO → Asunción race week",
        "core_read": (
            "Denver → Asunción is a long-haul with no nonstop option; realistic routings connect through Panama City "
            "(PTY, Copa), Lima (LIM), São Paulo (GRU) or Miami (MIA), typically 16–22 h door to door with an overnight leg. "
            "VERIFY THE TIME ZONE: the athlete guide lists Paraguay as GMT-4, but Paraguay moved to permanent UTC-3 in "
            "October 2024, which would make it +3 h from Denver (MDT), not +2 h. Confirm before booking race-morning logistics. "
            "The venue is 15 km from the official hotel across the Héroes del Chaco bridge and is CLOSED except on "
            "familiarization day and race day, so there is no course access midweek. Official LOC transfers run to a fixed "
            "timetable — women's race-day call time at the hotel is 04:30."
        ),
        "stats": [
            ("Primary route",    "DEN → ASU",        "1–2 stops (PTY/LIM/GRU/MIA)"),
            ("Time zones",       "+2h or +3h",       "Guide says UTC-4; verify UTC-3"),
            ("Airport Transfer", "ASU → Esplendor",  "~6 km / 15 min by car"),
            ("Hotel → Venue",    "~15 km / 25-35 min", "Traffic dependent, via bridge"),
            ("Race-day call",    "04:30 (W) / 06:35 (M)", "At Hotel Esplendor"),
        ],
        "hotel_title": "Host Hotel / Official Accommodation",
        "hotel_bullets": [
            "Official hotel: Esplendor by Wyndham Asunción, in the corporate centre of Asunción. Booking codes and event rates listed as TBD in the guide.",
            "Hotel is ~15 km / 25–35 min from the Delta venue — the longest hotel-to-venue transfer of the current preview set; budget for traffic.",
            "Reference hospital: Sanatorio Santa Bárbara (~4 km / 10–15 min from the hotel), +595 21 299 388. Venue medical is free; private hospital care is athlete-paid, so carry insurance.",
            "Non-official training: COP Olympic Aquatic Centre (swim, Tue-Fri 07:00–11:00) and COP athletics track (Tue-Fri 14:00–16:00). Costanera Norte is the only recommended road ride/run.",
        ],
        "food_title": "Food / Money / Practical Notes",
        "food_bullets": [
            "Shopping del Sol (~1 km) and Paseo La Galería (~1.5 km) sit next to the hotel: supermarkets, food courts, pharmacies, ATMs and exchange offices.",
            "Reference prices: food-court meal Gs. 35–60k (USD 5–8), mid-range restaurant Gs. 70–120k (USD 9–16), bottled water Gs. 5–8k.",
            "Bottled water is recommended for visitors. Cash is guaraníes only for taxis and small purchases; cards work at hotels, malls and ride apps.",
            "Ride-hailing (Bolt cheapest, plus Uber and local MUV) is well covered: ~USD 2–5 in-city, USD 8–12 airport run. Electricity 220 V 50 Hz, plug type C (European two-pin).",
            "LOC advises AGAINST riding on open roads in the city — bike safety is explicitly the athlete's own responsibility outside the venue.",
        ],
    },
    "weihai": {
        "origin":          "Denver, CO → Weihai race week",
        "core_read": (
            "The longest trip on the 2026 calendar. Denver → Weihai (WEH) needs two or three legs, typically DEN to a "
            "US west-coast gateway, then a trans-Pacific into Beijing/Shanghai/Seoul, then a domestic hop or the "
            "high-speed rail into Weihai. Budget 24–30 h door to door and plan 4–5 days on the ground: CST is UTC+8, "
            "14 hours ahead of Denver (MDT), which is close to maximum circadian disruption. A CHINESE VISA IS "
            "REQUIRED for most nationalities — request the invitation letter early via Ms. Ye Sinan "
            "(yesinan@ctsa.org.cn). LOC transfers are free from the airport, rail stations and port between 25–30 Aug, "
            "but ONLY if arrival details are submitted by 15 August; miss that deadline and there is no pickup."
        ),
        "stats": [
            ("Primary route",    "DEN → WEH",        "2–3 legs via PEK/PVG/ICN"),
            ("Time zones",       "+14h",             "CST UTC+8 vs Denver MDT"),
            ("Airport Transfer", "WEH → Hotel",      "46.2 km / ~50 min"),
            ("Rail alternative", "Weihai North HSR", "8.4 km / ~18 min to hotel"),
            ("Hotel → Venue",    "0.8 km",           "Walking distance"),
        ],
        "hotel_title": "Host Hotel / Accommodation",
        "hotel_bullets": [
            "Official hotel: Weihaiwei Hotel, No. 82 Haigang Road, Huancui District — 0.8 km from the main venue, an easy walk. 498 CNY/night single or twin, breakfast included.",
            "Briefing and press conference are at the Brigh Radiance Hotel, ~500 m from the official hotel (briefing Fri 28 Aug 16:00–16:30).",
            "Amenities: free Wi-Fi, fitness centre, free laundry on the 2nd floor, two bottles of water and fruit per room daily, free room-service delivery.",
            "Athlete/coach set meals 198–498 CNY per person in Liugong Island Hall (2/F). Pasta Party Fri 28 Aug 18:00–20:00, vegetarian and halal available.",
            "Weihai Municipal Hospital is ~1,600 m from the venue. Venue medical is free; hospital treatment is athlete-paid, so carry insurance.",
        ],
        "food_title": "China-Specific Practical Notes",
        "food_bullets": [
            "POWER BANKS: mainland domestic flights require the 3C/CCC logo — non-certified banks are routinely confiscated, including at transit security on a domestic connection. Under 100 Wh, carry-on only, never checked.",
            "Payment is QR-dominant (WeChat Pay / Alipay); link an international card before travelling and carry ¥500–1,000 cash for edge cases.",
            "Electricity 220 V / 50 Hz. Plug types A (US flat), I (AU) and some C (EU) — most hotels have universal sockets, but bring an adapter.",
            "Pool training available 27–28 Aug at Weihai Swimming Pool (25 m, 4 lanes reserved, free) with a free LOC shuttle from the hotel at 07:30/08:30/13:30/14:30.",
            "Bike and run training is on open public roads at the athlete's own risk; only the race venue gets traffic control, and only on familiarization and race day.",
        ],
    },
}


def add_travel_load_slide(prs: Presentation, venue: str) -> bool:
    """Travel & arrival logistics — DEN→venue routing, hotel, food."""
    travel = TRAVEL_PROFILES.get(venue.lower().strip())
    if not travel:
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Travel & Arrival Load", travel.get("origin", venue))

    # ── Core travel read (rounded card top) ──────────────────────────────────
    read_top = Inches(1.2)
    read_h   = Inches(1.05)
    read_bx = slide.shapes.add_shape(1, Inches(0.3), read_top, Inches(12.73), read_h)
    read_bx.fill.solid()
    read_bx.fill.fore_color.rgb = RGBColor(0xE9, 0xEF, 0xF9)
    read_bx.line.color.rgb = MID_GRAY
    _add_textbox(slide, "Core travel read",
                 Inches(0.4), read_top + Inches(0.07),
                 Inches(12.5), Inches(0.3),
                 font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _add_textbox(slide, travel.get("core_read", ""),
                 Inches(0.5), read_top + Inches(0.38),
                 Inches(12.3), Inches(0.62),
                 font_size=10.5, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # ── Stat cards row ───────────────────────────────────────────────────────
    stats = travel.get("stats", [])
    if stats:
        cards_top = read_top + read_h + Inches(0.2)
        n = min(len(stats), 5)
        card_total_w_in = 12.73
        gap_in = 0.15
        card_w_in = (card_total_w_in - gap_in * (n - 1)) / n
        for i, (label, main, sub) in enumerate(stats[:n]):
            cl = Inches(0.3 + (card_w_in + gap_in) * i)
            cw = Inches(card_w_in)
            bx = slide.shapes.add_shape(1, cl, cards_top, cw, Inches(0.95))
            bx.fill.solid()
            bx.fill.fore_color.rgb = WHITE
            bx.line.color.rgb = MID_GRAY
            _add_textbox(slide, label, cl + Inches(0.08), cards_top + Inches(0.05),
                         cw - Inches(0.16), Inches(0.22),
                         font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            _add_textbox(slide, main, cl + Inches(0.08), cards_top + Inches(0.28),
                         cw - Inches(0.16), Inches(0.32),
                         font_size=14, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
            if sub:
                _add_textbox(slide, sub, cl + Inches(0.08), cards_top + Inches(0.62),
                             cw - Inches(0.16), Inches(0.28),
                             font_size=9, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # ── Two panels (hotel + food) ────────────────────────────────────────────
    panels_top = read_top + read_h + Inches(1.35)
    panel_h    = Inches(2.6)
    panel_w    = Inches(6.2)
    panels = [
        (Inches(0.3),              travel.get("hotel_title", ""),  travel.get("hotel_bullets", [])),
        (Inches(0.3) + panel_w + Inches(0.3), travel.get("food_title", ""), travel.get("food_bullets", [])),
    ]
    for left, title, bullets in panels:
        bx = slide.shapes.add_shape(1, left, panels_top, panel_w, panel_h)
        bx.fill.solid()
        bx.fill.fore_color.rgb = WHITE
        bx.line.color.rgb = MID_GRAY
        _add_textbox(slide, title,
                     left + Inches(0.15), panels_top + Inches(0.12),
                     panel_w - Inches(0.3), Inches(0.32),
                     font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        for i, b in enumerate(bullets[:6]):
            _add_textbox(slide, f"•  {b}",
                         left + Inches(0.18), panels_top + Inches(0.5 + i * 0.34),
                         panel_w - Inches(0.36), Inches(0.36),
                         font_size=10, color=DARK_GRAY)
    return True


# ── Venue preview narrative + feature bullets (Key Course Differentiators) ──
# Used to populate the venue intro slide. When a venue has no historical race
# results, this is the primary way an athlete gets a feel for the course.
VENUE_PREVIEW: dict[str, dict] = {
    "karlovy vary": {
        "narrative": (
            "WTCS Karlovy Vary returns to the Czech spa town for its second year at Championship "
            "Series level (World Cup 2017–2024, WTCS from 2025). This is a standard-distance, "
            "two-transition race: a 1.5 km two-lap swim in Lake Rolava with an Australian exit, a "
            "6.5 km approach into 7 brutal 4.9 km bike laps (40.8 km, ~20% of every lap climbing, "
            "part cobbled), then a 4-lap 10 km run through the city centre from Theatre Square. "
            "The bike is the story here — lapped athletes are pulled, and the front group is "
            "usually shredded before the run even starts."
        ),
        "features": [
            "1.5 km swim (2 x 750 m) in Lake Rolava — pontoon start, Australian exit, island navigation. Wetsuit ruling has flipped year to year.",
            "40.8 km bike: 6.5 km approach + 7 x 4.9 km laps, ~20% ascent every lap, 80% asphalt / 20% cobbles. Lapped athletes REMOVED.",
            "10 km run (4 x 2.5 km) through the spa-town centre — mixed asphalt/pavement, aid twice per lap.",
            "Two transitions (T1 Lake Rolava to T2 Theatre Square) and split start times: Women 10:00, Men 15:00.",
        ],
        "format_label":   "Standard Triathlon (WTCS)",
        "race_info_url":  "https://events.triathlon.org/2026-wtcs-karlovy-vary/race-info",
        "race_info_text": "Race Info | 2026 WTCS Karlovy Vary",
    },
    "rio de janeiro": {
        "narrative": (
            "The World Triathlon Cup returns to Copacabana Beach — the 2016 Olympic triathlon "
            "venue — for a sprint-distance race on August 2. A beach-start 750 m ocean swim in "
            "Atlantic surf feeds a dead-flat, fast 20.8 km bike along Avenida Atlântica and a "
            "2-lap, 5 km run down the iconic Copacabana promenade. Warm 22–24 °C water means a "
            "non-wetsuit swim; the flat swim-and-bike put the whole race on the run. Morning "
            "starts (Women 07:00, Men 09:00) in Rio's mild August winter."
        ),
        "features": [
            "Beach-start 750 m ocean swim off Copacabana — Atlantic surf entry, non-wetsuit (22–24 °C).",
            "Flat & fast 20.8 km bike (6 laps) along the Avenida Atlântica seafront — big packs, high speed.",
            "5 km run (2 laps) on the iconic Copacabana promenade — flat and fast; this is where it's decided.",
            "2016 Olympic venue; morning starts (W 07:00 / M 09:00) in mild Rio winter conditions.",
        ],
        "format_label":   "Sprint Triathlon (World Cup)",
        "race_info_url":  "https://events.triathlon.org/2026-world-triathlon-cup-rio-de-janeiro/race-info",
        "race_info_text": "Race Info | 2026 World Triathlon Cup Rio de Janeiro",
    },
    "london": {
        "narrative": (
            "WTCS London moves to the ExCeL in the Royal Docks — a sprint-distance, "
            "stadium-style race with an afternoon start (Women 14:30, Men 16:15). The swim is "
            "a single 750 m loop in the enclosed Royal Victoria Dock, followed by a dead-flat, "
            "fast 20.64 km bike (7 short out-and-back laps on Royal Albert Way) and a 4.92 km run "
            "that weaves through the ExCeL venue hall on each of 3 laps. Cool UK-summer conditions "
            "and a flat, fast layout mean this race is decided on the run."
        ),
        "features": [
            "750 m single-loop swim in the enclosed Royal Victoria Dock — flat freshwater, 315 m to the first buoy, pontoon start.",
            "Flat & fast bike: 7 × 2.64 km out-and-back laps (20.64 km) on Royal Albert Way — pure power, big lead packs.",
            "4.92 km run, 3 laps partly indoors through the ExCeL hall — atmospheric and very fast.",
            "Afternoon start (W 14:30 / M 16:15) — unusual for WTCS; weather + water are at the day's warmest.",
        ],
        "format_label":   "Sprint Triathlon (WTCS)",
        "race_info_url":  "https://events.triathlon.org/2026-wtcs-london/race-info",
        "race_info_text": "Race Info | 2026 WTCS London",
    },
    "montreal": {
        "narrative": (
            "World Triathlon Para Series Montréal returns to Parc Jean-Drapeau for the 2026 "
            "edition. The race uses the venerable Olympic Basin for a 750 m swim, then heads to "
            "the Circuit Gilles-Villeneuve — the F1 track — for a flat, fast 21.5 km bike on a "
            "closed circuit. The run mixes the CGV and the path alongside the rowing basin. "
            "Late June in Montréal can run hot and humid; with 23 °C water, expect a non-wetsuit "
            "race and a power-driven bike."
        ),
        "features": [
            "Olympic Basin swim — sheltered, ~23 °C fresh water, in-water pontoon start. Non-wetsuit.",
            "Bike on Circuit Gilles-Villeneuve (F1 track): pancake-flat 4.1 km clockwise loop × 5 = 21.5 km.",
            "Run on CGV + Olympic Basin path: flat 2.5 km loop × 2 with five aid-station access points.",
            "12 distinct Para sport-class starts (PTWC, PTVI, PTS2–5) staged 07:00 → 08:30 across both genders.",
        ],
        "format_label":   "Sprint Para Triathlon",
        "race_info_url":  "https://events.triathlon.org/2026-world-triathlon-para-series-montreal",
        "race_info_text": "Race Info | 2026 World Triathlon Para Series Montréal",
    },
    "quiberon": {
        "narrative": (
            "Set on the southern tip of Brittany's Quiberon peninsula, this will be the first "
            "edition of WTCS Quiberon. Expect a cool-water, likely wetsuit sprint race with a "
            "decently long beach start. The bike is largely flat but has multiple sections "
            "along the coast with potential for exposure to coastal winds. The run is also a "
            "simple out-and-back on asphalt that should lead to fast splits."
        ),
        "features": [
            "Beach Start 750 m Atlantic Swim. Cool water, likely wetsuit-legal.",
            "Flat bike course. 4 laps of 5.5 km will make the distance slightly longer than standard. "
            "No real elevation but potential for exposed wind could come into play.",
            "Mostly asphalt out-and-back run. Simple 2-lap course.",
            "First race at this venue.",
        ],
        "format_label":   "Sprint Triathlon",
        "race_info_url":  "https://events.triathlon.org/2026-world-triathlon-championship-series-quiberon",
        "race_info_text": "Race Info | 2026 World Triathlon Championship Series Quiberon",
    },
    "huatulco": {
        "narrative": (
            "Huatulco's World Cup returns to the Bahía de Santa Cruz on Mexico's Pacific coast. "
            "A staple Tier-2 World Cup with strong historical depth in the database, the course "
            "rewards swimmers off the beach start, riders who can survive the short Vialidad 5 "
            "climb, and runners who can hold pace in tropical heat at the 08:30 men's start."
        ),
        "features": [
            "Bahía de Santa Cruz swim — warm water (~30 °C), beach start, 750 m single loop.",
            "5 km bike loop × 4 (20 km) with a punchy Vialidad 5 climb (max 18.4%) every lap.",
            "Hilly 2.5 km run loop × 2 — +55 m gain/lap, max 17.3% on Camino a Santa Cruz.",
            "Tropical heat is the dominant tactical variable. Pre-cool, sponge plan mandatory.",
        ],
        "format_label":   "Sprint Triathlon",
        "race_info_url":  "https://events.triathlon.org/2026-world-triathlon-cup-huatulco",
        "race_info_text": "Race Info | 2026 World Triathlon Cup Huatulco",
    },
    "hamburg": {
        "narrative": (
            "WTCS Hamburg is one of the sport's flagship city-center races, set in the heart of Hamburg "
            "along the Binnenalster and the historic Rathausmarkt. The sprint format — 750 m Alster "
            "swim, 6 × 3.3 km technical urban bike, 2 × 5 km flat run — rewards athletes who "
            "can handle high-frequency cornering on the bike and sustain race pace on fresh legs off "
            "a large pack. Hamburg has hosted elite triathlon since 2008 and consistently delivers "
            "fast, tactical racing in front of one of the largest crowds on the WTCS calendar."
        ),
        "features": [
            "Sheltered 750 m Binnenalster swim — flat calm fresh water, pontoon start, non-wetsuit typical in July.",
            "6 × 3.3 km urban bike (19.8 km) — flat but highly technical with multiple turns; large packs form.",
            "2 × 5 km flat run through city-center streets — fast finishing conditions, cobblestone sections.",
            "Elite Men 13:15 / Elite Women 15:00 Saturday July 11. Mixed Relay Sunday July 12.",
        ],
        "format_label":   "Sprint Triathlon (WTCS)",
        "race_info_url":  "https://events.triathlon.org/2026-wtcs-hamburg/race-info",
        "race_info_text": "Race Info | 2026 World Triathlon Championship Series Hamburg",
    },
    "edmonton": {
        "narrative": (
            "The 2026 World Triathlon Cup Edmonton returns elite racing to William Hawrelak Park, newly renovated "
            "ahead of the 2027 World Triathlon Multisport Championships — this race is the first elite test of the "
            "rebuilt course. Edmonton has a long elite history at this tier and above: WTS races in 2018 and 2019, "
            "and the 2021 World Triathlon Championship Finals, where Kristian Blummenfelt and Flora Duffy clinched "
            "world titles on this same lake. The last two editions (2023-24) stepped down to Americas Triathlon Cup "
            "level, so this World Cup marks a return to a stronger international field. Format is standard sprint: "
            "750 m lake swim, 20.4 km bike (3 laps, valley climb each lap), 5 km run (3 laps). Elite Men go off at "
            "14:00 and Elite Women at 16:00 on Saturday July 18; Mixed Relay follows Sunday July 19."
        ),
        "features": [
            "750 m Hawrelak Park lake swim — sheltered fresh water, single clockwise loop, ~19 °C (optional wetsuit).",
            "20.4 km bike — 3 × 6.8 km laps with a repeated river-valley climb/descent (+85 m/-87 m per lap).",
            "5 km run — 3 × 1.66 km laps on paved park paths, moderate heat risk (July highs ~23 °C).",
            "Elite Men 14:00 / Elite Women 16:00, Saturday July 18. Mixed Relay 14:30, Sunday July 19.",
        ],
        "format_label":   "Sprint Triathlon (World Cup)",
        "race_info_url":  "https://events.triathlon.org/2026-world-triathlon-cup-edmonton",
        "race_info_text": "Race Info | 2026 World Triathlon Cup Edmonton",
    },
    "asuncion": {
        "narrative": (
            "The 2026 World Triathlon Cup Asunción is the first World Triathlon elite race ever staged at this venue, "
            "and Paraguay's first World Cup. It runs on La Isla del Delta in Nueva Asunción, across the Héroes del Chaco "
            "bridge from the capital. Two things make it unusual on the 2026 calendar. First, it is a STANDARD-distance "
            "World Cup (1500 m / 40 km / 10 km) at a time when most Cups are sprints, so it rewards a different athlete "
            "profile. Second, the course is flat in all three disciplines, on a closed avenue with long straights and few "
            "turns, which points to a large intact bike pack and a run-off finish. There is no historical race data at "
            "this venue, so every read here comes from the athlete guide and course flyers rather than past results."
        ),
        "features": [
            "STANDARD distance (1500 m / 40 km / 10 km) — not a sprint. Different pacing and fuelling plan from most 2026 Cups.",
            "Flat across all three legs: 2 × 750 m lagoon swim, 6 × 6.66 km closed-avenue bike, 4 × 2.35 km + 600 m run.",
            "Wetsuit call is live: 18–22 °C forecast straddles the 20 °C elite cut-off. Pack both.",
            "August is Paraguayan winter — cool morning starts (Women 07:30, Men 10:30), both finishing before the daily max. Cold-weather warm-up kit matters more than cooling kit.",
        ],
        "format_label":   "Standard Distance (World Cup)",
        "race_info_url":  "https://events.triathlon.org/2026-world-triathlon-cup-asuncion/race-info",
        "race_info_text": "Race Info | 2026 World Triathlon Cup Asunción",
        # Used for max-points / tier labelling when the event is not yet in the DB.
        # WT event_id 195159 (2026-08-09) — ingest to unlock startlist-based slides.
        "event_name":     "2026 World Triathlon Cup Asuncion",
        "cat_name":       "World Cup",
    },
    "weihai": {
        "narrative": (
            "WTCS Weihai returns to the Wave to Wonderland venue on China's Shandong peninsula for the 2026 "
            "Championship Series. Weihai is one of the most established venues on the calendar, hosting elite racing "
            "since 2011 and stepping up from World Cup to WTCS in 2024. The format is standard distance: a 2-lap "
            "1.5 km sea swim, a lap-heavy 8 × 5 km bike, and a flat 4 × 2.5 km seafront run. Two things separate this "
            "edition from recent history. First, the date moves to 29 August, roughly four weeks earlier than the "
            "2024 and 2025 editions, which points at 2023-style heat rather than the milder late-September conditions "
            "of the last two years. Second, and more important tactically, the men start at 09:00 and the women at "
            "12:00, into peak midday sun. WT's recorded race-day temperature has come in about 4 °C higher for the "
            "women in every measured edition, and since ambient air only warms ~1 °C across those three hours here, "
            "that gap is essentially solar load on an unshaded seafront course."
        ),
        "features": [
            "STANDARD distance (1.5 km / 40 km / 10 km) at WTCS tier — 1000 points to the winner.",
            "Non-wetsuit is effectively certain: water has measured 24-26 °C and been ruled wetsuit-forbidden in all four recorded editions.",
            "8 × 5 km bike with a turnaround every lap — the most lap-dense, surge-heavy bike course on the preview set.",
            "Heat is asymmetric: men 09:00, women 12:00 into peak sun. Recorded race temps run ~4 °C higher for the women.",
        ],
        "format_label":   "Standard Distance (WTCS)",
        "race_info_url":  "https://events.triathlon.org/2026-wtcs-weihai/race-info",
        "race_info_text": "Race Info | 2026 WTCS Weihai",
        # Tier metadata for max-points labelling when the event is not yet ingested.
        # WT event_id 195151 (2026-08-29); startlist already present in program_entries.
        "event_name":     "2026 World Triathlon Championship Series Weihai",
        "cat_name":       "World Championship Series",
    },
}


# ── Race-week schedule data ───────────────────────────────────────────────────
# Each row: (time_label, activity_text, is_highlighted). Highlighted rows render
# in red (used for start times and mandatory meetings).
EVENT_SCHEDULES: dict[str, dict] = {
    "karlovy vary": {
        "title":      "2026 WTCS Karlovy Vary — Race Week",
        "date_range": "September 12 – 13, 2026",
        "venue_note": "Lake Rolava (T1) + Theatre Square (T2), Karlovy Vary (CEST, UTC+2)",
        "race_starts": [("Elite Women", "10:00"), ("Elite Men", "15:00")],
        "days": [
            ("Sat • Sep 12", "Familiarisation & Briefing", [
                ("07:30 – 09:30", "Athletes Lounge open — Lake Rolava", False),
                ("07:45 – 08:45", "Swim course familiarisation — Lake Rolava", False),
                ("09:00",         "Bike escorted familiarisation — approach + 2 laps from Rolava", False),
                ("16:00 – 16:30", "★ Mandatory Elite Athletes briefing — Spa Hotel Thermal", True),
            ]),
            ("Sun • Sep 13 — Women", "WTCS RACE DAY (AM)", [
                ("08:00 – 08:30", "Elite Women — TA2 check-in (Theatre Square)", False),
                ("08:30 – 09:30", "Elite Women — Athletes Lounge check-in (re-entry to 09:50)", False),
                ("08:45 – 09:45", "Elite Women — TA1 check-in (Lake Rolava)", False),
                ("09:00 – 09:45", "Elite Women — Swim warm-up", False),
                ("09:50 – 09:59", "Elite Women — Line-up and presentation", False),
                ("10:00",         "★ ELITE WOMEN RACE START", True),
                ("12:10 / 12:25", "Women awards + Chocolate ceremony — Theatre Square", False),
            ]),
            ("Sun • Sep 13 — Men", "WTCS RACE DAY (PM)", [
                ("13:00 – 13:30", "Elite Men — TA2 check-in (Theatre Square)", False),
                ("13:30 – 14:30", "Elite Men — Athletes Lounge check-in", False),
                ("13:45 – 14:45", "Elite Men — TA1 check-in (Lake Rolava)", False),
                ("14:00 – 14:45", "Elite Men — Swim warm-up", False),
                ("14:50 – 14:59", "Elite Men — Line-up and presentation", False),
                ("15:00",         "★ ELITE MEN RACE START", True),
                ("17:10 / 17:25", "Men awards + Chocolate ceremony — Theatre Square", False),
            ]),
        ],
    },
    "rio de janeiro": {
        "title":      "2026 World Triathlon Cup Rio de Janeiro — Race Week",
        "date_range": "August 1 – 2, 2026",
        "venue_note": "Triathlon Arena, Copacabana Beach, Av Atlântica (BRT, UTC-3)",
        "race_starts": [("Elite Women", "07:00"), ("Elite Men", "09:00")],
        "days": [
            ("Sat • Aug 1", "Familiarization & Briefing", [
                ("08:00 – 08:45", "Bike course familiarization — Triathlon Arena", False),
                ("09:00 – 09:45", "Swim course familiarization — Triathlon Arena", False),
                ("16:00 – 17:00", "★ Mandatory Elite Athletes' briefing — Othon Palace", True),
                ("Post-briefing", "Race pack + credentials distribution", False),
            ]),
            ("Sun • Aug 2", "WORLD CUP RACE DAY", [
                ("05:30 – 06:30", "Elite Women — Athletes' lounge open", False),
                ("05:30 – 06:45", "Elite Women — Transition open", False),
                ("06:20 – 06:45", "Elite Women — Swim course open for warm-up", False),
                ("06:50 – 06:58", "Elite Women — Presentation", False),
                ("07:00",         "★ ELITE WOMEN RACE START", True),
                ("06:45 – 08:45", "Elite Men — Athletes' lounge open", False),
                ("08:00 – 08:45", "Elite Men — Transition + swim warm-up", False),
                ("08:50 – 08:58", "Elite Men — Presentation", False),
                ("09:00",         "★ ELITE MEN RACE START", True),
                ("10:30",         "Elite Men & Women medal ceremony", False),
            ]),
        ],
    },
    "london": {
        "title":      "2026 WTCS London — Race Week",
        "date_range": "July 24 – 25, 2026",
        "venue_note": "ExCeL London, Royal Docks (BST, UTC+1)",
        "race_starts": [("Elite Women", "14:30"), ("Elite Men", "16:15")],
        "days": [
            ("Fri • July 24", "Familiarization & Briefing", [
                ("13:30 – 15:30", "Athlete lounge open — ExCeL South Halls", False),
                ("14:00 – 14:30", "Swim course familiarization — Dockside (via Elite Transition)", False),
                ("15:30 – 16:00", "Athlete sign-in for race briefing — South Halls 23/24", False),
                ("16:00 – 16:30", "★ Mandatory Elite Athletes' race briefing — South Halls 23/24", True),
                ("16:30 – 17:00", "Race pack + accreditation distribution", False),
                ("16:30 – 17:00", "WTCS Team Medical Meeting", False),
            ]),
            ("Sat • July 25", "WTCS SPRINT RACE DAY", [
                ("08:30",         "Bike course familiarization check-in — Elite Transition", False),
                ("09:00 – 09:20", "Bike course familiarization (escorted)", False),
                ("12:30 – 14:00", "Elite Women — Athletes' lounge check-in", False),
                ("13:45 – 14:15", "Elite Women — Transition check-in + swim warm-up", False),
                ("14:23",         "Elite Women — Athlete introductions", False),
                ("14:30",         "★ ELITE WOMEN WTCS START", True),
                ("14:45 – 15:45", "Elite Men — Athletes' lounge check-in", False),
                ("15:30 – 16:00", "Elite Men — Transition check-in + swim warm-up", False),
                ("16:08",         "Elite Men — Athlete introductions", False),
                ("16:15",         "★ ELITE MEN WTCS START", True),
                ("17:15",         "Elite Men awards ceremony", False),
                ("17:30 – 17:40", "Elite Women awards (immediately after men)", False),
            ]),
        ],
    },
    "montreal": {
        "title":      "2026 World Triathlon Para Series Montréal — Race Week",
        "date_range": "June 24 – 27, 2026",
        "venue_note": "Parc Jean-Drapeau, Notre-Dame Island, Montréal (EDT, UTC-4)",
        "race_starts": [("PTWC start", "07:00"), ("PTS/PTVI window", "08:00")],
        "days": [
            ("Thu • June 25", "Open Training", [
                ("11:00 – 15:00", "Swim training — Aquatic Complex (50 m pool, 4 lanes)", False),
                ("14:00 – 17:00", "Gym training — Athletes' Quarter gym", False),
                ("(Wed Jun 24)",   "Pool training 14:00–18:00 + gym 14:00–17:00", False),
            ]),
            ("Fri • June 26", "Familiarization & Briefing", [
                ("07:45",         "Departures from Host Hotel", False),
                ("08:45 – 09:15", "PTS/PTVI Bike Fam  •  PTWC Run Fam — Transition", False),
                ("09:30 – 10:00", "PTS/PTVI Run Fam  •  PTWC Bike Fam — Transition", False),
                ("10:15 – 11:00", "Swim Familiarisation (ALL) — Swim Start", False),
                ("10:15 – 10:45", "Coaches Onsite Meeting — Transition Area", False),
                ("10:30 – 12:00", "Equipment verification (ALL)", False),
                ("12:00 – 12:45", "★ WTPS Pre-race briefing (ALL) — Athletes' Quarter", True),
                ("12:45 – 13:30", "Race Package Distribution", False),
                ("13:45",         "Departures from venue", False),
            ]),
            ("Sat • June 27", "PARA RACE DAY", [
                ("05:00 / 05:30", "Hotel departures — PTWC 05:00 • PTS/PTVI 05:30", False),
                ("05:30 / 06:15", "Athletes' Lounge opens — PTWC 05:30 • PTS/PTVI 06:15", False),
                ("06:15 – 06:45", "PTWC swim warm-up (on course) + transition check-in", False),
                ("06:53",         "PTWC Athletes' Introduction", False),
                ("07:00 / 07:03", "★ PTWC1 Men 07:00  •  PTWC2 Men 07:03", True),
                ("07:04 / 07:08", "★ PTWC1 Women 07:04  •  PTWC2 Women 07:08", True),
                ("07:45",         "★ PTS5 Men start", True),
                ("07:50 / 07:53", "★ PTVI1 Men 07:50  •  PTVI2/3 Men 07:53", True),
                ("07:54 / 07:57", "★ PTVI1 Women 07:54  •  PTVI2/3 Women 07:57", True),
                ("08:10 / 08:15", "★ PTS5 Women 08:10  •  PTS2/3/4 Women 08:15", True),
                ("08:25 / 08:30", "★ PTS4 Men 08:25  •  PTS2/3 Men 08:30", True),
                ("10:15 / 12:00", "Venue departures — 1st bus 10:15  •  2nd 12:00", False),
                ("10:30 – 11:30", "WTPS Medal Ceremony — Podium", False),
            ]),
        ],
    },
    "quiberon": {
        "title":      "2026 WTCS Quiberon — Race Week",
        "date_range": "June 18 – 21, 2026",
        "venue_note": "Espace Louison Bobet, Bd René Cassin, Quiberon (CEST, UTC+2)",
        "race_starts": [("Men start", "10:00"), ("Women start", "12:00")],
        "days": [
            ("Friday • June 19", "Familiarization & Briefing", [
                ("07:00 – 10:00", "Pool training — Neptune Swimming Pool", False),
                ("09:00 – 09:30", "Bike familiarization — 2 laps from transition (escorted)", False),
                ("10:30",         "Swim familiarization — beach", False),
                ("14:00 – 17:00", "Pool training — Neptune", False),
                ("16:00 – 16:30", "★ Mandatory Elite Athletes' briefing — Espace Louison Bobet", True),
                ("16:30 – 17:00", "Elite team medical meeting — LOC Office", False),
                ("Post-briefing", "Race package distribution", False),
                ("19:00",         "Run course measurement (coaches) — Finish Area, Bd René Cassin", False),
            ]),
            ("Saturday • June 20", "WTCS SPRINT RACE DAY", [
                ("08:30 – 09:30", "WTCS Men — Athletes' lounge check-in", False),
                ("09:00 – 09:45", "WTCS Men — Transition + swim warm-up", False),
                ("09:50",         "WTCS Men — Athletes' introduction", False),
                ("10:00",         "★ WTCS MEN START", True),
                ("11:05",         "WTCS Men medals ceremony", False),
                ("10:30 – 11:30", "WTCS Women — Athletes' lounge check-in", False),
                ("11:00 – 11:45", "WTCS Women — Transition + swim warm-up", False),
                ("11:50",         "WTCS Women — Athletes' introduction", False),
                ("12:00",         "★ WTCS WOMEN START", True),
                ("13:10",         "WTCS Women medals ceremony", False),
                ("13:30 – 14:00", "Mixed Relay — Team Declaration", False),
                ("20:00 – 20:30", "Mixed Relay — Team Managers' Meeting (Elite Athletes Area)", False),
            ]),
            ("Sunday • June 21", "Mixed Relay Day", [
                ("07:00 – 10:00", "Pool training — Neptune", False),
                ("12:30 – 12:45", "Mixed Relay — Final Team Declaration", False),
                ("15:30 – 16:30", "Mixed Relay — Athletes' lounge check-in", False),
                ("16:00 – 16:45", "Mixed Relay — Transition + swim warm-up", False),
                ("16:50",         "Mixed Relay — Athletes' introduction", False),
                ("17:00",         "★ MIXED RELAY START", True),
                ("18:30",         "Mixed Relay medals ceremony", False),
            ]),
        ],
    },
    "huatulco": {
        "title":      "2026 World Triathlon Cup Huatulco — Race Week",
        "date_range": "June 12 – 14, 2026",
        "venue_note": "Plaza Santa Cruz Huatulco, Mexico",
        "race_starts": [("Female start", "06:30"), ("Male start", "08:30")],
        "days": [
            ("Friday • June 12", "Familiarization Day", [
                ("10:00",         "Bike familiarization — 3 laps from Bahía Santa Cruz", False),
                ("10:30",         "Run familiarization", False),
                ("11:00 – 11:45", "Swim familiarization — Santa Cruz", False),
            ]),
            ("Saturday • June 13", "Briefing & Packet Pickup", [
                ("16:00 – 16:30", "★ Mandatory Elite athlete meeting — Hotel Biniguenda", True),
                ("16:30 – 17:00", "Packet pickup — Hotel Biniguenda", False),
            ]),
            ("Sunday • June 14", "RACE DAY", [
                ("05:30", "Elite Female lounge opens", False),
                ("06:30", "★ ELITE FEMALE START", True),
                ("07:45", "Elite Male lounge opens", False),
                ("08:30", "★ ELITE MALE START", True),
                ("10:00", "Awards ceremony", False),
            ]),
        ],
    },
    "hamburg": {
        "title":      "2026 World Triathlon Championship Series Hamburg — Race Week",
        "date_range": "July 9 – 12, 2026",
        "venue_note": "Rathausmarkt / Binnenalster, Hamburg, Germany (CEST, UTC+2)",
        "race_date":  "2026-07-11",
        "race_starts": [("Elite Men", "13:15"), ("Elite Women", "15:00")],
        "days": [
            ("Thu • July 9", "Athlete Check-In / Open Training", [
                ("TBD",           "Athlete check-in and registration — TBD per athlete guide", False),
                ("TBD",           "Swim / bike / run open training windows — TBD per athlete guide", False),
            ]),
            ("Fri • July 10", "Familiarization & Briefing", [
                ("TBD",           "Course familiarization — swim, bike, run circuits", False),
                ("TBD",           "★ Mandatory Elite athlete briefing", True),
                ("TBD",           "Race package distribution", False),
            ]),
            ("Sat • July 11", "ELITE SPRINT RACE DAY", [
                ("11:30 – 12:45", "Elite Men — athletes' lounge check-in + warm-up", False),
                ("12:15 – 13:00", "Elite Men — swim warm-up + transition check", False),
                ("13:10",         "Elite Men — athletes' introduction", False),
                ("13:15",         "★ ELITE MEN START", True),
                ("~14:06",        "Elite Men — projected finish", False),
                ("14:15",         "Elite Men medals ceremony", False),
                ("13:45 – 14:45", "Elite Women — athletes' lounge check-in + warm-up", False),
                ("14:15 – 14:55", "Elite Women — swim warm-up + transition check", False),
                ("14:50",         "Elite Women — athletes' introduction", False),
                ("15:00",         "★ ELITE WOMEN START", True),
                ("~15:55",        "Elite Women — projected finish", False),
                ("16:05",         "Elite Women medals ceremony", False),
                ("Evening",       "Mixed Relay team declaration", False),
            ]),
            ("Sun • July 12", "Mixed Relay Day", [
                ("TBD",           "Mixed Relay team managers' meeting", False),
                ("17:30",         "★ MIXED RELAY START", True),
                ("~18:45",        "Mixed Relay projected finish", False),
                ("Post-finish",   "Mixed Relay medals ceremony", False),
            ]),
        ],
    },
    "edmonton": {
        "title":      "2026 World Triathlon Cup Edmonton — Race Week",
        "date_range": "July 15 – 19, 2026",
        "venue_note": "William Hawrelak Park, 9330 Groat Rd NW, Edmonton, AB (MDT, UTC-6)",
        "race_date":  "2026-07-18",
        "race_starts": [("Elite Men", "14:00"), ("Elite Women", "16:00")],
        "days": [
            ("Wed-Thu • July 15-16", "Open Swim Training", [
                ("10:30 – 12:30", "Swim Training — Kinsmen Sport Centre (Wed)", False),
                ("18:00 – 20:00", "Swim Training — Kinsmen Sport Centre (Wed)", False),
                ("11:45 – 13:45", "Swim Training — Kinsmen Sport Centre (Thu)", False),
                ("18:00 – 20:00", "Swim Training — Kinsmen Sport Centre (Thu)", False),
            ]),
            ("Fri • July 17", "Familiarization & Briefing", [
                ("7:45 – 9:45",   "Swim Training — Kinsmen Sport Centre", False),
                ("13:00 – 14:00", "Bike Familiarization — Hawrelak Park (police-escorted)", False),
                ("14:00 – 15:00", "Swim Familiarization — Hawrelak Park lake", False),
                ("16:00 – 16:45", "★ Elite & Mixed Relay Athletes' Briefing — Heritage Amphitheatre", True),
                ("16:45 – 17:15", "Elite Race Package Distribution — Heritage Amphitheatre", False),
            ]),
            ("Sat • July 18", "ELITE SPRINT RACE DAY", [
                ("12:30 – 13:30", "Elite Athlete Lounge open — Men", False),
                ("13:00 – 13:45", "Transition check-in + swim warm-up — Men", False),
                ("13:50 – 13:58", "Elite Men introductions", False),
                ("14:00",         "★ ELITE MEN — SPRINT TRIATHLON START", True),
                ("15:15",         "Elite Men medals ceremony", False),
                ("14:30 – 15:30", "Elite Athlete Lounge open — Women", False),
                ("15:00 – 15:45", "Transition check-in + swim warm-up — Women", False),
                ("15:30",         "Mixed Relay — online team declaration", False),
                ("15:50 – 15:58", "Elite Women introductions", False),
                ("16:00",         "★ ELITE WOMEN — SPRINT TRIATHLON START", True),
                ("17:15",         "Elite Women medals ceremony", False),
                ("17:30 – 18:00", "Mixed Relay — team declaration (Athletes' Lounge)", False),
            ]),
            ("Sun • July 19", "Mixed Relay Day", [
                ("12:30 – 12:45", "Mixed Relay final team declaration", False),
                ("13:00 – 14:00", "Elite Mixed Relay athlete lounge open", False),
                ("13:30 – 14:15", "Transition check-in + swim warm-up", False),
                ("14:20 – 14:28", "Mixed Relay team introductions", False),
                ("14:30",         "★ ELITE MIXED RELAY START", True),
                ("16:00",         "Elite Mixed Relay medals ceremony", False),
            ]),
        ],
    },
    "asuncion": {
        "title":      "2026 World Triathlon Cup Asunción — Race Week",
        "date_range": "August 4 – 9, 2026",
        "venue_note": "La Isla del Delta, Nueva Asunción (guide states UTC-4; verify — Paraguay moved to permanent UTC-3 in 2024)",
        "race_date":  "2026-08-09",
        "race_starts": [("Elite Women", "07:30"), ("Elite Men", "10:30")],
        "days": [
            ("Tue-Fri • Aug 4-7", "Non-Official Training", [
                ("07:00 – 11:00", "Swim — COP Olympic Aquatic Centre", False),
                ("14:00 – 16:00", "Run — COP Athletics Track", False),
                ("(any time)",    "Bike/run — Costanera Norte is the only recommended road option", False),
                ("—",             "⚠ Race venue CLOSED all week: no course access until Saturday", False),
            ]),
            ("Sat • Aug 8", "Familiarization & Briefing", [
                ("07:00 – 07:15", "Call time at Hotel Esplendor", False),
                ("08:35 – 08:40", "Departure from hotel  →  arrive El Delta 09:20-09:30", False),
                ("10:00 – 10:30", "Bike course familiarization (closed circuit)", False),
                ("10:45 – 11:30", "Swim course familiarization (marked area)", False),
                ("13:40 – 13:50", "Departure from venue  →  arrive hotel 14:45-14:55", False),
                ("16:00 – 17:00", "★ Mandatory Athlete Briefing — COP", True),
            ]),
            ("Sun • Aug 9", "ELITE WOMEN — STANDARD", [
                ("04:30",         "Call time at Hotel Esplendor (Women)", False),
                ("05:30",         "Departure from hotel  →  arrive El Delta 06:00", False),
                ("06:00 – 06:45", "Athletes' Lounge opens — check-in", False),
                ("06:10 – 07:00", "Transition area check-in", False),
                ("06:45 – 07:15", "Swim warm-up", False),
                ("07:20",         "Athletes' introduction", False),
                ("07:30",         "★ ELITE WOMEN START — 1500 / 40 / 10", True),
                ("~09:30",        "Projected finish", False),
            ]),
            ("Sun • Aug 9", "ELITE MEN — STANDARD", [
                ("06:35",         "Call time at Hotel Esplendor (Men)", False),
                ("07:50 – 07:55", "Departure from hotel  →  arrive El Delta 08:00-08:20", False),
                ("09:00 – 09:45", "Athletes' Lounge opens — check-in", False),
                ("09:10 – 10:00", "Transition area check-in", False),
                ("09:45 – 10:15", "Swim warm-up", False),
                ("10:20",         "Athletes' introduction", False),
                ("10:30",         "★ ELITE MEN START — 1500 / 40 / 10", True),
                ("12:30",         "Award ceremony — Elite Women & Elite Men", False),
            ]),
        ],
    },
    "weihai": {
        "title":      "2026 WTCS Weihai — Race Week",
        "date_range": "August 27 – 30, 2026",
        "venue_note": "Wave to Wonderland, Huancui District, Weihai, Shandong (CST, UTC+8)",
        "race_date":  "2026-08-29",
        "race_starts": [("Elite Men", "09:00"), ("Elite Women", "12:00")],
        "days": [
            ("Thu • Aug 27", "Arrival & Pool Training", [
                ("09:00 – 18:00", "Information desk + bike mechanic — Weihaiwei Hotel", False),
                ("08:00 – 10:00", "Pool training — Weihai Swimming Pool (25 m, 4 lanes)", False),
                ("14:00 – 16:00", "Pool training — afternoon session", False),
                ("07:30 / 13:30", "Free LOC shuttle hotel → pool (also 08:30 / 14:30)", False),
            ]),
            ("Fri • Aug 28", "Familiarization & Briefing", [
                ("09:30 – 10:00", "Elite bike course familiarization — 3 laps, officials lead, NO overtaking", False),
                ("10:00 – 11:00", "Elite swim course familiarization — all buoys set, boats + lifeguards", False),
                ("10:15 – 10:45", "Coaches' feedback on FOP to Technical Delegate", False),
                ("14:30 – 15:00", "Pre-race press conference — Brigh Radiance Hotel 12/F", False),
                ("16:00 – 16:30", "★ Elite Athletes' Briefing — Brigh Radiance Hotel 14/F", True),
                ("16:30 – 17:00", "Team medical meeting + race package distribution", False),
                ("18:00 – 20:00", "Pasta Party — Weihaiwei Hotel 2/F (veg + halal available)", False),
            ]),
            ("Sat • Aug 29", "ELITE MEN — 09:00", [
                ("07:30 – 08:30", "Elite Men athletes' lounge check-in", False),
                ("07:45 – 08:45", "Elite Men transition zone check-in", False),
                ("08:15 – 08:45", "Elite Men swim warm-up", False),
                ("08:30 – 08:45", "Opening ceremony — main venue", False),
                ("08:45 / 08:50", "Elite Men line-up (lounge) / introduction (swim exit)", False),
                ("09:00",         "★ ELITE MEN START — 1.5 / 40 / 10", True),
                ("10:55",         "Elite Men awards ceremony", False),
            ]),
            ("Sat • Aug 29", "ELITE WOMEN — 12:00", [
                ("10:30 – 11:30", "Elite Women athletes' lounge check-in", False),
                ("11:00 – 11:45", "Elite Women transition zone check-in", False),
                ("11:15 – 11:45", "Elite Women swim warm-up", False),
                ("11:15 – 11:45", "Elite coaches run course measurement", False),
                ("11:45 / 11:50", "Elite Women line-up / introduction", False),
                ("12:00",         "★ ELITE WOMEN START — hottest part of the day", True),
                ("14:10",         "Elite Women awards ceremony", False),
            ]),
        ],
    },
}


def build_bike_profile_chart(profile: dict) -> io.BytesIO:
    """Synthesize a stylised elevation profile over the full bike distance.

    The chart is a representative single-loop wave repeated `loops` times,
    scaled to match the published gain_per_lap_m. It is intentionally schematic
    — exact contours require a GPS-recorded FIT/GPX. The shape uses one big
    climb + descent + short rolling section per loop, matching the bike course
    map for Huatulco (and most short-circuit World Cups).
    """
    import numpy as _np

    loops      = int(profile.get("loops", 1))
    loop_km    = float(profile.get("loop_km", 5.0))
    gain_lap   = float(profile.get("gain_per_lap_m", 80.0))
    loss_lap   = float(profile.get("loss_per_lap_m", gain_lap))

    pts_per_loop = 400
    xs_loop = _np.linspace(0, loop_km, pts_per_loop, endpoint=False)
    # Asymmetric wave: short steep climb (first 40%) then a longer descent (40-90%)
    # then a small bump back to baseline (90-100%). Scale to gain_lap.
    def _loop_shape(x):
        t = x / loop_km
        y = _np.where(
            t < 0.4,
            (t / 0.4) ** 1.2,
            _np.where(
                t < 0.9,
                1.0 - ((t - 0.4) / 0.5) ** 1.1,
                0.15 * _np.sin((t - 0.9) / 0.1 * _np.pi),
            ),
        )
        return y
    base = _loop_shape(xs_loop)
    base = base - base.min()
    base = base / max(base.max(), 1e-6) * gain_lap

    xs_full = _np.concatenate([xs_loop + i * loop_km for i in range(loops)])
    ys_full = _np.tile(base, loops)

    fig, ax = plt.subplots(figsize=(11.0, 2.4), dpi=150)
    ax.fill_between(xs_full, ys_full, 0, color="#4472C4", alpha=0.35, linewidth=0)
    ax.plot(xs_full, ys_full, color="#002060", linewidth=1.4)

    for i in range(1, loops):
        ax.axvline(i * loop_km, color="#C00000", linestyle="--", linewidth=0.9, alpha=0.7)
        ax.text(i * loop_km, ys_full.max() * 1.02, f"Lap {i + 1}",
                ha="center", va="bottom", fontsize=8, color="#C00000")
    ax.text(loop_km / 2, ys_full.max() * 1.02, "Lap 1",
            ha="center", va="bottom", fontsize=8, color="#C00000")

    ax.set_xlim(0, loops * loop_km)
    ax.set_ylim(0, ys_full.max() * 1.18)
    ax.set_xlabel("Distance (km)", fontsize=9)
    ax.set_ylabel("Relative elevation (m)", fontsize=9)
    ax.set_title(f"Bike Course — {loops} × {loop_km:.0f} km loop ({loops * loop_km:.0f} km total)",
                 fontsize=11, fontweight="bold", color="#002060", pad=10)
    ax.tick_params(labelsize=8)
    ax.grid(axis="y", alpha=0.2)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    plt.tight_layout()
    return fig_to_image(fig)


def add_bike_course_profile_slide(prs: Presentation, venue: str, all_rows: list[dict]):
    """Dedicated bike course slide — published course specs + synthesized profile + tactical notes.

    Only added when BIKE_COURSE_PROFILES has data for this venue. Optionally embeds
    a course-map image from ppt files/support/{venue_lc}_bike_course.png if present.
    """
    profile = BIKE_COURSE_PROFILES.get(venue.lower().strip())
    if not profile:
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Bike Course Profile", venue)

    # ── Left: course map image (if present at expected path) ──────────────────
    map_path = os.path.join(REPO_ROOT, "ppt files", "support",
                            f"{venue.lower().strip()}_bike_course.png")
    map_left  = Inches(0.3)
    map_top   = Inches(1.25)
    map_w     = Inches(5.6)
    map_h     = Inches(3.6)

    if os.path.exists(map_path):
        slide.shapes.add_picture(map_path, map_left, map_top, map_w, map_h)
    else:
        ph = slide.shapes.add_shape(1, map_left, map_top, map_w, map_h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_GRAY
        ph.line.color.rgb = MID_GRAY
        _add_textbox(slide,
                     f"[Drop course map here]\n\n"
                     f"Save the organiser course map as:\n"
                     f"ppt files/support/{venue.lower().strip()}_bike_course.png\n\n"
                     f"It will be embedded automatically on the next run.",
                     map_left + Inches(0.2), map_top + Inches(1.2),
                     map_w - Inches(0.4), Inches(1.4),
                     font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Source caption under map
    _add_textbox(slide, f"Source: {profile.get('source', '—')}",
                 map_left, map_top + map_h + Inches(0.05),
                 map_w, Inches(0.25),
                 font_size=9, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # ── Right: stats card ─────────────────────────────────────────────────────
    stats_left = Inches(6.1)
    stats_top  = Inches(1.25)
    stats_w    = Inches(6.9)

    # Big distance card
    hero_h = Inches(0.9)
    hero = slide.shapes.add_shape(1, stats_left, stats_top, stats_w, hero_h)
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()
    _add_textbox(slide,
                 f"{profile['total_km']:g} km   •   {profile['loops']} × {profile['loop_km']:g} km loops",
                 stats_left, stats_top + Inches(0.12),
                 stats_w, Inches(0.4),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, profile.get("surface", ""),
                 stats_left, stats_top + Inches(0.54),
                 stats_w, Inches(0.32),
                 font_size=11, color=RGBColor(0xAA, 0xBB, 0xDD),
                 align=PP_ALIGN.CENTER, italic=True)

    # Stats grid (3 cols × 2 rows). Tolerate None values for venues with no
    # published elevation data (e.g. flat coastal courses).
    grid_top = stats_top + hero_h + Inches(0.15)
    grid_h   = Inches(1.5)
    cell_w   = stats_w / 3
    g  = profile.get("gain_per_lap_m")
    l_ = profile.get("loss_per_lap_m")
    avg_g  = profile.get("avg_grade_pct")
    max_p  = profile.get("max_grade_pos")
    max_n  = profile.get("max_grade_neg")
    wind   = profile.get("wind")           # optional — e.g. "Coastal (~20 km/h)"
    loops_ = profile.get("loops") or 1
    loop_km_val = profile.get("loop_km")
    # Loops cell prefers exact lap distance ("4 × 5.5 km") over loose rounding
    loops_str = (f"{loops_} × {loop_km_val:g} km"
                 if loop_km_val is not None else f"{loops_}")
    # If wind is set, it bumps Max Descent out of the grid (more actionable for
    # flat coastal courses like Quiberon/Montreal).
    middle_row_cell = (("Wind", wind) if wind
                       else ("Max Descent",
                             f"{max_n:.1f}%" if max_n is not None else "—"))
    grid_cells = [
        ("Elevation / Lap", f"+{g:.1f} m / {l_:.1f} m" if (g is not None and l_ is not None) else "Flat (no data)"),
        ("Total Climb",     f"~{g * loops_:.0f} m" if g is not None else "—"),
        ("Avg Grade",       f"{avg_g:.1f}%" if avg_g is not None else "—"),
        ("Max Climb Grade", f"+{max_p:.1f}%" if max_p is not None else "—"),
        middle_row_cell,
        ("Loops",           loops_str),
    ]
    for idx, (label, value) in enumerate(grid_cells):
        col = idx % 3
        row = idx // 3
        cl = stats_left + cell_w * col + Inches(0.05)
        cw = cell_w - Inches(0.1)
        ct = grid_top + (Inches(0.78) * row)
        bx = slide.shapes.add_shape(1, cl, ct, cw, Inches(0.7))
        bx.fill.solid()
        bx.fill.fore_color.rgb = LIGHT_GRAY
        bx.line.color.rgb = MID_GRAY
        _add_textbox(slide, label, cl + Inches(0.05), ct + Inches(0.04),
                     cw - Inches(0.1), Inches(0.22),
                     font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_textbox(slide, value, cl + Inches(0.05), ct + Inches(0.26),
                     cw - Inches(0.1), Inches(0.38),
                     font_size=15, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Tactical notes
    notes_top = grid_top + grid_h + Inches(0.25)
    notes_bar = slide.shapes.add_shape(1, stats_left, notes_top, stats_w, Inches(0.32))
    notes_bar.fill.solid()
    notes_bar.fill.fore_color.rgb = RED
    notes_bar.line.fill.background()
    _add_textbox(slide, "Tactical Implications",
                 stats_left + Inches(0.1), notes_top + Inches(0.03),
                 stats_w - Inches(0.2), Inches(0.26),
                 font_size=12, bold=True, color=WHITE)

    body_top = notes_top + Inches(0.4)
    for i, note in enumerate(profile.get("key_features", [])):
        _add_textbox(slide, f"•  {note}",
                     stats_left + Inches(0.1),
                     body_top + Inches(0.32 * i),
                     stats_w - Inches(0.2), Inches(0.34),
                     font_size=10.5, color=DARK_GRAY)

    # ── Bottom: synthesized elevation profile chart (skip for flat courses) ───
    has_elev = profile.get("gain_per_lap_m") is not None
    if not has_elev:
        chart_left   = Inches(0.3)
        chart_top    = Inches(5.6)
        chart_width  = Inches(12.73)
        chart_height = Inches(1.55)
        ph = slide.shapes.add_shape(1, chart_left, chart_top, chart_width, chart_height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_GRAY
        ph.line.color.rgb = MID_GRAY
        _add_textbox(slide,
                     "Flat course — published profile shows no significant climbs.\n"
                     "No elevation chart synthesized; if a race GPS FIT becomes available, "
                     "real contours will replace this block.",
                     chart_left, chart_top + Inches(0.5),
                     chart_width, Inches(0.6),
                     font_size=11, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    else:
        try:
            chart_buf = build_bike_profile_chart(profile)
            chart_left   = Inches(0.3)
            chart_top    = Inches(5.45)
            chart_width  = Inches(12.73)
            chart_height = Inches(1.85)
            slide.shapes.add_picture(chart_buf, chart_left, chart_top, chart_width, chart_height)
            _add_textbox(slide,
                         "Stylised — exact contours from a GPS-recorded race FIT would refine this.",
                         chart_left, chart_top + chart_height - Inches(0.05),
                         chart_width, Inches(0.22),
                         font_size=8.5, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
        except Exception as exc:
            _add_textbox(slide, f"[Profile chart could not be rendered: {exc}]",
                         Inches(0.3), Inches(5.5), Inches(12.73), Inches(0.4),
                         font_size=10, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    return True


def _add_discipline_profile_slide(prs: Presentation, venue: str, discipline: str,
                                  profile: dict, accent_color: RGBColor,
                                  map_filename_suffix: str) -> bool:
    """Shared layout for Swim / Run course profile slides.

    Left panel: course map image (if dropped at ppt files/support/{venue}_{suffix}.png)
                or a styled placeholder.
    Right panel: hero card (distance / laps / format), 2×3 stat grid, key features,
                and a 'Missing from race brief' note when relevant.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"{discipline} Course Profile", venue)

    # ── Left: course map slot ────────────────────────────────────────────────
    map_path = os.path.join(REPO_ROOT, "ppt files", "support",
                            f"{venue.lower().strip()}_{map_filename_suffix}.png")
    map_left = Inches(0.3)
    map_top  = Inches(1.25)
    map_w    = Inches(5.6)
    map_h    = Inches(4.2)

    if os.path.exists(map_path):
        slide.shapes.add_picture(map_path, map_left, map_top, map_w, map_h)
    else:
        ph = slide.shapes.add_shape(1, map_left, map_top, map_w, map_h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_GRAY
        ph.line.color.rgb = MID_GRAY
        _add_textbox(slide,
                     f"[Drop {discipline.lower()} course map here]\n\n"
                     f"Save the organiser map as:\n"
                     f"ppt files/support/{venue.lower().strip()}_{map_filename_suffix}.png\n\n"
                     f"It will be embedded automatically on the next run.",
                     map_left + Inches(0.2), map_top + Inches(1.5),
                     map_w - Inches(0.4), Inches(1.4),
                     font_size=10, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    _add_textbox(slide, f"Source: {profile.get('source', '—')}",
                 map_left, map_top + map_h + Inches(0.05),
                 map_w, Inches(0.25),
                 font_size=9, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

    # ── Right: hero + stats grid + features ──────────────────────────────────
    right_left = Inches(6.1)
    right_top  = Inches(1.25)
    right_w    = Inches(6.9)

    hero_h = Inches(0.9)
    hero = slide.shapes.add_shape(1, right_left, right_top, right_w, hero_h)
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    total_km = profile.get("total_km")
    laps     = profile.get("laps")
    loop_km  = profile.get("loop_km")

    def _fmt_distance(km: float | None) -> str:
        if not km:
            return "TBD"
        if km < 1:
            return f"{int(round(km * 1000))} m"
        return f"{km:.1f} km"

    hero_main = _fmt_distance(total_km) if total_km else "Distance TBD"
    if laps and laps > 1:
        per_lap = _fmt_distance(loop_km) if loop_km else _fmt_distance(total_km / laps if total_km else None)
        hero_main += f"   •   {laps} × {per_lap}"
    elif laps == 1:
        hero_main += "   •   single lap"
    _add_textbox(slide, hero_main,
                 right_left, right_top + Inches(0.12),
                 right_w, Inches(0.4),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    subtitle = profile.get("format") or profile.get("surface") or ""
    _add_textbox(slide, subtitle,
                 right_left, right_top + Inches(0.54),
                 right_w, Inches(0.32),
                 font_size=11, color=RGBColor(0xAA, 0xBB, 0xDD),
                 align=PP_ALIGN.CENTER, italic=True)

    # Stat grid — 3 cols × 2 rows. Cells unknown to this discipline render '—'.
    grid_top = right_top + hero_h + Inches(0.15)
    cell_w   = right_w / 3
    if discipline.lower() == "swim":
        water_temp = profile.get("water_temp_c")
        water_src  = profile.get("water_temp_source", "unavailable")
        # Wetsuit decision: use live/climatology SST if available, else fall back
        # to the published expected range from the profile, else "TBD".
        # Elite rules: mandatory <16 °C, optional 16–20 °C, forbidden >22 °C.
        expected_range = profile.get("expected_water_temp_range_c")
        ref_temp = water_temp
        if ref_temp is None and expected_range:
            ref_temp = sum(expected_range) / 2
        def _wetsuit_band(t: float) -> int:
            """Rule band index so a published range that crosses a threshold is
            reported as uncertain instead of collapsing to a midpoint verdict."""
            if t >= 22:  return 3   # forbidden
            if t >= 20:  return 2   # likely non-wetsuit
            if t >= 16:  return 1   # optional
            return 0                # mandatory

        straddles = (water_temp is None and expected_range
                     and _wetsuit_band(expected_range[0]) != _wetsuit_band(expected_range[1]))

        if ref_temp is None:
            wetsuit = "TBD — depends on day-of water temp"
        elif straddles:
            wetsuit = (f"Straddles cut-off "
                       f"({expected_range[0]:.0f}–{expected_range[1]:.0f} °C)")
        elif ref_temp >= 22:
            wetsuit = "Forbidden (> 22 °C)"
        elif ref_temp >= 20:
            wetsuit = f"Likely non-wetsuit ({ref_temp:.0f} °C)"
        elif ref_temp >= 16:
            wetsuit = f"Wetsuit optional ({ref_temp:.0f} °C)"
        else:
            wetsuit = f"Wetsuit mandatory (< 16 °C)"

        if water_temp is not None:
            tag = {"forecast": "live", "climatology": "7-yr avg"}.get(water_src, "")
            water_label = f"{water_temp:.1f} °C ({tag})" if tag else f"{water_temp:.1f} °C"
        elif expected_range:
            water_label = f"{expected_range[0]:.1f}–{expected_range[1]:.1f} °C (expected)"
        else:
            water_label = "Awaiting marine data"
        layout     = profile.get("layout") or profile.get("format", "—")
        grid_cells = [
            ("Distance",      _fmt_distance(total_km)),
            ("Layout",        layout.split(" — ")[0] if " — " in (layout or "") else (layout or "—")),
            ("Start",         (profile.get("start_type") or "—").split(" at ")[0]),
            ("Water Temp",    water_label),
            ("Wetsuit",       wetsuit),
            ("Lap Count",     f"{laps}" if laps else "—"),
        ]
    else:  # run
        gain_lap = profile.get("gain_per_lap_m")
        max_pos  = profile.get("max_grade_pos")
        max_neg  = profile.get("max_grade_neg")
        avg_g    = profile.get("avg_grade_pct")
        # Heat risk: prefer explicit profile field; default to "—" so it isn't
        # falsely flagged HIGH for cool venues. Forecast slide is the
        # authoritative heat-stress source.
        heat_risk = profile.get("heat_risk", "—")
        grid_cells = [
            ("Distance",      _fmt_distance(total_km)),
            ("Lap Structure", f"{laps} × {_fmt_distance(loop_km)}" if (laps and loop_km) else (f"{laps} laps" if laps else "Not published")),
            ("Total Climb",   f"~{gain_lap * laps:.0f} m" if (gain_lap and laps) else (f"+{gain_lap:.0f} m / lap" if gain_lap else "Not published")),
            ("Avg Grade",     f"{avg_g:.1f}%" if avg_g else "—"),
            ("Max ▲ / ▼",    f"+{max_pos:.1f}% / {max_neg:.1f}%" if (max_pos and max_neg) else (f"+{max_pos:.1f}%" if max_pos else "Not published")),
            ("Heat Risk",     heat_risk),
        ]

    for idx, (label, value) in enumerate(grid_cells):
        col = idx % 3
        row = idx // 3
        cl = right_left + cell_w * col + Inches(0.05)
        cw = cell_w - Inches(0.1)
        ct = grid_top + (Inches(0.78) * row)
        bx = slide.shapes.add_shape(1, cl, ct, cw, Inches(0.7))
        bx.fill.solid()
        bx.fill.fore_color.rgb = LIGHT_GRAY
        bx.line.color.rgb = MID_GRAY
        _add_textbox(slide, label, cl + Inches(0.05), ct + Inches(0.04),
                     cw - Inches(0.1), Inches(0.22),
                     font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_textbox(slide, value, cl + Inches(0.05), ct + Inches(0.26),
                     cw - Inches(0.1), Inches(0.38),
                     font_size=14, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Key features panel
    notes_top = grid_top + Inches(1.65)
    notes_bar = slide.shapes.add_shape(1, right_left, notes_top, right_w, Inches(0.32))
    notes_bar.fill.solid()
    notes_bar.fill.fore_color.rgb = accent_color
    notes_bar.line.fill.background()
    _add_textbox(slide, "Tactical Notes",
                 right_left + Inches(0.1), notes_top + Inches(0.03),
                 right_w - Inches(0.2), Inches(0.26),
                 font_size=12, bold=True, color=WHITE)

    body_top = notes_top + Inches(0.4)
    for i, note in enumerate(profile.get("key_features", [])):
        _add_textbox(slide, f"•  {note}",
                     right_left + Inches(0.1),
                     body_top + Inches(0.32 * i),
                     right_w - Inches(0.2), Inches(0.34),
                     font_size=10.5, color=DARK_GRAY)

    # ── Bottom: 'Missing from race brief' caveat strip ───────────────────────
    missing = profile.get("missing") or []
    if missing:
        miss_top = Inches(6.6)
        miss_bar = slide.shapes.add_shape(1, Inches(0.3), miss_top, Inches(12.73), Inches(0.32))
        miss_bar.fill.solid()
        miss_bar.fill.fore_color.rgb = LIGHT_GRAY
        miss_bar.line.color.rgb = MID_GRAY
        text = "Not in race brief: " + "  •  ".join(missing)
        _add_textbox(slide, text,
                     Inches(0.4), miss_top + Inches(0.05),
                     Inches(12.5), Inches(0.25),
                     font_size=9.5, italic=True, color=DARK_GRAY, align=PP_ALIGN.LEFT)

    return True


def add_swim_course_profile_slide(prs: Presentation, venue: str,
                                  water_temp_c: float | None = None,
                                  water_temp_source: str = "unavailable") -> bool:
    profile = SWIM_COURSE_PROFILES.get(venue.lower().strip())
    if not profile:
        return False
    profile = dict(profile)
    if water_temp_c is not None:
        profile["water_temp_c"] = water_temp_c
        profile["water_temp_source"] = water_temp_source
    return _add_discipline_profile_slide(
        prs, venue, "Swim", profile, RED, "swim_course"
    )


def add_run_course_profile_slide(prs: Presentation, venue: str) -> bool:
    profile = RUN_COURSE_PROFILES.get(venue.lower().strip())
    if not profile:
        return False
    return _add_discipline_profile_slide(
        prs, venue, "Run", profile, RED, "run_course"
    )


def add_race_week_schedule_slide(prs: Presentation, venue: str) -> bool:
    """Race-week timetable. Three day-columns side by side."""
    schedule = EVENT_SCHEDULES.get(venue.lower().strip())
    if not schedule:
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Race Week Schedule", venue)

    # Sub-header bar with date range + venue
    sub_top = Inches(1.2)
    sub_bar = slide.shapes.add_shape(1, Inches(0.3), sub_top, Inches(12.73), Inches(0.42))
    sub_bar.fill.solid()
    sub_bar.fill.fore_color.rgb = LIGHT_GRAY
    sub_bar.line.color.rgb = MID_GRAY
    _add_textbox(slide,
                 f"{schedule['date_range']}   •   {schedule.get('venue_note', '')}",
                 Inches(0.4), sub_top + Inches(0.07),
                 Inches(12.5), Inches(0.3),
                 font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    days = schedule["days"]
    n_days = len(days)
    col_top    = Inches(1.85)
    gap_in     = 0.25
    total_w_in = 12.73 - gap_in * (n_days - 1)
    col_w_in   = total_w_in / n_days
    col_w      = Inches(col_w_in)
    col_gap    = Inches(gap_in)
    col_left0  = Inches(0.3)

    for di, (day_label, day_subtitle, items) in enumerate(days):
        cl = col_left0 + (col_w + col_gap) * di

        # Day header
        header_h = Inches(0.5)
        header = slide.shapes.add_shape(1, cl, col_top, col_w, header_h)
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()
        _add_textbox(slide, day_label,
                     cl + Inches(0.1), col_top + Inches(0.04),
                     col_w - Inches(0.2), Inches(0.22),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, day_subtitle,
                     cl + Inches(0.1), col_top + Inches(0.26),
                     col_w - Inches(0.2), Inches(0.22),
                     font_size=9.5, color=RGBColor(0xAA, 0xBB, 0xDD),
                     italic=True, align=PP_ALIGN.CENTER)

        # Items table (time | activity)
        items_top = col_top + header_h + Inches(0.1)
        n_rows = len(items)
        tbl_h = Inches(0.55 * n_rows)
        tbl_shape = slide.shapes.add_table(n_rows, 2, cl, items_top, col_w, tbl_h)
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(col_w_in * 0.36)
        tbl.columns[1].width = Inches(col_w_in * 0.64)
        for ri, (tm, act, is_hi) in enumerate(items):
            bg     = RGBColor(0xFD, 0xEC, 0xEC) if is_hi else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
            fg     = RED if is_hi else DARK_GRAY
            _set_cell(tbl.cell(ri, 0), tm,
                      font_size=10, bold=is_hi, color=fg, bg_color=bg, align=PP_ALIGN.CENTER)
            _set_cell(tbl.cell(ri, 1), act,
                      font_size=10, bold=is_hi, color=fg, bg_color=bg, align=PP_ALIGN.LEFT)

    # Footer caption — venue/timezone pulled from schedule.venue_note when present
    tz_note = schedule.get("venue_note") or f"local to {venue}"
    _add_textbox(slide,
                 f"Times {tz_note}. Schedule confirmed from World Triathlon race-info page / "
                 f"athlete guide; subject to organiser updates.",
                 Inches(0.3), Inches(6.95), Inches(12.73), Inches(0.3),
                 font_size=9, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    return True


def _heat_band(apparent_c: float | None) -> tuple[str, RGBColor]:
    """Heat-stress band based on apparent temperature."""
    if apparent_c is None:
        return ("UNKNOWN", MID_GRAY)
    if apparent_c < 24:    return ("LOW",      RGBColor(0x1B, 0x7F, 0x3A))
    if apparent_c < 28:    return ("MODERATE", RGBColor(0xE6, 0xA8, 0x17))
    if apparent_c < 32:    return ("HIGH",     RGBColor(0xE6, 0x6A, 0x00))
    return                  ("EXTREME",        RGBColor(0xC0, 0x00, 0x00))


def add_race_day_forecast_slide(prs: Presentation, venue: str,
                                coords: tuple | None,
                                race_date: str | None,
                                race_starts: list[tuple[str, str]] | None = None) -> bool:
    """Race-day forecast slide. Pulls the forward forecast for race_date if it is
    inside Open-Meteo's 16-day window; otherwise falls back to a 7-year climatology
    averaged across the same calendar day."""
    if coords is None or race_date is None:
        return False
    lat, lon = coords

    forecast = fetch_openmeteo_forecast(lat, lon, race_date)
    mode = "forecast"
    if not forecast:
        forecast = fetch_openmeteo_climatology(lat, lon, race_date)
        mode = "climatology"
    if not forecast:
        return False

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Race-Day Weather Outlook", venue)

    mode_label = ("Forward forecast — Open-Meteo (≤16 days)"
                  if mode == "forecast"
                  else "Historical climatology — same calendar day, 7-year mean")
    sub_top = Inches(1.2)
    sub_bar = slide.shapes.add_shape(1, Inches(0.3), sub_top, Inches(12.73), Inches(0.42))
    sub_bar.fill.solid()
    sub_bar.fill.fore_color.rgb = LIGHT_GRAY
    sub_bar.line.color.rgb = MID_GRAY
    _add_textbox(slide,
                 f"{race_date}   •   {mode_label}",
                 Inches(0.4), sub_top + Inches(0.07),
                 Inches(12.5), Inches(0.3),
                 font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Hourly chart (top-left) ───────────────────────────────────────────────
    times = forecast["time"]
    temps = forecast["temperature_2m"]
    apparent = forecast.get("apparent_temperature") or [None] * len(times)
    humidity = forecast["relative_humidity_2m"]
    wind     = forecast["wind_speed_10m"]
    uv       = forecast.get("uv_index") or [None] * len(times)

    fig, ax1 = plt.subplots(figsize=(8.4, 3.6), dpi=150)
    x = list(range(len(times)))
    ax1.plot(x, temps,    color="#C00000", linewidth=2.2, marker="o", markersize=4, label="Air temp (°C)")
    ax1.plot(x, apparent, color="#E66A00", linewidth=1.6, linestyle="--", marker="x",
             markersize=4, label="Apparent (°C)", alpha=0.85)
    ax1.set_xlabel("Hour (local)", fontsize=10)
    ax1.set_ylabel("Temperature (°C)", color="#C00000", fontsize=10)
    ax1.tick_params(axis="y", labelcolor="#C00000")
    ax1.set_xticks(x)
    ax1.set_xticklabels(times, rotation=45, ha="right", fontsize=8)

    ax2 = ax1.twinx()
    ax2.plot(x, humidity, color="#002060", linewidth=2.0, marker="s", markersize=4, label="Humidity (%)")
    ax2.set_ylabel("Humidity (%)", color="#002060", fontsize=10)
    ax2.tick_params(axis="y", labelcolor="#002060")
    ax2.set_ylim(0, 100)

    # Race-start vertical markers
    if race_starts:
        for label, hhmm in race_starts:
            if hhmm in times:
                xi = times.index(hhmm)
                ax1.axvline(xi, color="black", linestyle=":", linewidth=1.5, alpha=0.7)
                ax1.text(xi, ax1.get_ylim()[1] * 0.98, label,
                         rotation=90, va="top", ha="right", fontsize=8.5,
                         color="black", fontweight="bold")

    ax1.grid(True, alpha=0.25)
    ax1.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)
    ax1.set_title("Hourly outlook — race window", fontsize=11, fontweight="bold", color="#002060")
    h1, l1 = ax1.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax1.legend(h1 + h2, l1 + l2, loc="lower right", fontsize=8, framealpha=0.9)
    plt.tight_layout()
    img_buf = fig_to_image(fig)
    slide.shapes.add_picture(img_buf, Inches(0.3), Inches(1.75), Inches(7.7), Inches(3.4))

    # ── Race-window summary stat cards (right side) ───────────────────────────
    right_left = Inches(8.2)
    right_w    = Inches(4.83)

    def _hour_idx(hhmm: str) -> int | None:
        """Find the time index closest to hhmm. Snaps half-hour starts (e.g.
        '06:30') to the hour bucket that contains them ('06:00')."""
        if hhmm in times:
            return times.index(hhmm)
        try:
            target_min = int(hhmm[:2]) * 60 + int(hhmm[3:5])
        except (ValueError, IndexError):
            return None
        best_i, best_d = None, None
        for i, t in enumerate(times):
            try:
                tm = int(t[:2]) * 60 + int(t[3:5])
            except (ValueError, IndexError):
                continue
            d = abs(tm - target_min)
            if best_d is None or d < best_d:
                best_d = d
                best_i = i
        return best_i if best_d is not None and best_d <= 60 else None

    def _race_window_avg(values: list, start_idx: int, end_idx: int) -> float | None:
        vals = [v for v in values[start_idx:end_idx + 1] if v is not None]
        return sum(vals) / len(vals) if vals else None

    cards: list[tuple[str, list[tuple[str, str, RGBColor]]]] = []
    for label, hhmm in (race_starts or []):
        idx = _hour_idx(hhmm)
        if idx is None:
            continue
        end = min(idx + 2, len(times) - 1)
        avg_temp = _race_window_avg(temps, idx, end)
        avg_app  = _race_window_avg(apparent, idx, end) if any(a is not None for a in apparent) else None
        avg_hum  = _race_window_avg(humidity, idx, end)
        avg_wind = _race_window_avg(wind, idx, end)
        peak_uv  = max((v for v in uv[idx:end + 1] if v is not None), default=None)
        heat_label, heat_color = _heat_band(avg_app if avg_app is not None else avg_temp)
        cards.append((label + f" — {hhmm}", [
            ("Temperature",    f"{avg_temp:.1f} °C" if avg_temp is not None else "—", DARK_GRAY),
            ("Apparent / Heat", f"{avg_app:.1f} °C" if avg_app is not None else "—", DARK_GRAY),
            ("Humidity",       f"{avg_hum:.0f}%" if avg_hum is not None else "—", DARK_GRAY),
            ("Wind",           f"{avg_wind:.1f} km/h" if avg_wind is not None else "—", DARK_GRAY),
            ("UV peak",        f"{peak_uv:.1f}" if peak_uv is not None else "—", DARK_GRAY),
            ("Heat Stress",    heat_label, heat_color),
        ]))

    card_top = Inches(1.75)
    for ci, (card_title, rows_) in enumerate(cards[:2]):
        ct = card_top + Inches(0.05 + ci * 1.75)
        hdr = slide.shapes.add_shape(1, right_left, ct, right_w, Inches(0.35))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = NAVY
        hdr.line.fill.background()
        _add_textbox(slide, card_title,
                     right_left + Inches(0.1), ct + Inches(0.04),
                     right_w - Inches(0.2), Inches(0.27),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        n_cols = 3
        cell_w = right_w / n_cols
        for ri, (k, v, color) in enumerate(rows_):
            col = ri % n_cols
            row = ri // n_cols
            cl = right_left + cell_w * col
            cw = cell_w - Inches(0.04)
            cyt = ct + Inches(0.4 + row * 0.62)
            bx = slide.shapes.add_shape(1, cl + Inches(0.02), cyt, cw, Inches(0.58))
            bx.fill.solid()
            bx.fill.fore_color.rgb = LIGHT_GRAY
            bx.line.color.rgb = MID_GRAY
            _add_textbox(slide, k, cl + Inches(0.04), cyt + Inches(0.03),
                         cw - Inches(0.04), Inches(0.2),
                         font_size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            _add_textbox(slide, v, cl + Inches(0.04), cyt + Inches(0.22),
                         cw - Inches(0.04), Inches(0.34),
                         font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ── Bottom strip: tactical recommendations ────────────────────────────────
    rec_top = Inches(5.3)
    rec_bar = slide.shapes.add_shape(1, Inches(0.3), rec_top, Inches(12.73), Inches(0.32))
    rec_bar.fill.solid()
    rec_bar.fill.fore_color.rgb = RED
    rec_bar.line.fill.background()
    _add_textbox(slide, "Hydration & Cooling Cues",
                 Inches(0.4), rec_top + Inches(0.03),
                 Inches(12.5), Inches(0.26),
                 font_size=12, bold=True, color=WHITE)

    # Build cues from the data
    cues: list[str] = []
    avg_window = _race_window_avg(apparent, 0, len(times) - 1) if any(a is not None for a in apparent) else None
    if avg_window is None:
        avg_window = _race_window_avg(temps, 0, len(times) - 1)
    if avg_window and avg_window >= 30:
        cues.append("Pre-cool 20 min before start (ice slurry, ice towel on neck); fluid 6–8 ml/kg in the hour pre-race")
        cues.append("Plan 2 ice-sock handoffs on the run — every other lap; sponges at every aid station on bike & run")
    elif avg_window and avg_window >= 25:
        cues.append("Standard heat protocol — pre-cool 10 min, sponges/ice at every aid station, target 600–800 ml/h on bike")
    else:
        cues.append("Cool to mild outlook — standard fluid plan (~500 ml/h), no special pre-cooling required")

    # Female vs male delta cue
    if len(cards) == 2:
        f_temp = next((float(v.split()[0]) for k, v, _ in cards[0][1] if k == "Temperature" and v != "—"), None)
        m_temp = next((float(v.split()[0]) for k, v, _ in cards[1][1] if k == "Temperature" and v != "—"), None)
        if f_temp is not None and m_temp is not None:
            delta = m_temp - f_temp
            if delta >= 2:
                # Pull the men's start time from its own card label ("Elite Men — 10:30")
                # rather than hard-coding a single venue's start.
                _m_start = cards[1][0].split("—")[-1].strip() if "—" in cards[1][0] else ""
                _when = f"the {_m_start} start" if _m_start else "the later start"
                cues.append(f"Men's race runs {delta:.1f} °C hotter than women's — heat-acclim work + heavier cooling kit for {_when}")
            elif delta >= 1:
                cues.append(f"Men's race ~{delta:.1f} °C warmer — add one extra cooling touchpoint vs. women's plan")
    # UV
    peak_uv_all = max((v for v in uv if v is not None), default=None)
    if peak_uv_all and peak_uv_all >= 8:
        cues.append(f"UV peaks at {peak_uv_all:.1f} (very high) — long-sleeve race kit + sunscreen reapply 30 min pre-start")

    # Wind — significant for coastal/exposed bike courses
    peak_wind = max((v for v in wind if v is not None), default=None)
    if peak_wind and peak_wind >= 25:
        cues.append(f"Wind peaks at {peak_wind:.0f} km/h — pack will split on exposed coastal sections; "
                    f"shallow front wheel + position discipline through crosswind segments")
    elif peak_wind and peak_wind >= 15:
        cues.append(f"Notable wind ({peak_wind:.0f} km/h) — drafting protection matters on exposed coastal bike sections; "
                    f"sight more frequently in chop")

    # Cold-water swim cue
    avg_temp_full = _race_window_avg(temps, 0, len(times) - 1)
    if avg_temp_full is not None and avg_temp_full < 16:
        cues.append("Cold air pre-race — extended warm-up + warm clothing to lounge; "
                    "wetsuit decision likely tilts mandatory if water tracks the same direction")

    body_top = rec_top + Inches(0.4)
    for i, cue in enumerate(cues[:4]):
        _add_textbox(slide, f"•  {cue}",
                     Inches(0.4), body_top + Inches(0.32 * i),
                     Inches(12.5), Inches(0.34),
                     font_size=10.5, color=DARK_GRAY)

    # Footnote
    _add_textbox(slide,
                 ("Data source: Open-Meteo. "
                  + ("Forecast refreshes every run. "
                     if mode == "forecast"
                     else "Forecast will replace climatology once race day is within 16 days. "))
                 + "Historical Race Conditions slide shows what actually happened in prior years.",
                 Inches(0.3), Inches(7.05), Inches(12.73), Inches(0.25),
                 font_size=9, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    return True


def add_race_overview_slide(
    prs: Presentation, gender: str, rows: list[dict], venue: str, content: dict
):
    """Race history paragraph + 3 Keys to Success for one gender."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Race Overview  —  Elite {gender.title()}", venue)

    # ── Left: Race history paragraph ───────────────────────────────────────────
    hist_bar = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2), Inches(7.0), Inches(0.35))
    hist_bar.fill.solid()
    hist_bar.fill.fore_color.rgb = NAVY
    hist_bar.line.fill.background()
    _add_textbox(slide, f"RACE HISTORY AT {venue.upper()}",
                 Inches(0.35), Inches(1.22), Inches(6.6), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    if rows:
        years = sorted(r["year"] for r in rows)
        year_range = f"{years[0]}–{years[-1]}"
        n = len(rows)
        winners = [r["winner_name"] for r in rows if r.get("winner_name")]
        winner_note = ""
        if winners:
            top = Counter(winners).most_common(1)[0]
            if top[1] > 1:
                winner_note = f" {top[0]} is the most successful athlete in this span with {top[1]} victories."
        cats = list(dict.fromkeys(
            r["cat_name"].replace("World Triathlon Championship Series", "WTCS")
            for r in rows if r.get("cat_name")
        ))
        cat_str = " / ".join(cats[:2])

        history_text = (
            f"The Elite {gender.title()} race at {venue} has been held {n} time{'s' if n > 1 else ''} "
            f"in this analysis window ({year_range}).{winner_note} "
            f"The event is staged as part of the {cat_str} calendar.\n\n"
            f"[PLACEHOLDER: Add 1–2 sentences on notable storylines, competitive history, "
            f"or what makes this event significant for Elite {gender.title()} athletes. "
            f"E.g. defending champions, record performances, dramatic finishes.]"
        )
    else:
        history_text = (
            f"[PLACEHOLDER: 3–4 sentences on the history of the Elite {gender.title()} race at {venue} — "
            "key moments, dominant athletes, course records, and the event's significance "
            "on the WTCS calendar.]"
        )

    _add_textbox(slide, history_text,
                 Inches(0.3), Inches(1.65), Inches(7.0), Inches(5.55),
                 font_size=12, color=DARK_GRAY)

    # ── Right: 3 Keys to Success ────────────────────────────────────────────────
    keys_bar = slide.shapes.add_shape(1, Inches(7.55), Inches(1.2), Inches(5.45), Inches(0.35))
    keys_bar.fill.solid()
    keys_bar.fill.fore_color.rgb = RED
    keys_bar.line.fill.background()
    _add_textbox(slide, "3 KEYS TO SUCCESS",
                 Inches(7.6), Inches(1.22), Inches(5.1), Inches(0.3),
                 font_size=12, bold=True, color=WHITE)

    keys = _derive_keys_to_success(rows)
    for i, key_text in enumerate(keys):
        top = Inches(1.7 + i * 1.85)
        # Numbered square badge
        badge = slide.shapes.add_shape(1, Inches(7.6), top, Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT

        _add_textbox(slide, key_text,
                     Inches(8.15), top - Inches(0.02), Inches(4.85), Inches(1.72),
                     font_size=11, color=DARK_GRAY)


# ── Deep-dive slides (single-race detailed analysis) ──────────────────────────

def add_pack_dynamics_slide(prs: Presentation, splits_df: pd.DataFrame,
                            gender: str, venue: str, year: int):
    """Three scatter plots: gap-to-leader vs placement at end of swim/bike/run."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Pack Dynamics — Elite {gender.title()} ({year})", venue)

    if splits_df is None or splits_df.empty:
        _add_textbox(slide, "Detailed split data unavailable.",
                     Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.5),
                     font_size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    data = compute_pack_scatter(splits_df)
    leg_titles = {"Swim": "Swim Exit", "Bike": "T2 (Bike Exit)", "Run": "Finish"}

    fig, axes = plt.subplots(1, 3, figsize=(15, 5.0))
    fig.patch.set_facecolor("white")

    for ax, leg in zip(axes, ["Swim", "Bike", "Run"]):
        d = data.get(leg, {"points": []})
        pts = d["points"]
        if not pts:
            ax.set_title(f"{leg_titles[leg]}: no data", fontsize=12)
            ax.set_facecolor("white")
            continue
        lead_pts  = [p for p in pts if p["in_lead_swim"]]
        other_pts = [p for p in pts if not p["in_lead_swim"]]

        if other_pts:
            ax.scatter([p["gap"] for p in other_pts],
                       [p["placement"] for p in other_pts],
                       s=70, c=C_GRAY, alpha=0.55, edgecolor="white", linewidth=0.8,
                       label="Other")
        if lead_pts:
            ax.scatter([p["gap"] for p in lead_pts],
                       [p["placement"] for p in lead_pts],
                       s=100, c=C_RED, alpha=0.9, edgecolor="white", linewidth=1.2,
                       label="Swim lead pack (≤15s)")

        for p in pts[:3]:
            last_name = str(p["name"]).split()[-1]
            ax.annotate(f"{p['placement']}. {last_name}",
                        (p["gap"], p["placement"]),
                        xytext=(6, 0), textcoords="offset points",
                        fontsize=8, va="center", color=C_NAVY, fontweight="bold")

        ax.set_title(f"At {leg_titles[leg]}", fontsize=12, pad=8, fontweight="bold")
        ax.set_xlabel("Gap to leader (s)", fontsize=10)
        ax.set_ylabel("Position", fontsize=10)
        ax.invert_yaxis()
        ax.grid(True, alpha=0.3, linestyle="--")
        ax.set_facecolor("white")
        if lead_pts and other_pts:
            ax.legend(fontsize=8, loc="lower right", framealpha=0.85)

    plt.tight_layout(pad=1.5)
    slide.shapes.add_picture(fig_to_image(fig),
                             Inches(0.2), Inches(1.3), Inches(12.93), Inches(5.4))

    _add_textbox(
        slide,
        "Each point = one athlete; position derives from elapsed time at the checkpoint. "
        "Red points are athletes who exited the swim within 15s of the leader — "
        "trace them across the three plots to see whether the early swim group held to the line.",
        Inches(0.3), Inches(6.78), Inches(12.73), Inches(0.55),
        font_size=10, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )


def add_top_splits_vs_season_slide(prs: Presentation, splits_df: pd.DataFrame,
                                   engine, gender: str, venue: str,
                                   year: int, prior_event_id: int,
                                   prior_prog_id: int, venue_date):
    """Top 3 splits per discipline at this race vs each athlete's WTCS 12-month best/avg."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Top Splits vs Season Norms — Elite {gender.title()} ({year})", venue)

    _add_textbox(
        slide,
        "Top 3 split times at this race vs each athlete's best and average split "
        "across WTCS Standard-distance races in the 12 months before the race.",
        Inches(0.3), Inches(1.18), Inches(12.73), Inches(0.4),
        font_size=10.5, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )

    if splits_df is None or splits_df.empty:
        _add_textbox(slide, "Detailed split data unavailable.",
                     Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.5),
                     font_size=14, color=DARK_GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    # Identify top 3 per discipline (finishers only — DNF rows excluded)
    finishers = splits_df[~splits_df["dnf"]].copy()
    top_by_leg: dict[str, list[dict]] = {}
    for leg_label, col in [("Swim", "swim_sec"), ("Bike", "bike_sec"), ("Run", "run_sec")]:
        col_clean = adjust_outlier(finishers[col].dropna()).reindex(finishers.index)
        valid = finishers[col_clean.notna()].copy()
        valid[col] = col_clean[col_clean.notna()]
        valid = valid.nsmallest(3, col)
        top_by_leg[leg_label] = [{"name": r["Name"], "split_sec": float(r[col])}
                                 for _, r in valid.iterrows()]

    all_names = list({p["name"] for plist in top_by_leg.values() for p in plist})
    name_to_aid = map_excel_names_to_athlete_ids(engine, prior_event_id, prior_prog_id, all_names)

    aids = list(name_to_aid.values())
    season_df = query_wtcs_season_splits(engine, aids, venue_date, gender)
    season_df = season_df[season_df.event_date < venue_date]

    norms_by_aid: dict[int, dict] = {}
    for aid, g in season_df.groupby("athlete_id"):
        norms_by_aid[int(aid)] = {
            "swim_best": g["swim_sec"].min(skipna=True) if g["swim_sec"].notna().any() else None,
            "swim_avg":  g["swim_sec"].mean(skipna=True) if g["swim_sec"].notna().any() else None,
            "bike_best": g["bike_sec"].min(skipna=True) if g["bike_sec"].notna().any() else None,
            "bike_avg":  g["bike_sec"].mean(skipna=True) if g["bike_sec"].notna().any() else None,
            "run_best":  g["run_sec"].min(skipna=True)  if g["run_sec"].notna().any()  else None,
            "run_avg":   g["run_sec"].mean(skipna=True) if g["run_sec"].notna().any()  else None,
            "n":         int(g.shape[0]),
        }

    panel_width = Inches(4.2)
    panel_left  = [Inches(0.3), Inches(4.55), Inches(8.8)]
    panel_top   = Inches(1.7)
    bar_color   = {"Swim": C_LIGHT_BLUE, "Bike": C_NAVY, "Run": C_RED}

    for left, leg_label in zip(panel_left, ["Swim", "Bike", "Run"]):
        bar = slide.shapes.add_shape(1, left, panel_top, panel_width, Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(bar_color[leg_label].lstrip("#"))
        bar.line.fill.background()
        _add_textbox(slide, leg_label.upper(),
                     left + Inches(0.1), panel_top + Inches(0.04),
                     panel_width - Inches(0.2), Inches(0.28),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tbl_top = panel_top + Inches(0.4)
        cols = [("Athlete", 1.65), ("Split", 0.72), ("Best", 0.72), ("Avg", 0.72), ("Δ Avg", 0.55)]
        n_rows = 1 + len(top_by_leg[leg_label])
        tbl_shape = slide.shapes.add_table(
            n_rows, len(cols), left, tbl_top, panel_width, Inches(0.4 * n_rows + 0.1)
        )
        tbl = tbl_shape.table
        for ci, (_, w) in enumerate(cols):
            tbl.columns[ci].width = Inches(w)
        for ci, (hdr, _) in enumerate(cols):
            _set_cell(tbl.cell(0, ci), hdr, bold=True, color=WHITE, bg_color=NAVY,
                      font_size=9)

        leg_key = leg_label.lower()
        for ri, entry in enumerate(top_by_leg[leg_label], start=1):
            bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
            aid = name_to_aid.get(entry["name"])
            norms = norms_by_aid.get(aid) if aid is not None else None
            best = norms.get(f"{leg_key}_best") if norms else None
            avg  = norms.get(f"{leg_key}_avg")  if norms else None
            split = entry["split_sec"]
            best_s = seconds_to_mmss(best) if best is not None and not pd.isna(best) else "—"
            avg_s  = seconds_to_mmss(avg)  if avg  is not None and not pd.isna(avg)  else "—"
            split_s = seconds_to_mmss(split)
            if avg is not None and not pd.isna(avg):
                delta = int(round(split - float(avg)))
                if delta < 0:
                    delta_s = f"{delta}s"
                    delta_color = RGBColor(0x1B, 0x7F, 0x3A)  # green
                elif delta > 0:
                    delta_s = f"+{delta}s"
                    delta_color = RGBColor(0xC0, 0x00, 0x00)  # red
                else:
                    delta_s = "0s"
                    delta_color = DARK_GRAY
            else:
                delta_s = "—"
                delta_color = DARK_GRAY
            last_first = str(entry["name"])
            parts = last_first.split()
            short_name = f"{parts[0][0]}. {parts[-1]}" if len(parts) > 1 else last_first
            _set_cell(tbl.cell(ri, 0), short_name, font_size=8.5, bg_color=bg, align=PP_ALIGN.LEFT)
            _set_cell(tbl.cell(ri, 1), split_s, font_size=8.5, bold=True, bg_color=bg)
            _set_cell(tbl.cell(ri, 2), best_s,  font_size=8.5, bg_color=bg)
            _set_cell(tbl.cell(ri, 3), avg_s,   font_size=8.5, bg_color=bg)
            _set_cell(tbl.cell(ri, 4), delta_s, font_size=8.5, bold=True, bg_color=bg,
                      color=delta_color)

    _add_textbox(
        slide,
        "Best / Avg = athlete's WTCS Standard-distance splits in the 12 months ending the day before this race. "
        "Δ Avg = this race's split minus the athlete's 12-month average; "
        "green = faster than average, red = slower.",
        Inches(0.3), Inches(6.4), Inches(12.73), Inches(0.55),
        font_size=9, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )


def add_who_to_watch_slide(prs: Presentation, engine, venue: str,
                           upcoming_event_id: int | None,
                           prior_men: dict | None,
                           prior_women: dict | None):
    """Single slide covering both genders: top 3 by ranking + returning prior-year podium."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, f"Who to Watch — {venue} Race Week", "")

    panel_w = Inches(6.2)
    gap     = Inches(0.3)
    left_m  = Inches(0.3)
    left_w  = left_m + panel_w + gap
    panel_top = Inches(1.3)

    for left, gender, prior in [(left_m, "Men", prior_men),
                                (left_w, "Women", prior_women)]:
        header = slide.shapes.add_shape(1, left, panel_top, panel_w, Inches(0.4))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()
        _add_textbox(slide, f"ELITE {gender.upper()}",
                     left + Inches(0.1), panel_top + Inches(0.05),
                     panel_w - Inches(0.2), Inches(0.32),
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rankings_df = query_top_by_world_ranking(engine, gender,
                                                 on_startlist_event_id=upcoming_event_id,
                                                 limit=10)
        rankings_label = "Top 10 Ranked on Startlist" if (upcoming_event_id and not rankings_df.empty) else "Top 10 by World Ranking"
        if rankings_df.empty and upcoming_event_id:
            rankings_df = query_top_by_world_ranking(engine, gender,
                                                     on_startlist_event_id=None, limit=10)
            rankings_label = "Top 10 by World Ranking (no startlist data yet)"

        sub_top = panel_top + Inches(0.5)
        sub_bar = slide.shapes.add_shape(1, left, sub_top, panel_w, Inches(0.3))
        sub_bar.fill.solid()
        sub_bar.fill.fore_color.rgb = RED
        sub_bar.line.fill.background()
        _add_textbox(slide, rankings_label,
                     left + Inches(0.1), sub_top + Inches(0.02),
                     panel_w - Inches(0.2), Inches(0.26),
                     font_size=11, bold=True, color=WHITE)

        tbl_top = sub_top + Inches(0.35)
        rank_cols = [("Rank", 0.8), ("Athlete", 3.6), ("Points", 1.6)]
        n_rows = 1 + max(len(rankings_df), 1)
        tbl_shape = slide.shapes.add_table(n_rows, len(rank_cols), left, tbl_top,
                                           panel_w, Inches(0.28 * n_rows + 0.1))
        tbl = tbl_shape.table
        for ci, (_, w) in enumerate(rank_cols):
            tbl.columns[ci].width = Inches(w)
        for ci, (hdr, _) in enumerate(rank_cols):
            _set_cell(tbl.cell(0, ci), hdr, bold=True, color=WHITE, bg_color=NAVY, font_size=10)
        if rankings_df.empty:
            _set_cell(tbl.cell(1, 0), "—", font_size=10, bg_color=WHITE)
            _set_cell(tbl.cell(1, 1), "No ranking data found", font_size=10, bg_color=WHITE,
                      align=PP_ALIGN.LEFT, italic=True)
            _set_cell(tbl.cell(1, 2), "—", font_size=10, bg_color=WHITE)
        else:
            for ri, r in enumerate(rankings_df.itertuples(index=False), start=1):
                bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
                pts = f"{int(r.total_points):,}" if r.total_points else "—"
                _set_cell(tbl.cell(ri, 0), str(int(r.rank_position)),
                          font_size=10, bold=True, bg_color=bg)
                _set_cell(tbl.cell(ri, 1), str(r.athlete_name),
                          font_size=10, bg_color=bg, align=PP_ALIGN.LEFT)
                _set_cell(tbl.cell(ri, 2), pts, font_size=10, bg_color=bg)

        # Prior-year podium block
        podium_top = tbl_top + Inches(0.28 * n_rows + 0.3)
        podium_bar = slide.shapes.add_shape(1, left, podium_top, panel_w, Inches(0.3))
        podium_bar.fill.solid()
        podium_bar.fill.fore_color.rgb = RED
        podium_bar.line.fill.background()
        podium_label = "Returning Podium Athletes on Startlist" if upcoming_event_id else f"Last Year's Podium at {venue}"
        _add_textbox(slide, podium_label,
                     left + Inches(0.1), podium_top + Inches(0.02),
                     panel_w - Inches(0.2), Inches(0.26),
                     font_size=11, bold=True, color=WHITE)

        body_top = podium_top + Inches(0.35)
        if upcoming_event_id and prior:
            podium_df = query_prior_podium_on_startlist(
                engine, prior["event_id"], prior["prog_id"], upcoming_event_id
            )
            if podium_df.empty:
                _add_textbox(slide,
                             "None of last year's podium are on the current startlist.",
                             left + Inches(0.1), body_top + Inches(0.1),
                             panel_w - Inches(0.2), Inches(0.5),
                             font_size=10.5, color=DARK_GRAY, italic=True)
            else:
                for ri, r in enumerate(podium_df.itertuples(index=False)):
                    txt = f"{int(r.position)}.  {r.athlete_full_name}   ({r.total_time})"
                    _add_textbox(slide, txt,
                                 left + Inches(0.15),
                                 body_top + Inches(0.05 + 0.32 * ri),
                                 panel_w - Inches(0.3), Inches(0.32),
                                 font_size=11, bold=(ri == 0), color=DARK_GRAY)
        elif prior:
            # No startlist available — just list last year's top 3
            sql = text("""
                SELECT position_sort AS position, athlete_full_name, total_time
                FROM race_results
                WHERE event_id = :eid AND prog_id = :pid
                  AND position_sort IS NOT NULL AND position_sort <= 3
                ORDER BY position_sort
            """)
            podium_df = pd.read_sql(sql, engine, params={"eid": prior["event_id"],
                                                         "pid": prior["prog_id"]})
            if podium_df.empty:
                _add_textbox(slide, "No prior-race podium data found.",
                             left + Inches(0.1), body_top + Inches(0.1),
                             panel_w - Inches(0.2), Inches(0.5),
                             font_size=10.5, color=DARK_GRAY, italic=True)
            else:
                for ri, r in enumerate(podium_df.itertuples(index=False)):
                    txt = f"{int(r.position)}.  {r.athlete_full_name}   ({r.total_time})"
                    _add_textbox(slide, txt,
                                 left + Inches(0.15),
                                 body_top + Inches(0.05 + 0.32 * ri),
                                 panel_w - Inches(0.3), Inches(0.32),
                                 font_size=11, bold=(ri == 0), color=DARK_GRAY)
        else:
            _add_textbox(slide, "No prior race history at this venue.",
                         left + Inches(0.1), body_top + Inches(0.1),
                         panel_w - Inches(0.2), Inches(0.5),
                         font_size=10.5, color=DARK_GRAY, italic=True)

    _add_textbox(
        slide,
        "Rankings = current World Triathlon rankings snapshot; startlist intersect uses program_entries when available.",
        Inches(0.3), Inches(7.05), Inches(12.73), Inches(0.25),
        font_size=9, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )


# ── New slides: Swim Threats / OQR / Rankings ────────────────────────────────

def _wtcs_pts(pos: int, win: float = 1000.0) -> float:
    return round(win * (0.925 ** (pos - 1)), 2)


def _tier_pts_scale(event_name: str | None,
                    cat_name: str | None = None) -> tuple[float, str]:
    """Return (winner base points, tier label) for the max-points scheme.

    Prefers the authoritative World Triathlon `cat_name` field; falls back to
    event-name keyword matching. Note WT brands World Cups as "World Triathlon
    Cup", so the naive `"world cup" in name` check MISSES them (there's a
    'triathlon' in between) — we match 'triathlon cup' explicitly.
    """
    c = (cat_name or "").lower()
    n = (event_name or "").lower()

    # 1. Authoritative cat_name (most specific first)
    if "championship series" in c or "wtcs" in c:
        return (1000.0, "WTCS")
    if "world cup" in c:
        return (500.0, "World Cup")
    if "continental championships" in c:
        return (400.0, "Continental Champs")
    if "continental cup" in c:
        return (250.0, "Continental Cup")
    if "major games" in c or "olympic" in c:
        return (1000.0, "Olympic Games")

    # 2. Event-name fallback (handles the "World Triathlon Cup" branding)
    if "world triathlon cup" in n or "triathlon cup" in n or "world cup" in n:
        return (500.0, "World Cup")
    if "continental" in n and "championship" in n:
        return (400.0, "Continental Champs")
    if "continental" in n or "americas triathlon cup" in n:
        return (250.0, "Continental Cup")
    if "olympic" in n:
        return (1000.0, "Olympic Games")
    return (1000.0, "WTCS")


def query_swim_threats(engine, prog_regex: str,
                       on_startlist_event_id: int | None,
                       since_years: int = 4) -> pd.DataFrame:
    today = date.today()
    since = f"{today.year - since_years}-{today.month:02d}-01"
    startlist_join = ""
    if on_startlist_event_id:
        startlist_join = f"""
        JOIN program_entries pe
          ON pe.athlete_id = pm.athlete_id
         AND pe.event_id = {int(on_startlist_event_id)}
         AND pe.is_active = TRUE
         AND pe.entry_type = 'start'"""
    sql = text(f"""
        SELECT pm.athlete_id, rr.athlete_full_name, a.country, e.event_date,
               e.prog_distance_category AS dist,
               pm.elapsedswim, pm.behindswim
        FROM position_metrics pm
        JOIN events e ON e.event_id = pm.event_id AND e.prog_id = pm.prog_id
        JOIN race_results rr
          ON rr.event_id = pm.event_id AND rr.prog_id = pm.prog_id
         AND rr.athlete_id = pm.athlete_id
        LEFT JOIN athlete a ON a.athlete_id = pm.athlete_id
        {startlist_join}
        WHERE e.event_date >= :since
          AND (e.cat_name ILIKE '%Championship Series%'
               OR e.cat_name ILIKE '%World Cup%'
               OR e.cat_name ILIKE '%Continental%'
               OR e.event_name ~* 'World Triathlon Championship (Series|Finals)')
          AND e.prog_name  ~* :prog
          AND rr.finish_status = 'FINISH'
          AND pm.elapsedswim IS NOT NULL
          -- Plausibility filter: exclude relay legs / super-sprint / zero-time
          -- contamination. No real 750 m swim is faster than ~6:40 (400 s);
          -- no real standard swim is slower than ~40 min (2400 s).
          AND pm.elapsedswim BETWEEN 400 AND 2400
          AND e.prog_distance_category IN ('sprint', 'standard')
        ORDER BY pm.athlete_id, e.event_date ASC
    """)
    return pd.read_sql(sql, engine, params={"since": since, "prog": prog_regex})


def _compute_swim_metrics(raw: pd.DataFrame, top_n: int = 6,
                          rank_by: str = "sprint") -> pd.DataFrame:
    """Per-athlete swim metrics with EWMA computed SEPARATELY for sprint vs
    standard distance (they aren't comparable — a 750 m and 1500 m swim differ
    by ~7 minutes). Ranks by the race's own distance (`rank_by`), falling back
    to the other distance when an athlete has no history at the race distance.

    Lead Pack % and Avg Gap are gap-based (distance-independent) so they stay
    pooled across all races.
    """
    LEAD_GAP = 15
    rank_by = (rank_by or "sprint").lower()
    other = "standard" if rank_by == "sprint" else "sprint"

    def _ewma(sub: pd.DataFrame):
        if sub.empty:
            return None
        return sub.sort_values("event_date")["elapsedswim"].ewm(
            span=5, min_periods=1).mean().iloc[-1]

    results = []
    for aid, grp in raw.groupby("athlete_id"):
        grp = grp.sort_values("event_date")
        if len(grp) < 3:  # need a few races so front-pack % isn't a 1-2 race fluke
            continue
        d = grp["dist"].astype(str).str.lower() if "dist" in grp else pd.Series(["standard"] * len(grp))
        spr = grp[d == "sprint"]
        std = grp[d == "standard"]
        ewma_sprint = _ewma(spr)
        ewma_std    = _ewma(std)
        avg_gap  = grp["behindswim"].mean()
        lead_pct = (grp["behindswim"] <= LEAD_GAP).mean() * 100
        country  = grp["country"].iloc[-1] if "country" in grp else None
        # Rank key = EWMA at the race's distance, else the other distance
        rank_ewma = (ewma_sprint if rank_by == "sprint" else ewma_std)
        if rank_ewma is None:
            rank_ewma = (ewma_std if rank_by == "sprint" else ewma_sprint)
        results.append({
            "athlete_id": aid,
            "athlete_name": grp["athlete_full_name"].iloc[-1],
            "country": country,
            "races": len(grp),
            "n_sprint": int(len(spr)),
            "n_std": int(len(std)),
            "ewma_sprint": ewma_sprint,
            "ewma_std": ewma_std,
            "rank_ewma": rank_ewma,
            "avg_gap": avg_gap,
            "lead_pct": lead_pct,
        })
    if not results:
        return pd.DataFrame()
    # Primary sort: front-pack rate (course/field-robust — being within 15 s of
    # the leader means the same on a fast or slow course). EWMA at the race
    # distance breaks ties.
    return (pd.DataFrame(results)
            .sort_values(["lead_pct", "rank_ewma"],
                         ascending=[False, True], na_position="last")
            .head(top_n)
            .reset_index(drop=True))


def add_top_swim_threats_slide(prs: Presentation, engine, venue: str,
                               upcoming_event_id: int | None,
                               race_distance: str | None = None):
    """Swim threats tables ranked by front-pack rate. Shows Elite Men/Women,
    plus U23 Men/Women only when the event has a U23 category. EWMA swim is
    shown separately for sprint and standard distance; ranking uses the race's
    own distance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Top Swim Threats", venue)

    # Determine the race distance (drives which EWMA column ranks the table)
    if not race_distance and upcoming_event_id:
        try:
            race_distance = pd.read_sql(text(
                "SELECT prog_distance_category FROM events "
                "WHERE event_id = :eid AND prog_distance_category IS NOT NULL LIMIT 1"
            ), engine, params={"eid": upcoming_event_id})["prog_distance_category"].iloc[0]
        except Exception:
            race_distance = None
    rank_by = (race_distance or "sprint").lower()
    rank_by = "sprint" if "sprint" in rank_by else ("standard" if "standard" in rank_by else "sprint")

    # ── Methodology description bar ───────────────────────────────────────────
    desc = slide.shapes.add_shape(1, Inches(0.3), Inches(1.18),
                                  Inches(12.73), Inches(0.35))
    desc.fill.solid()
    desc.fill.fore_color.rgb = LIGHT_GRAY
    desc.line.color.rgb = MID_GRAY
    desc.line.width = Pt(1)
    desc.shadow.inherit = False
    _add_textbox(
        slide,
        f"Ranked by Lead Pack % (within 15 s of swim leader — robust to course speed)   •   "
        f"EWMA swim split (span-5) shown per distance; {rank_by.title()} = this race   •   "
        f"WTCS / World Cup / Continental, last 4 yrs",
        Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.27),
        font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Only show U23 panels if the upcoming event actually has a U23 category.
    # For Elite-only events (e.g. most World Cups), the U23 distinction is
    # meaningless — show just Elite Men/Women and fill the space with more names.
    has_u23 = False
    if upcoming_event_id:
        try:
            has_u23 = bool(pd.read_sql(text(
                "SELECT 1 FROM events WHERE event_id = :eid "
                "AND prog_name ILIKE 'U23%' LIMIT 1"
            ), engine, params={"eid": upcoming_event_id}).shape[0])
        except Exception:
            has_u23 = False

    if has_u23:
        panels = [
            ("ELITE MEN",   r"Elite Men(?!.*Women)", Inches(0.3), Inches(1.65)),
            ("ELITE WOMEN", r"Elite Women",           Inches(6.7), Inches(1.65)),
            ("U23 MEN",     r"U23 Men(?!.*Women)",    Inches(0.3), Inches(4.45)),
            ("U23 WOMEN",   r"U23 Women",             Inches(6.7), Inches(4.45)),
        ]
        panel_top_n = 7
    else:
        panels = [
            ("ELITE MEN",   r"Elite Men(?!.*Women)", Inches(0.3), Inches(1.65)),
            ("ELITE WOMEN", r"Elite Women",           Inches(6.7), Inches(1.65)),
        ]
        panel_top_n = 14   # only two panels — use the vertical space for more names
    panel_w = Inches(6.35)
    # Highlight whichever distance is the race distance in the header
    spr_hdr = "EWMA Spr ▸" if rank_by == "sprint" else "EWMA Spr"
    std_hdr = "EWMA Std ▸" if rank_by == "standard" else "EWMA Std"
    col_spec = [
        ("#",        0.374),
        ("Athlete",  1.950),
        ("Lead %",   0.900),
        (spr_hdr,    1.040),
        (std_hdr,    1.040),
        ("Spr/Std",  1.046),
    ]

    def _fmt_swim(v):
        if v is None or (isinstance(v, float) and pd.isna(v)):
            return "—"
        mins, secs = divmod(int(v), 60)
        return f"{mins}:{secs:02d}"

    for label, prog_regex, left, top in panels:
        hdr = slide.shapes.add_shape(1, left, top, panel_w, Inches(0.3))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = NAVY
        hdr.line.fill.background()
        _add_textbox(slide, label, left, top + Inches(0.03),
                     panel_w, Inches(0.24),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        raw = query_swim_threats(engine, prog_regex, upcoming_event_id)
        df  = _compute_swim_metrics(raw, top_n=panel_top_n, rank_by=rank_by)

        tbl_top = top + Inches(0.35)
        n_rows  = 1 + max(len(df), 1)
        tbl = slide.shapes.add_table(
            n_rows, len(col_spec),
            left, tbl_top, panel_w, Inches(0.28 * n_rows)
        ).table
        for ci, (_, w) in enumerate(col_spec):
            tbl.columns[ci].width = Inches(w)
        for ci, (h, _) in enumerate(col_spec):
            _set_cell(tbl.cell(0, ci), h, bold=True, color=WHITE, bg_color=NAVY,
                      font_size=8.5, align=PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER)

        if df.empty:
            _set_cell(tbl.cell(1, 1), "No data available", font_size=8.5,
                      align=PP_ALIGN.LEFT, italic=True)
        else:
            for ri, r in enumerate(df.itertuples(index=False), start=1):
                bg     = LIGHT_GRAY if ri % 2 == 0 else WHITE
                is_usa = (r.country or "") == "United States"
                # Bold whichever EWMA column is the race distance
                spr_bold = (rank_by == "sprint")
                std_bold = (rank_by == "standard")
                _set_cell(tbl.cell(ri, 0), str(ri), font_size=8.5, bold=True,
                          color=NAVY, bg_color=bg)
                _set_cell(tbl.cell(ri, 1), r.athlete_name, font_size=8.5,
                          bold=is_usa, color=(RED if is_usa else DARK_GRAY),
                          bg_color=bg, align=PP_ALIGN.LEFT)
                _set_cell(tbl.cell(ri, 2), f"{r.lead_pct:.0f}%", font_size=8.5,
                          bold=True, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 3), _fmt_swim(r.ewma_sprint), font_size=8.5,
                          bold=spr_bold, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 4), _fmt_swim(r.ewma_std), font_size=8.5,
                          bold=std_bold, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 5), f"{int(r.n_sprint)}/{int(r.n_std)}",
                          font_size=8.5, color=MID_GRAY, bg_color=bg)


def query_usa_world_rankings(engine, gender: str,
                             on_startlist_event_id: int | None) -> pd.DataFrame:
    """USA athletes from world rankings (cat 13=Men, 14=Women), filtered to startlist if given."""
    cat_id = 13 if gender.lower().startswith("m") else 14
    if on_startlist_event_id:
        sql = text("""
            WITH latest AS (
              SELECT MAX(retrieved_at) AS dt
                FROM athlete_rankings WHERE ranking_cat_id = :cat
            )
            SELECT DISTINCT ar.athlete_id, ar.athlete_name,
                   ar.rank_position, ar.total_points,
                   ar.events_current_period, ar.events_previous_period
            FROM athlete_rankings ar
            JOIN latest l ON ar.retrieved_at = l.dt
            JOIN athlete a ON a.athlete_id = ar.athlete_id
            JOIN program_entries pe
              ON pe.athlete_id = ar.athlete_id
             AND pe.event_id = :eid
             AND pe.is_active = TRUE
             AND pe.entry_type = 'start'
            WHERE ar.ranking_cat_id = :cat
              AND a.country = 'United States'
            ORDER BY ar.rank_position ASC
        """)
        return pd.read_sql(sql, engine,
                           params={"cat": cat_id, "eid": on_startlist_event_id})
    sql = text("""
        SELECT ar.athlete_id, ar.athlete_name, ar.rank_position, ar.total_points,
               ar.events_current_period, ar.events_previous_period
        FROM athlete_rankings ar
        JOIN athlete a ON a.athlete_id = ar.athlete_id
        WHERE ar.ranking_cat_id = :cat
          AND ar.retrieved_at = (
                SELECT MAX(retrieved_at) FROM athlete_rankings
                 WHERE ranking_cat_id = :cat)
          AND a.country = 'United States'
        ORDER BY ar.rank_position ASC
    """)
    return pd.read_sql(sql, engine, params={"cat": cat_id})


def add_usa_world_rankings_slide(prs: Presentation, engine, venue: str,
                                 upcoming_event_id: int | None,
                                 event_name: str | None = None,
                                 cat_name: str | None = None):
    """Max-points header bar + USA world-ranking tables (cat 13/14)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "USA World Rankings Snapshot", venue)

    win_pts, tier_label = _tier_pts_scale(event_name, cat_name)

    # Latest ranking snapshot date for the context bar
    try:
        as_of = pd.read_sql(
            text("SELECT MAX(retrieved_at)::date AS d FROM athlete_rankings "
                 "WHERE ranking_cat_id IN (13, 14)"), engine)["d"].iloc[0]
    except Exception:
        as_of = None
    as_of_str = f"As of {as_of}   •   " if as_of else ""

    # ── Context description bar ────────────────────────────────────────────────
    desc = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2),
                                  Inches(12.73), Inches(0.42))
    desc.fill.solid()
    desc.fill.fore_color.rgb = LIGHT_GRAY
    desc.line.color.rgb = MID_GRAY
    desc.line.width = Pt(1)
    desc.shadow.inherit = False
    _add_textbox(
        slide,
        f"{as_of_str}World Triathlon Rankings   •   Heading into {venue}"
        f"   •   Race tier: {tier_label}",
        Inches(0.4), Inches(1.28), Inches(12.5), Inches(0.28),
        font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Max-points banner ──────────────────────────────────────────────────────
    banner = slide.shapes.add_shape(1, Inches(0.3), Inches(1.75),
                                    Inches(12.73), Inches(0.32))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RED
    banner.line.fill.background()
    banner.shadow.inherit = False
    _add_textbox(slide, "Max Points Available at this Race",
                 Inches(0.4), Inches(1.78), Inches(12.5), Inches(0.27),
                 font_size=12, bold=True, color=WHITE)

    # ── Max-points KPI cards (kept colour-coded per preference) ────────────────
    kpi_positions = [1,   3,    5,     10,    15,    20]
    kpi_labels    = ["Win","Podium","Top 5","Top 10","Top 15","Top 20"]
    OQR_PACE = 302
    card_w   = Inches(1.972)
    card_step = Inches(2.152)
    cx = Inches(0.3)
    for pos, lbl in zip(kpi_positions, kpi_labels):
        pts_val  = int(_wtcs_pts(pos, win_pts))
        on_pace  = pts_val >= OQR_PACE
        clr = RED if pos == 1 else (NAVY if on_pace else MID_GRAY)
        card = slide.shapes.add_shape(1, cx, Inches(2.17), card_w, Inches(0.55))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = clr
        card.line.width = Pt(1.5)
        card.shadow.inherit = False
        _add_textbox(slide, lbl,
                     cx, Inches(2.2), card_w, Inches(0.2),
                     font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add_textbox(slide, f"{pts_val:,} pts",
                     cx, Inches(2.39), card_w, Inches(0.32),
                     font_size=14, bold=True, color=clr, align=PP_ALIGN.CENTER)
        cx += card_step

    # ── Men / Women ranking tables ────────────────────────────────────────────
    rank_cols = [("Rank", 0.857), ("Athlete", 2.769), ("Points", 1.253),
                 ("Curr", 0.725), ("Prev", 0.725)]
    panels = [
        ("ELITE MEN",   "men",   Inches(0.3)),
        ("ELITE WOMEN", "women", Inches(6.7)),
    ]
    panel_w = Inches(6.33)
    MAX_ROWS = 15

    for label, gender, left in panels:
        hdr = slide.shapes.add_shape(1, left, Inches(2.85), panel_w, Inches(0.32))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = NAVY
        hdr.line.fill.background()
        _add_textbox(slide, label, left, Inches(2.88),
                     panel_w, Inches(0.26),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Always show the full USA team picture (all top-ranked USA athletes),
        # not just those on this race's startlist — this slide is a team-wide
        # world-ranking snapshot, independent of the specific event field.
        df = query_usa_world_rankings(engine, gender, None)
        df = df.head(MAX_ROWS)

        tbl_top = Inches(3.21)
        n_rows  = 1 + max(len(df), 1)
        tbl = slide.shapes.add_table(
            n_rows, len(rank_cols),
            left, tbl_top, panel_w, Inches(0.245 * n_rows)
        ).table
        for ci, (_, w) in enumerate(rank_cols):
            tbl.columns[ci].width = Inches(w)
        for ci, (h, _) in enumerate(rank_cols):
            _set_cell(tbl.cell(0, ci), h, bold=True, color=WHITE, bg_color=NAVY,
                      font_size=9,
                      align=PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER)

        if df.empty:
            _set_cell(tbl.cell(1, 0), "—", font_size=8.5)
            _set_cell(tbl.cell(1, 1), "No USA athletes found", font_size=8.5,
                      align=PP_ALIGN.LEFT, italic=True)
            for ci in range(2, len(rank_cols)):
                _set_cell(tbl.cell(1, ci), "—", font_size=8.5)
        else:
            for ri, r in enumerate(df.itertuples(index=False), start=1):
                bg   = LIGHT_GRAY if ri % 2 == 0 else WHITE
                pts  = f"{int(r.total_points):,}" if r.total_points else "—"
                curr = str(int(r.events_current_period))  if r.events_current_period  else "—"
                prev = str(int(r.events_previous_period)) if r.events_previous_period else "—"
                _set_cell(tbl.cell(ri, 0), str(int(r.rank_position)),
                          font_size=8.5, bold=True, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 1), str(r.athlete_name),
                          font_size=8.5, bg_color=bg, align=PP_ALIGN.LEFT, color=DARK_GRAY)
                _set_cell(tbl.cell(ri, 2), pts,  font_size=8.5, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 3), curr, font_size=8.5, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 4), prev, font_size=8.5, color=DARK_GRAY, bg_color=bg)


def query_usa_oqr_athletes(engine, gender: str) -> pd.DataFrame:
    """USA athletes in LA 2028 OQR rankings (cat 11=Men, 12=Women)."""
    cat_id = 11 if gender.lower().startswith("m") else 12
    sql = text("""
        SELECT ar.athlete_id, ar.athlete_name, ar.rank_position, ar.total_points,
               ar.events_current_period, ar.events_previous_period
        FROM athlete_rankings ar
        JOIN athlete a ON a.athlete_id = ar.athlete_id
        WHERE ar.ranking_cat_id = :cat
          AND ar.retrieved_at = (
                SELECT MAX(retrieved_at) FROM athlete_rankings
                 WHERE ranking_cat_id = :cat)
          AND a.country = 'United States'
        ORDER BY ar.rank_position ASC
    """)
    return pd.read_sql(sql, engine, params={"cat": cat_id})


def query_oqr_on_pace_counts(engine, athlete_ids: list[int],
                              cat_id: int, pace: float) -> pd.DataFrame:
    """Per-athlete count of included OQR events scoring >= pace threshold."""
    if not athlete_ids:
        return pd.DataFrame(columns=["athlete_id", "on_pace_events"])
    sql = text("""
        WITH latest AS (
          SELECT MAX(retrieved_at) AS dt
            FROM athlete_ranking_breakdown WHERE ranking_cat_id = :cat
        )
        SELECT b.athlete_id,
               COUNT(*) FILTER (WHERE b.points >= :pace AND b.included = TRUE)
                 AS on_pace_events
        FROM athlete_ranking_breakdown b
        JOIN latest l ON b.retrieved_at = l.dt
        WHERE b.athlete_id = ANY(:ids)
          AND b.ranking_cat_id = :cat
        GROUP BY b.athlete_id
    """)
    return pd.read_sql(sql, engine,
                       params={"cat": cat_id, "ids": athlete_ids, "pace": pace})


def add_oqr_pace_slide(prs: Presentation, engine, venue: str):
    """LA 2028 OQR on-pace analysis: tier table + USA Men/Women OQR tables."""
    OQR_CUTOFF = 3626.69
    OQR_CAP    = 12
    PACE       = OQR_CUTOFF / OQR_CAP   # ~= 302.22

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "Olympic Qualification Pace", venue)

    # -- Pace banner (two lines) --
    banner = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2),
                                    Inches(12.73), Inches(0.727))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()
    banner.shadow.inherit = False
    _add_textbox(
        slide,
        f"Paris 2024 OQR cutoff (30th): {OQR_CUTOFF:,.2f} pts   ÷   "
        f"{OQR_CAP}-event cap =   {PACE:.0f} pts / event on pace",
        Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.36),
        font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(
        slide,
        f"An athlete who averages ~{PACE:.0f} pts per event across their best "
        f"{OQR_CAP} events is on pace to clear the LA 2028 qualifying threshold.",
        Inches(0.4), Inches(1.595), Inches(12.5), Inches(0.36),
        font_size=11, italic=True, color=SUBTITLE_BLUE, align=PP_ALIGN.CENTER)

    # -- Section bar --
    sec = slide.shapes.add_shape(1, Inches(0.3), Inches(1.982),
                                 Inches(12.73), Inches(0.32))
    sec.fill.solid()
    sec.fill.fore_color.rgb = RED
    sec.line.fill.background()
    sec.shadow.inherit = False
    _add_textbox(slide, f"FINISH POSITION AT EACH TIER THAT HITS {PACE:.0f}-PT PACE",
                 Inches(0.4), Inches(2.012), Inches(12.5), Inches(0.26),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # -- Tier reference table (with Verdict) --
    tier_rows = [
        ("World Championship Series",           1000.0),
        ("World Cup",                            500.0),
        ("Continental Championships (Americas)", 590.0),
        ("Continental Cup (Americas)",           280.0),
    ]
    t_cols = [("Tier", 4.811), ("Win pts", 1.504),
              ("Last on-pace pos", 2.406), ("First off-pace pos (pts)", 3.007),
              ("Verdict", 1.002)]
    t_top = Inches(2.402)
    tbl   = slide.shapes.add_table(
        len(tier_rows) + 1, len(t_cols),
        Inches(0.3), t_top, Inches(12.73),
        Inches(0.3 * (len(tier_rows) + 1))
    ).table
    for ci, (_, w) in enumerate(t_cols):
        tbl.columns[ci].width = Inches(w)
    for ci, (h, _) in enumerate(t_cols):
        _set_cell(tbl.cell(0, ci), h, bold=True, color=WHITE, bg_color=NAVY,
                  font_size=9.5, align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

    for ri, (tier_name, win) in enumerate(tier_rows, start=1):
        on_pos = off_pos = None
        for p in range(1, 61):
            pts = win * (0.925 ** (p - 1))
            if pts >= PACE:
                on_pos = p
            elif on_pos and not off_pos:
                off_pos = p
                break
        bg      = LIGHT_GRAY if ri % 2 == 0 else WHITE
        on_str  = f"#{on_pos}" if on_pos else f"None (win = {int(win)})"
        off_str = (f"#{off_pos} ({win*(0.925**(off_pos-1)):.0f} pts)"
                   if off_pos else "—")
        verdict = f"Top {on_pos}" if on_pos else "Not achievable"
        v_clr   = GREEN if on_pos else RED
        _set_cell(tbl.cell(ri, 0), tier_name, font_size=9.5, bg_color=bg,
                  align=PP_ALIGN.LEFT)
        _set_cell(tbl.cell(ri, 1), f"{int(win)}", font_size=9.5, bold=True, bg_color=bg)
        _set_cell(tbl.cell(ri, 2), on_str, font_size=9.5, bold=True, bg_color=bg)
        _set_cell(tbl.cell(ri, 3), off_str, font_size=9.5, bg_color=bg)
        _set_cell(tbl.cell(ri, 4), verdict, font_size=9.5, bold=True,
                  color=v_clr, bg_color=bg)

    # -- USA OQR tables (Men / Women) --
    oqr_cols = [("Rk", 0.45), ("Athlete", 2.1), ("Total", 0.8), ("Races", 0.55),
                ("On", 0.45), ("Off", 0.45), ("Pace", 0.75), ("Proj @12", 0.8)]
    oqr_top  = Inches(3.95)
    panel_w  = Inches(6.35)
    MAX_ROWS = 11

    for label, gender, cat_id, left in [
        ("USA MEN",   "men",   11, Inches(0.3)),
        ("USA WOMEN", "women", 12, Inches(6.7)),
    ]:
        hdr = slide.shapes.add_shape(1, left, oqr_top, panel_w, Inches(0.3))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = NAVY
        hdr.line.fill.background()
        _add_textbox(slide, label, left, oqr_top + Inches(0.03),
                     panel_w, Inches(0.24),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        athletes = query_usa_oqr_athletes(engine, gender)
        pace_df  = query_oqr_on_pace_counts(
            engine, athletes["athlete_id"].tolist(), cat_id, PACE)
        if not pace_df.empty:
            athletes = athletes.merge(pace_df, on="athlete_id", how="left")
            athletes["on_pace_events"] = athletes["on_pace_events"].fillna(0).astype(int)
        else:
            athletes["on_pace_events"] = 0
        athletes = athletes.head(MAX_ROWS)
        fs = 8.5 if len(athletes) <= 10 else 7.5

        tbl_top2 = oqr_top + Inches(0.35)
        n_rows   = 1 + max(len(athletes), 1)
        tbl2 = slide.shapes.add_table(
            n_rows, len(oqr_cols),
            left, tbl_top2, panel_w,
            Inches(0.22 * n_rows)
        ).table
        for ci, (_, w) in enumerate(oqr_cols):
            tbl2.columns[ci].width = Inches(w)
        for ci, (h, _) in enumerate(oqr_cols):
            _set_cell(tbl2.cell(0, ci), h, bold=True, color=WHITE, bg_color=NAVY,
                      font_size=8.5,
                      align=PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER)

        if athletes.empty:
            _set_cell(tbl2.cell(1, 1), "No USA athletes in OQR",
                      font_size=fs, align=PP_ALIGN.LEFT, italic=True)
        else:
            for ri, r in enumerate(athletes.itertuples(index=False), start=1):
                bg     = LIGHT_GRAY if ri % 2 == 0 else WHITE
                total  = r.total_points or 0
                races  = (r.events_current_period or 0) + (r.events_previous_period or 0)
                on_n   = r.on_pace_events
                off_n  = max(races - on_n, 0)
                pace_r = (total / races) if races > 0 else 0
                proj   = round(total + pace_r * max(OQR_CAP - races, 0))
                proj_str  = f"{proj:,}" if races > 0 else "—"
                proj_clr  = GREEN if proj >= OQR_CUTOFF else DARK_GRAY
                _set_cell(tbl2.cell(ri, 0), str(int(r.rank_position)),
                          font_size=fs, bold=True, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl2.cell(ri, 1), str(r.athlete_name),
                          font_size=fs, bg_color=bg, align=PP_ALIGN.LEFT, color=DARK_GRAY)
                _set_cell(tbl2.cell(ri, 2), f"{int(total):,}",
                          font_size=fs, bold=True, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl2.cell(ri, 3), str(races),
                          font_size=fs, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl2.cell(ri, 4), str(on_n),
                          font_size=fs, bold=True, color=GREEN, bg_color=bg)
                _set_cell(tbl2.cell(ri, 5), str(off_n),
                          font_size=fs, bold=True, color=RED, bg_color=bg)
                _set_cell(tbl2.cell(ri, 6), str(int(pace_r)) if races > 0 else "—",
                          font_size=fs, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl2.cell(ri, 7), proj_str,
                          font_size=fs, bold=True, color=proj_clr, bg_color=bg)

    _add_textbox(
        slide,
        "On = events scoring ≥ pace.  Off = events < pace.  "
        "Pace = total breakdown pts ÷ events raced.  "
        "Proj @12 = current total + (avg pts × remaining events to cap); green if ≥ Paris cutoff.  "
        f"Cutoff estimate: Paris 2024 = {OQR_CUTOFF:,.0f} pts / {OQR_CAP} events = {PACE:.0f} pts/event.  "
        "Tier positions use a geometric 0.925 points model.  "
        "Source: World Triathlon rankings database.",
        Inches(0.3), Inches(7.02), Inches(12.73), Inches(0.3),
        font_size=8, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )


def add_usa_oqr_snapshot_slide(prs: Presentation, engine, venue: str):
    """USA athletes in LA 2028 OQR standings (cat 11=Men, 12=Women)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, "USA Olympic Qualification Snapshot", venue)

    # Latest OQR snapshot date for the context bar
    try:
        as_of = pd.read_sql(
            text("SELECT MAX(retrieved_at)::date AS d FROM athlete_rankings "
                 "WHERE ranking_cat_id IN (11, 12)"), engine)["d"].iloc[0]
    except Exception:
        as_of = None
    as_of_str = f"As of {as_of}   •   " if as_of else ""

    desc = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2),
                                  Inches(12.73), Inches(0.42))
    desc.fill.solid()
    desc.fill.fore_color.rgb = LIGHT_GRAY
    desc.line.color.rgb = MID_GRAY
    desc.line.width = Pt(1)
    desc.shadow.inherit = False
    _add_textbox(
        slide,
        f"{as_of_str}Olympic Qualification Rankings (LA 2028 cycle)"
        f"   •   Heading into {venue}",
        Inches(0.4), Inches(1.28), Inches(12.5), Inches(0.28),
        font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    tbl_cols = [("Rank", 0.857), ("Athlete", 2.769),
                ("Points", 1.243), ("Second", 0.735), ("First", 0.725)]
    panels = [
        ("ELITE MEN",   "men",   Inches(0.3)),
        ("ELITE WOMEN", "women", Inches(6.7)),
    ]
    panel_w   = Inches(6.33)
    panel_top = Inches(1.85)
    MAX_ROWS  = 12

    for label, gender, left in panels:
        hdr = slide.shapes.add_shape(1, left, panel_top, panel_w, Inches(0.32))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = NAVY
        hdr.line.fill.background()
        _add_textbox(slide, label, left, panel_top + Inches(0.03),
                     panel_w, Inches(0.26),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        df = query_usa_oqr_athletes(engine, gender).head(MAX_ROWS)

        tbl_top = panel_top + Inches(0.36)
        n_rows  = 1 + max(len(df), 1)
        tbl = slide.shapes.add_table(
            n_rows, len(tbl_cols),
            left, tbl_top, panel_w,
            Inches(0.245 * n_rows)
        ).table
        for ci, (_, w) in enumerate(tbl_cols):
            tbl.columns[ci].width = Inches(w)
        for ci, (h, _) in enumerate(tbl_cols):
            _set_cell(tbl.cell(0, ci), h, bold=True, color=WHITE, bg_color=NAVY,
                      font_size=9,
                      align=PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER)

        if df.empty:
            _set_cell(tbl.cell(1, 0), "—", font_size=8.5)
            _set_cell(tbl.cell(1, 1), "No USA athletes in OQR",
                      font_size=8.5, align=PP_ALIGN.LEFT, italic=True)
            for ci in range(2, len(tbl_cols)):
                _set_cell(tbl.cell(1, ci), "—", font_size=8.5)
        else:
            for ri, r in enumerate(df.itertuples(index=False), start=1):
                bg   = LIGHT_GRAY if ri % 2 == 0 else WHITE
                pts  = f"{int(r.total_points):,}" if r.total_points else "—"
                curr = str(int(r.events_current_period))  if r.events_current_period  else "0"
                prev = str(int(r.events_previous_period)) if r.events_previous_period else "0"
                _set_cell(tbl.cell(ri, 0), str(int(r.rank_position)),
                          font_size=8.5, bold=True, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 1), str(r.athlete_name),
                          font_size=8.5, bg_color=bg, align=PP_ALIGN.LEFT, color=DARK_GRAY)
                _set_cell(tbl.cell(ri, 2), pts,  font_size=8.5, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 3), curr, font_size=8.5, color=DARK_GRAY, bg_color=bg)
                _set_cell(tbl.cell(ri, 4), prev, font_size=8.5, color=DARK_GRAY, bg_color=bg)

    _add_textbox(
        slide,
        "Rankings: latest weekly snapshot.  Second / First = events counted in the "
        "second vs first scoring period of the LA 2028 OQR cycle.  "
        "Because the cycle is still in its first year, most results appear under First.  "
        "Source: World Triathlon rankings database.",
        Inches(0.3), Inches(7.02), Inches(12.73), Inches(0.3),
        font_size=8, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER,
    )


# ── Main ────────────────────────────────────────────────────────────────────────

def _autodiscover_detailed_splits(venue: str) -> str | None:
    """Look in data/ for 'Detailed results <Venue> <YYYY>.xlsx' (most recent year)."""
    data_dir = os.path.join(REPO_ROOT, "data")
    if not os.path.isdir(data_dir):
        return None
    pattern = re.compile(rf"^Detailed results {re.escape(venue)} (\d{{4}})\.xlsx$",
                         re.IGNORECASE)
    matches = []
    for fn in os.listdir(data_dir):
        m = pattern.match(fn)
        if m:
            matches.append((int(m.group(1)), os.path.join(data_dir, fn)))
    if not matches:
        return None
    matches.sort(reverse=True)
    return matches[0][1]


def main():
    parser = argparse.ArgumentParser(description="Generate venue historical race analysis PPT")
    parser.add_argument("--venue", required=True, help="Venue name (e.g. 'Yokohama')")
    parser.add_argument("--years", type=int, default=8, help="Years to look back (default: 8)")
    parser.add_argument("--gender", choices=["men", "women", "both"], default="both")
    parser.add_argument("--output", default=None, help="Output .pptx filename (optional)")
    parser.add_argument("--deep-dive", action="store_true",
                        help="Force the single-race deep-dive section (auto-enabled when N≤1 prior race)")
    parser.add_argument("--detailed-splits", default=None,
                        help="Path to detailed per-lap splits Excel (auto-discovered in data/ if omitted). Implies --deep-dive.")
    parser.add_argument("--upcoming-event-id", type=int, default=None,
                        help="Event ID of the upcoming race for 'Who to Watch' (auto-detected if omitted)")
    parser.add_argument("--preview-only", action="store_true",
                        help="Suppress Elite gender-section slides (e.g. for Para Series venues where "
                             "historical Elite race data is not relevant to the upcoming race).")
    parser.add_argument("--para", action="store_true",
                        help="Query Para Series events instead of (or in addition to) Elite events. "
                             "Adds Para-specific results summary and run analysis slides.")
    parser.add_argument("--para-class", default=None,
                        help="Comma-separated Para sport classes to show, e.g. 'PTS5,PTVI,PTWC'. "
                             "When omitted all classes found at the venue are included.")
    parser.add_argument("--include-elite", action="store_true",
                        help="In --para mode, also append the Elite gender-section slides. "
                             "By default Para mode is Para-focused and suppresses Elite sections "
                             "when Para data is present.")
    args = parser.parse_args()

    fname = args.output or f"{args.venue.replace(' ', '_')}_{date.today()}.pptx"
    if not os.path.isabs(fname):
        if os.path.dirname(fname):
            # Relative path already includes a directory — just ensure it exists
            os.makedirs(os.path.dirname(fname), exist_ok=True)
        else:
            # Bare filename — place it in the default output directory
            os.makedirs(DEFAULT_OUTPUT_DIR, exist_ok=True)
            fname = os.path.join(DEFAULT_OUTPUT_DIR, fname)

    print(f"Connecting to database...")
    engine = get_engine()

    # Parse Para class filter
    para_classes: list[str] | None = None
    if args.para_class:
        para_classes = [c.strip() for c in args.para_class.split(",") if c.strip()]

    print(f"Querying events at '{args.venue}' (past {args.years} years)...")
    events_df = query_venue_events(engine, args.venue, args.years, args.gender)

    # Para mode: also query Para events
    para_events_df = pd.DataFrame()
    if args.para:
        para_events_df = query_venue_para_events(engine, args.venue, args.years, para_classes)
        if not para_events_df.empty:
            print(f"Found {len(para_events_df)} Para program(s) at '{args.venue}':")
            for cls in sorted(para_events_df.prog_name.unique()):
                count = len(para_events_df[para_events_df.prog_name == cls])
                print(f"  {cls}: {count} editions")
        else:
            print(f"No Para events found for '{args.venue}'")

    if events_df.empty and para_events_df.empty:
        # No prior race history. Continue in preview-only mode iff we have
        # static profile data + an upcoming event at this venue; otherwise abort.
        venue_key = args.venue.lower().strip()
        has_profile = (venue_key in BIKE_COURSE_PROFILES
                       or venue_key in SWIM_COURSE_PROFILES
                       or venue_key in RUN_COURSE_PROFILES
                       or venue_key in EVENT_SCHEDULES)
        if not has_profile:
            print(f"No events found for venue '{args.venue}' and no preview profile data.")
            sys.exit(1)
        print(f"No prior race history for '{args.venue}' — building preview-only deck "
              f"from course profile data.")
    elif not events_df.empty:
        print(f"Found {len(events_df)} Elite program(s):")
        for _, r in events_df.iterrows():
            print(f"  {r.event_date}  {r.event_name}  ({r.prog_name})")

    men_df   = events_df[events_df.prog_name.str.contains("Elite Men",   case=False, na=False) &
                         ~events_df.prog_name.str.contains("Women",      case=False, na=False)]
    women_df = events_df[events_df.prog_name.str.contains("Elite Women", case=False, na=False)]

    print("\nCollecting race data...")
    men_data   = collect_race_data(engine, men_df)   if not men_df.empty   else []
    women_data = collect_race_data(engine, women_df) if not women_df.empty else []

    # Set the weather sampling window from this venue's race start times so the
    # window covers the actual race (afternoon races included), then geocode +
    # enrich all rows with Open-Meteo weather + AQ.
    _sched = EVENT_SCHEDULES.get(args.venue.lower().strip())
    _starts = _sched.get("race_starts") if _sched else None
    win = set_weather_window(_starts)
    print(f"Geocoding '{args.venue}' and fetching Open-Meteo data... "
          f"(weather window {win[0]:02d}:00–{win[1]:02d}:00 local)")
    coords = geocode_venue(args.venue)
    enrich_rows_with_openmeteo(men_data,   coords)
    enrich_rows_with_openmeteo(women_data, coords)

    print("Fetching venue content (WT API + URL construction)...")
    venue_content = fetch_venue_content(args.venue, events_df)

    # Build deduped venue-level rows for the environmental slide
    seen_eids: set = set()
    unique_env_rows: list[dict] = []
    for row in sorted(men_data + women_data, key=lambda r: r["year"]):
        if row["event_id"] not in seen_eids:
            seen_eids.add(row["event_id"])
            unique_env_rows.append(row)

    if args.para and not para_events_df.empty and not unique_env_rows:
        # For Para-only decks with no Elite data: build minimal env rows from
        # para event dates so the weather + env slides still render.
        for _, ev_row in para_events_df.drop_duplicates("event_id").iterrows():
            if int(ev_row.event_id) not in seen_eids:
                seen_eids.add(int(ev_row.event_id))
                unique_env_rows.append({
                    "event_id": int(ev_row.event_id),
                    "year": pd.to_datetime(ev_row.event_date).year,
                    "date": ev_row.event_date,
                    "event_name": ev_row.event_name,
                    "temp_air": ev_row.temperature_air,
                    "temp_water": ev_row.temperature_water,
                    "wind_kmh": ev_row.wind_speed_kmh, "wind_raw": ev_row.wind,
                    "humidity": None, "precip": None,
                    "pm25": None, "aqi": None, "uv_index": None,
                })
        unique_env_rows.sort(key=lambda r: r["year"])
        enrich_rows_with_openmeteo(unique_env_rows, coords)

    # ── Deep-dive trigger ─────────────────────────────────────────────────────
    n_men   = len(men_data)
    n_women = len(women_data)
    auto_trigger = max(n_men, n_women) <= 1
    deep_dive = args.deep_dive or bool(args.detailed_splits) or auto_trigger
    if deep_dive:
        reason = "explicit flag" if (args.deep_dive or args.detailed_splits) else "N<=1 prior race per gender"
        print(f"Deep-dive mode ON (reason: {reason})")

    detailed: dict[str, pd.DataFrame] = {}
    if deep_dive:
        xlsx_path = args.detailed_splits or _autodiscover_detailed_splits(args.venue)
        if xlsx_path and os.path.exists(xlsx_path):
            print(f"Loading detailed splits: {xlsx_path}")
            try:
                detailed = load_detailed_splits(xlsx_path)
            except Exception as exc:
                print(f"  Warning: failed to load detailed splits ({exc}); deep-dive split slides will be skipped")
                detailed = {}
        else:
            print("  No detailed-splits Excel found (looked in data/); split-based deep-dive slides will be skipped")

    # Upcoming-event lookup runs regardless of deep-dive — Who to Watch is useful for any venue
    # where we have prior race history.
    upcoming = None
    if args.upcoming_event_id:
        # Look up the real event row so event_name + cat_name are populated
        # (needed for correct tier / max-points detection downstream).
        eid = int(args.upcoming_event_id)
        _row = pd.read_sql(text(
            "SELECT event_id, MIN(event_date) AS event_date, "
            "MIN(event_name) AS event_name, MIN(cat_name) AS cat_name "
            "FROM events WHERE event_id = :eid GROUP BY event_id"
        ), engine, params={"eid": eid})
        if not _row.empty:
            upcoming = _row.iloc[0].to_dict()
            print(f"  Upcoming event (CLI override): {upcoming['event_name']} "
                  f"({upcoming['event_date']}) — cat={upcoming.get('cat_name')}")
        else:
            # Startlist exists in program_entries but the event row was never
            # ingested. Fall back to VENUE_PREVIEW metadata so tier / max-points
            # detection still resolves instead of reading a synthetic placeholder.
            _pv = VENUE_PREVIEW.get(args.venue.lower().strip(), {})
            upcoming = {"event_id": eid, "event_date": None,
                        "event_name": _pv.get("event_name") or f"Upcoming event {eid}",
                        "cat_name": _pv.get("cat_name")}
            print(f"  event_id={eid} not in events table — using startlist from "
                  f"program_entries with VENUE_PREVIEW tier metadata "
                  f"(cat={upcoming['cat_name']})")
    else:
        upcoming = find_upcoming_event(engine, args.venue)
        if upcoming:
            print(f"  Upcoming event detected: {upcoming['event_name']} ({upcoming['event_date']}) "
                  f"— event_id={upcoming['event_id']}, cat={upcoming.get('cat_name')}")
        else:
            print("  No upcoming event found at this venue; 'Who to Watch' will use rankings-only mode")

    # Prior-race anchor per gender (for athlete_id lookup + podium queries)
    prior_men   = men_data[-1]   if men_data   else None
    prior_women = women_data[-1] if women_data else None

    print("Building PowerPoint...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    all_data = men_data + women_data
    title_label = ("Para Series Analysis"
                   if (args.para and not para_events_df.empty)
                   else "Elite Analysis")
    add_title_slide(prs, args.venue, args.years, analysis_label=title_label)
    if add_race_week_schedule_slide(prs, args.venue):
        print(f"  Race-week schedule slide added (venue: {args.venue})")
    add_course_differentiators_slide(prs, args.venue, all_data, events_df, venue_content)
    add_course_map_slide(prs, args.venue, all_data, venue_content)

    # Sea-surface temp for race day (forecast if in window, else prior-years avg).
    # Prefer the DB event date; fall back to EVENT_SCHEDULES["race_date"] so preview
    # venues whose event is not yet ingested still get race-day weather.
    race_date_str = (str(pd.to_datetime(upcoming["event_date"]).date())
                     if upcoming and upcoming.get("event_date") else None)
    if not race_date_str:
        _sched_rd = EVENT_SCHEDULES.get(args.venue.lower().strip(), {}).get("race_date")
        if _sched_rd:
            race_date_str = str(_sched_rd)
            print(f"  Using EVENT_SCHEDULES race_date fallback: {race_date_str}")
    sst_c, sst_src = (None, "unavailable")
    if coords and race_date_str:
        sst_c, sst_src = fetch_sst_for_race_day(coords[0], coords[1], race_date_str)
        if sst_c is not None:
            print(f"  Sea-surface temp for {race_date_str}: {sst_c:.1f} °C ({sst_src})")

    if add_swim_course_profile_slide(prs, args.venue, water_temp_c=sst_c,
                                     water_temp_source=sst_src):
        print(f"  Swim course profile slide added (venue: {args.venue})")
    if add_bike_course_profile_slide(prs, args.venue, all_data):
        print(f"  Bike course profile slide added (venue: {args.venue})")
    if add_run_course_profile_slide(prs, args.venue):
        print(f"  Run course profile slide added (venue: {args.venue})")

    # Race-day forecast — race-start times pulled from EVENT_SCHEDULES per venue.
    venue_schedule = EVENT_SCHEDULES.get(args.venue.lower().strip())
    race_starts = venue_schedule.get("race_starts") if venue_schedule else None
    if race_date_str and add_race_day_forecast_slide(prs, args.venue, coords,
                                                    race_date_str, race_starts):
        print(f"  Race-day forecast slide added (race_date={race_date_str})")

    if add_travel_load_slide(prs, args.venue):
        print(f"  Travel & arrival slide added (venue: {args.venue})")

    add_environmental_risk_slide(prs, args.venue, unique_env_rows)

    # ── Para-specific historical analysis (per-gender sections) ───────────────
    if args.para and not para_events_df.empty:
        print(f"  Adding Para Series historical sections ({len(para_events_df)} programs)...")
        for gender_label in ("Men", "Women"):
            progs = _para_gender_progs(para_events_df, gender_label, para_classes)
            if not progs:
                continue
            print(f"    Para {gender_label}: {len(progs)} sport class(es)")
            add_section_divider(prs, f"Para {gender_label}", prefix="")
            add_para_results_summary_slide(prs, para_events_df, engine, args.venue,
                                           gender_label, para_classes)
            add_para_split_trend_slide(prs, para_events_df, engine, args.venue,
                                       gender_label, para_classes,
                                       split="total", split_label="Winning Time")

    # Preview venues (no prior races, upcoming event not yet in the DB) still get
    # the ranking-driven slides — they query athlete_rankings, not the event. Tier
    # labels fall back to VENUE_PREVIEW's event_name/cat_name so the max-points
    # scale is right even with no DB event to read it from.
    _preview_meta = VENUE_PREVIEW.get(args.venue.lower().strip(), {})
    _tier_event_name = (upcoming.get("event_name") if upcoming
                        else _preview_meta.get("event_name"))
    _tier_cat_name   = (upcoming.get("cat_name") if upcoming
                        else _preview_meta.get("cat_name"))

    if prior_men or prior_women or upcoming or _preview_meta:
        add_who_to_watch_slide(
            prs, engine, args.venue,
            upcoming_event_id=upcoming["event_id"] if upcoming else None,
            prior_men={"event_id": prior_men["event_id"], "prog_id": prior_men["prog_id"]} if prior_men else None,
            prior_women={"event_id": prior_women["event_id"], "prog_id": prior_women["prog_id"]} if prior_women else None,
        )
        add_top_swim_threats_slide(prs, engine, args.venue,
                                   upcoming_event_id=upcoming["event_id"] if upcoming else None)
        add_usa_world_rankings_slide(prs, engine, args.venue,
                                     upcoming_event_id=upcoming["event_id"] if upcoming else None,
                                     event_name=_tier_event_name,
                                     cat_name=_tier_cat_name)
        add_oqr_pace_slide(prs, engine, args.venue)
        add_usa_oqr_snapshot_slide(prs, engine, args.venue)

    # Para mode is Para-focused: suppress Elite gender sections unless the user
    # opts in with --include-elite (or there is simply no Para data to show).
    para_focus = args.para and not para_events_df.empty and not args.include_elite
    preview_only = (args.preview_only
                    or para_focus
                    or (args.para and not men_data and not women_data))

    gender_sections = []
    if preview_only:
        print("  Preview-only mode — skipping Elite gender-section slides")
    else:
        if men_data   and args.gender in ("men",   "both"): gender_sections.append(("Men",   men_data))
        if women_data and args.gender in ("women", "both"): gender_sections.append(("Women", women_data))

    for gender_label, data in gender_sections:
        print(f"  Building {gender_label} slides ({len(data)} races)...")
        add_section_divider(prs, gender_label)
        add_race_overview_slide(prs, gender_label, data, args.venue, venue_content)
        add_overview_slide(prs, data, gender_label, args.venue)

        # Deep-dive slides (single-race detail) — insert after Results Summary
        if deep_dive:
            gkey = gender_label.lower()
            splits_df = detailed.get(gkey)
            prior = data[-1]
            if splits_df is not None and not splits_df.empty:
                add_pack_dynamics_slide(prs, splits_df, gender_label, args.venue, prior["year"])
                add_top_splits_vs_season_slide(
                    prs, splits_df, engine, gender_label, args.venue,
                    prior["year"], prior["event_id"], prior["prog_id"], prior["date"],
                )

        add_swim_slide(prs, data, gender_label, args.venue)
        add_bike_evolution_slide(prs, data, gender_label, args.venue)
        add_run_slide(prs, data, gender_label, args.venue)
        add_position_times_slide(prs, data, gender_label, args.venue)
        add_weather_slide(prs, data, gender_label, args.venue)

    prs.save(fname)
    print(f"\nDone! Saved to: {fname}")


if __name__ == "__main__":
    main()
